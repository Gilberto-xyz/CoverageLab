
# --- START OF FILE coverage_g_mod_v11_Optimized.py ---
print("Iniciando script de coberturas...")

# Bibliotecas necesarias
# --------------------------------------------------------------------------------------------------
from matplotlib import pyplot as plt
import dataframe_image as dfi
import pandas as pd
import numpy as np
import warnings
import matplotlib
matplotlib.use("Agg")   
import os
import re
import io
from matplotlib import pyplot as plt
from datetime import datetime as dt, timedelta
from scipy.stats import pearsonr
from pptx import Presentation
from pptx.util import Inches
from openpyxl.utils import get_column_letter
from tqdm import tqdm
import colorama
from colorama import Fore, Style
import matplotlib.ticker as mtick
from matplotlib import pyplot as plt
from matplotlib.dates import MonthLocator, DateFormatter
import matplotlib.style # Import specific module for context
from rich.progress import Progress, BarColumn, TextColumn, TimeElapsedColumn, TimeRemainingColumn, SpinnerColumn
from PIL import Image, ImageOps

# Instalar las bibliotecas necesarias si no están instaladas
# pip install pandas numpy matplotlib openpyxl tqdm colorama rich dataframe_image scipy python-pptx


# --- Configuraciones Globales ---
colorama.init(autoreset=True)
pd.set_option('future.no_silent_downcasting', True)
pd.set_option('mode.chained_assignment', None)
warnings.filterwarnings('ignore')

# --- Constantes ---
PPT_LAYOUT_INDEX = 21 # Layout usado para las diapositivas de gráficos
DEFAULT_POP_COVERAGE = "100%"
EXCEL_TEMP_FILENAME = 'file_temp_coverage.xlsx' # Nombre temporal para el Excel

# Nombres de Columnas estándar
COL_DATA = 'Data'
COL_SELL_IN = 'Sell_in'
COL_SELL_OUT = 'Sell_out'
COL_PENET = 'Penet'
COL_COMPRA_MEDIA = 'Compra_Media'
COL_COMPRA_OCA = 'Compra_por_Oca'
COL_FREQ = 'Freq'
COL_BUYERS = 'Buyers'
COL_SELL_IN_SIM = 'Sell_in_sim'
COL_ACUM_SELL_OUT = 'Acum_Sell_out'
COL_ACUM_SELL_IN = 'Acum_Sell_in'
COL_ANO = 'Ano'
COL_TRI = 'Tri'
COL_SEM = 'Sem'

# Colores para gráficos
COLOR_KANTAR_LINE = '#2C3E50'        # Gris azulado oscuro
COLOR_SELLIN_LINE = '#D4AC0D'        # Dorado
COLOR_KANTAR_BAR_VAR = '#7F8C8D'    # Gris medio
COLOR_SELLIN_BAR_VAR = '#F1C40F'    # Dorado más claro
COLOR_KANTAR_EDGE_VAR = '#2C3E50'
COLOR_SELLIN_EDGE_VAR = '#B7950B'
COLOR_COBERTURA_BAR = '#D9D9D9'      # Gris claro (igual que antes)
COLOR_PENETRACION_BAR = '#FFC000'    # Naranja/Amarillo (igual que antes)
COLOR_SELLIN_TREND_LINE = '#F8E092'  # Amarillo pálido (igual que antes)
COLOR_SELLOUT_TREND_LINE = 'black'   # Negro (igual que antes)
COLOR_POS_LABEL = '#1e8449'          # Verde oscuro
COLOR_NEG_LABEL = '#8b0000'          # Rojo oscuro
COLOR_POS_LABEL_ALT = '#27ae60'       # Verde más brillante
COLOR_NEG_LABEL_ALT = '#c0392b'       # Rojo más brillante

# --- Datos Embebidos (Categorías) ---
# Conservar todo en el script pero en formato CSV embebido reducirá el ruido visual
# y seguirá siendo ligero.
CATEGORIES_CSV_DATA = """cod,cest,cat
ALCB,Bebidas,Bebidas Alcohólicas
BEER,Bebidas,Cervezas
CARB,Bebidas,Bebidas Gaseosas
CWAT,Bebidas,Agua Gasificada
COCW,Bebidas,Água de Coco
COFF,Bebidas,Café-Consolidado de Café
CRBE,Bebidas,Cross Category (Bebidas)
ENDR,Bebidas,Bebidas Energéticas
FLBE,Bebidas,Bebidas Saborizadas Sin Gas
GCOF,Bebidas,Café Tostado y Molido
HJUI,Bebidas,Jugos Caseros
ITEA,Bebidas,Té Helado
ICOF,Bebidas,Café Instantáneo-Café Sucedáneo
JUNE,Bebidas,Jugos y Nectares
VEJU,Bebidas,Zumos de Vegetales
WATE,Bebidas,Agua Natural
CSDW,Bebidas,Gaseosas + Aguas
MXCM,Bebidas,Mixta Café+Malta
MXDG,Bebidas,Mixta Dolce Gusto-Mixta Té Helado + Café + Modificadores
MXJM,Bebidas,Mixta Jugos y Leches
MXJS,Bebidas,Mixta Jugos Líquidos + Bebidas de Soja
MXTC,Bebidas,Mixta Té+Café
JUIC,Bebidas,Jugos Liquidos-Jugos Polvo
PWDJ,Bebidas,Refrescos en Polvo-Jugos - Bebidas Instantáneas En Polvo - Jugos Polvo
RFDR,Bebidas,Bebidas Refrescantes
RTDJ,Bebidas,Refrescos Líquidos-Jugos Líquidos
RTEA,Bebidas,Té Líquido - Listo para Tomar
SOYB,Bebidas,Bebidas de Soja
SPDR,Bebidas,Bebidas Isotónicas
TEAA,Bebidas,Té e Infusiones-Te-Infusión Hierbas
YERB,Bebidas,Yerba Mate
BUTT,Lacteos,Manteca
CHEE,Lacteos,Queso Fresco y para Untar
CMLK,Lacteos,Leche Condensada
CRCH,Lacteos,Queso Untable
DYOG,Lacteos,Yoghurt p-beber
EMLK,Lacteos,Leche Culinaria-Leche Evaporada
FRMM,Lacteos,Leche Fermentada
FMLK,Lacteos,Leche Líquida Saborizada-Leche Líquida Con Sabor
FRMK,Lacteos,Fórmulas Infantiles
LQDM,Lacteos,Leche Líquida
LLFM,Lacteos,Leche Larga Vida
MARG,Lacteos,Margarina
MCHE,Lacteos,Queso Fundido
MKCR,Lacteos,Crema de Leche
MXDI,Lacteos,Mixta Lácteos-Postre+Leches+Yogurt
MXMI,Lacteos,Mixta Leches
MXYD,Lacteos,Mixta Yoghurt+Postres
PTSS,Lacteos,Petit Suisse
PWDM,Lacteos,Leche en Polvo
SYOG,Lacteos,Yoghurt p-comer
MILK,Lacteos,Leche-Leche Líquida Blanca - Leche Liq. Natural
YOGH,Lacteos,Yoghurt
CLOT,Ropas y Calzados,Ropas
FOOT,Ropas y Calzados,Calzados
SOCK,Ropas y Calzados,Medias-Calcetines
AREP,Alimentos,Arepas
BCER,Alimentos,Cereales Infantiles
BABF,Alimentos,Nutrición Infantil-Colados y Picados
BEAN,Alimentos,Frijoles
BISC,Alimentos,Galletas
BOUI,Alimentos,Caldos-Caldos y Sazonadores
BREA,Alimentos,Pan
BRCR,Alimentos,Apanados-Empanizadores
BRDC,Alimentos,Empanados
CERE,Alimentos,Cereales-Cereales Desayuno-Avenas y Cereales
BURG,Alimentos,Hamburguesas
CCMX,Alimentos,Mezclas Listas para Tortas-Preparados Base Harina Trigo
CAKE,Alimentos,Queques-Ponques Industrializados
FISH,Alimentos,Conservas De Pescado
CFAV,Alimentos,Conservas de Frutas y Verduras
CRML,Alimentos,Dulce de Leche-Manjar
CMLC,Alimentos,Alfajores
CBAR,Alimentos,Barras de Cereal
CHCK,Alimentos,Pollo
CHOC,Alimentos,Chocolate
COCO,Alimentos,Chocolate de Taza-Achocolatados - Cocoas
COLS,Alimentos,Salsas Frías
COMP,Alimentos,Compotas
SPIC,Alimentos,Condimentos y Especias
CKCH,Alimentos,Chocolate de Mesa
COIL,Alimentos,Aceite-Aceites Comestibles
CSAU,Alimentos,Salsas Listas-Salsas Caseras Envasadas
CNML,Alimentos,Grano - Harina y Masa de Maíz
CNST,Alimentos,Fécula de Maíz
CNFL,Alimentos,Harina De Maíz
CAID,Alimentos,Ayudantes Culinarios
DESS,Alimentos,Postres Preparados
DHAM,Alimentos,Jamón Endiablado
DFNS,Alimentos,Semillas y Frutos Secos
EBRE,Alimentos,Pan de Pascua
EEGG,Alimentos,Huevos de Páscua
EGGS,Alimentos,Huevos
FLSS,Alimentos,Flash Cecinas
FLOU,Alimentos,Harinas
MEAT,Alimentos,Carne Fresca
FRDS,Alimentos,Platos Listos Congelados
FRFO,Alimentos,Alimentos Congelados
HAMS,Alimentos,Jamones
HCER,Alimentos,Cereales Calientes-Cereales Precocidos
HOTS,Alimentos,Salsas Picantes
ICEC,Alimentos,Helados
IBRE,Alimentos,Pan Industrializado
IMPO,Alimentos,Puré Instantáneo
INOO,Alimentos,Fideos Instantáneos
JAMS,Alimentos,Mermeladas
KETC,Alimentos,Ketchup
LJDR,Alimentos,Jugo de Limon Adereso
MALT,Alimentos,Maltas
SEAS,Alimentos,Adobos - Sazonadores
MAYO,Alimentos,Mayonesa
MEAT,Alimentos,Cárnicos # Segunda ocurrencia de MEAT
MLKM,Alimentos,Modificadores de Leche-Saborizadores p-leche
MXCO,Alimentos,Mixta Cereales Infantiles+Avenas
MXBS,Alimentos,Mixta Caldos + Saborizantes
MXSB,Alimentos,Mixta Caldos + Sopas
MXCH,Alimentos,Mixta Cereales + Cereales Calientes
MXCC,Alimentos,Mixta Chocolate + Manjar
MXSN,Alimentos,Galletas - snacks y mini tostadas
COBT,Alimentos,Aceites + Mantecas
COCF,Alimentos,Aceites + Conservas De Pescado
CABB,Alimentos,Ayudantes Culinarios + Bolsa de Hornear
MXEC,Alimentos,Mixta Huevos de Páscua + Chocolates
MXDP,Alimentos,Mixta Platos Listos Congelados + Pasta
MXFR,Alimentos,Mixta Platos Congelados y Listos para Comer
MXFM,Alimentos,Mixta Alimentos Congelados + Margarina
MXMC,Alimentos,Mixta Modificadores + Cocoa
MXPS,Alimentos,Mixta Pastas
MXSO,Alimentos,Mixta Sopas+Cremas+Ramen
MXSP,Alimentos,Mixta Margarina + Mayonesa + Queso Crema
MXSW,Alimentos,Mixta Azúcar+Endulzantes
MUST,Alimentos,Mostaza
NDCR,Alimentos,Sustitutos de Crema
NOOD,Alimentos,Fideos
NUGG,Alimentos,Nuggets
OAFL,Alimentos,Avena en hojuelas-liquidas
OLIV,Alimentos,Aceitunas
PANC,Alimentos,Tortilla
PANE,Alimentos,Panetón
PAST,Alimentos,Pastas
PSAU,Alimentos,Salsas para Pasta
PNOU,Alimentos,Turrón de maní
PORK,Alimentos,Carne Porcina
PPMX,Alimentos,Postres en Polvo-Postres para Preparar - Horneables-Gelificables
PWSM,Alimentos,Leche de Soya en Polvo
PCCE,Alimentos,Cereales Precocidos
DOUG,Alimentos,Masas Frescas-Tapas Empanadas y Tarta
PPIZ,Alimentos,Pre-Pizzas
REFR,Alimentos,Meriendas listas
RICE,Alimentos,Arroz
RBIS,Alimentos,Galletas de Arroz
RTEB,Alimentos,Frijoles Procesados
RTEM,Alimentos,Pratos Prontos - Comidas Listas
SDRE,Alimentos,Aderezos para Ensalada
SALT,Alimentos,Sal
SLTC,Alimentos,Galletas Saladas-Galletas No Dulce
SARD,Alimentos,Sardina Envasada
SAUS,Alimentos,Cecinas
SCHN,Alimentos,Milanesas
SNAC,Alimentos,Snacks
SNOO,Alimentos,Fideos Sopa
SOUP,Alimentos,Sopas-Sopas Cremas
SOYS,Alimentos,Siyau
SPAG,Alimentos,Tallarines-Spaguetti
SPCH,Alimentos,Chocolate para Untar
SUGA,Alimentos,Azucar
SWCO,Alimentos,Galletas Dulces
SWSP,Alimentos,Untables Dulces
SWEE,Alimentos,Endulzantes
TOAS,Alimentos,Torradas - Tostadas
TOMA,Alimentos,Salsas de Tomate
TUNA,Alimentos,Atún Envasado
VMLK,Alimentos,Leche Vegetal
WFLO,Alimentos,Harinas de trigo
AIRC,Cuidado del Hogar,Ambientadores-Desodorante Ambiental
BARS,Cuidado del Hogar,Jabón en Barra-Jabón de lavar
BLEA,Cuidado del Hogar,Cloro-Lavandinas-Lejías-Blanqueadores
CBLK,Cuidado del Hogar,Pastillas para Inodoro
CGLO,Cuidado del Hogar,Guantes de látex
CLSP,Cuidado del Hogar,Esponjas de Limpieza-Esponjas y paños
CLTO,Cuidado del Hogar,Utensilios de Limpieza
FILT,Cuidado del Hogar,Filtros de Café
CRHC,Cuidado del Hogar,Cross Category (Limpiadores Domesticos)
CRLA,Cuidado del Hogar,Cross Category (Lavandería)
CRPA,Cuidado del Hogar,Cross Category (Productos de Papel)
DISH,Cuidado del Hogar,Lavavajillas-Lavaplatos - Lavalozas mano
DPAC,Cuidado del Hogar,Empaques domésticos-Bolsas plásticas-Plástico Adherente-Papel encerado-Papel aluminio
DRUB,Cuidado del Hogar,Destapacañerias
FBRF,Cuidado del Hogar,Perfumantes para Ropa-Perfumes para Ropa
FWAX,Cuidado del Hogar,Cera p-pisos
FDEO,Cuidado del Hogar,Desodorante para Pies
FRNP,Cuidado del Hogar,Lustramuebles
GBBG,Cuidado del Hogar,Bolsas de Basura
GCLE,Cuidado del Hogar,Limpiadores verdes
CLEA,Cuidado del Hogar,Limpiadores-Limpiadores y Desinfectantes
INSE,Cuidado del Hogar,Insecticidas-Raticidas
KITT,Cuidado del Hogar,Toallas de papel-Papel Toalla - Toallas de Cocina - Rollos Absorbentes de Papel
LAUN,Cuidado del Hogar,Detergentes para ropa
LSTA,Cuidado del Hogar,Apresto
MXBC,Cuidado del Hogar,Mixta Pastillas para Inodoro + Limpiadores
MXHC,Cuidado del Hogar,Mixta Home Care-Cloro-Limpiadores-Ceras-Ambientadores
MXCB,Cuidado del Hogar,Mixta Limpiadores + Cloro
MXLB,Cuidado del Hogar,Mixta Detergentes + Cloro
MXLD,Cuidado del Hogar,Mixta Detergentes + Lavavajillas
CRTO,Cuidado del Hogar,Pañitos + Papel Higienico
NAPK,Cuidado del Hogar,Servilletas
PLWF,Cuidado del Hogar,Film plastico e papel aluminio
SCOU,Cuidado del Hogar,Esponjas de Acero
SOFT,Cuidado del Hogar,Suavizantes de Ropa
STRM,Cuidado del Hogar,Quitamanchas-Desmanchadores
TOIP,Cuidado del Hogar,Papel Higiénico
WIPE,Cuidado del Hogar,Paños de Limpieza
ANLG,OTC,Analgésicos-Painkillers
FSUP,OTC,Suplementos alimentares
GMED,OTC,Gastrointestinales-Efervescentes
VITA,OTC,Vitaminas y Calcio
nan,Otros,Categoría Desconocida
BATT,Otros,Pilas-Baterías
CGAS,Otros,Combustible Gas
PFHH,Otros,Panel Financiero de Hogares
PFIN,Otros,Panel Financiero de Hogares
INKC,Otros,Cartuchos de Tintas
PETF,Otros,Alimento para Mascota-Alim.p - perro - gato
TELE,Otros,Telecomunicaciones - Convergencia
TILL,Otros,Tickets - Till Rolls
TOBA,Otros,Tabaco - Cigarrillos
ADIP,Cuidado Personal,Incontinencia de Adultos
BSHM,Cuidado Personal,Shampoo Infantil
RAZO,Cuidado Personal,Maquinas de Afeitar
BDCR,Cuidado Personal,Cremas Corporales
CWIP,Cuidado Personal,Paños Húmedos
COMB,Cuidado Personal,Cremas para Peinar
COND,Cuidado Personal,Acondicionador-Bálsamo
CRHY,Cuidado Personal,Cross Category (Higiene)
CRPC,Cuidado Personal,Cross Category (Personal Care)
DEOD,Cuidado Personal,Desodorantes
DIAP,Cuidado Personal,Pañales-Pañales Desechables
FCCR,Cuidado Personal,Cremas Faciales
FTIS,Cuidado Personal,Pañuelos Faciales
FEMI,Cuidado Personal,Protección Femenina-Toallas Femeninas
FRAG,Cuidado Personal,Fragancias
HAIR,Cuidado Personal,Cuidado del Cabello-Hair Care
HRCO,Cuidado Personal,Tintes para el Cabello-Tintes - Tintura - Tintes y Coloración para el cabello
HREM,C cuidado Personal,Depilación
HRST,Cuidado Personal,Alisadores para el Cabello
HSTY,Cuidado Personal,Fijadores para el Cabello-Modeladores-Gel-Fijadores para el cabello
HRTR,Cuidado Personal,Tratamientos para el Cabello
LINI,Cuidado Personal,Óleo Calcáreo
MAKE,Cuidado Personal,Maquillaje-Cosméticos
MEDS,Cuidado Personal,Jabón Medicinal
CRDT,Cuidado Personal,Pañitos + Pañales # Nota: CRDT duplicado - verificar
MXMH,Cuidado Personal,Mixta Make Up+Tinturas
MOWA,Cuidado Personal,Enjuague Bucal-Refrescante Bucal
ORAL,Cuidado Personal,Cuidado Bucal
SPAD,Cuidado Personal,Protectores Femeninos
STOW,Cuidado Personal,Toallas Femininas
SHAM,Cuidado Personal,Shampoo
SHAV,Cuidado Personal,Afeitado-Crema afeitar-Loción de afeitar-Pord. Antes del afeitado
SKCR,Cuidado Personal,Cremas Faciales y Corporales-Cremas de Belleza - Cremas Cuerp y Faciales
SUNP,Cuidado Personal,Protección Solar
TALC,Cuidado Personal,Talcos-Talco para pies
TAMP,Cuidado Personal,Tampones Femeninos
TOIL,Cuidado Personal,Jabón de Tocador
TOOB,Cuidado Personal,Cepillos Dentales
TOOT,Cuidado Personal,Pastas Dentales
BAGS,Material Escolar,Morrales y MAletas Escoalres
CLPC,Material Escolar,Lapices de Colores
GRPC,Material Escolar,Lapices De Grafito
MRKR,Material Escolar,Marcadores
NTBK,Material Escolar,Cuadernos
SCHS,Material Escolar,Útiles Escolares
CSTD,Diversos,Estudio de Categorías
CORP,Diversos,Corporativa
CROS,Diversos,Cross Category
CRBA,Diversos,Cross Category (Bebés)
CRBR,Diversos,Cross Category (Desayuno)-Yogurt - Cereal - Pan - Queso
CRDT,Diversos,Cross Category (Diet y Light) # Segunda ocurrencia de CRDT
CRDF,Diversos,Cross Category (Alimentos Secos)
CRFO,Diversos,Cross Category (Alimentos)
CRSA,Diversos,Cross Category (Salsas)-Mayonesas-Ketchup - Salsas Frías
CRSN,Diversos,Cross Category (Snacks) # Nota: CRSN duplicado- verificar
DEMO,Diversos,Demo
FLSH,Diversos,Flash
HLVW,Diversos,Holistic View
COCP,Diversos,Mezcla para café instantaneo y crema no láctea
CRSN,Diversos,Mezclas nutricionales y suplementos # Segunda ocurrencia de CRSN
MULT,Diversos,Consolidado-Multicategory
PCHK,Diversos,Pantry Check
STCK,Diversos,Inventario
MIHC,Diversos,Leche y Cereales Calientes-Cereales Precocidos y Leche Líquida Blanca
"""

# --- Datos Estáticos (País, Cobertura Poblacional) ---
pais = pd.DataFrame(
    {
      'cod': [10, 54, 91, 55, 12, 56, 57, 93, 52, 51],
      'pais': ['LatAm', 'Argentina', 'Bolivia', 'Brasil', 'CAM', 'Chile', 'Colombia', 'Ecuador', 'Mexico', 'Peru']
    }
)

pop_coverage = {
    "Argentina": "90%", "Bolivia": "60%", "Brasil": "82%", "CAM": "63%",
    "Chile": "78%", "Colombia": "65%", "Costa Rica": "70%", "Ecuador": "55%",
    "El Salvador": "70%", "Guatemala": "70%", "Honduras": "70%", "Mexico": "64%",
    "Nicaragua": "70%", "Panamá": "70%", "Perú": "66%", "RD": "63.29%"
    # Añadir más si es necesario
}

# --- Función para cargar categorías ---
def load_categories():
    """Carga el catálogo de categorías desde el string embebido."""
    try:
        categories_file = io.StringIO(CATEGORIES_CSV_DATA)
        df = pd.read_csv(categories_file, dtype={'cod': str}).set_index('cod')
        if df.index.duplicated().any():
            duplicates = df.index[df.index.duplicated()].unique().tolist()
            print(
                f"{Fore.YELLOW}Advertencia: Se encontraron códigos de categoría duplicados en los datos embebidos: {duplicates}. Se usará la última entrada encontrada para cada código."
            )
        print(Fore.GREEN + "Datos de categorías cargados correctamente desde el script.")
        return df
    except Exception as e:
        print(f"{Fore.RED}{Style.BRIGHT}Error Crítico al cargar datos de categorías desde el string embebido: {e}")
        exit()

# --- Variables Globales y Funciones de Utilidad ---
SELECTIONS = {} # Para guardar las respuestas interactivas

def clear_and_print_summary():
    """Limpia la terminal y muestra un resumen de las selecciones del usuario."""
    os.system('cls' if os.name == 'nt' else 'clear') # Compatible con Windows y Linux/Mac
    print(Fore.CYAN + Style.BRIGHT + "Resumen de opciones seleccionadas:")
    if 'Excel' in SELECTIONS:
        print(Fore.BLUE + "Archivo Excel: " + Fore.YELLOW + f"{SELECTIONS['Excel']}")
    if 'Cobertura' in SELECTIONS:
        print(Fore.BLUE + "Tipo de cobertura: " + Fore.YELLOW + f"{SELECTIONS['Cobertura']}")
    if 'Razón' in SELECTIONS:
        print(Fore.BLUE + "Razón de Cobertura: " + Fore.YELLOW + f"{SELECTIONS['Razón']}")
    print("\n" + "-"*50 + "\n")

def calc_var1(df, coluna, p):
    """
    Calcula variaciones vs período anterior (Y-1) en Python.

    Args:
        df (pd.DataFrame): DataFrame con los datos.
        coluna (str): Nombre de la columna a calcular (e.g., COL_SELL_OUT).
        p (int): Pipeline (shift para Sell_in).

    Returns:
        list: Lista con variaciones [Anual, Semestral, Trimestral].
              Retorna NaN para cálculos imposibles (datos insuficientes).
    """
    n_rows = len(df)
    variations = []

    # Anual (12 vs 12 meses)
    if n_rows >= 24 + p:
        current_sum = df[coluna][n_rows-12-p : n_rows-p].sum() if p != 0 else df[coluna][-12:].sum()
        previous_sum = df[coluna][n_rows-24-p : n_rows-12-p].sum() if p!= 0 else df[coluna][-24:-12].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    # Semestral (6 vs 6 meses)
    if n_rows >= 12 + p:
        current_sum = df[coluna][n_rows-6-p : n_rows-p].sum() if p != 0 else df[coluna][-6:].sum()
        previous_sum = df[coluna][n_rows-12-p : n_rows-6-p].sum() if p!= 0 else df[coluna][-12:-6].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    # Trimestral (3 vs 3 meses)
    if n_rows >= 6 + p:
        current_sum = df[coluna][n_rows-3-p : n_rows-p].sum() if p != 0 else df[coluna][-3:].sum()
        previous_sum = df[coluna][n_rows-6-p : n_rows-3-p].sum() if p!= 0 else df[coluna][-6:-3].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    return variations


def calc_var2(df, coluna, p):
    """
    Calcula variaciones vs período retrasado (Y-2) en Python.

    Args:
        df (pd.DataFrame): DataFrame con los datos.
        coluna (str): Nombre de la columna a calcular (e.g., COL_SELL_OUT).
        p (int): Pipeline (shift para Sell_in).

    Returns:
        list: Lista con variaciones [Anual, Semestral, Trimestral].
              Retorna NaN para cálculos imposibles (datos insuficientes).
    """
    n_rows = len(df)
    variations = []

    # Anual (12 meses actuales vs 12 meses de hace 2 años)
    if n_rows >= 36 + p:
        current_sum = df[coluna][n_rows-12-p : n_rows-p].sum() if p != 0 else df[coluna][-12:].sum()
        previous_sum = df[coluna][n_rows-36-p : n_rows-24-p].sum() if p!= 0 else df[coluna][-36:-24].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    # Semestral (6 meses actuales vs 6 meses de hace 2 años) - CORREGIDO
    if n_rows >= 30 + p: # Necesitamos 6 actuales + 24 para ir 2 años atrás
        current_sum = df[coluna][n_rows-6-p : n_rows-p].sum() if p != 0 else df[coluna][-6:].sum()
        previous_sum = df[coluna][n_rows-30-p : n_rows-24-p].sum() if p!= 0 else df[coluna][-30:-24].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    # Trimestral (3 meses actuales vs 3 meses de hace 2 años) - CORREGIDO
    if n_rows >= 27 + p: # Necesitamos 3 actuales + 24 para ir 2 años atrás
        current_sum = df[coluna][n_rows-3-p : n_rows-p].sum() if p != 0 else df[coluna][-3:].sum()
        previous_sum = df[coluna][n_rows-27-p : n_rows-24-p].sum() if p!= 0 else df[coluna][-27:-24].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    return variations


def escalona(df_to_scale):
    """
    Desplaza los valores de cada columna hacia abajo, rellenando con NaN al principio.
    Se utiliza para alinear datos en fórmulas de Excel para cálculos de cobertura.

    Args:
        df_to_scale (pd.DataFrame): DataFrame cuyas columnas serán escalonadas.
    """
    for col in df_to_scale.columns:
        col_idx = df_to_scale.columns.get_loc(col)
        values = list(df_to_scale[col].values)
        # Invierte, trunca desde el inicio según índice, rellena, invierte de nuevo
        scaled_values = (values[::-1][col_idx:] + [np.nan]*col_idx)[::-1]
        df_to_scale[col] = scaled_values

def razao_cov():
    """Solicita interactivamente la razón de la cobertura y la devuelve."""
    print(Fore.CYAN + "\nPregunta: ¿Cuál es la razón de la cobertura?")
    print(Fore.MAGENTA + "Opciones:")
    print(Fore.MAGENTA + "1 - Actualización periódica por contrato")
    print(Fore.MAGENTA + "2 - Conocer nivel de cobertura o pipeline")
    print(Fore.MAGENTA + "3 - Tendencias Contrarias")
    print(Fore.MAGENTA + "4 - Renovación de contrato")
    print(Fore.MAGENTA + "5 - Otras")

    razones = {
        '1': "Actualización periódica por contrato",
        '2': "Conocer nivel de cobertura o pipeline",
        '3': "Tendencias Contrarias",
        '4': "Renovación de contrato",
        '5': "Otras"
    }
    eleccion = input(Fore.GREEN + "Elija el número de la opción (1-5): ")
    razon_seleccionada = razones.get(eleccion, "Otras") # Default a 'Otras'
    SELECTIONS['Razón'] = razon_seleccionada
    clear_and_print_summary()
    return razon_seleccionada

def tipo_cobertura():
    """Solicita interactivamente el tipo de cobertura (Absoluta o Relativa)."""
    print(Fore.CYAN + "\nPregunta: ¿Qué tipo de cobertura se calculará?")
    print(Fore.MAGENTA + "Opciones:")
    print(Fore.MAGENTA + "1 - Cobertura Absoluta")
    print(Fore.MAGENTA + "2 - Cobertura Relativa")
    tipos = {'1': "Absoluta", '2': "relativa"}
    eleccion = input(Fore.GREEN + "Elija 1 o 2: ")
    tipo_seleccionado = tipos.get(eleccion, "Absoluta") # Default a 'Absoluta'
    SELECTIONS['Cobertura'] = tipo_seleccionado
    clear_and_print_summary()
    return tipo_seleccionado

def tipo_eje_tendencia():
    """
    Pregunta al usuario si desea el gráfico de tendencia en un solo eje o doble eje (Kantar en eje secundario).
    """
    print(Fore.CYAN + "\n¿Desea el gráfico de tendencia con doble eje?")
    print(Fore.MAGENTA + "1 - Un solo eje (Sell-in y Kantar juntos)")
    print(Fore.MAGENTA + "2 - Doble eje (Kantar en eje secundario)")
    opciones = {'1': "simple", '2': "doble"}
    eleccion = input(Fore.GREEN + "Elija 1 o 2: ")
    tipo_eje = opciones.get(eleccion, "simple")
    SELECTIONS['Eje tendencia'] = tipo_eje
    clear_and_print_summary()
    return tipo_eje

def load_and_preprocess_sheet(excel_file_obj, sheet_name):
    """
    Carga una hoja del archivo Excel, la preprocesa (renombra, limpia, fechas)
    y devuelve el DataFrame procesado y la unidad de medida.

    Args:
        excel_file_obj (pd.ExcelFile): Objeto ExcelFile abierto.
        sheet_name (str): Nombre de la hoja a procesar.

    Returns:
        tuple: (pd.DataFrame, str) - El DataFrame procesado y la unidad de medida.
               Retorna (None, None) si hay un error al cargar o procesar.
    """
    try:
        df_sheet = excel_file_obj.parse(sheet_name)

        # Validar estructura mínima esperada (al menos 2 filas, 8 columnas)
        if df_sheet.shape[0] < 2 or df_sheet.shape[1] < 8:
             print(f"{Fore.YELLOW}Advertencia: La hoja '{sheet_name}' tiene una estructura inesperada ({df_sheet.shape[0]} filas, {df_sheet.shape[1]} columnas). Se omitirá.")
             return None, None

        # Obtiene la 'unidad' o 'medida' de la primera fila, columna 2 (índice 1)
        measure = str(df_sheet.iat[0, 1]).replace('Weighted', '').strip()

        # Renombra las columnas al formato estándar
        df_sheet.columns = [COL_DATA, COL_SELL_OUT, COL_PENET, COL_COMPRA_MEDIA, COL_COMPRA_OCA, COL_FREQ, COL_BUYERS, COL_SELL_IN] + list(df_sheet.columns[8:]) # Mantiene columnas extra si existen
        df_sheet = df_sheet.loc[:, [COL_DATA, COL_SELL_IN, COL_SELL_OUT, COL_COMPRA_MEDIA, COL_COMPRA_OCA, COL_FREQ, COL_PENET, COL_BUYERS]] # Reordena y selecciona

        # Elimina la primera fila (encabezados repetidos) y resetea el índice
        df_sheet = df_sheet.iloc[1:].reset_index(drop=True)

        # Convierte la columna "Data" a tipo datetime
        # Maneja posibles errores de formato o valores nulos
        original_dates = df_sheet[COL_DATA].copy() # Guardar original por si falla
        try:
            # Intenta convertir primero todos los que sean strings
            is_string = df_sheet[COL_DATA].apply(lambda x: isinstance(x, str))
            if is_string.any():
                 # Intenta formato específico primero, maneja errores individuales
                 df_sheet.loc[is_string, COL_DATA] = df_sheet.loc[is_string, COL_DATA].apply(
                     lambda x: dt.strptime(x, '%b-%y  ') if isinstance(x, str) and re.match(r'\w{3}-\d{2}\s{2}', x) else x
                 )
            # Convierte el resto (o los ya convertidos) a datetime
            df_sheet[COL_DATA] = pd.to_datetime(df_sheet[COL_DATA], errors='coerce')
        except Exception as e:
             print(f"{Fore.YELLOW}Advertencia: Problema al convertir fechas en hoja '{sheet_name}'. Error: {e}. Se usará la columna original si es posible.")
             df_sheet[COL_DATA] = pd.to_datetime(original_dates, errors='coerce') # Reintentar con la original

        # Eliminar filas donde la fecha no se pudo convertir (NaT)
        initial_rows = len(df_sheet)
        df_sheet.dropna(subset=[COL_DATA], inplace=True)
        if len(df_sheet) < initial_rows:
            print(f"{Fore.YELLOW}Advertencia: Se eliminaron {initial_rows - len(df_sheet)} filas de la hoja '{sheet_name}' por fechas inválidas.")

        if df_sheet.empty:
            print(f"{Fore.YELLOW}Advertencia: La hoja '{sheet_name}' está vacía o no contiene fechas válidas después del preprocesamiento. Se omitirá.")
            return None, None

        # Asegurar tipos numéricos (intentar convertir, rellenar NaN con 0 si falla)
        numeric_cols = [COL_SELL_IN, COL_SELL_OUT, COL_COMPRA_MEDIA, COL_COMPRA_OCA, COL_FREQ, COL_PENET, COL_BUYERS]
        for col in numeric_cols:
            df_sheet[col] = pd.to_numeric(df_sheet[col], errors='coerce').fillna(0)

        # Añade columnas de Año, Trimestre, Semestre
        df_sheet[COL_ANO] = df_sheet[COL_DATA].dt.year
        df_sheet[COL_TRI] = df_sheet[COL_DATA].dt.quarter
        df_sheet[COL_SEM] = (df_sheet[COL_DATA].dt.month - 1) // 6 + 1
        df_sheet[COL_DATA] = df_sheet[COL_DATA].dt.date # Convertir a solo fecha al final

        return df_sheet, measure

    except Exception as e:
        print(f"{Fore.RED}Error crítico al cargar o preprocesar la hoja '{sheet_name}': {e}")
        return None, None


# --- Funciones de Generación de Gráficos ---

def generar_grafico_evolucion_mensual(df_graf, pipeline_meses=0):
    """
    Genera un gráfico de evolución mensual de Kantar vs Sell-in con variación interanual.

    Args:
        df_graf (pd.DataFrame): DataFrame con datos mensuales (col 'Data' debe ser datetime).
        pipeline_meses (int): Número de meses de pipeline para desplazar Sell-in.

    Returns:
        matplotlib.figure.Figure: Figura de matplotlib con el gráfico, o None si no hay datos.
    """
    if df_graf is None or df_graf.empty or len(df_graf) < 24: # Necesita al menos 24 meses para var YOY
        print(f"{Fore.YELLOW}Advertencia: No se puede generar gráfico de evolución mensual. Datos insuficientes (se requieren >= 24 meses).")
        return None

    # Usar contexto de estilo para evitar afectar otros gráficos
    with matplotlib.style.context('seaborn-v0_8-whitegrid'):
        df_plot = df_graf.copy()
        df_plot[COL_DATA] = pd.to_datetime(df_plot[COL_DATA]) # Asegurar datetime

        # Si hay pipeline, desplazar Sell-in y guardar original si es necesario
        if pipeline_meses > 0:
            # df_plot["Sell_in_original"] = df_plot[COL_SELL_IN].copy() # Descomentar si se necesita el original
            df_plot[COL_SELL_IN] = df_plot[COL_SELL_IN].shift(pipeline_meses)

        # Calcular sumas móviles y variaciones interanuales
        df_plot["Kantar_12m"] = df_plot[COL_SELL_OUT].rolling(12).sum()
        df_plot["Sellin_12m"] = df_plot[COL_SELL_IN].rolling(12).sum()
        df_plot["Kantar_yoy"] = ((df_plot["Kantar_12m"] / df_plot["Kantar_12m"].shift(12)) - 1) * 100
        df_plot["Sellin_yoy"] = ((df_plot["Sellin_12m"] / df_plot["Sellin_12m"].shift(12)) - 1) * 100

        # Filtrar NaNs resultantes de rolling/shift
        df_plot = df_plot.dropna(subset=["Kantar_yoy", "Sellin_yoy"]).copy()

        if df_plot.empty:
            print(f"{Fore.YELLOW}Advertencia: No quedan datos para el gráfico de evolución después de calcular YOY.")
            return None

        # Crear figura y ejes con márgenes personalizados
        fig = plt.figure(figsize=(16.5, 8), dpi=100) # Ajustar tamaño si es necesario
        left_margin, right_margin, bottom_margin, top_margin = 0.08, 0.92, 0.18, 0.90
        ax1 = fig.add_axes([left_margin, bottom_margin, right_margin-left_margin, top_margin-bottom_margin])
        ax2 = ax1.twinx()

        # Eje primario (Líneas)
        sellin_label = f"{COL_SELL_IN} (Mensual)" + (f" - P:{pipeline_meses}" if pipeline_meses > 0 else "")
        ax1.plot(df_plot[COL_DATA], df_plot[COL_SELL_OUT], color=COLOR_KANTAR_LINE, marker="o", linewidth=2, markersize=5, label=f"{COL_SELL_OUT} (Mensual)")
        ax1.plot(df_plot[COL_DATA], df_plot[COL_SELL_IN], color=COLOR_SELLIN_LINE, marker="o", linewidth=2, markersize=5, label=sellin_label)
        ax1.set_ylabel("Volumen Mensual", fontsize=11, labelpad=15)
        ax1.tick_params(axis='y', labelsize=9)
        ax1.set_ylim(bottom=0)
        ax1.grid(axis='y', linestyle='--', alpha=0.4)

        # Eje secundario (Barras de Variación YOY)
        width = 8
        offset = 4
        ax2.bar(df_plot[COL_DATA] - pd.DateOffset(days=offset), df_plot["Kantar_yoy"], width=width, color=COLOR_KANTAR_BAR_VAR, edgecolor=COLOR_KANTAR_EDGE_VAR, alpha=0.7, label="% Var Kantar")
        ax2.bar(df_plot[COL_DATA] + pd.DateOffset(days=offset), df_plot["Sellin_yoy"], width=width, color=COLOR_SELLIN_BAR_VAR, edgecolor=COLOR_SELLIN_EDGE_VAR, alpha=0.7, label="% Var Sell-in")
        ax2.set_ylabel("Variación Interanual (%)", fontsize=11, labelpad=15)
        ax2.yaxis.set_major_formatter(mtick.PercentFormatter(decimals=0))
        ax2.tick_params(axis='y', labelsize=9)
        ax2.axhline(y=0, color='gray', linestyle='-', alpha=0.5, linewidth=0.8)

        # Etiquetas en barras
        for _, row in df_plot.iterrows():
            for col_yoy, x_offset, color_pos, color_neg in [("Kantar_yoy", -offset, COLOR_POS_LABEL, COLOR_NEG_LABEL),
                                                             ("Sellin_yoy", offset, COLOR_POS_LABEL_ALT, COLOR_NEG_LABEL_ALT)]:
                if not pd.isna(row[col_yoy]):
                    valor = row[col_yoy]
                    pos_vert = valor + 1 if valor >= 0 else valor - 1
                    va_align = "bottom" if valor >= 0 else "top"
                    color_etiq = color_pos if valor >= 0 else color_neg
                    ax2.text(row[COL_DATA] + pd.Timedelta(days=x_offset), pos_vert, f"{valor:.1f}%",
                             ha="center", va=va_align, fontsize=7, fontweight='bold', color=color_etiq) # Añadir etiquetas, con un decimal y color según el signo

        # Ajustar límites eje Y secundario para dar espacio a etiquetas
        y2_min, y2_max = ax2.get_ylim()
        padding = max(abs(y2_min), abs(y2_max)) * 0.15 # 15% padding
        ax2.set_ylim(y2_min - padding, y2_max + padding*2) # Más espacio arriba

        # Formato Eje X (Fechas)
        fechas_validas = df_plot[COL_DATA] # Obtener fechas válidas para el eje X (no extras)
        ax1.set_xlim([fechas_validas.min(), fechas_validas.max()])
        ax1.xaxis.set_major_locator(MonthLocator(interval=1)) # Ajustar intervalo dinámicamente
        ax1.xaxis.set_major_formatter(DateFormatter('%b-%y'))
        ax1.tick_params(axis='x', rotation=45, labelsize=8)

        # Título y Leyenda
        # titulo = "Evolución Mensual y Variación " + (f" (Pipeline: {pipeline_meses})" if pipeline_meses > 0 else "")
        # fig.suptitle(titulo, fontsize=16, fontweight='bold', y=top_margin + 0.05) # Título de la figura
        lines1, labels1 = ax1.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax2.legend(lines1 + lines2, labels1 + labels2, loc="upper left", bbox_to_anchor=(0.01, 0.98), fontsize=9, frameon=True, framealpha=0.8)

        # No usar tight_layout con add_axes, márgenes manuales ya aplicados
        # fig.tight_layout(rect=[0, 0, 1, 0.95]) # Ajustar rect si el título se solapa
  
        return fig

def generar_grafico_cobertura(slide, marca_clean, pipeline, df_cov_pipe, df_pen_pipe, lang_idx, cov_type_str, labels_dict):
    """
    Genera el gráfico de barras de Cobertura vs Penetración y lo añade al slide.

    Args:
        slide: Objeto slide de pptx.
        marca_clean (str): Nombre de la marca sin prefijo de pipeline.
        pipeline (int): Número de pipeline.
        df_cov_pipe (pd.Series): Serie con datos de cobertura para el pipeline.
        df_pen_pipe (pd.Series): Serie con datos de penetración para el mismo período.
        lang_idx (int): Índice de idioma (1=PT, 2=ES).
        cov_type_str (str): "Absoluta" o "relativa".
        labels_dict (dict): Diccionario con etiquetas de idioma.
    """
    if df_cov_pipe.empty or df_pen_pipe.empty:
         print(f"{Fore.YELLOW}Advertencia: Datos insuficientes para gráfico Cobertura/Penetración (Marca: {marca_clean}, P:{pipeline}).")
         return

    # Crear figura y ejes
    fig_cov, ax_cov = plt.subplots(figsize=(12, 4.25), dpi=100) # Usar subplots

    # Datos y etiquetas del eje X
    cov_data = df_cov_pipe.values
    pen_data = df_pen_pipe.values
    x_labels = [d.strftime('%m-%y') for d in df_cov_pipe.index] # Asegurar formato mes-año corto
    x_pos = np.arange(len(x_labels))
    

    # Barras
    bar_width = 0.35
    
    # Desplazamiento para separar las barras
    offset = bar_width / 2

    # Penetración a la izquierda, cobertura a la derecha
    rects2 = ax_cov.bar(x_pos - offset/1.2, pen_data, bar_width,
                        label=labels_dict.get((lang_idx, 'Graf cob Penet Men'), 'Penetración Mensual'),
                        color=COLOR_PENETRACION_BAR, edgecolor='black', zorder=1)
    rects1 = ax_cov.bar(x_pos + offset, cov_data, bar_width,
                        label=f"{coverage_label}",
                        color=COLOR_COBERTURA_BAR, edgecolor='black', linewidth=2, zorder=2, alpha=0.85)

    # Etiquetas de datos sobre las barras
    for rect_group, offset in [(rects2, 0), (rects1, 0)]:
        for i, rect in enumerate(rect_group):
            height = rect.get_height()
            # Añadir etiqueta solo si la altura es > 0 para evitar solapamiento en cero
            if height > 0.1: # Umbral pequeño
                # Decidir color de fondo de etiqueta
                bbox_props = dict(facecolor='#F2F2F2', edgecolor='black', boxstyle='round,pad=0.3') # Blanco por defecto
                if rect_group == rects1: # Cobertura
                    if i % 12 == (len(rect_group) % 12 -1) : # Último mes de cada año móvil (ajustado)
                         bbox_props['facecolor'] = '#A6A6A6' # Gris
                         bbox_props['edgecolor'] = 'black'
                else: # Penetración
                    bbox_props['facecolor'] = '#FDEAD9' # Amarillo pálido

                ax_cov.annotate(f'{height:.1f}',
                                xy=(rect.get_x() + rect.get_width() / 2, height),
                                xytext=(0, 3),  # 3 puntos de offset vertical
                                textcoords="offset points",
                                ha='center', va='bottom', fontsize=8, bbox=bbox_props)

    # Configuración del gráfico
    ax_cov.set_ylabel(f"{coverage_label} | {labels_dict.get((lang_idx, 'Graf cob Penet Men'), 'Penetración Mensual')}", fontsize=9)
    title_key = 'Titulo Cob'
    ax_cov.set_title(f"{labels_dict.get((lang_idx, title_key), 'Cobertura Año Móvil')} | {marca_clean} Pipeline {pipeline}", size=16)
    ax_cov.set_xticks(x_pos)
    ax_cov.set_xticklabels(x_labels, rotation=30, ha='right', fontsize=9)
    ax_cov.legend(loc='lower center', bbox_to_anchor=(0.5, -0.30), # Ajustar posición leyenda
                  frameon=False, prop={'size': 11}, ncol=2)
    ax_cov.grid(axis='y', linestyle='--', alpha=0.6)
    ax_cov.set_axisbelow(True)
    ax_cov.spines['top'].set_visible(False)
    ax_cov.spines['right'].set_visible(False)
    ax_cov.spines['left'].set_visible(False)
    ax_cov.set_ylim(bottom=0, top=max(np.nanmax(cov_data) if cov_data.size > 0 else 0,
                                       np.nanmax(pen_data) if pen_data.size > 0 else 0) * 1.15) # Añadir padding arriba
    ax_cov.margins(x=0)

    # Guardar y añadir a PPT
    plt.tight_layout() # Ajustar layout
    img_stream = io.BytesIO()
    fig_cov.savefig(img_stream, format='png', bbox_inches='tight', pad_inches=0.1, transparent=True)
    img_stream.seek(0)

    # --- AGREGAR CONTORNO NEGRO ---
    img_pil = Image.open(img_stream)
    bordered = ImageOps.expand(img_pil, border=1, fill='black')  # 1 píxel de borde negro
    img_stream_bordered = io.BytesIO()
    bordered.save(img_stream_bordered, format='PNG')
    img_stream_bordered.seek(0)

    slide.shapes.add_picture(img_stream_bordered, Inches(0.5), Inches(2.0), height=Inches(4.2)) # Ajustar posición/tamaño
    plt.close(fig_cov) # Cerrar figura para liberar memoria

def generar_grafico_tendencia(slide, marca_clean, pipeline, df_plot, lang_idx, labels_dict, doble_eje=False):
    """
    Genera el gráfico de líneas de Tendencia (Sell-in vs Sell-out) y lo añade al slide.
    Si doble_eje=True, Kantar (Sell-out) va en eje secundario.
    """
    if df_plot is None or df_plot.empty or pipeline >= len(df_plot):
         print(f"{Fore.YELLOW}Advertencia: Datos insuficientes para gráfico de Tendencia (Marca: {marca_clean}, P:{pipeline}).")
         return

    fig_trend, ax_trend = plt.subplots(figsize=(13, 5), dpi=100)

    sell_out_data = df_plot[COL_SELL_OUT].iloc[pipeline:].values
    sell_in_data = df_plot[COL_SELL_IN].iloc[:len(df_plot)-pipeline].values
    x_labels = df_plot[COL_DATA].iloc[pipeline:].values

    if len(sell_out_data) != len(sell_in_data):
         print(f"{Fore.RED}Error: Discrepancia de longitud en datos de tendencia para {marca_clean} P:{pipeline}.")
         plt.close(fig_trend)
         return

    if doble_eje:
        ax2 = ax_trend.twinx()
        lns1 = ax_trend.plot(x_labels, sell_in_data, color=COLOR_SELLIN_TREND_LINE, linewidth=4, label=f'{COL_SELL_IN} (P:{pipeline})')
        lns2 = ax2.plot(x_labels, sell_out_data, color=COLOR_SELLOUT_TREND_LINE, linewidth=4, label=COL_SELL_OUT)
        ax_trend.set_ylabel(f'{COL_SELL_IN}', color=COLOR_SELLIN_TREND_LINE, fontsize=11)
        ax2.set_ylabel(f'{COL_SELL_OUT}', color=COLOR_SELLOUT_TREND_LINE, fontsize=11)
        # --- CORRECCIÓN: Configurar ambos ejes para empezar desde 0 ---
        ax_trend.set_ylim(bottom=0)
        ax2.set_ylim(bottom=0)
        lns = lns1 + lns2
        labs = [l.get_label() for l in lns]
        ax2.legend(lns, labs, loc='lower center', bbox_to_anchor=(0.5, -0.28), frameon=False, prop={'size': 11}, ncol=2)
    else:
        lns1 = ax_trend.plot(x_labels, sell_in_data, color=COLOR_SELLIN_TREND_LINE, linewidth=4, label=f'{COL_SELL_IN} (P:{pipeline})')
        lns2 = ax_trend.plot(x_labels, sell_out_data, color=COLOR_SELLOUT_TREND_LINE, linewidth=4, label=COL_SELL_OUT)
        ax_trend.set_ylabel(f'{COL_SELL_IN} / {COL_SELL_OUT}', color='black', fontsize=11)
        ax_trend.set_ylim(bottom=0)
        lns = lns1 + lns2
        labs = [l.get_label() for l in lns]
        ax_trend.legend(lns, labs, loc='lower center', bbox_to_anchor=(0.5, -0.28), frameon=False, prop={'size': 11}, ncol=2)

    ax_trend.tick_params(axis='x', rotation=30, labelsize=9)
    for label in ax_trend.get_xticklabels():
        label.set_ha('right')
    ax_trend.grid(axis='y', linestyle='--', alpha=0.6)
    ax_trend.spines['top'].set_visible(False)
    ax_trend.spines['right'].set_visible(False)
    ax_trend.set_title(f"{labels_dict.get((lang_idx, 'Titulo Vol'), 'Tendencia en Volumen')} | {marca_clean} P:{pipeline}", size=17)

    plt.tight_layout()
    img_stream = io.BytesIO()
    fig_trend.savefig(img_stream, format='png', bbox_inches='tight', pad_inches=0.1, transparent=True)
    img_stream.seek(0)
    img_pil = Image.open(img_stream)
    bordered = ImageOps.expand(img_pil, border=2, fill='black')
    img_stream_bordered = io.BytesIO()
    bordered.save(img_stream_bordered, format='PNG')
    img_stream_bordered.seek(0)
    slide.shapes.add_picture(img_stream_bordered, Inches(0.5), Inches(1.8), height=Inches(4.5))
    plt.close(fig_trend)
    

# --------------------------------------------------------------------------------------------------
# INICIO DEL SCRIPT PRINCIPAL
# --------------------------------------------------------------------------------------------------

# --- Configuración de directorio ---
root_dir = os.path.dirname(os.path.abspath(__file__))
os.chdir(root_dir)

# --- Selección de Archivo y Opciones ---

excel_list = [f for f in os.listdir(root_dir) if f.endswith('.xlsx') and not f.startswith('~$') and f != EXCEL_TEMP_FILENAME] # Evitar temporales

if not excel_list:
    print(f"{Fore.RED}{Style.BRIGHT}Error: No se encontraron archivos .xlsx en la carpeta: {root_dir}")
    exit()

print(Fore.CYAN + "Archivos Excel (.xlsx) encontrados:")
for i, archivo in enumerate(excel_list, start=1):
    print(Fore.MAGENTA + f"{i}. {archivo}")

while True:
    try:
        opcion = int(input(Fore.GREEN + f"Seleccione el número de archivo a procesar (1-{len(excel_list)}): "))
        if 1 <= opcion <= len(excel_list):
            excel_file_name = excel_list[opcion - 1]
            SELECTIONS['Excel'] = excel_file_name
            clear_and_print_summary()
            break
        else:
            print(Fore.YELLOW + "Opción inválida. Intente de nuevo.")
    except ValueError:
        print(Fore.YELLOW + "Entrada inválida. Ingrese un número.")

categ = load_categories()  # Cargar categorías después de seleccionar el archivo

cov_type = tipo_cobertura() # Preguntar tipo de cobertura
razon_cobertura = razao_cov() # Preguntar razón
tipo_eje_tend = tipo_eje_tendencia() # Preguntar tipo de eje para tendencia

# --- Procesamiento del Archivo Excel Seleccionado ---
excel_file_path = os.path.join(root_dir, excel_file_name)
try:
    excel_file_obj = pd.ExcelFile(excel_file_path)
    marcas = excel_file_obj.sheet_names
except FileNotFoundError:
    print(f"{Fore.RED}{Style.BRIGHT}Error: No se encontró el archivo seleccionado: {excel_file_path}")
    exit()
except Exception as e:
     print(f"{Fore.RED}{Style.BRIGHT}Error al abrir el archivo Excel '{excel_file_name}': {e}")
     exit()

# --- Obtener Metadatos del Nombre de Archivo ---
try:
    parts = os.path.splitext(excel_file_name)[0].split('_')
    if len(parts) < 3:
        raise ValueError("El nombre de archivo no contiene suficientes partes separadas por '_' (se esperan pais_cat_fab).")

    country_code_str = parts[0]
    category_code = parts[1]
    fabricante = parts[2]

    country_code = int(country_code_str)
    pais_nombre = pais.loc[pais.cod == country_code, 'pais'].iloc[0]
    cesta_nombre = categ.loc[category_code, 'cest']
    categoria_nombre = categ.loc[category_code, 'cat']

except (IndexError, ValueError, KeyError) as e:
    print(f"{Fore.RED}{Style.BRIGHT}Error al procesar metadatos del nombre de archivo '{excel_file_name}': {e}")
    print(f"{Fore.YELLOW}Asegúrese que el nombre siga el formato 'CodigoPais_CodigoCategoria_Fabricante.xlsx'.")

    # Usar country_code_str (que se define antes de la posible falla de int()) para los mensajes
    country_code_str_defined = 'country_code_str' in locals()
    category_code_defined = 'category_code' in locals()

    if isinstance(e, IndexError):
        print(f"{Fore.YELLOW}El nombre del archivo no tiene las partes esperadas (país_categoría_fabricante).")
    elif isinstance(e, ValueError):
        # Si country_code_str se definió, úsalo en el mensaje
        if country_code_str_defined:
            print(f"{Fore.YELLOW}El código de país '{country_code_str}' extraído del nombre de archivo no es un número válido.")
        else:
            print(f"{Fore.YELLOW}No se pudo extraer o convertir el código de país a número.")
    elif isinstance(e, KeyError):
        # Si el error es de clave, verificar si fue por categoría o país (usando str)
        if category_code_defined and category_code not in categ.index:
             print(f"{Fore.YELLOW}Verifique que el código de categoría '{category_code}' esté en los datos embebidos.")
        elif country_code_str_defined:
             # No podemos verificar directamente 'country_code not in pais['cod']' si falló int()
             # pero podemos indicar que el código extraído podría ser el problema.
             print(f"{Fore.YELLOW}Verifique que el código de país '{country_code_str}' exista en la definición de países o que la categoría sea correcta.")
        else:
             print(f"{Fore.YELLOW}Ocurrió un error al buscar un código (país o categoría).")
    else:
        print(f"{Fore.YELLOW}Ocurrió un error inesperado al procesar los metadatos.")

    exit()

coverage_label = "Cobertura Absoluta" if cov_type == "Absoluta" else "Cobertura Relativa"
ref_month_year = "" # Se actualizará en el bucle con la última fecha

# --- (1) CREACIÓN DEL TEMPLATE EN EXCEL ---
print(Fore.CYAN + "\nGenerando archivo Excel temporal...")
excel_temp_path = os.path.join(root_dir, EXCEL_TEMP_FILENAME)
try:
    with pd.ExcelWriter(excel_temp_path) as writer:
        # Recorrer cada hoja (marca) del archivo
        for marca_sheet_name in tqdm(marcas, desc="Procesando Hojas Excel"):

            # 1.1) Carga y preprocesa la hoja usando la función refactorizada
            df_marca, measure_unit = load_and_preprocess_sheet(excel_file_obj, marca_sheet_name)

            # Si la carga falló, continuar con la siguiente hoja
            if df_marca is None:
                continue

            # Guardar número original de filas de datos para fórmulas Excel
            original_data_rows = len(df_marca)
            if original_data_rows < 12:
                print(f"{Fore.YELLOW}Advertencia: Hoja '{marca_sheet_name}' tiene < 12 meses de datos ({original_data_rows}). Algunos cálculos de Excel pueden fallar o dar NaN.")
                # Continuar de todos modos, pero con precaución

            # Actualizar fecha de referencia global (usará la de la última hoja procesada con éxito)
            ref_month_year = df_marca[COL_DATA].iloc[-1].strftime('%m-%y')

            # --- 1.5) Creación de columnas con fórmulas Excel ---
            df_excel = df_marca.copy() # Trabajar sobre una copia para Excel
            # Hacer los índices basados en 1 y añadir offset de header (fila 1)
            excel_row_offset = 2

            # Sell_in_sim (Ejemplo - ajustable manualmente en Excel si se necesita)
            # La fórmula asume que Sell_in está en la columna B
            df_excel[COL_SELL_IN_SIM] = [f"=B{r}" for r in range(excel_row_offset, original_data_rows + excel_row_offset)] + [np.nan] * (len(df_excel) - original_data_rows)

            # Acumulados (MAT - Moving Annual Total) - comienzan desde la fila 12 de datos
            # Las fórmulas asumen Sell_out en C y Sell_in_sim en L
            for i in range(11, original_data_rows):
                row_excel = i + excel_row_offset
                df_excel.loc[i, COL_ACUM_SELL_OUT] = f"=SUM(C{row_excel - 11}:C{row_excel})"
                df_excel.loc[i, COL_ACUM_SELL_IN] = f"=SUM(L{row_excel - 11}:L{row_excel})" # Usa Sell_in_sim (L)

            # --- 1.6) Cálculo de coberturas (pipeline 0 a 6) en Excel ---
            pop_value_str = pop_coverage.get(pais_nombre, DEFAULT_POP_COVERAGE)
            cov_formulas_list = []
            max_rows_excel = original_data_rows + excel_row_offset -1 # Última fila con datos en Excel

            for r_idx in range(original_data_rows): # Iterar sobre índices de datos (0 a N-1)
                excel_current_row = r_idx + excel_row_offset
                row_formulas = {}
                if r_idx >= 11: # Cobertura solo se calcula desde el mes 12
                    for p in range(7): # Pipelines P0 a P6
                        # Fila de Excel para el numerador (Acum_Sell_in) - con pipeline
                        num_row_excel = excel_current_row + p
                        # Fila de Excel para el denominador (Acum_Sell_out) - sin pipeline
                        den_row_excel = excel_current_row

                        # Verificar que las filas referenciadas existan
                        if num_row_excel <= max_rows_excel and den_row_excel <= max_rows_excel:
                            # La fórmula asume Acum_Sell_in en N y Acum_Sell_out en M
                            base_formula = f"M{den_row_excel}/N{num_row_excel}*100"
                            if cov_type == "relativa":
                                # CORRECCIÓN: Cambiar formato de porcentaje y usar NA()
                                pop_value_decimal = float(pop_value_str.replace("%", "")) / 100
                                formula = f"=IFERROR(({base_formula})/{pop_value_decimal},NA())"
                            else:
                                formula = f"=IFERROR({base_formula},NA())"
                            row_formulas[f'P{p}'] = formula
                        else:
                             row_formulas[f'P{p}'] = np.nan # O "" o NA()
                else:
                    # Rellenar con NaN para las primeras 11 filas
                    for p in range(7):
                         row_formulas[f'P{p}'] = np.nan

                cov_formulas_list.append(row_formulas)

            df_cov_excel = pd.DataFrame(cov_formulas_list, index=df_excel.index[:original_data_rows])

            # Escalonar las columnas de cobertura
            df_cov_excel_scaled = df_cov_excel.copy()
            escalona(df_cov_excel_scaled) # Escalonar la copia

            # --- 1.7 & 1.8) Cálculo de variaciones (Y-1 e Y-2) en Excel --- (REEMPLAZADO)
            var_summary_list = []
            n_data = original_data_rows  # Número de filas con datos
            last_row_excel = n_data + excel_row_offset - 1

            # Definir los periodos y desplazamientos
            periods = [
                ("Anual", 12),
                ("Semestral", 6),
                ("Trimestral", 3)
            ]

            # Y-1
            for period_type, months in periods:
                end_curr = last_row_excel
                start_curr = end_curr - months + 1
                end_prev = end_curr - months
                start_prev = start_curr - months
                period_label = f"{period_type} vs Y-1"
                row_vars = {'Tipo': period_type, 'Periodo': period_label}
                # Kantar (columna C)
                if start_curr >= excel_row_offset and start_prev >= excel_row_offset:
                    formula_kantar = f"=IFERROR(SUM(C{start_curr}:C{end_curr})/SUM(C{start_prev}:C{end_prev})-1,NA())"
                    row_vars['Kantar'] = formula_kantar
                else:
                    row_vars['Kantar'] = "-"
                # Cliente (columna L) para P0-P6
                for p in range(7):
                    end_curr_p = end_curr - p
                    start_curr_p = end_curr_p - months + 1
                    end_prev_p = end_curr_p - months
                    start_prev_p = start_curr_p - months
                    if start_curr_p >= excel_row_offset and start_prev_p >= excel_row_offset:
                        formula_cliente = f"=IFERROR(SUM(L{start_curr_p}:L{end_curr_p})/SUM(L{start_prev_p}:L{end_prev_p})-1,NA())"
                        row_vars[f'Cliente P{p}'] = formula_cliente
                    else:
                        row_vars[f'Cliente P{p}'] = "-"
                var_summary_list.append(row_vars)

            # Y-2
            for period_type, months in periods:
                end_curr = last_row_excel
                start_curr = end_curr - months + 1
                end_prev = end_curr - 24  # 24 meses atrás
                start_prev = start_curr - 24
                period_label = f"{period_type} vs Y-2"
                row_vars = {'Tipo': period_type, 'Periodo': period_label}
                # Kantar (columna C)
                if start_curr >= excel_row_offset and start_prev >= excel_row_offset:
                    formula_kantar = f"=IFERROR(SUM(C{start_curr}:C{end_curr})/SUM(C{start_prev}:C{end_prev})-1,NA())"
                    row_vars['Kantar'] = formula_kantar
                else:
                    row_vars['Kantar'] = "-"
                # Cliente (columna L) para P0-P6
                for p in range(7):
                    end_curr_p = end_curr - p
                    start_curr_p = end_curr_p - months + 1
                    end_prev_p = end_curr_p - 24
                    start_prev_p = start_curr_p - 24
                    if start_curr_p >= excel_row_offset and start_prev_p >= excel_row_offset:
                        formula_cliente = f"=IFERROR(SUM(L{start_curr_p}:L{end_curr_p})/SUM(L{start_prev_p}:L{end_prev_p})-1,NA())"
                        row_vars[f'Cliente P{p}'] = formula_cliente
                    else:
                        row_vars[f'Cliente P{p}'] = "-"
                var_summary_list.append(row_vars)

            df_variations_excel = pd.DataFrame(var_summary_list)


            # --- 1.9) Cálculo de correlaciones en Excel (MAT) ---
            correl_formulas = {'Correl': 'Anual'}
            if n_data >= 12:
                 start_row_corr = last_row_excel - 11
                 end_row_corr = last_row_excel
                 # Corrección: Correlación debe ser entre series originales (C vs L con pipeline) no acumuladas
                 sell_out_range = f"C{start_row_corr}:C{end_row_corr}"

                 for p in range(7): # Incluir P0 ahora
                     sell_in_range_p = f"L{start_row_corr-p}:L{end_row_corr-p}"
                     if start_row_corr-p >= excel_row_offset:
                         # CORRECCIÓN: Usar IFERROR y NA()
                         correl_formulas[f'P{p}'] = f"=IFERROR(CORREL({sell_out_range},{sell_in_range_p}),NA())"
                     else:
                         correl_formulas[f'P{p}'] = np.nan
            else:
                 for p in range(7): correl_formulas[f'P{p}'] = np.nan # O NA()

            df_correlations_excel = pd.DataFrame([correl_formulas])


            # --- 1.10) Promedio de Penetración y Buyers (MAT) en Excel ---
            avg_formulas = []
            # MAT Actual
            if n_data >= 12:
                 start_avg_curr = last_row_excel - 11
                 end_avg_curr = last_row_excel
                 # Asume Penet en G, Buyers en H
                 avg_formulas.append({'Media': 'Penet MAT Actual', 'Valor': f"=AVERAGE(G{start_avg_curr}:G{end_avg_curr})"})
                 avg_formulas.append({'Media': 'Buyers MAT Actual', 'Valor': f"=AVERAGE(H{start_avg_curr}:H{end_avg_curr})"})
            else:
                 avg_formulas.append({'Media': 'Penet MAT Actual', 'Valor': f"=AVERAGE(G{excel_row_offset}:G{last_row_excel})"}) # Promedio de lo disponible
                 avg_formulas.append({'Media': 'Buyers MAT Actual', 'Valor': f"=AVERAGE(H{excel_row_offset}:H{last_row_excel})"})

            # MAT Anterior
            if n_data >= 24:
                 start_avg_prev = last_row_excel - 23
                 end_avg_prev = last_row_excel - 12
                 avg_formulas.append({'Media': 'Penet MAT Anterior', 'Valor': f"=AVERAGE(G{start_avg_prev}:G{end_avg_prev})"})
                 avg_formulas.append({'Media': 'Buyers MAT Anterior', 'Valor': f"=AVERAGE(H{start_avg_prev}:H{end_avg_prev})"})
            else:
                 avg_formulas.append({'Media': 'Penet MAT Anterior', 'Valor': np.nan}) # O NA()
                 avg_formulas.append({'Media': 'Buyers MAT Anterior', 'Valor': np.nan})

            df_averages_excel = pd.DataFrame(avg_formulas)


            # --- 1.11) Calcular Estabilidad en Excel ---
            # Diferencia entre última cobertura y cobertura de hace 12 meses
            estab_data = {"Estabilidad": "Estabilidad"}
            # Asume Cobertura P0-P6 en columnas O a U (después de escalonar)
            coverage_start_col_letter = 'O'
            coverage_start_col_idx = 15 # Col O es la 15

            last_data_row_idx = original_data_rows -1 # Índice base 0

            for p in range(7):
                 col_letter = get_column_letter(coverage_start_col_idx + p)
                 row_last_cov = last_row_excel - p
                 row_prev_cov = row_last_cov - 12

                 # Verificar si las filas son válidas y si hay suficientes datos
                 if row_last_cov >= excel_row_offset and row_prev_cov >= excel_row_offset and (original_data_rows >= 23+p):
                     # CORRECCIÓN: Usar IFERROR y NA()
                     formula = f"=IFERROR({col_letter}{row_last_cov}-{col_letter}{row_prev_cov},NA())"
                     estab_data[f'P{p}'] = formula
                 else:
                     estab_data[f'P{p}'] = np.nan







            df_stability_excel = pd.DataFrame([estab_data])

            # --- 1.12) Ensamblar DataFrame final para Excel ---
            # Unir datos originales con coberturas escalonadas
            df_excel_final = pd.concat([df_excel, df_cov_excel_scaled], axis=1)

            # Crear la sección de resumen (Variaciones, Promedios, Correlación, Estabilidad)
            # Añadir filas vacías y reorganizar
            df_variations_excel['spacer1'] = np.nan
            df_averages_excel['spacer2'] = np.nan
            df_correlations_excel['spacer3'] = np.nan

            # Aplanar las tablas de resumen para concatenarlas horizontalmente
            summary_part1 = df_variations_excel.T.reset_index().T # Variaciones
            summary_part2 = df_averages_excel.T.reset_index().T   # Promedios
            summary_part3 = df_correlations_excel.T.reset_index().T # Correlaciones
            summary_part4 = df_stability_excel.T.reset_index().T  # Estabilidad

            # Crear un DataFrame vacío con el número correcto de columnas para alinear
            max_cols = df_excel_final.shape[1]
            summary_placeholder = pd.DataFrame(np.nan, index=range(max(len(summary_part1), len(summary_part2), len(summary_part3), len(summary_part4))), columns=df_excel_final.columns)

            # Rellenar el placeholder (esto requiere manejo cuidadoso de índices y columnas)
            # Simplificación: Crear el df_excel_summary_part como antes y concatenar al final
            df_excel_summary_part = pd.concat([df_variations_excel.reset_index(drop=True),
                                              df_averages_excel.reset_index(drop=True),
                                              df_correlations_excel.reset_index(drop=True),
                                              df_stability_excel.reset_index(drop=True)], axis=1)

            # Añadir fila vacía de separación
            df_excel_final.loc[len(df_excel_final)] = [np.nan] * len(df_excel_final.columns)

            # Añadir nombres de columnas del resumen como cabecera
            summary_header = pd.DataFrame([df_excel_summary_part.columns], columns=df_excel_summary_part.columns)
            df_excel_summary_part_with_header = pd.concat([summary_header, df_excel_summary_part], ignore_index=True)

            # Ajustar columnas del resumen para que coincidan con el df principal y concatenar
            # --- INICIO CAMBIO ---
            # Si el número de columnas no coincide, agrega columnas vacías
            n_main_cols = df_excel_final.shape[1]
            n_summary_cols = df_excel_summary_part_with_header.shape[1]
            if n_summary_cols < n_main_cols:
                # Agrega columnas vacías al resumen
                for i in range(n_summary_cols, n_main_cols):
                    df_excel_summary_part_with_header[f'empty_{i}'] = np.nan
            elif n_summary_cols > n_main_cols:
                # Si el resumen tiene más columnas, recórtalas
                df_excel_summary_part_with_header = df_excel_summary_part_with_header.iloc[:, :n_main_cols]
            # Ahora reasigna los nombres de columnas
            df_excel_summary_part_with_header.columns = df_excel_final.columns
            # --- FIN CAMBIO ---

            df_excel_final = pd.concat([df_excel_final, df_excel_summary_part_with_header], ignore_index=True)

            # --- 1.13) Exportar a la hoja de Excel ---
            df_excel_final.to_excel(writer, sheet_name=marca_sheet_name, index=False)

    print(Fore.GREEN + f"Archivo Excel temporal '{EXCEL_TEMP_FILENAME}' generado.")

except Exception as e:
    print(f"{Fore.RED}{Style.BRIGHT}Error crítico durante la generación del archivo Excel: {e}")
    if os.path.exists(excel_temp_path):
         os.remove(excel_temp_path) # Limpiar si falla
    exit()

# --- 1.14) Renombrar y mover archivo Excel final ---
if not ref_month_year:
     print(f"{Fore.RED}{Style.BRIGHT}No se pudo determinar la fecha de referencia. No se puede renombrar el archivo Excel.")
     if os.path.exists(excel_temp_path):
         os.remove(excel_temp_path)
     exit()

nombre_base_archivo = f"{pais_nombre}-{categoria_nombre}-{fabricante}-{ref_month_year}_{coverage_label}"
carpeta_salida = os.path.join(root_dir, nombre_base_archivo) # Carpeta con el mismo nombre base

if not os.path.exists(carpeta_salida):
    try:
        os.makedirs(carpeta_salida)
        print(Fore.BLUE + "Carpeta de salida creada: " + Fore.YELLOW + f"{carpeta_salida}")
    except OSError as e:
        print(f"{Fore.RED}Error al crear carpeta de salida '{carpeta_salida}': {e}")
        if os.path.exists(excel_temp_path): os.remove(excel_temp_path)
        exit()
else:
    print(Fore.BLUE + "Carpeta de salida ya existe: " + Fore.YELLOW + f"{carpeta_salida}")

nombre_template_final = f"{nombre_base_archivo}.xlsx"
ruta_template_final = os.path.join(carpeta_salida, nombre_template_final)

try:
    if os.path.exists(ruta_template_final):
        print(Fore.YELLOW + f"Archivo Excel '{nombre_template_final}' ya existe. Se sobrescribirá.")
        os.remove(ruta_template_final)
    os.rename(excel_temp_path, ruta_template_final)
    print(Fore.GREEN + "Archivo Excel final guardado en: " + Fore.YELLOW + f"{ruta_template_final}")
except Exception as e:
    print(f"{Fore.RED}Error al mover/renombrar archivo Excel final: {e}")
    if os.path.exists(excel_temp_path): os.remove(excel_temp_path) # Limpiar temporal si falla el renombrado
    exit()


# --------------------------------------------------------------------------------------------------
# (2) CREACIÓN DE PPT CON GRÁFICOS
# --------------------------------------------------------------------------------------------------
print(Fore.CYAN + "\nGenerando presentación PowerPoint...")
ppt = Presentation('Modelo_Revision.pptx')  # Cargar plantilla PPT

# Definir etiquetas de idioma (simplificado)
lang_index = 1 if pais_nombre == 'Brasil' else 2 # 1 -> PT / 2 -> ES
labels = {
    (1, 'S1'): 'Construa sua história',
    (1, 'Summary'): ['Marca', 'Pipeline', 'Penetração Média Mensal', fabricante, 'Kantar Worldpanel',
                      f'Cobertura {(dt.strptime(ref_month_year, "%m-%y") - timedelta(days=365)).strftime("%b-%y")}', # Formato mes-año
                      f'Cobertura {dt.strptime(ref_month_year, "%m-%y").strftime("%b-%y")}', 'Estabilidad'],
    (1, 'Graf cob Penet Men'): 'Penetração Mensal',
    (1, 'Titulo Cob'): 'Cobertura em Ano Móvel',
    (1, 'Var'): 'com',
    (1, 'Titulo Vol'): 'Tendência em Volumen',

    (2, 'S1'): '-', # Placeholder para español
    (2, 'Summary'): ['Marca', 'Pipeline', 'Penetración Media Mensual', f'%VAR {fabricante}', '% VAR Kantar Worldpanel',
                      f'Cobertura {(dt.strptime(ref_month_year, "%m-%y") - timedelta(days=365)).strftime("%b-%y")}',
                      f'Cobertura {dt.strptime(ref_month_year, "%m-%y").strftime("%b-%y")}', 'Estabilidad'],
    (2, 'Graf cob Penet Men'): 'Penetración Mensual',
    (2, 'Titulo Cob'): 'Cobertura en Año Móvil', # Corregido Móbil->Móvil
    (2, 'Var'): 'con',
    (2, 'Titulo Vol'): 'Tendencia en Volumen'
}

# DataFrames para resumen final y banco de coberturas
df_summary_ppt = pd.DataFrame(columns=labels[(lang_index, 'Summary')])
df_coverage_bank = pd.DataFrame(
    columns=['Periodo', 'Fabricante', 'Categoria', 'Marca', 'Cesta', 'Panel', 'Unidade', 'Razon', 'Pais', 'Ampliacion',
             'Penet Media Ano Mov Atual', 'Penet Media Ano Mov Anterior',
             'Raw Buyers Media Ano Mov Atual', 'Pipeline',
             'Cobertura Año Mov Actual', 'Cobertura Año Mov Anterior',
             '%VAR Cliente', '% VAR Kantar', 'Misma Tendencia']
)

# --- Bucle principal para generar diapositivas ---
total_slides_to_generate = 0
for marca_sheet_name in marcas:
    df_marca_ppt, _ = load_and_preprocess_sheet(excel_file_obj, marca_sheet_name)
    if df_marca_ppt is None: continue

    # Determinar pipelines a procesar para esta marca
    match = re.match(r"(?i)^p([0-6])_", marca_sheet_name)
    if match:
        pipelines_to_run = [int(match.group(1))]
    else:
        pipelines_to_run = list(range(7)) # P0 a P6

    # Contar cuántos slides se generarán (aprox 3 por pipeline + 1 resumen)
    n_slides_marca = len(pipelines_to_run) * (2 + (1 if len(df_marca_ppt) >= 24 else 0)) # Cobertura, Tendencia, [Evolución]
    total_slides_to_generate += n_slides_marca

print(f"Total slides a generar (estimado): {total_slides_to_generate}") # Añade esto para depurar

# --- INICIO CAMBIO: Usar rich Progress ---
progress = Progress(
    SpinnerColumn(),
    TextColumn("[progress.description]{task.description}"),
    BarColumn(),
    TextColumn("[progress.percentage]{task.percentage:>3.0f}%"),
    TextColumn("{task.completed}/{task.total}"),
    TimeElapsedColumn(),
    TimeRemainingColumn(),
    transient=True
)
progress_task = None
# --- FIN CAMBIO ---

# Bucle para cada marca
with progress:
    progress_task = progress.add_task("Creando Diapositivas PPT", total=total_slides_to_generate + 1)
    for marca_sheet_name in marcas:

        # 2.1) Cargar y preprocesar datos para PPT
        df_marca_ppt, measure_unit_ppt = load_and_preprocess_sheet(excel_file_obj, marca_sheet_name)
        if df_marca_ppt is None:
            continue # Saltar si la hoja no se pudo cargar/procesar

        marca_nombre_limpio = re.sub(r"(?i)^p[0-6]_", "", marca_sheet_name)

        # Determinar pipelines a procesar para esta marca
        match = re.match(r"(?i)^p([0-6])_", marca_sheet_name)
        if match:
            pipelines_to_run = [int(match.group(1))]
        else:
            pipelines_to_run = list(range(7)) # P0 a P6

        # 2.2) Cálculo de coberturas en Python (rolling 12 meses)
        acum_sell_out_py = df_marca_ppt[COL_SELL_OUT].rolling(window=12, min_periods=12).sum()
        acum_sell_in_py = df_marca_ppt[COL_SELL_IN].rolling(window=12, min_periods=12).sum()
        acum_sell_out_py.index = df_marca_ppt[COL_DATA]
        acum_sell_in_py.index = df_marca_ppt[COL_DATA]

        df_coverage_py = pd.DataFrame(index=acum_sell_out_py.index)
        for p in range(7):
            sell_in_shifted = df_marca_ppt[COL_SELL_IN].shift(p)
            acum_sell_in_shifted_py = sell_in_shifted.rolling(window=12, min_periods=12).sum()
            acum_sell_in_shifted_py.index = df_marca_ppt[COL_DATA]
            coverage_p = (acum_sell_out_py / acum_sell_in_shifted_py) * 100
            df_coverage_py[f'P{p}'] = coverage_p

        pop_val_num = float(pop_coverage.get(pais_nombre, DEFAULT_POP_COVERAGE).replace("%", "")) / 100.0
        if cov_type == "relativa" and pop_val_num > 0:
            df_coverage_py = df_coverage_py / pop_val_num

        df_coverage_py = df_coverage_py.round(1)

        # 2.3) Cálculo de variaciones en Python
        df_variations_py = pd.DataFrame(columns=['Tipo', 'Periodo', 'Kantar'] + [f'Cliente P{p}' for p in range(7)])
        period_types = ["Anual", "Semestral", "Trimestral"]
        # Y-1
        kantar_vars_y1 = calc_var1(df_marca_ppt, COL_SELL_OUT, 0)
        cliente_vars_y1 = {p: calc_var1(df_marca_ppt, COL_SELL_IN, p) for p in range(7)}
        for i, p_type in enumerate(period_types):
            row = {'Tipo': p_type, 'Periodo': f'{p_type} vs Y-1', 'Kantar': kantar_vars_y1[i]}
            for p in range(7): row[f'Cliente P{p}'] = cliente_vars_y1[p][i]
            df_variations_py.loc[len(df_variations_py)] = row
        # Y-2
        kantar_vars_y2 = calc_var2(df_marca_ppt, COL_SELL_OUT, 0)
        cliente_vars_y2 = {p: calc_var2(df_marca_ppt, COL_SELL_IN, p) for p in range(7)}
        for i, p_type in enumerate(period_types):
            row = {'Tipo': p_type, 'Periodo': f'{p_type} vs Y-2', 'Kantar': kantar_vars_y2[i]}
            for p in range(7): row[f'Cliente P{p}'] = cliente_vars_y2[p][i]
            df_variations_py.loc[len(df_variations_py)] = row

        # 2.4) Correlaciones (Pearson entre Sell-Out y Sell-In shifteado, últimos 12m)
        correlations_py = {"Periodo": "Correlación MAT"}
        if len(df_marca_ppt) >= 12:
            sell_out_last12 = df_marca_ppt[COL_SELL_OUT].iloc[-12:]
            for p in range(7):
                if len(df_marca_ppt) >= 12 + p:
                    sell_in_last12_shifted = df_marca_ppt[COL_SELL_IN].iloc[-(12+p):-p] if p > 0 else df_marca_ppt[COL_SELL_IN].iloc[-12:]
                    if len(sell_out_last12) == 12 and len(sell_in_last12_shifted) == 12 and \
                        sell_out_last12.notna().all() and sell_in_last12_shifted.notna().all() and \
                        np.isfinite(sell_out_last12).all() and np.isfinite(sell_in_last12_shifted).all() and \
                        sell_out_last12.std() > 0 and sell_in_last12_shifted.std() > 0:
                        corr_val, _ = pearsonr(sell_out_last12, sell_in_last12_shifted)
                        correlations_py[f'P{p}'] = corr_val
                    else:
                        correlations_py[f'P{p}'] = np.nan
                else:
                    correlations_py[f'P{p}'] = np.nan
        else:
            for p in range(7): correlations_py[f'P{p}'] = np.nan # O NA()

        df_correlations_py = pd.DataFrame([correlations_py])


        # 2.5) Promedios MAT (Penetración, Buyers)
        averages_py = {}
        n_data_ppt = len(df_marca_ppt)
        if n_data_ppt >= 12:
            averages_py['Penet_MAT_Actual'] = df_marca_ppt[COL_PENET].iloc[-12:].mean()
            averages_py['Buyers_MAT_Actual'] = df_marca_ppt[COL_BUYERS].iloc[-12:].mean()
        else:
            averages_py['Penet_MAT_Actual'] = df_marca_ppt[COL_PENET].mean()
            averages_py['Buyers_MAT_Actual'] = df_marca_ppt[COL_BUYERS].mean()

        if n_data_ppt >= 24:
            averages_py['Penet_MAT_Anterior'] = df_marca_ppt[COL_PENET].iloc[-24:-12].mean()
            averages_py['Buyers_MAT_Anterior'] = df_marca_ppt[COL_BUYERS].iloc[-24:-12].mean()
        else:
            averages_py['Penet_MAT_Anterior'] = np.nan
            averages_py['Buyers_MAT_Anterior'] = np.nan

        # 2.6) Preparar DF para gráficos de tendencia (fecha como string)
        df_trend_plot = df_marca_ppt[[COL_DATA, COL_SELL_IN, COL_SELL_OUT]].copy()
        df_trend_plot[COL_DATA] = df_trend_plot[COL_DATA].apply(lambda x: x.strftime('%m-%y'))

        # --- Llenar Banco de Coberturas y Resumen PPT ---
        for p in pipelines_to_run:
            coverage_series = df_coverage_py[f'P{p}'].dropna()
            if not coverage_series.empty:
                coverage_actual = coverage_series.iloc[-1]
                coverage_anterior = coverage_series.iloc[-13] if len(coverage_series) >= 13 else np.nan
            else:
                coverage_actual = np.nan
                coverage_anterior = np.nan

            var_cliente_anual_y1 = df_variations_py.loc[df_variations_py['Tipo'] == 'Anual', f'Cliente P{p}'].iloc[0]
            var_kantar_anual_y1 = df_variations_py.loc[df_variations_py['Tipo'] == 'Anual', 'Kantar'].iloc[0]

            tendencia_alineada = "NO"
            if pd.notna(var_cliente_anual_y1) and pd.notna(var_kantar_anual_y1):
                if (var_cliente_anual_y1 * var_kantar_anual_y1) > 0:
                    tendencia_alineada = "SI"
                elif var_cliente_anual_y1 == 0 and var_kantar_anual_y1 == 0:
                    tendencia_alineada = "SI"

            banco_row = {
                'Periodo': dt.strptime(ref_month_year, '%m-%y').strftime('%b-%y'),
                'Fabricante': fabricante,
                'Categoria': categoria_nombre,
                'Marca': marca_nombre_limpio,
                'Cesta': cesta_nombre,
                'Panel': 'PNC',
                'Unidade': measure_unit_ppt,
                'Razon': razon_cobertura,
                'Pais': pais_nombre,
                'Ampliacion': 'SI',
                'Penet Media Ano Mov Atual': round(averages_py.get('Penet_MAT_Actual', 0), 1),
                'Penet Media Ano Mov Anterior': round(averages_py.get('Penet_MAT_Anterior', 0), 1),
                'Raw Buyers Media Ano Mov Atual': round(averages_py.get('Buyers_MAT_Actual', 0), 1),
                'Pipeline': p,
                'Cobertura Año Mov Actual': round(coverage_actual, 1) if pd.notna(coverage_actual) else 0,
                'Cobertura Año Mov Anterior': round(coverage_anterior, 1) if pd.notna(coverage_anterior) else 0,
                '%VAR Cliente': round(var_cliente_anual_y1 * 100, 1) if pd.notna(var_cliente_anual_y1) else 0,
                '% VAR Kantar': round(var_kantar_anual_y1 * 100, 1) if pd.notna(var_kantar_anual_y1) else 0,
                'Misma Tendencia': tendencia_alineada
            }
            df_coverage_bank.loc[len(df_coverage_bank)] = banco_row

            estabilidad = round(coverage_actual - coverage_anterior, 1) if pd.notna(coverage_actual) and pd.notna(coverage_anterior) else np.nan
            summary_row = {
                labels[(lang_index, 'Summary')][0]: marca_nombre_limpio,
                labels[(lang_index, 'Summary')][1]: p,
                labels[(lang_index, 'Summary')][2]: f"{averages_py.get('Penet_MAT_Actual', 0):.1f}%",
                labels[(lang_index, 'Summary')][3]: f"{var_cliente_anual_y1*100:.1f}%" if pd.notna(var_cliente_anual_y1) else "0.0%",
                labels[(lang_index, 'Summary')][4]: f"{var_kantar_anual_y1*100:.1f}%" if pd.notna(var_kantar_anual_y1) else "0.0%",
                labels[(lang_index, 'Summary')][5]: f"{coverage_anterior:.1f}" if pd.notna(coverage_anterior) else "0.0",
                labels[(lang_index, 'Summary')][6]: f"{coverage_actual:.1f}" if pd.notna(coverage_actual) else "0.0",
                labels[(lang_index, 'Summary')][7]: f"{estabilidad:.1f}" if pd.notna(estabilidad) else "0.0"
            }
             
            df_summary_ppt.loc[len(df_summary_ppt)] = pd.Series(summary_row)

        # --- 2.9) Generar Diapositivas ---
        for p in pipelines_to_run:
            # A) Slide de Cobertura y Penetración
            slide_cov = ppt.slides.add_slide(ppt.slide_layouts[PPT_LAYOUT_INDEX])
            # Añadir título al slide
            tx_title_cov = slide_cov.shapes.title # Usar placeholder de título si existe
            if tx_title_cov is None: # Si no hay placeholder de título, añadir textbox
                 tx_title_cov = slide_cov.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)).text_frame # Ajustar tamaño/posición
            else:
                 tx_title_cov = tx_title_cov.text_frame # Acceder al text_frame si existe
            t_cov = tx_title_cov.paragraphs[0]
            t_cov.text = f"{marca_nombre_limpio} - Pipeline {p}"
            t_cov.font.bold = True
            t_cov.font.size = Inches(0.3)

            # Extraer datos para el gráfico P
            cov_series_p = df_coverage_py[f'P{p}'].dropna()
            # Asegurar que la penetración tenga el mismo índice que la cobertura
            pen_series_p = df_marca_ppt.set_index(COL_DATA)[COL_PENET].loc[cov_series_p.index]

            generar_grafico_cobertura(slide_cov, marca_nombre_limpio, p,
                                  cov_series_p, pen_series_p,
                                  lang_index, cov_type, labels)

            # B) Añadir tabla de Variación MAT Y-1 al slide de cobertura
            var_cliente_p_mat = df_variations_py.loc[0, f'Cliente P{p}'] # Fila 0 es Anual Y-1
            var_kantar_mat = df_variations_py.loc[0, 'Kantar']

            var_mat_data = {
             " ": [f"VAR % MAT ({ref_month_year})"], # Fila 0, columna 'Tipo'/'Periodo'
             f"{fabricante} {labels[(lang_index, 'Var')]} Pipeline {p}": [f"{var_cliente_p_mat*100:.1f}%" if pd.notna(var_cliente_p_mat) else "-"],
             "Kantar Worldpanel": [f"{var_kantar_mat*100:.1f}%" if pd.notna(var_kantar_mat) else "-"]
            }
            df_var_mat_table = pd.DataFrame(var_mat_data)

            try:
             # Exportar tabla como imagen
             styler = df_var_mat_table.style
             styler.set_table_styles([
                 {'selector': '*',  'props': [('font-size', '10pt'), ('font-family','Calibri'), ('color','black'), ('border-style','solid'), ('border-width','1px'), ('text-align', 'center')]},
                 {'selector': 'th', 'props': [('background-color','lightgray'), ('font-weight', 'bold')]},
                 {'selector': 'td', 'props': [('padding', '2px 4px')]} # Añadir padding
             ]).hide(axis="index") # Ocultar índice numérico

             table_stream = io.BytesIO()
             # Ajustar dpi y table_conversion según sea necesario
             dfi.export(styler, table_stream, table_conversion='matplotlib', dpi=200)
             table_stream.seek(0)
             # Añadir imagen de la tabla al slide
             #  slide_cov.shapes.add_picture(table_stream, Inches(0.5), Inches(1.1), height=Inches(0.6)) # Ajustar posición
             img_pil = Image.open(table_stream)
             bordered = ImageOps.expand(img_pil, border=2, fill='black')
             img_stream_bordered = io.BytesIO()
             bordered.save(img_stream_bordered, format='PNG')
             img_stream_bordered.seek(0)
             slide_cov.shapes.add_picture(img_stream_bordered, Inches(0.5), Inches(1.1), height=Inches(0.6))

            except Exception as e:
                print(f"{Fore.YELLOW}Advertencia: No se pudo generar la tabla de variación MAT para {marca_nombre_limpio} P{p}. Error: {e}")

            # C) Slide de Tendencia (Sell-in vs Sell-out)
            slide_trend = ppt.slides.add_slide(ppt.slide_layouts[PPT_LAYOUT_INDEX])
            # Añadir título
            tx_title_trend = slide_trend.shapes.title # Usar placeholder de título
            if tx_title_trend is None:
                 tx_title_trend = slide_trend.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)).text_frame
            else:
                 tx_title_trend = tx_title_trend.text_frame
            t_trend = tx_title_trend.paragraphs[0]
            t_trend.text = f"{marca_nombre_limpio} - Pipeline {p}"
            t_trend.font.bold = True
            t_trend.font.size = Inches(0.3)

            generar_grafico_tendencia(
    slide_trend, marca_nombre_limpio, p,
    df_trend_plot, lang_index, labels,
    doble_eje=(tipo_eje_tend == "doble")
)

            # D) Slide de Evolución Mensual y Variación YOY (si hay datos)
            if len(df_marca_ppt) >= 24:
                df_evol_plot = df_marca_ppt[[COL_DATA, COL_SELL_IN, COL_SELL_OUT]].copy()
                df_evol_plot[COL_DATA] = pd.to_datetime(df_evol_plot[COL_DATA])

                fig_evol = generar_grafico_evolucion_mensual(df_evol_plot, p)

                if fig_evol is not None:
                    slide_evol = ppt.slides.add_slide(ppt.slide_layouts[PPT_LAYOUT_INDEX])
                    # Añadir título
                    tx_title_evol = slide_evol.shapes.title
                    if tx_title_evol is None:
                         tx_title_evol = slide_evol.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)).text_frame
                    else:
                         tx_title_evol = tx_title_evol.text_frame
                    t_evol = tx_title_evol.paragraphs[0]
                    t_evol.text = f"{marca_nombre_limpio} - Pipeline {p} - Evolución Mensual y Variación"
                    t_evol.font.bold = True
                    t_evol.font.size = Inches(0.3)

                    # Guardar gráfico en memoria y añadir al slide
                    img_stream_evol = io.BytesIO()
                    fig_evol.savefig(img_stream_evol, format='png', bbox_inches='tight', pad_inches=0.1, transparent=True)
                    img_stream_evol.seek(0)
                    slide_evol.shapes.add_picture(img_stream_evol, Inches(0.4), Inches(1.0), height=Inches(5.5)) # Ajustar posición y tamaño gráfico variacion mensual
                    plt.close(fig_evol)
                    progress.update(progress_task, advance=1)
            progress.update(progress_task, advance=1)
            progress.update(progress_task, advance=1)

# --- 2.10) Creación del slide "Summary" ---
print(Fore.CYAN + "\nAgregando slide de resumen...")
slide_summary = ppt.slides.add_slide(ppt.slide_layouts[PPT_LAYOUT_INDEX])
# Título
tx_title_summary = slide_summary.shapes.title
if tx_title_summary is None:
     tx_title_summary = slide_summary.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)).text_frame
else:
     tx_title_summary = tx_title_summary.text_frame
t_summary = tx_title_summary.paragraphs[0]
t_summary.text = f"Summary - {pais_nombre} {categoria_nombre} - {coverage_label}"
t_summary.font.bold = True
t_summary.font.size = Inches(0.35)

# Añadir texto S1 si es necesario
tx_s1 = slide_summary.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5)) # Posición inferior
tf_s1 = tx_s1.text_frame
t_s1 = tf_s1.add_paragraph()
t_s1.text = labels.get((lang_index, 'S1'), '')
t_s1.font.size = Inches(0.18)

# --- AGREGAR CUADRO DE TEXTO PARA COMENTARIOS ---
# Puedes ajustar la posición y tamaño según lo necesites
comentarios_box = slide_summary.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(8.5), Inches(0.7))
comentarios_frame = comentarios_box.text_frame
comentarios_frame.word_wrap = True
comentarios_frame.auto_size = True
comentarios_frame.text = "Comentarios:"  # Texto inicial, el usuario puede editarlo después en PowerPoint


# Crear tabla de resumen como imagen
if not df_summary_ppt.empty:
    try:
        # Aplicar estilo a la tabla resumen
        styler_summary = df_summary_ppt.style
        styler_summary.set_table_styles([
            {'selector': '*', 'props': [('font-size', '9pt'), ('font-family','Calibri'), ('color', 'black'), ('border-style', 'solid'), ('border-width', '1px'), ('text-align', 'center')]},
            {'selector': 'th', 'props': [('background-color', '#D9E1F2'), ('font-weight', 'bold'), ('border-width', '1px'), ('padding', '3px 5px')]}, # Cabecera azul claro
            {'selector': 'td', 'props': [('border-width', '1px'), ('padding', '3px 5px')]}
        ]).hide(axis="index")

        summary_table_stream = io.BytesIO()
        dfi.export(styler_summary, summary_table_stream, table_conversion='matplotlib', dpi=250) # Aumentar dpi si es necesario
        summary_table_stream.seek(0)

        # Añadir imagen de la tabla al slide (ajustar tamaño y posición)
        # Calcular altura necesaria basada en número de filas
        n_rows_summary = len(df_summary_ppt)
        img_height_inch = min(5.5, 0.3 + n_rows_summary * 0.25) # Ajustar multiplicador según tamaño de fuente/padding
        img_width_inch = 9.0 # Ancho

        # slide_summary.shapes.add_picture(summary_table_stream, Inches(0.5), Inches(1.0), width=Inches(img_width_inch) , height=Inches(img_height_inch))
        img_pil = Image.open(summary_table_stream)
        bordered = ImageOps.expand(img_pil, border=2, fill='black')
        img_stream_bordered = io.BytesIO()
        bordered.save(img_stream_bordered, format='PNG')
        img_stream_bordered.seek(0)
        slide_summary.shapes.add_picture(img_stream_bordered, Inches(0.5), Inches(1.0), width=Inches(img_width_inch), height=Inches(img_height_inch))

    except Exception as e:
        print(f"{Fore.YELLOW}Advertencia: No se pudo generar la tabla resumen en el PPT. Error: {e}")
else:
    print(f"{Fore.YELLOW}Advertencia: No hay datos para generar la tabla resumen en el PPT.")

# --- 2.11) Guardar PPT y Banco ---
try:
    # Mover slide de resumen a la segunda posición (índice 1)
    if len(ppt.slides) > 1:
        summary_slide_xml = ppt.slides._sldIdLst[-1] # El último slide añadido es el resumen
        ppt.slides._sldIdLst.insert(1, summary_slide_xml) # Insertar en posición 1 (segundo slide)

    nombre_ppt_final = f"{nombre_base_archivo}.pptx"
    ruta_ppt_final = os.path.join(carpeta_salida, nombre_ppt_final)
    ppt.save(ruta_ppt_final)
    print(Fore.GREEN + "Presentación PowerPoint guardada en: " + Fore.YELLOW + f"{ruta_ppt_final}")
except Exception as e:
    print(f"{Fore.RED}{Style.BRIGHT}Error al guardar la presentación PowerPoint: {e}")

# --- (3) Guardado del banco de coberturas ---
try:
    nombre_banco_final = f"Banco_{fabricante}_{categoria_nombre}_{pais_nombre}_{ref_month_year}_{coverage_label}.xlsx"
    ruta_banco_final = os.path.join(carpeta_salida, nombre_banco_final)
    df_coverage_bank.to_excel(ruta_banco_final, index=False)
    print(Fore.GREEN + "Banco de coberturas guardado en: " + Fore.YELLOW + f"{ruta_banco_final}")
except Exception as e:
    print(f"{Fore.RED}{Style.BRIGHT}Error al guardar el banco de coberturas: {e}")

print(f"\n{Fore.CYAN}{Style.BRIGHT}--- Proceso completado ---")

# --- END OF FILE ---