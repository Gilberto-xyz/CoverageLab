"""Coverage Studio Ultra

Refactor mantenible de coverage_studio.py. Conserva la logica original pero
organiza la generacion de reportes en componentes reutilizables
"""

from __future__ import annotations

import io
import math
import os
import posixpath
import re
import shutil
import sys
import tempfile
import threading
import time
import unicodedata
import uuid
import zipfile
from collections import OrderedDict
from dataclasses import dataclass, field
from datetime import datetime
from typing import Dict, Iterable, List, Optional, Sequence, Tuple, Callable, Set
from calendar import month_abbr
import xml.etree.ElementTree as ET

import colorama
from colorama import Back, Fore, Style
from rich.console import Console
from rich.panel import Panel

colorama.init(autoreset=True)
console = Console()

BRAND_EXCEPTION_REASONS: Dict[str, Set[str]] = {}

EXCEPTION_STYLES: Dict[str, Dict[str, str]] = {
    "zero_dash": {
        "brand_color": Fore.YELLOW,
        "message": "contiene 0s en los en algunos meses, graficando con exepcion",
        "summary_tag": "0/-",
    },
    "negative": {
        "brand_color": Fore.YELLOW,
        "message": "contiene valores negativos, graficando con exepcion",
        "summary_tag": "neg",
    },
}

SUMMARY_EXTRA_MONTHS_ENV_KEYS: Tuple[str, ...] = ("AUTO_EXTEA", "AUTO_EXTRA_MONTHS")
SUMMARY_EXTRA_MONTHS_MODE_ENV_KEYS: Tuple[str, ...] = ("AUTO_EXTEA_MODE", "AUTO_EXTRA_MONTHS_MODE")
VARIATIONS_BOX_STYLE_ENV_KEYS: Tuple[str, ...] = ("AUTO_VAR_BOX_STYLE", "AUTO_VAR_STYLE")
COVERAGE_SLIDE_VARIANT_ENV_KEYS: Tuple[str, ...] = ("AUTO_COV_SLIDE", "AUTO_COV_SLIDE_STYLE")
EVOLUTION_SLIDE_VARIANT_ENV_KEYS: Tuple[str, ...] = ("AUTO_EVO_SLIDE", "AUTO_EVO_SLIDE_STYLE")
TREND_GRANULARITY_ENV_KEYS: Tuple[str, ...] = ("AUTO_TREND_MODE", "AUTO_TREND_GRANULARITY")
MONTH_TOKEN_TO_NUMBER: Dict[str, int] = {
    "ene": 1, "enero": 1, "jan": 1, "janeiro": 1, "january": 1,
    "feb": 2, "febrero": 2, "fev": 2, "fevereiro": 2, "february": 2,
    "mar": 3, "marzo": 3, "marco": 3, "march": 3,
    "abr": 4, "abril": 4, "apr": 4, "april": 4,
    "may": 5, "mayo": 5, "maio": 5,
    "jun": 6, "junio": 6, "junho": 6, "june": 6,
    "jul": 7, "julio": 7, "julho": 7, "july": 7,
    "ago": 8, "agosto": 8, "aug": 8, "august": 8,
    "sep": 9, "sept": 9, "set": 9, "septiembre": 9, "setiembre": 9, "setembro": 9, "september": 9,
    "oct": 10, "octubre": 10, "out": 10, "outubro": 10, "october": 10,
    "nov": 11, "noviembre": 11, "novembro": 11, "november": 11,
    "dic": 12, "diciembre": 12, "dez": 12, "dezembro": 12, "dec": 12, "december": 12,
}

ANSI_RESET = "\033[0m"
PREVIEW_COPY_SUFFIX_RE = re.compile(
    r"^(?P<brand>.+?)(?P<suffix>\s*[-‑–—−‒]\s*(?:copia|copy|preview|previa?|borrador|draft).*)?$",
    re.IGNORECASE,
)
# Umbral mínimo para que los nombres resaltados sean legibles en fondos oscuros
# o temas de terminal con bajo contraste.
MIN_READABLE_LUMINANCE = 170.0
TERMINAL_BRAND_COLOR_SEQUENCE: Tuple[Tuple[int, int, int], ...] = (
    (31, 119, 180),   # Azul
    (214, 39, 40),    # Rojo
    (44, 160, 44),    # Verde
    (255, 127, 14),   # Naranja
    (148, 103, 189),  # Morado
    (23, 190, 207),   # Cian
    (140, 86, 75),    # Marron
    (227, 119, 194),  # Rosa
    (188, 189, 34),   # Oliva
    (63, 81, 181),    # Indigo
    (0, 173, 181),    # Turquesa
    (255, 64, 129),   # Magenta
    (57, 59, 121),    # Azul oscuro
    (82, 84, 163),    # Azul violeta
    (107, 110, 207),  # Lavanda intensa
    (99, 121, 57),    # Verde musgo
    (140, 162, 82),   # Verde oliva claro
    (140, 109, 49),   # Cafe
    (189, 158, 57),   # Mostaza
    (132, 60, 57),    # Rojo ladrillo
    (173, 73, 74),    # Coral oscuro
    (123, 65, 115),   # Purpura profundo
    (165, 81, 148),   # Fucsia oscuro
    (230, 85, 13),    # Naranja quemado
)
TEMPLATE_TAB_COLOR_SEQUENCE: List[str] = [
    "#1F77B4", "#FF7F0E", "#2CA02C", "#D62728", "#9467BD", "#8C564B", "#E377C2", "#7F7F7F", "#BCBD22", "#17BECF",
    "#AEC7E8", "#FFBB78", "#98DF8A", "#FF9896", "#C5B0D5", "#C49C94", "#F7B6D2", "#C7C7C7", "#DBDB8D", "#9EDAE5",
    "#393B79", "#5254A3", "#6B6ECF", "#9C9EDE", "#637939", "#8CA252", "#B5CF6B", "#CEDB9C", "#8C6D31", "#BD9E39",
    "#E7BA52", "#E7CB94", "#843C39", "#AD494A", "#D6616B", "#E7969C", "#7B4173", "#A55194", "#CE6DBD", "#DE9ED6",
    "#3182BD", "#6BAED6", "#9ECAE1", "#C6DBEF", "#E6550D", "#FD8D3C", "#FDAE6B", "#FDD0A2", "#31A354", "#74C476",
    "#A1D99B", "#C7E9C0", "#756BB1", "#9E9AC8", "#BCBDDC", "#DADAEB", "#636363", "#969696", "#BDBDBD", "#D9D9D9",
    "#393E46", "#00ADB5", "#FF5722", "#795548", "#607D8B", "#8BC34A", "#CDDC39", "#FFC107", "#FF4081", "#3F51B5",
]
PPTX_PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
PPTX_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PPTX_PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PPTX_DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PPTX_SECTION_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
PPTX_SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"
SECTION_LABEL_ALIASES: Dict[str, str] = {
    "laun": "Laundry",
    "mayo": "Mayonesa",
    "cond": "Acondicionador",
    "shampoo": "Shampoo",
    "deos": "Deos",
    "clean": "Clean",
    "hair": "Hair",
    "bar": "Bar",
    "liquido": "Liquido",
    "fe": "FE",
    "sc": "SC",
}
PRIMARY_SECTION_ALIASES: Dict[str, str] = {
    "clean": "Clean",
    "fe": "FE",
    "laun": "Laundry",
    "laundry": "Laundry",
    "mayo": "Mayonesa",
    "mayonesa": "Mayonesa",
    "deos": "Deos",
    "deod": "Deos",
    "hair": "Hair",
    "sc": "SC",
}


def normalize_brand_key(brand: str) -> str:
    """Normaliza fabricante para mapearlo siempre al mismo color."""
    normalized = unicodedata.normalize("NFKD", str(brand or ""))
    normalized = "".join(ch for ch in normalized if not unicodedata.combining(ch))
    normalized = re.sub(r"\s+", " ", normalized).strip().lower()
    return normalized


def normalize_section_title(value: str) -> str:
    """Normaliza una etiqueta para usarla como nombre de seccion."""
    normalized = unicodedata.normalize("NFKD", str(value or ""))
    normalized = "".join(ch for ch in normalized if not unicodedata.combining(ch))
    normalized = normalized.replace("_", " ").replace(".", " ")
    normalized = re.sub(r"\s+", " ", normalized).strip()
    if not normalized:
        return ""
    lower_key = normalized.lower()
    return SECTION_LABEL_ALIASES.get(lower_key, normalized)


def is_total_group_sheet(sheet_name: str) -> bool:
    """Detecta hojas que representan un total/subgrupo y sirven como ancla de seccion."""
    cleaned = _clean_brand_name_from_sheet(sheet_name)
    normalized = normalize_brand_key(cleaned)
    return normalized.startswith("t ") or normalized.startswith("t.") or normalized.startswith("total ")


def extract_total_group_tokens(sheet_name: str) -> List[str]:
    """Extrae tokens significativos de una hoja total para clasificar su nivel de agrupacion."""
    cleaned = _clean_brand_name_from_sheet(sheet_name)
    normalized = unicodedata.normalize("NFKD", cleaned)
    normalized = "".join(ch for ch in normalized if not unicodedata.combining(ch))
    return [
        token.lower()
        for token in re.split(r"[^A-Za-z0-9']+", normalized)
        if token
    ]


def extract_total_group_display_tokens(sheet_name: str) -> List[str]:
    """Extrae tokens de una hoja total conservando capitalizacion para mostrar."""
    cleaned = _clean_brand_name_from_sheet(sheet_name)
    normalized = unicodedata.normalize("NFKD", cleaned)
    normalized = "".join(ch for ch in normalized if not unicodedata.combining(ch))
    return [
        token
        for token in re.split(r"[^A-Za-z0-9']+", normalized)
        if token
    ]


def derive_primary_section_title_from_total_sheet(sheet_name: str) -> str:
    """Devuelve la seccion principal si la hoja total representa un grupo raiz."""
    generic_tokens = {"t", "total", "ul"}
    for token in extract_total_group_tokens(sheet_name):
        if token in generic_tokens:
            continue
        mapped = PRIMARY_SECTION_ALIASES.get(token)
        if mapped:
            return mapped
    return ""


def derive_section_title_from_total_sheet(sheet_name: str) -> str:
    """Deriva el nombre de la seccion principal a partir de una hoja total."""
    cleaned = _clean_brand_name_from_sheet(sheet_name)
    tokens = extract_total_group_display_tokens(sheet_name)
    generic_tokens = {"t", "total"}
    for token in reversed(tokens):
        token_key = token.lower()
        if token_key in generic_tokens:
            continue
        if len(token_key) <= 2 and token_key.isalpha():
            continue
        return normalize_section_title(token)
    return normalize_section_title(cleaned)


def _section_membership_key(value: str) -> str:
    """Normaliza una hoja/seccion a su marca base para decidir herencia."""
    cleaned = _clean_brand_name_from_sheet(value)
    cleaned = re.sub(r"\([^)]*\)", "", cleaned)
    cleaned = re.sub(r"(?i)^\s*(?:t|total)\s*[.\-_\s]+", "", cleaned).strip()
    return normalize_brand_key(normalize_section_title(cleaned))


def _section_exact_title_key(value: str) -> str:
    """Normaliza una hoja/seccion conservando el detalle entre parentesis."""
    cleaned = _clean_brand_name_from_sheet(value)
    cleaned = re.sub(r"(?i)^\s*(?:t|total)\s*[.\-_\s]+", "", cleaned).strip()
    return normalize_brand_key(normalize_section_title(cleaned))


def _section_has_detail(value: str) -> bool:
    """Detecta titulos con detalle explicito, por ejemplo 'Marca (Segmento)'."""
    return bool(re.search(r"\([^)]*\)", str(value or "")))


def should_inherit_current_section(sheet_name: str, section_title: Optional[str]) -> bool:
    """Indica si una hoja debe heredar la seccion activa sin perder subsegmentos."""
    if not section_title:
        return False
    current_exact_key = _section_exact_title_key(section_title)
    sheet_exact_key = _section_exact_title_key(sheet_name)
    if current_exact_key and current_exact_key == sheet_exact_key:
        return True
    if _section_has_detail(section_title):
        return False
    return sheet_belongs_to_section(sheet_name, section_title)


def sheet_belongs_to_section(sheet_name: str, section_title: Optional[str]) -> bool:
    """Indica si una hoja debe heredar la seccion abierta actualmente."""
    current_key = _section_membership_key(section_title or "")
    if not current_key:
        return False
    sheet_key = _section_membership_key(sheet_name)
    return bool(sheet_key and sheet_key == current_key)


def build_section_title_for_sheet(sheet_name: str, current_group: Optional[str]) -> Tuple[str, Optional[str]]:
    """Resuelve la seccion aplicable a una hoja y actualiza el grupo actual si corresponde."""
    if is_total_group_sheet(sheet_name):
        primary_group_title = derive_primary_section_title_from_total_sheet(sheet_name)
        if primary_group_title:
            return primary_group_title, primary_group_title
        group_title = derive_section_title_from_total_sheet(sheet_name)
        if current_group and (
            should_inherit_current_section(sheet_name, current_group)
            or "(" not in str(sheet_name or "")
        ):
            return current_group, current_group
        return group_title, group_title
    brand_title = normalize_section_title(_clean_brand_name_from_sheet(sheet_name))
    if current_group and should_inherit_current_section(sheet_name, current_group):
        return current_group, current_group
    return brand_title, brand_title


def build_metadata_group_for_sheet(sheet_name: str, current_group: Optional[str], fabricante: str = "") -> Tuple[str, Optional[str]]:
    """Resuelve el grupo semantico para metadata del banco con reglas por fabricante."""
    manufacturer_key = _resolve_mult_manufacturer_key(fabricante) if fabricante else None
    if manufacturer_key:
        rule_config = MULT_METADATA_RULES.get(manufacturer_key, {})
        category_rules = rule_config.get("category_rules", ())
        if rule_config.get("category_source") == "section":
            # Solo promovemos hojas "ancla" que realmente correspondan a una
            # categoria; totales de submarca como "T.Lux" deben heredar grupo.
            if is_total_group_sheet(sheet_name):
                candidate_titles = [
                    derive_primary_section_title_from_total_sheet(sheet_name),
                    derive_section_title_from_total_sheet(sheet_name),
                ]
                for candidate_title in candidate_titles:
                    if candidate_title and _match_metadata_rule(candidate_title, category_rules):
                        return candidate_title, candidate_title
    brand_title = normalize_section_title(_clean_brand_name_from_sheet(sheet_name))
    return current_group or brand_title, current_group


def register_section_slide_range(
    section_slide_map: Dict[str, List[int]],
    section_title: str,
    start_idx: int,
    count: int,
) -> None:
    """Agrega un rango de slides a una seccion conservando el orden."""
    if count <= 0:
        return
    title = normalize_section_title(section_title)
    if not title:
        return
    bucket = section_slide_map.setdefault(title, [])
    known = set(bucket)
    for slide_idx in range(start_idx, start_idx + count):
        if slide_idx not in known:
            bucket.append(slide_idx)
            known.add(slide_idx)


def apply_powerpoint_sections(pptx_path: str, section_slide_map: Dict[str, List[int]]) -> None:
    """Inyecta secciones de PowerPoint dentro del .pptx final."""
    if not pptx_path or not section_slide_map:
        return
    ET.register_namespace("a", PPTX_DRAWING_NS)
    ET.register_namespace("r", PPTX_REL_NS)
    ET.register_namespace("p", PPTX_PRESENTATION_NS)
    ET.register_namespace("p14", PPTX_SECTION_NS)
    with zipfile.ZipFile(pptx_path, "r") as src_zip:
        root = ET.fromstring(src_zip.read("ppt/presentation.xml"))
    sld_id_list = root.find(f"{{{PPTX_PRESENTATION_NS}}}sldIdLst")
    if sld_id_list is None:
        return
    slide_ids = [
        int(node.attrib["id"])
        for node in sld_id_list.findall(f"{{{PPTX_PRESENTATION_NS}}}sldId")
        if node.attrib.get("id")
    ]
    valid_sections: List[Tuple[str, List[int]]] = []
    for title, slide_indexes in section_slide_map.items():
        slide_id_values: List[int] = []
        seen_slide_ids: Set[int] = set()
        for slide_idx in slide_indexes:
            if 0 <= slide_idx < len(slide_ids):
                slide_id = slide_ids[slide_idx]
                if slide_id not in seen_slide_ids:
                    slide_id_values.append(slide_id)
                    seen_slide_ids.add(slide_id)
        if slide_id_values:
            valid_sections.append((title, slide_id_values))
    if not valid_sections:
        return
    ext_lst = root.find(f"{{{PPTX_PRESENTATION_NS}}}extLst")
    if ext_lst is None:
        ext_lst = ET.SubElement(root, f"{{{PPTX_PRESENTATION_NS}}}extLst")
    for ext in list(ext_lst.findall(f"{{{PPTX_PRESENTATION_NS}}}ext")):
        if ext.attrib.get("uri") == PPTX_SECTION_EXT_URI:
            ext_lst.remove(ext)
    sections_ext = ET.SubElement(
        ext_lst,
        f"{{{PPTX_PRESENTATION_NS}}}ext",
        {"uri": PPTX_SECTION_EXT_URI},
    )
    section_lst = ET.SubElement(sections_ext, f"{{{PPTX_SECTION_NS}}}sectionLst")
    for title, slide_id_values in valid_sections:
        section = ET.SubElement(
            section_lst,
            f"{{{PPTX_SECTION_NS}}}section",
            {"name": title, "id": "{" + str(uuid.uuid4()).upper() + "}"},
        )
        section_slide_list = ET.SubElement(section, f"{{{PPTX_SECTION_NS}}}sldIdLst")
        for slide_id in slide_id_values:
            ET.SubElement(
                section_slide_list,
                f"{{{PPTX_SECTION_NS}}}sldId",
                {"id": str(slide_id)},
            )
    updated_xml = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    tmp_path: Optional[str] = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx", dir=os.path.dirname(pptx_path)) as tmp_file:
            tmp_path = tmp_file.name
        with zipfile.ZipFile(pptx_path, "r") as read_zip, zipfile.ZipFile(tmp_path, "w") as write_zip:
            for entry in read_zip.infolist():
                payload = updated_xml if entry.filename == "ppt/presentation.xml" else read_zip.read(entry.filename)
                write_zip.writestr(entry, payload)
        os.replace(tmp_path, pptx_path)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


def _make_table_border_xml(side: str, *, visible: bool, width: int = 12700, color_hex: str = "000000") -> ET.Element:
    line_width = str(int(width if visible else 0))
    line = ET.Element(
        f"{{{PPTX_DRAWING_NS}}}{side}",
        {"w": line_width, "cap": "flat", "cmpd": "sng", "algn": "ctr"},
    )
    solid_fill = ET.SubElement(line, f"{{{PPTX_DRAWING_NS}}}solidFill")
    if visible:
        ET.SubElement(solid_fill, f"{{{PPTX_DRAWING_NS}}}srgbClr", {"val": color_hex})
    else:
        color = ET.SubElement(solid_fill, f"{{{PPTX_DRAWING_NS}}}prstClr", {"val": "black"})
        ET.SubElement(color, f"{{{PPTX_DRAWING_NS}}}alpha", {"val": "0"})
    ET.SubElement(line, f"{{{PPTX_DRAWING_NS}}}prstDash", {"val": "solid"})
    ET.SubElement(line, f"{{{PPTX_DRAWING_NS}}}round")
    ET.SubElement(line, f"{{{PPTX_DRAWING_NS}}}headEnd", {"type": "none", "w": "med", "len": "med"})
    ET.SubElement(line, f"{{{PPTX_DRAWING_NS}}}tailEnd", {"type": "none", "w": "med", "len": "med"})
    return line


def _force_text_size_xml(container_node: ET.Element, font_size_points: int) -> None:
    """Fija el tamano de texto en runs XML para evitar defaults ambiguos en PowerPoint Web."""
    size_value = str(int(font_size_points) * 100)
    for paragraph_node in container_node.findall(f".//{{{PPTX_DRAWING_NS}}}p"):
        for run_node in paragraph_node.findall(f"{{{PPTX_DRAWING_NS}}}r"):
            run_props = run_node.find(f"{{{PPTX_DRAWING_NS}}}rPr")
            if run_props is None:
                run_props = ET.Element(f"{{{PPTX_DRAWING_NS}}}rPr")
                run_node.insert(0, run_props)
            run_props.set("sz", size_value)
        end_props = paragraph_node.find(f"{{{PPTX_DRAWING_NS}}}endParaRPr")
        if end_props is not None:
            end_props.set("sz", size_value)


def apply_summary_table_border_style_in_pptx(pptx_path: str) -> None:
    """Evita que PowerPoint Web pinte contornos negros en la tabla summary."""
    if not pptx_path or not os.path.exists(pptx_path):
        return
    ET.register_namespace("a", PPTX_DRAWING_NS)
    ET.register_namespace("r", PPTX_REL_NS)
    ET.register_namespace("p", PPTX_PRESENTATION_NS)

    updated_parts: Dict[str, bytes] = {}
    with zipfile.ZipFile(pptx_path, "r") as src_zip:
        for entry in src_zip.infolist():
            if not re.match(r"^ppt/slides/slide\d+\.xml$", entry.filename):
                continue
            slide_xml = src_zip.read(entry.filename)
            slide_root = ET.fromstring(slide_xml)
            slide_changed = False
            for tbl_node in slide_root.findall(f".//{{{PPTX_DRAWING_NS}}}tbl"):
                rows = tbl_node.findall(f"{{{PPTX_DRAWING_NS}}}tr")
                if not rows:
                    continue
                header_texts = [
                    "".join(text_node.text or "" for text_node in cell_node.findall(f".//{{{PPTX_DRAWING_NS}}}t")).strip()
                    for cell_node in rows[0].findall(f"{{{PPTX_DRAWING_NS}}}tc")
                ]
                normalized_headers = {re.sub(r"\s+", " ", header).strip().lower() for header in header_texts}
                is_summary_table = (
                    any(header in normalized_headers for header in ("fabricante/marca", "manufacturer/brand"))
                    and "pipeline" in normalized_headers
                    and any("worldpanel by numerator" in header for header in normalized_headers)
                    and any(header.startswith(("cobertura ", "coverage ")) for header in normalized_headers)
                )
                if not is_summary_table:
                    continue

                tbl_pr = tbl_node.find(f"{{{PPTX_DRAWING_NS}}}tblPr")
                if tbl_pr is None:
                    tbl_pr = ET.Element(f"{{{PPTX_DRAWING_NS}}}tblPr")
                    tbl_node.insert(0, tbl_pr)
                for style_node in list(tbl_pr.findall(f"{{{PPTX_DRAWING_NS}}}tableStyleId")):
                    tbl_pr.remove(style_node)
                for attr_name in ("firstRow", "bandRow", "lastRow", "firstCol", "lastCol", "bandCol"):
                    tbl_pr.attrib.pop(attr_name, None)

                row_count = len(rows)
                for row_idx, row_node in enumerate(rows):
                    cells = row_node.findall(f"{{{PPTX_DRAWING_NS}}}tc")
                    col_count = len(cells)
                    for col_idx, cell_node in enumerate(cells):
                        tc_pr = cell_node.find(f"{{{PPTX_DRAWING_NS}}}tcPr")
                        if tc_pr is None:
                            tc_pr = ET.SubElement(cell_node, f"{{{PPTX_DRAWING_NS}}}tcPr")
                        for child in list(tc_pr):
                            local_name = child.tag.rsplit("}", 1)[-1]
                            if local_name in ("lnL", "lnR", "lnT", "lnB"):
                                tc_pr.remove(child)
                        side_visibility = {
                            "lnL": False,
                            "lnR": col_idx < col_count - 1,
                            "lnT": False,
                            "lnB": row_idx > 0 and row_idx < row_count - 1,
                        }
                        insert_index = 0
                        for side, visible in side_visibility.items():
                            tc_pr.insert(
                                insert_index,
                                _make_table_border_xml(side, visible=visible, width=int(5715), color_hex="FFFFFF"),
                            )
                            insert_index += 1
                slide_changed = True
            if slide_changed:
                updated_parts[entry.filename] = ET.tostring(slide_root, encoding="utf-8", xml_declaration=True)

    if not updated_parts:
        return

    tmp_path: Optional[str] = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx", dir=os.path.dirname(pptx_path)) as tmp_file:
            tmp_path = tmp_file.name
        with zipfile.ZipFile(pptx_path, "r") as read_zip, zipfile.ZipFile(tmp_path, "w") as write_zip:
            for entry in read_zip.infolist():
                payload = updated_parts.get(entry.filename)
                write_zip.writestr(entry, payload if payload is not None else read_zip.read(entry.filename))
        os.replace(tmp_path, pptx_path)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


def apply_variation_table_internal_borders_in_pptx(pptx_path: str) -> None:
    """Oculta el perimetro de tablas VAR % MAT y conserva solo bordes internos."""
    if not pptx_path or not os.path.exists(pptx_path):
        return
    ET.register_namespace("a", PPTX_DRAWING_NS)
    ET.register_namespace("r", PPTX_REL_NS)
    ET.register_namespace("p", PPTX_PRESENTATION_NS)

    updated_parts: Dict[str, bytes] = {}
    with zipfile.ZipFile(pptx_path, "r") as src_zip:
        for entry in src_zip.infolist():
            if not re.match(r"^ppt/slides/slide\d+\.xml$", entry.filename):
                continue
            slide_xml = src_zip.read(entry.filename)
            slide_root = ET.fromstring(slide_xml)
            slide_changed = False
            for tbl_node in slide_root.findall(f".//{{{PPTX_DRAWING_NS}}}tbl"):
                table_text = " ".join(
                    text_node.text or ""
                    for text_node in tbl_node.findall(f".//{{{PPTX_DRAWING_NS}}}t")
                )
                if "VAR % MAT" not in table_text:
                    continue

                tbl_pr = tbl_node.find(f"{{{PPTX_DRAWING_NS}}}tblPr")
                if tbl_pr is None:
                    tbl_pr = ET.Element(f"{{{PPTX_DRAWING_NS}}}tblPr")
                    tbl_node.insert(0, tbl_pr)
                for style_node in list(tbl_pr.findall(f"{{{PPTX_DRAWING_NS}}}tableStyleId")):
                    tbl_pr.remove(style_node)
                for attr_name in ("firstRow", "bandRow", "lastRow", "firstCol", "lastCol", "bandCol"):
                    tbl_pr.attrib.pop(attr_name, None)

                rows = tbl_node.findall(f"{{{PPTX_DRAWING_NS}}}tr")
                row_count = len(rows)
                for row_idx, row_node in enumerate(rows):
                    cells = row_node.findall(f"{{{PPTX_DRAWING_NS}}}tc")
                    col_count = len(cells)
                    for col_idx, cell_node in enumerate(cells):
                        tc_pr = cell_node.find(f"{{{PPTX_DRAWING_NS}}}tcPr")
                        if tc_pr is None:
                            tc_pr = ET.SubElement(cell_node, f"{{{PPTX_DRAWING_NS}}}tcPr")
                        for side in ("lnL", "lnR", "lnT", "lnB"):
                            existing = tc_pr.find(f"{{{PPTX_DRAWING_NS}}}{side}")
                            if existing is not None:
                                tc_pr.remove(existing)
                        side_visibility = {
                            "lnL": col_idx > 0,
                            "lnR": col_idx < col_count - 1,
                            "lnT": row_idx > 0,
                            "lnB": row_idx < row_count - 1,
                        }
                        insert_index = 0
                        for side, visible in side_visibility.items():
                            tc_pr.insert(insert_index, _make_table_border_xml(side, visible=visible))
                            insert_index += 1
                slide_changed = True
            if slide_changed:
                updated_parts[entry.filename] = ET.tostring(slide_root, encoding="utf-8", xml_declaration=True)

    if not updated_parts:
        return

    tmp_path: Optional[str] = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx", dir=os.path.dirname(pptx_path)) as tmp_file:
            tmp_path = tmp_file.name
        with zipfile.ZipFile(pptx_path, "r") as read_zip, zipfile.ZipFile(tmp_path, "w") as write_zip:
            for entry in read_zip.infolist():
                payload = updated_parts.get(entry.filename)
                write_zip.writestr(entry, payload if payload is not None else read_zip.read(entry.filename))
        os.replace(tmp_path, pptx_path)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


def apply_trend_variation_table_transparent_style_in_pptx(pptx_path: str) -> None:
    """Deja transparente la tabla de variaciones del slide de tendencia."""
    if not pptx_path or not os.path.exists(pptx_path):
        return
    ET.register_namespace("a", PPTX_DRAWING_NS)
    ET.register_namespace("r", PPTX_REL_NS)
    ET.register_namespace("p", PPTX_PRESENTATION_NS)

    updated_parts: Dict[str, bytes] = {}
    with zipfile.ZipFile(pptx_path, "r") as src_zip:
        for entry in src_zip.infolist():
            if not re.match(r"^ppt/slides/slide\d+\.xml$", entry.filename):
                continue
            slide_xml = src_zip.read(entry.filename)
            slide_root = ET.fromstring(slide_xml)
            slide_changed = False
            for tbl_node in slide_root.findall(f".//{{{PPTX_DRAWING_NS}}}tbl"):
                rows = tbl_node.findall(f"{{{PPTX_DRAWING_NS}}}tr")
                if not rows:
                    continue
                header_texts = [
                    "".join(text_node.text or "" for text_node in cell_node.findall(f".//{{{PPTX_DRAWING_NS}}}t")).strip()
                    for cell_node in rows[0].findall(f"{{{PPTX_DRAWING_NS}}}tc")
                ]
                normalized_headers = {re.sub(r"\s+", " ", header).strip().lower() for header in header_texts}
                is_trend_variation_table = (
                    "tipo" in normalized_headers
                    and "periodo" in normalized_headers
                    and any("wp by numerator" in header for header in normalized_headers)
                )
                if not is_trend_variation_table:
                    continue

                tbl_pr = tbl_node.find(f"{{{PPTX_DRAWING_NS}}}tblPr")
                if tbl_pr is None:
                    tbl_pr = ET.Element(f"{{{PPTX_DRAWING_NS}}}tblPr")
                    tbl_node.insert(0, tbl_pr)
                for style_node in list(tbl_pr.findall(f"{{{PPTX_DRAWING_NS}}}tableStyleId")):
                    tbl_pr.remove(style_node)
                for attr_name in ("firstRow", "bandRow", "lastRow", "firstCol", "lastCol", "bandCol"):
                    tbl_pr.attrib.pop(attr_name, None)

                for row_idx, row_node in enumerate(rows):
                    row_font_size = 8 if row_idx == 0 or len(rows) > 5 else 9
                    for cell_node in row_node.findall(f"{{{PPTX_DRAWING_NS}}}tc"):
                        _force_text_size_xml(cell_node, row_font_size)
                        tc_pr = cell_node.find(f"{{{PPTX_DRAWING_NS}}}tcPr")
                        if tc_pr is None:
                            tc_pr = ET.SubElement(cell_node, f"{{{PPTX_DRAWING_NS}}}tcPr")
                        for child in list(tc_pr):
                            local_name = child.tag.rsplit("}", 1)[-1]
                            if local_name in ("lnL", "lnR", "lnT", "lnB"):
                                tc_pr.remove(child)
                        insert_index = 0
                        for side in ("lnL", "lnR", "lnT", "lnB"):
                            tc_pr.insert(insert_index, _make_table_border_xml(side, visible=False))
                            insert_index += 1
                slide_changed = True
            if slide_changed:
                updated_parts[entry.filename] = ET.tostring(slide_root, encoding="utf-8", xml_declaration=True)

    if not updated_parts:
        return

    tmp_path: Optional[str] = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx", dir=os.path.dirname(pptx_path)) as tmp_file:
            tmp_path = tmp_file.name
        with zipfile.ZipFile(pptx_path, "r") as read_zip, zipfile.ZipFile(tmp_path, "w") as write_zip:
            for entry in read_zip.infolist():
                payload = updated_parts.get(entry.filename)
                write_zip.writestr(entry, payload if payload is not None else read_zip.read(entry.filename))
        os.replace(tmp_path, pptx_path)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


def apply_table_grid_widths_in_pptx(
    pptx_path: str,
    *,
    slide_index: int,
    header_row: Sequence[str],
    column_widths: Sequence[int],
) -> None:
    """Reescribe el tblGrid de una tabla para fijar anchos persistentes en el .pptx.

    python-pptx deja la tabla correcta en memoria, pero en algunos casos el save()
    reserializa el tblGrid con anchos distintos. Este post-proceso fuerza el ancho
    final que realmente termina viendo PowerPoint.
    """
    if not pptx_path or slide_index < 0 or not header_row or not column_widths:
        return
    ET.register_namespace("a", PPTX_DRAWING_NS)
    ET.register_namespace("r", PPTX_REL_NS)
    ET.register_namespace("p", PPTX_PRESENTATION_NS)
    header_texts = [str(value).strip() for value in header_row]
    updated_slide_xml: Optional[bytes] = None
    slide_part_path: Optional[str] = None
    with zipfile.ZipFile(pptx_path, "r") as src_zip:
        presentation_root = ET.fromstring(src_zip.read("ppt/presentation.xml"))
        slide_id_list = presentation_root.find(f"{{{PPTX_PRESENTATION_NS}}}sldIdLst")
        if slide_id_list is None:
            return
        slide_nodes = slide_id_list.findall(f"{{{PPTX_PRESENTATION_NS}}}sldId")
        if not (0 <= slide_index < len(slide_nodes)):
            return
        slide_rel_id = slide_nodes[slide_index].attrib.get(f"{{{PPTX_REL_NS}}}id")
        if not slide_rel_id:
            return
        rels_root = ET.fromstring(src_zip.read("ppt/_rels/presentation.xml.rels"))
        slide_target: Optional[str] = None
        for rel_node in rels_root.findall(f"{{{PPTX_PACKAGE_REL_NS}}}Relationship"):
            if rel_node.attrib.get("Id") == slide_rel_id:
                slide_target = rel_node.attrib.get("Target")
                break
        if not slide_target:
            return
        slide_part_path = posixpath.normpath(posixpath.join("ppt", slide_target.lstrip("/")))
        slide_root = ET.fromstring(src_zip.read(slide_part_path))
        for tbl_node in slide_root.findall(f".//{{{PPTX_DRAWING_NS}}}tbl"):
            first_row = tbl_node.find(f"{{{PPTX_DRAWING_NS}}}tr")
            tbl_grid = tbl_node.find(f"{{{PPTX_DRAWING_NS}}}tblGrid")
            if first_row is None or tbl_grid is None:
                continue
            current_headers: List[str] = []
            for cell_node in first_row.findall(f"{{{PPTX_DRAWING_NS}}}tc"):
                text_fragments = [
                    text_node.text or ""
                    for text_node in cell_node.findall(f".//{{{PPTX_DRAWING_NS}}}t")
                ]
                current_headers.append("".join(text_fragments).strip())
            if current_headers != header_texts:
                continue
            grid_cols = tbl_grid.findall(f"{{{PPTX_DRAWING_NS}}}gridCol")
            if len(grid_cols) != len(column_widths):
                continue
            for grid_col, width in zip(grid_cols, column_widths):
                grid_col.set("w", str(int(width)))
            updated_slide_xml = ET.tostring(slide_root, encoding="utf-8", xml_declaration=True)
            break
    if not slide_part_path or updated_slide_xml is None:
        return
    tmp_path: Optional[str] = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx", dir=os.path.dirname(pptx_path)) as tmp_file:
            tmp_path = tmp_file.name
        with zipfile.ZipFile(pptx_path, "r") as read_zip, zipfile.ZipFile(tmp_path, "w") as write_zip:
            for entry in read_zip.infolist():
                payload = updated_slide_xml if entry.filename == slide_part_path else read_zip.read(entry.filename)
                write_zip.writestr(entry, payload)
        os.replace(tmp_path, pptx_path)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


def parse_filename_brand(filename: str) -> Tuple[str, str, str, str]:
    """Separa nombre de archivo en prefijo, marca, sufijo y extensión."""
    base, ext = os.path.splitext(filename)
    parts = base.split("_", 2)
    if len(parts) < 3:
        return "", "", "", ext
    prefix = f"{parts[0]}_{parts[1]}_"
    raw_brand_segment = parts[2].strip()
    if not raw_brand_segment:
        return prefix, "", "", ext
    # Conserva sufijos como "- copia"/"- preview" fuera del texto coloreado para
    # que el color represente únicamente al fabricante.
    match = PREVIEW_COPY_SUFFIX_RE.match(raw_brand_segment)
    if not match:
        return prefix, raw_brand_segment, "", ext
    brand = (match.group("brand") or "").strip()
    suffix = match.group("suffix") or ""
    return prefix, brand, suffix, ext


def ansi_truecolor(text: str, rgb: Tuple[int, int, int]) -> str:
    """Envuelve texto con ANSI 24-bit."""
    r, g, b = rgb
    return f"\033[38;2;{r};{g};{b}m{text}{ANSI_RESET}"


SCENARIO_AUTO = "AUTO"
SCENARIO_AUTO_DUAL_AXIS = "AUTO_DOBLE_EJE"
SCENARIO_AUTO_OPTIMAL_PIPELINE = "AUTO_PIPELINE_OPTIMO"
SCENARIO_PG_GLOBAL_EN = "PG_GLOBAL_EN"
SCENARIO_NATURA_BR = "NATURA_BR"
SCENARIO_PG_COLOR = (64, 105, 205)
SCENARIO_NATURA_COLOR = (255, 105, 19)


def relative_luminance(rgb: Tuple[int, int, int]) -> float:
    """Calcula luminancia relativa simple para evitar colores oscuros."""
    r, g, b = rgb
    return (0.2126 * r) + (0.7152 * g) + (0.0722 * b)


def lift_color_to_min_luminance(
    rgb: Tuple[int, int, int], min_luminance: float = MIN_READABLE_LUMINANCE
) -> Tuple[int, int, int]:
    """Aclara un color mezclándolo con blanco hasta alcanzar luminancia mínima."""
    lum = relative_luminance(rgb)
    if lum >= min_luminance:
        return rgb
    r, g, b = rgb
    # Mezcla lineal con blanco; evita colores apagados/obscuros en terminal.
    mix = min(0.7, max(0.0, (min_luminance - lum) / 255.0))
    r = int(round(r + (255 - r) * mix))
    g = int(round(g + (255 - g) * mix))
    b = int(round(b + (255 - b) * mix))
    return (r, g, b)

def _hex_to_tab_argb(hex_color: str) -> str:
    """Convierte #RRGGBB a AARRGGBB para color de pestaña de Excel."""
    value = str(hex_color or "").strip()
    if value.startswith("#"):
        value = value[1:]
    if len(value) != 6:
        return "FF808080"
    return f"FF{value.upper()}"

def assign_brand_tab_color(brand_label: str, brand_color_lookup: Dict[str, str]) -> str:
    """
    Asigna color consistente por marca usando la misma lógica de paleta del generador de informe:
    primer color libre no usado, fallback cíclico.
    """
    key = normalize_brand_key(brand_label)
    if not key:
        return TEMPLATE_TAB_COLOR_SEQUENCE[0]
    existing = brand_color_lookup.get(key)
    if existing:
        return existing
    used_colors = {color for color in brand_color_lookup.values() if color}
    for candidate in TEMPLATE_TAB_COLOR_SEQUENCE:
        if candidate not in used_colors:
            brand_color_lookup[key] = candidate
            return candidate
    fallback = TEMPLATE_TAB_COLOR_SEQUENCE[len(brand_color_lookup) % len(TEMPLATE_TAB_COLOR_SEQUENCE)]
    brand_color_lookup[key] = fallback
    return fallback

def apply_template_tab_colors(xlsx_path: str, marcas: Sequence[str]) -> None:
    """Pinta pestañas del template por marca para facilitar lectura de agrupado."""
    from openpyxl import load_workbook as _load_wb_tabs
    wb_tabs = _load_wb_tabs(xlsx_path)
    color_lookup: Dict[str, str] = {}
    for brand in marcas:
        if brand not in wb_tabs.sheetnames:
            continue
        color_hex = assign_brand_tab_color(brand, color_lookup)
        ws = wb_tabs[brand]
        ws.sheet_properties.tabColor = _hex_to_tab_argb(color_hex)
    wb_tabs.save(xlsx_path)

def normalize_variations_box_style(raw_value: Optional[str]) -> str:
    """Normaliza el estilo del cuadro de variaciones (classic | pretty)."""
    val = (raw_value or "").strip().lower()
    if not val:
        return "classic"
    if val in {"pretty", "bonito", "nuevo", "nice", "card", "cards", "2"}:
        return "pretty"
    if val in {"classic", "clasico", "clásico", "tabla", "1"}:
        return "classic"
    return "classic"


def normalize_coverage_slide_variant(raw_value: Optional[str]) -> str:
    """Normaliza el modo del slide de cobertura (classic | complemented | pg)."""
    val = (raw_value or "").strip().lower()
    if not val:
        return "classic"
    if val in {"pg", "p&g", "p & g", "procter", "procter & gamble", "procter and gamble", "3"}:
        return "pg"
    if val in {"complemented", "complementado", "complement", "penetracion", "penetración", "penetration", "2"}:
        return "complemented"
    if val in {"classic", "clasico", "clásico", "variacion", "variación", "var", "1"}:
        return "classic"
    return "classic"


def coverage_slide_variant_label(variant: str) -> str:
    normalized = normalize_coverage_slide_variant(variant)
    if normalized == "complemented":
        return "Complementado"
    if normalized == "pg":
        return "P&G"
    return "Clasico"

def normalize_evolution_slide_variant(raw_value: Optional[str]) -> str:
    """Normaliza el modo del slide de Evolucion mensual y variacion (classic | simple)."""
    val = (raw_value or "").strip().lower()
    if not val:
        return "classic"
    if val in {"simple", "basico", "basica", "basic", "1"}:
        return "simple"
    if val in {"classic", "clasico", "clásico", "avanzado", "advanced", "2"}:
        return "classic"
    return "classic"


def normalize_trend_granularity(raw_value: Optional[str]) -> str:
    """Normaliza la periodicidad del gráfico de tendencia (monthly | quarterly)."""
    val = (raw_value or "").strip().lower()
    if not val:
        return "monthly"
    if val in {"quarterly", "quarter", "trimestral", "trimestre", "quarterly_3m", "q", "2"}:
        return "quarterly"
    if val in {"monthly", "month", "mensual", "mes", "m", "1"}:
        return "monthly"
    return "monthly"


def trend_granularity_label(granularity: str) -> str:
    normalized = normalize_trend_granularity(granularity)
    return "Trimestral" if normalized == "quarterly" else "Mensual"

def _register_brand_exception(marca_label: Optional[str], reason: str) -> None:
    normalized = (marca_label or "N/D").strip() or "N/D"
    reason_set = BRAND_EXCEPTION_REASONS.setdefault(normalized, set())
    if reason in reason_set:
        return
    reason_set.add(reason)
    style = EXCEPTION_STYLES.get(reason, EXCEPTION_STYLES["zero_dash"])
    # Mensaje en rojo con el nombre de la marca en amarillo al inicio
    print(f"{Fore.RED}{Fore.YELLOW}{normalized}{Fore.RED} {style['message']}")


def notify_zero_months_exception(marca_label: Optional[str]) -> None:
    _register_brand_exception(marca_label, "zero_dash")


def notify_negative_values_exception(marca_label: Optional[str]) -> None:
    _register_brand_exception(marca_label, "negative")


def notify_buyers_threshold(marca_label: Optional[str], buyers_value: Optional[float], threshold: float = 200) -> None:
    if buyers_value is None:
        return
    try:
        if pd.isna(buyers_value):
            return
        buyers_num = float(buyers_value)
    except Exception:
        return
    normalized = (marca_label or "N/D").strip() or "N/D"
    buyers_display = f"{buyers_num:.0f}"
    if buyers_num < threshold:
        print(Fore.RED + f"{normalized} cuenta con {buyers_display} compradores promedio, tener precaución")
    else:
        print(Fore.GREEN + f"{normalized} si cuenta con al menos {int(threshold)} compradores")


def report_zero_months_exceptions() -> None:
    if not BRAND_EXCEPTION_REASONS:
        return
    print(f"{Fore.RED}Marcas con excepción detectada:")
    for marca in sorted(BRAND_EXCEPTION_REASONS):
        tags = "/".join(
            sorted(
                {
                    EXCEPTION_STYLES.get(reason, {}).get("summary_tag", reason)
                    for reason in BRAND_EXCEPTION_REASONS[marca]
                }
            )
        )
        print(f"{Fore.RED}- {Fore.YELLOW}{marca}{Fore.RED} [{tags}]")
    BRAND_EXCEPTION_REASONS.clear()


def detect_brand_data_issues(df_marca: "pd.DataFrame", window: int = 12) -> Set[str]:
    issues: Set[str] = set()
    if df_marca is None or df_marca.empty:
        return issues
    cols_to_check = [COL_SELL_IN, COL_SELL_OUT]
    tail_df = df_marca.tail(window) if window > 0 else df_marca
    for col in cols_to_check:
        series = tail_df[col]
        str_series = series.astype(str).str.strip()
        if str_series.eq("-").any():
            issues.add("zero_dash")
        numeric = pd.to_numeric(series, errors="coerce")
        if (numeric == 0).any():
            issues.add("zero_dash")
        if (numeric < 0).any():
            issues.add("negative")
    return issues

CATEGORIES_CSV_DATA = """cod,cest,cat
ALCB,Bebidas,Bebidas Alcoholicas
BEER,Bebidas,Cervezas
CARB,Bebidas,Bebidas Gaseosas
CWAT,Bebidas,Agua Gasificada
COCW,Bebidas,Agua de Coco
COFF,Bebidas,Cafe-Consolidado de Cafe
CRBE,Bebidas,Cross Category (Bebidas)
ENDR,Bebidas,Bebidas Energeticas
FLBE,Bebidas,Bebidas Saborizadas Sin Gas
GCOF,Bebidas,Cafe Tostado y Molido
HJUI,Bebidas,Jugos Caseros
ITEA,Bebidas,Te Helado
ICOF,Bebidas,Cafe Instantaneo-Cafe Sucedaneo
JUNE,Bebidas,Jugos y Nectares
VEJU,Bebidas,Zumos de Vegetales
WATE,Bebidas,Agua Natural
CSDW,Bebidas,Gaseosas + Aguas
MXCM,Bebidas,Mixta Cafe+Malta
MXDG,Bebidas,Mixta Dolce Gusto-Mixta Te Helado + Cafe + Modificadores
MXJM,Bebidas,Mixta Jugos y Leches
MXJS,Bebidas,Mixta Jugos Liquidos + Bebidas de Soja
MXTC,Bebidas,Mixta Te+Cafe
JUIC,Bebidas,Jugos Liquidos-Jugos Polvo
PWDJ,Bebidas,Refrescos en Polvo-Jugos - Bebidas Instantaneas En Polvo - Jugos Polvo
RFDR,Bebidas,Bebidas Refrescantes
RTDJ,Bebidas,Refrescos Liquidos-Jugos Liquidos
RTEA,Bebidas,Te Liquido - Listo para Tomar
SOYB,Bebidas,Bebidas de Soja
SPDR,Bebidas,Bebidas Isotonicas
TEAA,Bebidas,Te e Infusiones-Te-Infusion Hierbas
YERB,Bebidas,Yerba Mate
BUTT,Lacteos,Manteca
CHEE,Lacteos,Queso Fresco y para Untar
CMLK,Lacteos,Leche Condensada
CRCH,Lacteos,Queso Untable
DYOG,Lacteos,Yoghurt p-beber
EMLK,Lacteos,Leche Culinaria-Leche Evaporada
FRMM,Lacteos,Leche Fermentada
FMLK,Lacteos,Leche Liquida Saborizada-Leche Liquida Con Sabor
FRMK,Lacteos,Formulas Infantiles
LQDM,Lacteos,Leche Liquida
LLFM,Lacteos,Leche Larga Vida
MARG,Lacteos,Margarina
MCHE,Lacteos,Queso Fundido
MKCR,Lacteos,Crema de Leche
MXDI,Lacteos,Mixta Lacteos-Postre+Leches+Yogurt
MXMI,Lacteos,Mixta Leches
MXYD,Lacteos,Mixta Yoghurt+Postres
PTSS,Lacteos,Petit Suisse
PWDM,Lacteos,Leche en Polvo
SYOG,Lacteos,Yoghurt p-comer
MILK,Lacteos,Leche-Leche Liquida Blanca - Leche Liq. Natural
YOGH,Lacteos,Yoghurt
CLOT,Ropas y Calzados,Ropas
FOOT,Ropas y Calzados,Calzados
SOCK,Ropas y Calzados,Medias-Calcetines
AREP,Alimentos,Arepas
BCER,Alimentos,Cereales Infantiles
BABF,Alimentos,Nutricion Infantil-Colados y Picados
BEAN,Alimentos,Frijoles
BISC,Alimentos,Galletas
BOUI,Alimentos,Caldos-Caldos y Sazonadores
BREA,Alimentos,Pan
BRCR,Alimentos,Apanados-Empanizadores
BRDC,Alimentos,Empanados
CERE,Alimentos,Cereales-Cereales Desayuno-Avenas y Cereales
BURG,Alimentos,Hamburguesas
CCMX,Alimentos,Mezclas Listas para Tortas-Preparados Base Harina Trigo
CAKE,Alimentos,Queques-Ponques Industrializados
FISH,Alimentos,Conservas De Pescado
CFAV,Alimentos,Conservas de Frutas y Verduras
CRML,Alimentos,Dulce de Leche-Manjar
CMLC,Alimentos,Alfajores
CBAR,Alimentos,Barras de Cereal
CHCK,Alimentos,Pollo
CHOC,Alimentos,Chocolate
COCO,Alimentos,Chocolate de Taza-Achocolatados - Cocoas
COLS,Alimentos,Salsas Frias
COMP,Alimentos,Compotas
SPIC,Alimentos,Condimentos y Especias
CKCH,Alimentos,Chocolate de Mesa
COIL,Alimentos,Aceite-Aceites Comestibles
CSAU,Alimentos,Salsas Listas-Salsas Caseras Envasadas
CNML,Alimentos,Grano- Harina y Masa de Maiz
CNST,Alimentos,Fecula de Maiz
CNFL,Alimentos,Harina De Maiz
CAID,Alimentos,Ayudantes Culinarios
DESS,Alimentos,Postres Preparados
DHAM,Alimentos,Jamon Endiablado
DFNS,Alimentos,Semillas y Frutos Secos
EBRE,Alimentos,Pan de Pascua
EEGG,Alimentos,Huevos de Pascua
EGGS,Alimentos,Huevos
FLSS,Alimentos,Flash Cecinas
FLOU,Alimentos,Harinas
MEAT,Alimentos,Carne Fresca
FRDS,Alimentos,Platos Listos Congelados
FRFO,Alimentos,Alimentos Congelados
HAMS,Alimentos,Jamones
HCER,Alimentos,Cereales Calientes-Cereales Precocidos
HOTS,Alimentos,Salsas Picantes
ICEC,Alimentos,Helados
IBRE,Alimentos,Pan Industrializado
IMPO,Alimentos,Pure Instantaneo
INOO,Alimentos,Fideos Instantaneos
JAMS,Alimentos,Mermeladas
KETC,Alimentos,Ketchup
LJDR,Alimentos,Jugo de Limon Adereso
MALT,Alimentos,Maltas
SEAS,Alimentos,Adobos - Sazonadores
MAYO,Alimentos,Mayonesa
MEAT,Alimentos,Carnicos
SNAG,Alimentos,Salchichas
MLKM,Alimentos,Modificadores de Leche-Saborizadores p-leche
MXCO,Alimentos,Mixta Cereales Infantiles+Avenas
MXBS,Alimentos,Mixta Caldos + Saborizantes
MXSB,Alimentos,Mixta Caldos + Sopas
MXCH,Alimentos,Mixta Cereales + Cereales Calientes
MXCC,Alimentos,Mixta Chocolate + Manjar
MXSN,Alimentos,Galletas - snacks y mini tostadas
COBT,Alimentos,Aceites + Mantecas
COCF,Alimentos,Aceites + Conservas De Pescado
CABB,Alimentos,Ayudantes Culinarios + Bolsa de Hornear
MXEC,Alimentos,Mixta Huevos de Pascua + Chocolates
MXDP,Alimentos,Mixta Platos Listos Congelados + Pasta
MXFR,Alimentos,Mixta Platos Congelados y Listos para Comer
MXFM,Alimentos,Mixta Alimentos Congelados + Margarina
MXMC,Alimentos,Mixta Modificadores + Cocoa
MXPS,Alimentos,Mixta Pastas
MXSO,Alimentos,Mixta Sopas+Cremas+Ramen
MXSP,Alimentos,Mixta Margarina + Mayonesa + Queso Crema
MXSW,Alimentos,Mixta Azucar+Endulzantes
MUST,Alimentos,Mostaza
NDCR,Alimentos,Sustitutos de Crema
NOOD,Alimentos,Fideos
NUGG,Alimentos,Nuggets
OAFL,Alimentos,Avena en hojuelas-liquidas
OLIV,Alimentos,Aceitunas
PANC,Alimentos,Tortilla
PANE,Alimentos,Paneton
PAST,Alimentos,Pastas
PSAU,Alimentos,Salsas para Pasta
PNOU,Alimentos,Turron de mani
PORK,Alimentos,Carne Porcina
PPMX,Alimentos,Postres en Polvo-Postres para Preparar - Horneables-Gelificables
PWSM,Alimentos,Leche de Soya en Polvo
PCCE,Alimentos,Cereales Precocidos
DOUG,Alimentos,Masas Frescas-Tapas Empanadas y Tarta
PPIZ,Alimentos,Pre-Pizzas
REFR,Alimentos,Meriendas listas
RICE,Alimentos,Arroz
RBIS,Alimentos,Galletas de Arroz
RTEB,Alimentos,Frijoles Procesados
RTEM,Alimentos,Pratos Prontos - Comidas Listas
SDRE,Alimentos,Aderezos para Ensalada
SALT,Alimentos,Sal
SLTC,Alimentos,Galletas Saladas-Galletas No Dulce
SARD,Alimentos,Sardina Envasada
SAUS,Alimentos,Cecinas
SCHN,Alimentos,Milanesas
SNAC,Alimentos,Snacks
SNOO,Alimentos,Fideos Sopa
SOUP,Alimentos,Sopas-Sopas Cremas
SOYS,Alimentos,Siyau
SPAG,Alimentos,Tallarines-Spaguetti
SPCH,Alimentos,Chocolate para Untar
SUGA,Alimentos,Azucar
SWCO,Alimentos,Galletas Dulces
SWSP,Alimentos,Untables Dulces
SWEE,Alimentos,Endulzantes
TOAS,Alimentos,Torradas - Tostadas
TOMA,Alimentos,Salsas de Tomate
TUNA,Alimentos,Atun Envasado
VMLK,Alimentos,Leche Vegetal
WFLO,Alimentos,Harinas de trigo
AIRC,Cuidado del Hogar,Ambientadores-Desodorante Ambiental
BARS,Cuidado del Hogar,Jabon en Barra-Jabon de lavar
BLEA,Cuidado del Hogar,Cloro-Lavandinas-Lejias-Blanqueadores
CBLK,Cuidado del Hogar,Pastillas para Inodoro
CGLO,Cuidado del Hogar,Guantes de latex
CLSP,Cuidado del Hogar,Esponjas de Limpieza-Esponjas y panos
CLTO,Cuidado del Hogar,Utensilios de Limpieza
FILT,Cuidado del Hogar,Filtros de Cafe
CRHC,Cuidado del Hogar,Cross Category (Limpiadores Domesticos)
CRLA,Cuidado del Hogar,Cross Category (Lavanderia)
CRPA,Cuidado del Hogar,Cross Category (Productos de Papel)
DISH,Cuidado del Hogar,Lavavajillas-Lavaplatos - Lavalozas mano
DPAC,Cuidado del Hogar,Empaques domesticos-Bolsas plasticas-Plastico Adherente-Papel encerado-Papel aluminio
DRUB,Cuidado del Hogar,Destapacanerias
FBRF,Cuidado del Hogar,Perfumantes para Ropa-Perfumes para Ropa
FWAX,Cuidado del Hogar,Cera p-pisos
FDEO,Cuidado del Hogar,Desodorante para Pies
FRNP,Cuidado del Hogar,Lustramuebles
GBBG,Cuidado del Hogar,Bolsas de Basura
GCLE,Cuidado del Hogar,Limpiadores verdes
CLEA,Cuidado del Hogar,Limpiadores-Limpiadores y Desinfectantes
INSE,Cuidado del Hogar,Insecticidas-Raticidas
KITT,Cuidado del Hogar,Toallas de papel-Papel Toalla - Toallas de Cocina - Rollos Absorbentes de Papel
LAUN,Cuidado del Hogar,Detergentes para ropa
LSTA,Cuidado del Hogar,Apresto
MXBC,Cuidado del Hogar,Mixta Pastillas para Inodoro + Limpiadores
MXHC,Cuidado del Hogar,Mixta Home Care-Cloro-Limpiadores-Ceras-Ambientadores
MXCB,Cuidado del Hogar,Mixta Limpiadores + Cloro
MXLB,Cuidado del Hogar,Mixta Detergentes + Cloro
MXLD,Cuidado del Hogar,Mixta Detergentes + Lavavajillas
CRTO,Cuidado del Hogar,Panitos + Papel Higienico
NAPK,Cuidado del Hogar,Servilletas
PLWF,Cuidado del Hogar,Film plastico e papel aluminio
SCOU,Cuidado del Hogar,Esponjas de Acero
SOFT,Cuidado del Hogar,Suavizantes de Ropa
STRM,Cuidado del Hogar,Quitamanchas-Desmanchadores
TOIP,Cuidado del Hogar,Papel Higienico
WIPE,Cuidado del Hogar,Panos de Limpieza
ANLG,OTC,Analgesicos-Painkillers
FSUP,OTC,Suplementos alimentares
GMED,OTC,Gastrointestinales-Efervescentes
VITA,OTC,Vitaminas y Calcio
nan,Otros,Categoria Desconocida
BATT,Otros,Pilas-Baterias
CGAS,Otros,Combustible Gas
PFHH,Otros,Panel Financiero de Hogares
PFIN,Otros,Panel Financiero de Hogares
INKC,Otros,Cartuchos de Tintas
PETF,Otros,Alimento para Mascota-Alim.p - perro - gato
TELE,Otros,Telecomunicaciones - Convergencia
TILL,Otros,Tickets - Till Rolls
TOBA,Otros,Tabaco - Cigarrillos
ADIP,Cuidado Personal,Incontinencia de Adultos
BSHM,Cuidado Personal,Shampoo Infantil
RAZO,Cuidado Personal,Maquinas de Afeitar
BDCR,Cuidado Personal,Cremas Corporales
CWIP,Cuidado Personal,Panos Humedos
COMB,Cuidado Personal,Cremas para Peinar
COND,Cuidado Personal,Acondicionador-Balsamo
CRHY,Cuidado Personal,Cross Category (Higiene)
CRPC,Cuidado Personal,Cross Category (Personal Care)
DEOD,Cuidado Personal,Desodorantes
DIAP,Cuidado Personal,Panales-Panales Desechables
FCCR,Cuidado Personal,Cremas Faciales
FTIS,Cuidado Personal,Panuelos Faciales
FEMI,Cuidado Personal,Proteccion Femenina
FRAG,Cuidado Personal,Fragancias
HAIR,Cuidado Personal,Cuidado del Cabello-Hair Care
HRCO,Cuidado Personal,Tintes para el Cabello-Tintes - Tintura - Tintes y Coloracion para el cabello
HREM,Cuidado Personal,Depilacion
HRST,Cuidado Personal,Alisadores para el Cabello
HSTY,Cuidado Personal,Fijadores para el Cabello-Modeladores-Gel-Fijadores para el cabello
HRTR,Cuidado Personal,Tratamientos para el Cabello
LINI,Cuidado Personal,Oleo Calcareo
MAKE,Cuidado Personal,Maquillaje-Cosmeticos
MEDS,Cuidado Personal,Jabon Medicinal
CRDT,Cuidado Personal,Panitos + Panales
MXMH,Cuidado Personal,Mixta Make Up+Tinturas
MOWA,Cuidado Personal,Enjuague Bucal-Refrescante Bucal
ORAL,Cuidado Personal,Cuidado Bucal
SPAD,Cuidado Personal,Protectores Femeninos
STOW,Cuidado Personal,Toallas Femininas
SHAM,Cuidado Personal,Shampoo
SHAV,Cuidado Personal,Afeitado-Crema afeitar-Locion de afeitar-Pord. Antes del afeitado
SKCR,Cuidado Personal,Cremas Faciales y Corporales-Cremas de Belleza - Cremas Cuerp y Faciales
SUNP,Cuidado Personal,Proteccion Solar
TALC,Cuidado Personal,Talcos-Talco para pies
TAMP,Cuidado Personal,Tampones Femeninos
TOIL,Cuidado Personal,Jabon de Tocador
TOOB,Cuidado Personal,Cepillos Dentales
TOOT,Cuidado Personal,Pastas Dentales
BAGS,Material Escolar,Morrales y MAletas Escoalres
CLPC,Material Escolar,Lapices de Colores
GRPC,Material Escolar,Lapices De Grafito
MRKR,Material Escolar,Marcadores
NTBK,Material Escolar,Cuadernos
SCHS,Material Escolar,Utiles Escolares
CSTD,Diversos,Estudio de Categorias
CORP,Diversos,Corporativa
CROS,Diversos,Cross Category
CRBA,Diversos,Cross Category (Bebes)
CRBR,Diversos,Cross Category (Desayuno)-Yogurt - Cereal - Pan y Queso
CRDT,Diversos,Cross Category (Diet y Light)
CRDF,Diversos,Cross Category (Alimentos Secos)
CRFO,Diversos,Cross Category (Alimentos)
CRCU,Diversos,Cross Category (Untables + Leche Condensada)
MXEV,Diversos,Cross Category (Leche Evaporada)
CRSA,Diversos,Cross Category (Salsas)-Mayonesas-Ketchup - Salsas Frias
CRSN,Diversos,Cross Category (Snacks)
DEMO,Diversos,Demo
FLSH,Diversos,Flash
HLVW,Diversos,Holistic View
COCP,Diversos,Mezcla para cafe instantaneo y crema no lactea
CRSN,Diversos,Mezclas nutricionales y suplementos
MULT,Diversos,Consolidado-Multicategory
PCHK,Diversos,Pantry Check
STCK,Diversos,Inventario
MIHC,Diversos,Leche y Cereales Calientes-Cereales Precocidos y Leche Liquida Blanca
FLWT,Alimentos,Agua Saborizada
"""
COUNTRY_MAP = {
    "10": "LatAm",
    "54": "Argentina",
    "91": "Bolivia",
    "55": "Brasil",
    "12": "CAM",
    "56": "Chile",
    "57": "Colombia",
    "93": "Ecuador",
    "52": "Mexico",
    "51": "Peru",
    "69": "Republica Dominicana",
    "62": "Guatemala",
    "63": "El Salvador",
    "64": "Honduras",
    "65": "Nicaragua",
    "66": "Costa Rica",
    "67": "Panamá",
}
CATEGORY_MAP: dict[str, str] = {}
for _line in CATEGORIES_CSV_DATA.splitlines()[1:]:
    _parts = _line.split(',')
    if len(_parts) >= 3:
        CATEGORY_MAP[_parts[0]] = _parts[2]
CATEGORY_CODE_SET: Set[str] = frozenset(
    str(code).strip().upper()
    for code in CATEGORY_MAP
    if str(code).strip()
)
CATEGORY_CODE_ALIASES: Dict[str, str] = {
    "CROSS": "CROS",
}
METADATA_RESOLUTION_CATEGORY_CODES: Set[str] = frozenset({"MULT", "CROS"})

PPT_LAYOUT_INDEX = 1
DEFAULT_POP_COVERAGE = "100%"
EXCEL_TEMP_FILENAME = "file_temp_coverage.xlsx"
POP_COVERAGE_MAP = {
    "Argentina": "90%",
    "Bolivia": "60%",
    "Brasil": "82%",
    "Chile": "78%",
    "Colombia": "65%",
    "Ecuador": "61%",
    "Mexico": "64%",
    "Peru": "66%",
    "CAM": "74%",
    "Costa Rica": "94%",
    "El Salvador": "86%",
    "Guatemala": "69%",
    "Honduras": "65%",
    "Nicaragua": "57%",
    "Panama": "92%",
    "Republica Dominicana": "63%",
}

COL_DATA = "Data"
COL_SELL_IN = "Sell_in"
COL_SELL_OUT = "Sell_out"
COL_PENET = "Penet"
COL_COMPRA_MEDIA = "Compra_Media"
COL_COMPRA_OCA = "Compra_por_Oca"
COL_FREQ = "Freq"
COL_BUYERS = "Buyers"
COL_SELL_IN_SIM = "Sell_in_sim"
COL_ACUM_SELL_OUT = "Acum_Sell_out"
COL_ACUM_SELL_IN = "Acum_Sell_in"
COL_ANO = "Ano"
COL_TRI = "Tri"
COL_SEM = "Sem"
COL_EVO_KANTAR_YOY = "% VAR WP by Numerator"
COL_EVO_SELLIN_YOY = "% VAR Sell-in (Cliente)"

VISIBLE_SELL_IN_LABEL = "Sell-in"
VISIBLE_SELL_OUT_LABEL = "Compras de Worldpanel"
VISIBLE_EXCEL_HEADER_MAP = {
    COL_SELL_IN: VISIBLE_SELL_IN_LABEL,
    COL_SELL_OUT: VISIBLE_SELL_OUT_LABEL,
    COL_SELL_IN_SIM: "Sell-in sim",
    COL_ACUM_SELL_IN: "Acum Sell-in",
    COL_ACUM_SELL_OUT: "Acum Compras de Worldpanel",
}

COLOR_KANTAR_LINE = "#2C3E50"
COLOR_SELLIN_LINE = "#D4AC0D"
COLOR_SELLOUT_LINE = "#1F618D"
COLOR_TENDENCIA_FILL = "#EBF5FB"
COLOR_COVERAGE_BAR = "#3498DB"
COLOR_PENET_LINE = "#E74C3C"
COLOR_KANTAR_BAR_VAR = '#7F8C8D'
COLOR_SELLIN_BAR_VAR = '#F1C40F'
COLOR_KANTAR_EDGE_VAR = '#2C3E50'
COLOR_SELLIN_EDGE_VAR = '#B7950B'
COLOR_COBERTURA_BAR = '#D9D9D9'
COLOR_PENETRACION_BAR = '#FFC000'
COLOR_POS_LABEL = '#1E8449'
COLOR_NEG_LABEL = '#8B0000'
COLOR_POS_LABEL_ALT = '#27AE60'
COLOR_NEG_LABEL_ALT = '#C0392B'
COLOR_SELLIN_TREND_LINE = "#D4AC0D"
COLOR_SELLOUT_TREND_LINE = "#2C3E50"
EXCEL_TREND_INITIAL_GAP_MONTHS = 6


def visible_sell_in_label() -> str:
    return VISIBLE_SELL_IN_LABEL


def visible_sell_out_label(lang_idx: int = 2) -> str:
    if lang_idx == 3:
        return "Worldpanel Purchases"
    if lang_idx == 1:
        return "Compras do Worldpanel"
    return VISIBLE_SELL_OUT_LABEL


def visible_monthly_label(base_label: str, lang_idx: int) -> str:
    suffix = "Monthly" if lang_idx == 3 else "Mensual"
    return f"{base_label} ({suffix})"


def visible_accum_sell_out_label(lang_idx: int) -> str:
    if lang_idx == 3:
        return "MAT Worldpanel Purchases"
    if lang_idx == 1:
        return "Acum Compras do Worldpanel"
    return "Acum Compras de Worldpanel"


def visible_accum_sell_in_label(lang_idx: int) -> str:
    return "MAT Sell-in" if lang_idx == 3 else "Acum Sell-in"


def evolution_mat_axis_label(lang_idx: int) -> str:
    return "MAT Volume" if lang_idx == 3 else "Volumen MAT"


def short_visible_sell_out_axis_label(lang_idx: int) -> str:
    return "WP Purchases" if lang_idx == 3 else "Compras WP"


def format_trend_axis_tick(value: object, _position: object = None) -> str:
    """Formatea marcas sin escala cientifica usando separadores de miles."""
    try:
        number = float(value)
    except (TypeError, ValueError):
        return str(value)
    if abs(number) < 0.0005:
        number = 0.0
    if number.is_integer():
        return f"{number:,.0f}"
    decimals = 3 if abs(number) < 1 else 2
    return f"{number:,.{decimals}f}".rstrip("0").rstrip(".")


def trend_axis_magnitude_exponent(values: Iterable[object]) -> int:
    """Devuelve una escala compacta de ingenieria desde miles."""
    finite_values: List[float] = []
    for value in values:
        try:
            number = float(value)
        except (TypeError, ValueError):
            continue
        if math.isfinite(number):
            finite_values.append(abs(number))
    max_abs = max(finite_values, default=0.0)
    if max_abs < 1_000:
        return 0
    return int(math.floor(math.log10(max_abs)) // 3 * 3)


def trend_axis_magnitude_label(exponent: int, lang_idx: int) -> str:
    labels = {
        1: {
            3: "milhares", 6: "milhões", 9: "bilhões", 12: "trilhões",
            15: "quadrilhões", 18: "quintilhões", 21: "sextilhões", 24: "septilhões",
        },
        2: {
            3: "miles", 6: "millones", 9: "miles de millones", 12: "billones",
            15: "miles de billones", 18: "trillones", 21: "miles de trillones", 24: "cuatrillones",
        },
        3: {
            3: "thousands", 6: "millions", 9: "billions", 12: "trillions",
            15: "quadrillions", 18: "quintillions", 21: "sextillions", 24: "septillions",
        },
    }
    return labels.get(lang_idx, labels[2]).get(exponent, f"10^{exponent}")


def trend_axis_magnitude_abbreviation(exponent: int) -> str:
    abbreviations = {
        3: "K",
        6: "M",
        9: "B",
        12: "T",
        15: "Q",
        18: "Qi",
        21: "Sx",
        24: "Sp",
    }
    return abbreviations.get(exponent, f"E{exponent}" if exponent else "")


def trend_axis_scale_text(exponent: int, lang_idx: int) -> str:
    if not exponent:
        return ""
    abbreviation = trend_axis_magnitude_abbreviation(exponent)
    return f"{abbreviation} ({trend_axis_magnitude_label(exponent, lang_idx)})"


def trend_axis_unit_text(exponent: int, lang_idx: int) -> str:
    if not exponent:
        return ""
    magnitude = trend_axis_magnitude_label(exponent, lang_idx)
    abbreviation = trend_axis_magnitude_abbreviation(exponent)
    return f"{abbreviation} = {magnitude}"


def trend_axis_title(metric_label: str, exponent: int, lang_idx: int) -> str:
    unit_text = trend_axis_unit_text(exponent, lang_idx)
    return f"{metric_label} ({unit_text})" if unit_text else metric_label


def build_trend_axis_formatter(_lang_idx: int, exponent: int = 0):
    """Acorta todas las magnitudes con su letra y las explica en el titulo."""
    divisor = float(10 ** exponent) if exponent else 1.0
    suffix = trend_axis_magnitude_abbreviation(exponent)

    def _format(value: object, position: object = None) -> str:
        try:
            scaled_value = float(value) / divisor
        except (TypeError, ValueError):
            return str(value)
        return f"{format_trend_axis_tick(scaled_value, position)}{suffix}"

    return mtick.FuncFormatter(_format)


def apply_trend_grid_style(axis: object, granularity: str) -> None:
    """Jerarquiza el grid horizontal y las guias mensuales de tendencia."""
    axis.set_axisbelow(True)
    axis.grid(axis="y", linestyle="--", linewidth=0.8, alpha=0.55)
    axis.grid(axis="x", visible=False)
    if normalize_trend_granularity(granularity) == "monthly":
        axis.grid(
            axis="x",
            visible=True,
            color="#AEB4BA",
            linestyle="-",
            linewidth=0.45,
            alpha=0.12,
        )

def _load_heavy_modules() -> None:
    """Carga en segundo plano las bibliotecas pesadas y datos estaticos."""
    try:
        global pd, np, dfi, plt, warnings, matplotlib, dt, timedelta, pearsonr
        global Presentation, Inches, get_column_letter, tqdm, mtick, MonthLocator
        global DateFormatter, matplotlib_style, Progress, BarColumn, TextColumn
        global TimeElapsedColumn, TimeRemainingColumn, SpinnerColumn, Image, ImageOps
        global RGBColor, Pt, MSO_SHAPE, MSO_VERTICAL_ANCHOR, pais, pop_coverage, OxmlElement, qn

        import dataframe_image as dfi
        import pandas as pd
        import numpy as np
        import warnings
        import matplotlib

        matplotlib.use("Agg")
        from matplotlib import pyplot as plt
        from datetime import datetime as dt, timedelta
        from scipy.stats import pearsonr
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_VERTICAL_ANCHOR
        from pptx.dml.color import RGBColor
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.oxml.ns import qn
        from openpyxl.utils import get_column_letter
        from openpyxl import load_workbook
        from openpyxl.formatting.rule import ColorScaleRule
        from tqdm import tqdm
        import matplotlib.ticker as mtick
        from matplotlib.dates import MonthLocator, DateFormatter
        import matplotlib.style as matplotlib_style
        from rich.progress import (
            Progress,
            BarColumn,
            TextColumn,
            TimeElapsedColumn,
            TimeRemainingColumn,
            SpinnerColumn,
        )
        from PIL import Image, ImageOps

        # This option is deprecated in pandas 3.x and removed in pandas 4.x.
        # Keep it only for older versions where it still applies.
        try:
            pandas_major = int(str(pd.__version__).split('.', 1)[0])
        except (TypeError, ValueError):
            pandas_major = 0
        if pandas_major < 3:
            pd.set_option('future.no_silent_downcasting', True)
        pd.set_option('mode.chained_assignment', None)
        warnings.filterwarnings('ignore')

        _codes = sorted((int(k), v) for k, v in COUNTRY_MAP.items())
        pais = pd.DataFrame({"cod": [c for c, _ in _codes], "pais": [v for _, v in _codes]})

        pop_coverage = dict(POP_COVERAGE_MAP)
    finally:
        LOADER_READY.set()


LOADER_READY = threading.Event()

def wait_for_heavy_modules() -> None:
    """Bloquea hasta que los módulos pesados hayan terminado de cargarse."""
    if not LOADER_READY.is_set():
        _loader_thread.join()

_loader_thread = threading.Thread(target=_load_heavy_modules)
_loader_thread.start()

SELECTIONS: Dict[str, str] = {}
ROUND_COVERAGE = False

def _normalize_lookup_text(value: object) -> str:
    text = str(value or "").strip().lower()
    normalized = unicodedata.normalize("NFKD", text)
    text = "".join(ch for ch in normalized if not unicodedata.combining(ch))
    return re.sub(r"\s+", " ", text).strip()

def _parse_percent_value(raw_value: object, fallback: float) -> float:
    try:
        return float(str(raw_value).replace("%", "").strip())
    except Exception:
        return fallback

DEFAULT_POP_COVERAGE_PERCENT = _parse_percent_value(DEFAULT_POP_COVERAGE, 100.0)
POP_COVERAGE_PERCENT_BY_COUNTRY_NORM: Dict[str, float] = {
    _normalize_lookup_text(country_name): _parse_percent_value(raw_percent, DEFAULT_POP_COVERAGE_PERCENT)
    for country_name, raw_percent in POP_COVERAGE_MAP.items()
}

def get_population_coverage_percent(country_name: str) -> float:
    target = _normalize_lookup_text(country_name)
    return POP_COVERAGE_PERCENT_BY_COUNTRY_NORM.get(target, DEFAULT_POP_COVERAGE_PERCENT)

def quick_file_metadata(filename: str) -> str:
    """Obtiene metadatos básicos del nombre de archivo."""
    base = os.path.splitext(filename)[0]
    parts = base.split('_')
    if len(parts) < 2:
        return ""
    country = COUNTRY_MAP.get(parts[0], "Desconocido")
    category_code = _normalize_category_code(parts[1])
    category = CATEGORY_MAP.get(category_code, "Categoria desconocida")
    return f"{country} - {category}"


def _normalize_category_code(category_code: object) -> str:
    normalized = str(category_code or "").strip().upper()
    return CATEGORY_CODE_ALIASES.get(normalized, normalized)


def _requires_metadata_category_resolution(category_code: object) -> bool:
    return _normalize_category_code(category_code) in METADATA_RESOLUTION_CATEGORY_CODES


def parse_input_filename_parts(excel_file_name: str) -> Tuple[str, str, str]:
    """Extrae codigo de pais, codigo de categoria y fabricante del nombre del archivo.

    Acepta un descriptor opcional entre categoria y fabricante, por ejemplo:
    52_CRCU_Untables_Nestle.xlsx.
    """
    parts = os.path.splitext(excel_file_name)[0].split('_')
    if len(parts) < 3:
        raise ValueError("El nombre de archivo no contiene suficientes partes (pais_categoria_fabricante)")
    return parts[0], _normalize_category_code(parts[1]), parts[-1]


def extract_input_filename_descriptor(excel_file_name: str) -> str:
    """Devuelve el descriptor opcional entre categoria y fabricante."""
    parts = os.path.splitext(excel_file_name)[0].split('_')
    if len(parts) <= 3:
        return ""
    return " ".join(part.strip() for part in parts[2:-1] if part.strip())


def sanitize_output_name_segment(value: object) -> str:
    """Limpia caracteres no validos para segmentos de nombres de archivo en Windows."""
    text = str(value or "").strip()
    text = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', " ", text)
    text = re.sub(r"\s+", " ", text).strip(" .")
    return text


def build_output_category_segment(categoria_nombre_corto: object, descriptor: object = "") -> str:
    base = sanitize_output_name_segment(categoria_nombre_corto)
    qualifier = sanitize_output_name_segment(descriptor)
    if qualifier and qualifier.casefold() != base.casefold():
        return f"{base} - {qualifier}"
    return base


def build_bounded_output_filename(directory: str, desired_filename: str, *, max_path_len: int = 240) -> str:
    """Recorta el nombre de archivo si la ruta completa se acerca al limite de Windows."""
    safe_filename = sanitize_output_name_segment(os.path.splitext(desired_filename)[0])
    ext = os.path.splitext(desired_filename)[1]
    filename = f"{safe_filename}{ext}"
    directory_abs = os.path.abspath(directory)
    full_path = os.path.join(directory_abs, filename)
    if len(full_path) <= max_path_len:
        return filename

    digest = uuid.uuid5(uuid.NAMESPACE_URL, filename).hex[:8]
    suffix = f"_{digest}{ext}"
    max_filename_len = max_path_len - len(directory_abs) - 1
    if max_filename_len <= len(suffix) + 8:
        compact_root = "output"
    else:
        compact_root = safe_filename[: max_filename_len - len(suffix)].rstrip(" .-_")
        if not compact_root:
            compact_root = "output"
    return f"{compact_root}{suffix}"


def build_category_short_name(categoria_nombre: object) -> str:
    """Genera una version corta de la categoria tomando el texto previo al primer guion."""
    try:
        dash_split = re.split(r"\s*[-‑–—−‒]\s*", str(categoria_nombre), maxsplit=1)
        categoria_corta = dash_split[0].strip() if dash_split else str(categoria_nombre).strip()
        if not categoria_corta:
            categoria_corta = str(categoria_nombre).strip()
        return categoria_corta
    except Exception:
        return str(categoria_nombre).strip()


@dataclass(frozen=True)
class SheetBankMetadata:
    pais_nombre: str
    cesta_nombre: str
    categoria_nombre: str
    categoria_nombre_corto: str
    categoria_codigo: str = ""


def _normalize_metadata_match_text(value: object) -> str:
    normalized = _normalize_lookup_text(value)
    return re.sub(r"[^a-z0-9]+", " ", normalized).strip()


@dataclass(frozen=True)
class MultSheetMetadataHints:
    clue_texts: Tuple[str, ...] = ()
    semantic_segments: Tuple[str, ...] = ()
    exact_category_codes: Tuple[str, ...] = ()
    opaque_tokens: Tuple[str, ...] = ()


@dataclass(frozen=True)
class MultCategoryResolution:
    category_code: Optional[str]
    source: str
    confidence: int


@dataclass
class SheetLoadCacheEntry:
    df_sheet: Optional["pd.DataFrame"]
    measure: Optional[str]
    metadata_source: Optional["pd.DataFrame"] = None
    metadata_hints: Optional[MultSheetMetadataHints] = None


@dataclass(frozen=True)
class PreparedMultSemanticSegment:
    raw_text: str
    normalized_text: str
    tokens: Tuple[str, ...]
    is_specific_phrase: bool


MULT_METADATA_TOKEN_STOPWORDS: Set[str] = {
    "a",
    "an",
    "and",
    "brand",
    "brands",
    "by",
    "categoria",
    "categorias",
    "category",
    "fabricante",
    "fabricantes",
    "for",
    "mercado",
    "mercado1",
    "t",
    "table",
    "total",
    "weighted",
    "with",
}
MULT_METADATA_OPAQUE_TOKENS_IGNORE: Set[str] = {"cg", "kw", "mc", "pnc", "ul"}
MULT_METADATA_OPAQUE_TOKEN_CATEGORY_MAP: Dict[str, str] = {
    # Siglas realmente opacas donde el archivo no despliega la categoria.
    "dw": "DISH",
}
_MULT_CATEGORY_PROFILE_CACHE: Dict[int, Dict[str, Dict[str, object]]] = {}
_MULT_CATEGORY_TOKEN_FREQ_CACHE: Dict[int, Dict[str, int]] = {}
_MULT_CATEGORY_TOKEN_INDEX_CACHE: Dict[int, Dict[str, Tuple[str, ...]]] = {}
_MULT_SEMANTIC_RESOLUTION_CACHE: Dict[Tuple[int, Tuple[str, ...]], Optional[MultCategoryResolution]] = {}


def _normalize_metadata_token(token: str) -> str:
    token = _normalize_metadata_match_text(token)
    if len(token) > 4 and token.endswith("es"):
        token = token[:-2]
    elif len(token) > 4 and token.endswith("s"):
        token = token[:-1]
    return token


def _tokenize_metadata_segment(value: object) -> Set[str]:
    normalized = _normalize_metadata_match_text(value)
    if not normalized:
        return set()
    tokens: Set[str] = set()
    for token in re.findall(r"[a-z0-9]+", normalized):
        if token.isdigit():
            continue
        reduced = _normalize_metadata_token(token)
        if reduced and reduced not in MULT_METADATA_TOKEN_STOPWORDS:
            tokens.add(reduced)
    return tokens


def _split_metadata_segments(value: object) -> List[str]:
    raw_text = str(value or "").strip()
    if not raw_text:
        return []

    segments: List[str] = []
    seen: Set[str] = set()
    raw_candidates = [raw_text]
    table_split = re.split(r"(?i)\s*-\s*table\s*-\s*", raw_text, maxsplit=1)
    if table_split:
        raw_candidates.append(table_split[0].strip())

    for candidate in raw_candidates:
        for segment in re.split(r"[=\\/|]+", candidate):
            clean_segment = re.sub(r"\s+", " ", str(segment or "").strip(" _-"))
            if not clean_segment:
                continue
            normalized_segment = _normalize_metadata_match_text(clean_segment)
            if not normalized_segment or normalized_segment in seen:
                continue
            segments.append(clean_segment)
            seen.add(normalized_segment)
    return segments


def _extract_exact_category_codes_from_text(value: object, valid_codes: Set[str]) -> List[str]:
    raw_text = str(value or "").strip()
    if not raw_text:
        return []

    candidates: List[str] = []
    seen: Set[str] = set()
    for match in re.finditer(r"(?i)\b([a-z]{4,5})(?=\d)", raw_text):
        candidate = _normalize_category_code(match.group(1))
        if candidate in valid_codes and candidate not in seen:
            candidates.append(candidate)
            seen.add(candidate)

    alpha_chunks = re.sub(r"\d+", " ", raw_text)
    for token in re.split(r"[^A-Za-z]+", alpha_chunks):
        candidate = _normalize_category_code(token)
        if candidate in valid_codes and candidate not in seen:
            candidates.append(candidate)
            seen.add(candidate)
    return candidates


def _extract_opaque_tokens_from_text(value: object, valid_codes: Set[str]) -> List[str]:
    raw_text = str(value or "").strip()
    if not raw_text:
        return []

    tokens: List[str] = []
    seen: Set[str] = set()
    alpha_chunks = re.sub(r"\d+", " ", raw_text)
    for token in re.split(r"[^A-Za-z]+", alpha_chunks):
        normalized_token = _normalize_metadata_match_text(token)
        if (
            2 <= len(normalized_token) <= 3
            and normalized_token not in MULT_METADATA_TOKEN_STOPWORDS
            and normalized_token not in MULT_METADATA_OPAQUE_TOKENS_IGNORE
            and normalized_token.upper() not in valid_codes
            and normalized_token not in seen
        ):
            tokens.append(normalized_token)
            seen.add(normalized_token)
    return tokens


def _find_sheet_table_row_idx(raw_sheet: "pd.DataFrame") -> Optional[int]:
    try:
        first_col = raw_sheet.iloc[:, 0].astype(str)
        table_mask = first_col.str.contains(r"\btable\b", flags=re.IGNORECASE, na=False)
        if table_mask.any():
            return int(table_mask[table_mask].index[0])
    except Exception:
        return None
    return None


def _build_sheet_metadata_source(
    raw_sheet: "pd.DataFrame",
    table_row_idx: Optional[int],
) -> Optional["pd.DataFrame"]:
    if raw_sheet is None or raw_sheet.empty:
        return None
    max_rows = min((table_row_idx + 1) if table_row_idx is not None else 3, len(raw_sheet.index))
    max_cols = min(len(raw_sheet.columns), 8)
    if max_rows <= 0 or max_cols <= 0:
        return None
    return raw_sheet.iloc[:max_rows, :max_cols].copy()


def _build_mult_category_profiles(categories_df: "pd.DataFrame") -> Dict[str, Dict[str, object]]:
    cache_key = id(categories_df)
    cached = _MULT_CATEGORY_PROFILE_CACHE.get(cache_key)
    if cached is not None:
        return cached

    profiles: Dict[str, Dict[str, object]] = {}
    token_freq: Dict[str, int] = {}
    token_index: Dict[str, Set[str]] = {}
    for category_code, row in categories_df.iterrows():
        category_code_str = _normalize_category_code(category_code)
        if not category_code_str or category_code_str in METADATA_RESOLUTION_CATEGORY_CODES:
            continue
        category_name = str(row.get("cat", "")).strip()
        short_name = build_category_short_name(category_name)
        basket_name = str(row.get("cest", "")).strip()
        profile_tokens: Set[str] = set()
        normalized_texts: List[str] = []
        for clue_text in (category_name, short_name, basket_name):
            profile_tokens.update(_tokenize_metadata_segment(clue_text))
            normalized_text = _normalize_metadata_match_text(clue_text)
            if normalized_text and normalized_text not in normalized_texts:
                normalized_texts.append(normalized_text)
        for token in profile_tokens:
            token_freq[token] = token_freq.get(token, 0) + 1
            token_index.setdefault(token, set()).add(category_code_str)
        profiles[category_code_str] = {
            "tokens": frozenset(profile_tokens),
            "texts": tuple(normalized_texts),
        }

    _MULT_CATEGORY_PROFILE_CACHE[cache_key] = profiles
    _MULT_CATEGORY_TOKEN_FREQ_CACHE[cache_key] = token_freq
    _MULT_CATEGORY_TOKEN_INDEX_CACHE[cache_key] = {
        token: tuple(sorted(category_codes))
        for token, category_codes in token_index.items()
    }
    return profiles


def _prepare_mult_semantic_segment(
    segment_text: str,
    token_frequencies: Dict[str, int],
) -> Optional[PreparedMultSemanticSegment]:
    segment_norm = _normalize_metadata_match_text(segment_text)
    segment_tokens = tuple(sorted(_tokenize_metadata_segment(segment_text)))
    if not segment_norm or not segment_tokens:
        return None
    segment_has_rare_token = any(token_frequencies.get(token, 99) <= 2 for token in segment_tokens)
    segment_is_specific_phrase = len(segment_tokens) >= 2 or segment_has_rare_token
    return PreparedMultSemanticSegment(
        raw_text=segment_text,
        normalized_text=segment_norm,
        tokens=segment_tokens,
        is_specific_phrase=segment_is_specific_phrase,
    )


def _score_mult_semantic_segment(
    segment: PreparedMultSemanticSegment,
    category_profile: Dict[str, object],
    token_frequencies: Dict[str, int],
) -> int:
    segment_norm = segment.normalized_text
    segment_tokens = segment.tokens
    if not segment_norm or not segment_tokens:
        return 0

    profile_tokens = set(category_profile.get("tokens", set()))
    profile_texts = tuple(category_profile.get("texts", ()))
    score = 0
    segment_is_specific_phrase = segment.is_specific_phrase

    if segment_norm in profile_texts and segment_is_specific_phrase:
        score += 130
    elif (
        len(segment_norm) >= 4
        and segment_is_specific_phrase
        and any(segment_norm in profile_text for profile_text in profile_texts)
    ):
        score += 100

    covered_tokens = 0
    token_strength = 0
    for segment_token in segment_tokens:
        best_token_score = 0
        for profile_token in profile_tokens:
            if segment_token == profile_token:
                freq_weight = max(1, 5 - token_frequencies.get(profile_token, 5))
                best_token_score = max(best_token_score, 3 * freq_weight)
            elif (
                len(segment_token) >= 4
                and len(profile_token) >= 4
                and (segment_token in profile_token or profile_token in segment_token)
            ):
                freq_weight = max(1, 5 - token_frequencies.get(profile_token, 5))
                best_token_score = max(best_token_score, 2 * freq_weight)
        if best_token_score:
            covered_tokens += 1
            token_strength += best_token_score

    score += token_strength * 9
    if covered_tokens == len(segment_tokens) and len(segment_tokens) >= 2:
        score += 25
    elif covered_tokens == len(segment_tokens) and len(segment_tokens) == 1:
        rare_token = next(iter(segment_tokens))
        if token_frequencies.get(rare_token, 99) <= 2:
            score += 15
    return score


def _resolve_mult_semantic_category(
    semantic_segments: Sequence[str],
    categories_df: "pd.DataFrame",
) -> Optional[MultCategoryResolution]:
    profiles = _build_mult_category_profiles(categories_df)
    cache_key = id(categories_df)
    token_frequencies = _MULT_CATEGORY_TOKEN_FREQ_CACHE.get(cache_key, {})
    token_index = _MULT_CATEGORY_TOKEN_INDEX_CACHE.get(cache_key, {})
    semantic_cache_key = (
        cache_key,
        tuple(
            dict.fromkeys(
                normalized_segment
                for normalized_segment in (
                    _normalize_metadata_match_text(segment_text)
                    for segment_text in semantic_segments
                )
                if normalized_segment
            )
        ),
    )
    if semantic_cache_key in _MULT_SEMANTIC_RESOLUTION_CACHE:
        return _MULT_SEMANTIC_RESOLUTION_CACHE[semantic_cache_key]

    prepared_segments: List[PreparedMultSemanticSegment] = []
    seen_segment_keys: Set[str] = set()
    for segment_text in semantic_segments:
        prepared_segment = _prepare_mult_semantic_segment(segment_text, token_frequencies)
        if (
            prepared_segment is None
            or prepared_segment.normalized_text in seen_segment_keys
        ):
            continue
        prepared_segments.append(prepared_segment)
        seen_segment_keys.add(prepared_segment.normalized_text)

    best_code: Optional[str] = None
    best_score = 0
    second_score = 0
    best_segment = ""

    for segment in prepared_segments:
        candidate_codes: Set[str] = set()
        for segment_token in segment.tokens:
            candidate_codes.update(token_index.get(segment_token, ()))
        if not candidate_codes and len(segment.normalized_text) >= 4:
            candidate_codes.update(
                category_code
                for category_code, profile in profiles.items()
                if any(
                    segment.normalized_text in profile_text or profile_text in segment.normalized_text
                    for profile_text in profile.get("texts", ())
                )
            )

        for category_code in candidate_codes:
            profile = profiles[category_code]
            score = _score_mult_semantic_segment(segment, profile, token_frequencies)
            if score > best_score:
                second_score = best_score
                best_score = score
                best_code = category_code
                best_segment = segment.raw_text
            elif score > second_score:
                second_score = score

    if best_code is None or best_score < 75:
        _MULT_SEMANTIC_RESOLUTION_CACHE[semantic_cache_key] = None
        return None
    if best_score == second_score and best_score < 110:
        _MULT_SEMANTIC_RESOLUTION_CACHE[semantic_cache_key] = None
        return None
    resolution = MultCategoryResolution(best_code, f"semantic:{best_segment}", best_score)
    _MULT_SEMANTIC_RESOLUTION_CACHE[semantic_cache_key] = resolution
    return resolution


def _resolve_mult_exact_category(exact_category_codes: Sequence[str]) -> Optional[MultCategoryResolution]:
    unique_codes = [code for code in dict.fromkeys(str(code).strip().upper() for code in exact_category_codes if str(code).strip())]
    if len(unique_codes) != 1:
        return None
    return MultCategoryResolution(unique_codes[0], "exact_code", 94)


def _resolve_mult_opaque_category(opaque_tokens: Sequence[str]) -> Optional[MultCategoryResolution]:
    mapped_codes = [
        MULT_METADATA_OPAQUE_TOKEN_CATEGORY_MAP[token]
        for token in dict.fromkeys(str(token).strip().lower() for token in opaque_tokens if str(token).strip())
        if token in MULT_METADATA_OPAQUE_TOKEN_CATEGORY_MAP
    ]
    unique_codes = list(dict.fromkeys(mapped_codes))
    if len(unique_codes) != 1:
        return None
    return MultCategoryResolution(unique_codes[0], "opaque_token", 96)


def _extract_sheet_metadata_hints(raw_sheet: "pd.DataFrame", sheet_name: str = "") -> MultSheetMetadataHints:
    clue_texts: List[str] = []
    semantic_segments: List[str] = []
    exact_category_codes: List[str] = []
    opaque_tokens: List[str] = []
    seen_clues: Set[str] = set()
    seen_segments: Set[str] = set()
    seen_codes: Set[str] = set()
    seen_opaque_tokens: Set[str] = set()
    valid_codes = CATEGORY_CODE_SET

    table_row_idx = _find_sheet_table_row_idx(raw_sheet)

    if table_row_idx is not None:
        max_rows = min(table_row_idx + 1, len(raw_sheet.index))
    else:
        max_rows = min(len(raw_sheet.index), 3)
    max_cols = min(len(raw_sheet.columns), 8)

    for row_idx in range(max_rows):
        row_max_cols = 1 if table_row_idx is not None and row_idx == table_row_idx else max_cols
        for col_idx in range(row_max_cols):
            cell_value = raw_sheet.iat[row_idx, col_idx]
            if pd.isna(cell_value):
                continue
            clue_text = str(cell_value).strip()
            clue_key = _normalize_metadata_match_text(clue_text)
            if not clue_key or clue_key in seen_clues:
                continue
            clue_texts.append(clue_text)
            seen_clues.add(clue_key)

            for segment_text in _split_metadata_segments(clue_text):
                segment_key = _normalize_metadata_match_text(segment_text)
                if (
                    segment_key
                    and segment_key not in seen_segments
                    and len(_tokenize_metadata_segment(segment_text)) >= 2
                ):
                    semantic_segments.append(segment_text)
                    seen_segments.add(segment_key)

            for category_code in _extract_exact_category_codes_from_text(clue_text, valid_codes):
                if category_code not in seen_codes:
                    exact_category_codes.append(category_code)
                    seen_codes.add(category_code)

            for opaque_token in _extract_opaque_tokens_from_text(clue_text, valid_codes):
                if opaque_token not in seen_opaque_tokens:
                    opaque_tokens.append(opaque_token)
                    seen_opaque_tokens.add(opaque_token)

    if sheet_name:
        extra_segments = [sheet_name, _clean_brand_name_from_sheet(sheet_name)]
        sheet_subcategory = extract_sheet_subcategory(sheet_name)
        if sheet_subcategory:
            extra_segments.append(sheet_subcategory)
        for extra_segment in extra_segments:
            for segment_text in _split_metadata_segments(extra_segment):
                segment_key = _normalize_metadata_match_text(segment_text)
                if segment_key and segment_key not in seen_segments:
                    semantic_segments.append(segment_text)
                    seen_segments.add(segment_key)
            for opaque_token in _extract_opaque_tokens_from_text(extra_segment, valid_codes):
                if opaque_token not in seen_opaque_tokens:
                    opaque_tokens.append(opaque_token)
                    seen_opaque_tokens.add(opaque_token)

    return MultSheetMetadataHints(
        clue_texts=tuple(clue_texts),
        semantic_segments=tuple(semantic_segments),
        exact_category_codes=tuple(exact_category_codes),
        opaque_tokens=tuple(opaque_tokens),
    )


SUBCATEGORY_CATALOG_TEXT = """
Absorvente de Olores
Aceite
Acondicionador
Aderezo de Mayonesa
Adulto
Afeitadora
Afeitadora Desechable
Ampollas Bebibles
Antiadherente
Anticaspa
Arroz con Leche
Avellana
Baño
Barquillo - Cono - Canasta
Barra
Base
Bata - Camison
Batido
Bebe
Bebidas con Fruta o Jugo
Bebidas Saborizadas Artificiales
Bebidas Saborizadas Envasadas
Bebidas Saborizadas Naturales
Bloqueador
Blusa
Botas
Brasieres - Corpiños
Brillo Labial
Bronceador
Cacahuate
Calzon - Pants
Calzon Menstrual
Calzonillos - Pantaleta
Camisa de Vestir
Camisa Sport
Camiseta
Camiseta Interior
Capsulas
Capsulas Blandas
Caramelos
Carne de Res y Puerco
Cartucho - Repuesto
Cepillo Electrico
Cepillo Manual
Ceras
Cereales
Chaleco
Chamarra
Chicles
Chocolate con Malvavisco
Chocolate Galleta
Cloro
Cobertura Confitado - Grajeado
Combo
Complemento Nutricional Polvo
Comprimidos - Pastillas
Con Alcohol
Con Gas
Con Jugo
Con Relleno
Concentrado
Condensada
Congelada
Conjunto Deportivo
Copa Menstrual
Corbata
Corrector Facial
Crema
Crema de Peinar
Crema Gel
Cubos
Cuerpo
De Sanitario
Delineador de Cejas
Delineador de Ojos
Delineador Labial
Desengrasantes
Deshidratada
Desinfectante
Desmanchadores
Desodorante - Aromatizante
Destilado de Agave
Donas
Dulce
Dulce Cremoso
Dulces
Elote
En Vaso
Ensamble
Envasados
Esmalte
Espuma
Faja
Falda Larga - Corta
Flan
Fondo
Formulada Polvo
Fresca - Refrigerada
Frutas
Gato
Gato Humedo
Gato Seco
Gel
Gelatina
Gomitas
Gomitas - Masticables
Gotas
Grano - Tostado Molido
Helado
Hoja de Afeitar
Infantil
Insecticida
Instantaneo - Soluble
Intensificador de Perfume
Jalea
Jarabe
Jarabe de Maiz
Jeans - Pantalon Mezclilla
Jegging
Juego de Brasier y Pantaleta
Jugo de Frutas Envasado
Jugo de Frutas Naturales
Jugos Puros
Labial
Lacteo Fermentado
Laminas
Leche Pasteurizada
Leche Ultra Pasteurizada UHT
Legging
Licuado
Ligeramente Gasificada
Light
Limpiamuebles
Limpiavidrios
Liquido
Liquido - Gel
Liquido - Jarabe
Listo para Beber
Maiz Pozolero
Mameluco
Manos
Mantequilla
Margarina
Mascarilla
Masticables - Gomitas
Mayonesa
Mermelada
Mezcal
Miel
Mixto
Mousse
Multiuso - Perfumado - Piso
Natilla
Natural
Nectar
Otros Postres Caseros
Otros Postres Industrializados
Paleta Helada
Paletas
Pan Artesanal Ambos
Pan Artesanal Dulce
Pan Artesanal Salado
Pan de Muerto
Pantalon Casual - Informal
Pantalon de Vestir
Pantiblusa
Pants - Legging Deportivo
Pañal
Papas
Papillas y Postres
Pastel Helado
Perro
Perro Humedo
Perro Seco
Pescado - Marisco
Pescadores Capri
Petit Suisse
Pies
Pijama
Playera tipo Golf - Polo
Playera Tradicional con Manga
Playera Tradicional sin Manga
Pollo
Polvo
Polvo - Granulado
Polvo Facial
Postre con Crema
Postre de Yogurt
Preparados
Producto Lacteo Pasteurizado
Producto Lacteo Ultra Pasteurizado UHT
Protectores de Cama
Protectores Diarios
Pulpas y Polvos
Regular
Repelente de Insectos
Rimel
Rosca de Reyes
Rubor
Saborizada
Saco Sport - Blazer
Saladas
Salado
Sandalias
Sandwich
Seca
Shampoo
Short y Bermuda de Mezclilla
Short y Bermuda de Vestir
Silicas
Sin Alcohol
Sin Gas
Solido
Sombra de Ojos
Sorpresa
Soya
Spray
Spray - Aerosol
Suavizante
Sudadera Deportiva
Sueter
Tableta Solida
Tabletas Efervescentes
Tampones
Tenis Casual
Tenis Deportivo
Tequila
Toallas - Protectores
Toallas Femeninas
Topico - Ungüento
Tops Casual - Elegante - Deportivo
Tortillas
Tostadas
Traje de Caballero
Traje Sastre con Falda
Traje Sastre con Pantalon
Tratamiento
Tratamiento de Uñas
Trusas - Boxes
Tunica
Unidades
Vegetales
Verduras
Vestido Casual
Vestido Noche - Coctel
Yoghurt Vegetal
Zapato Casual
Zapato Formal
"""

SUBCATEGORY_CATALOG: Tuple[str, ...] = tuple(
    line.strip() for line in SUBCATEGORY_CATALOG_TEXT.splitlines() if line.strip()
)
SUBCATEGORY_CANONICAL_BY_NORM: Dict[str, str] = {
    _normalize_lookup_text(subcategory): subcategory for subcategory in SUBCATEGORY_CATALOG
}


MULT_MANUFACTURER_ALIASES: Dict[str, Tuple[str, ...]] = {
    "unilever": ("unilever",),
    "colgate": ("colgate",),
}

# Libreria declarativa para casos MULT. Para ampliar soporte, agrega aqui el
# fabricante y define si categoria/pais se resuelven por seccion o por marca.
MULT_METADATA_RULES: Dict[str, Dict[str, object]] = {
    "unilever": {
        "category_source": "section",
        "brand_category_rules": (
            (r"\bfab\s*clean\b|\bfabclean\b", "LAUN"),
            (r"\bamaciantes?\b", "SOFT"),
            (r"\bpos\s+shampoo\b|\bposshampoo\b", "COND"),
            (r"\bclean\b", "CLEA"),
        ),
        "category_rules": (
            (r"\bclean\b", "CLEA"),
            (r"\bamaciantes?\b", "SOFT"),
            (r"\bfe\b", "SOFT"),
            (r"\bfab\s*clean\b|\bfabclean\b", "LAUN"),
            (r"\blaundry\b", "LAUN"),
            (r"\bmayonesa\b", "MAYO"),
            (r"\bsc\b", "TOIL"),
            (r"\bbar\b", "TOIL"),
            (r"\bliquido\b", "TOIL"),
            (r"\bdeos\b", "DEOD"),
            (r"\bhair\b", "HAIR"),
            (r"\bpos\s+shampoo\b|\bposshampoo\b", "COND"),
            (r"\bshampoo\b", "SHAM"),
            (r"\b(cond|acondicionador)\b", "COND"),
        ),
    },
    "colgate": {
        "category_source": "brand",
        "category_rules": (
            (r"\bjabon de tocador\b", "TOIL"),
            (r"\bcrema dental\b", "TOOT"),
            (r"\bsuavizante\b", "SOFT"),
            (r"\blimpiadores\b", "CLEA"),
            (r"\blavavajillas\b", "DISH"),
        ),
        "country_source": "brand",
        "country_rules": (
            (r"\bguatemala\b", "Guatemala"),
            (r"\bel salvador\b", "El Salvador"),
            (r"\bhonduras\b", "Honduras"),
            (r"\bnicaragua\b", "Nicaragua"),
            (r"\bcosta rica\b", "Costa Rica"),
            (r"\bpanama\b", "Panama"),
        ),
    },
}


def _resolve_mult_manufacturer_key(fabricante: str) -> Optional[str]:
    fabricante_norm = _normalize_metadata_match_text(fabricante)
    for manufacturer_key, aliases in MULT_MANUFACTURER_ALIASES.items():
        for alias in aliases:
            if alias in fabricante_norm:
                return manufacturer_key
    return None


def _lookup_category_metadata(category_code: str, categories_df: "pd.DataFrame") -> Tuple[str, str, str]:
    category_code = _normalize_category_code(category_code)
    if category_code not in categories_df.index:
        raise ValueError(f"El codigo de categoria '{category_code}' no esta en el catalogo")
    cesta_nombre = str(categories_df.loc[category_code, 'cest']).strip()
    categoria_nombre = str(categories_df.loc[category_code, 'cat']).strip()
    categoria_corta = build_category_short_name(categoria_nombre)
    return cesta_nombre, categoria_nombre, categoria_corta


def _match_metadata_rule(source_value: object, rules: Sequence[Tuple[str, str]]) -> Optional[str]:
    normalized_source = _normalize_metadata_match_text(source_value)
    if not normalized_source:
        return None
    for pattern, resolved_value in rules:
        if re.search(pattern, normalized_source):
            return resolved_value
    return None


def _resolve_rule_based_mult_category(
    manufacturer_key: Optional[str],
    marca_nombre_limpio: str,
    section_title: Optional[str],
) -> Optional[MultCategoryResolution]:
    if not manufacturer_key:
        return None

    rule_config = MULT_METADATA_RULES.get(manufacturer_key, {})
    brand_category_override = _match_metadata_rule(
        marca_nombre_limpio,
        rule_config.get("brand_category_rules", ()),
    )
    if brand_category_override:
        return MultCategoryResolution(brand_category_override, "manufacturer_brand_rule", 300)

    category_sources: List[object] = []
    if rule_config.get("category_source") == "section":
        category_sources.extend([section_title, marca_nombre_limpio])
    else:
        category_sources.extend([marca_nombre_limpio, section_title])

    for source_value in category_sources:
        category_override_code = _match_metadata_rule(source_value, rule_config.get("category_rules", ()))
        if category_override_code:
            return MultCategoryResolution(category_override_code, "manufacturer_rule", 70)
    return None


def resolve_sheet_bank_metadata(
    category_code: str,
    fabricante: str,
    marca_nombre_limpio: str,
    section_title: Optional[str],
    categories_df: "pd.DataFrame",
    default_pais_nombre: str,
    default_cesta_nombre: str,
    default_categoria_nombre: str,
    default_categoria_nombre_corto: str,
    sheet_metadata_hints: Optional[MultSheetMetadataHints] = None,
    inherited_category_code: Optional[str] = None,
) -> SheetBankMetadata:
    """Resuelve metadata del banco a nivel hoja para escenarios MULT/CROSS."""
    normalized_category_code = _normalize_category_code(category_code)
    metadata = SheetBankMetadata(
        pais_nombre=default_pais_nombre,
        cesta_nombre=default_cesta_nombre,
        categoria_nombre=default_categoria_nombre,
        categoria_nombre_corto=default_categoria_nombre_corto,
        categoria_codigo=normalized_category_code,
    )
    if not _requires_metadata_category_resolution(normalized_category_code):
        return metadata

    manufacturer_key = _resolve_mult_manufacturer_key(fabricante)
    sheet_metadata_hints = sheet_metadata_hints or MultSheetMetadataHints()
    semantic_segments: List[str] = list(dict.fromkeys(sheet_metadata_hints.semantic_segments))
    if section_title:
        semantic_segments.append(section_title)
    if marca_nombre_limpio:
        semantic_segments.append(marca_nombre_limpio)
    semantic_segments = list(dict.fromkeys(segment for segment in semantic_segments if str(segment or "").strip()))

    exact_resolution = _resolve_mult_exact_category(sheet_metadata_hints.exact_category_codes)
    selected_resolution = _resolve_mult_opaque_category(sheet_metadata_hints.opaque_tokens)
    rule_resolution = _resolve_rule_based_mult_category(manufacturer_key, marca_nombre_limpio, section_title)

    if selected_resolution is None:
        resolution_candidates: List[Tuple[int, int, MultCategoryResolution]] = []
        if semantic_segments:
            semantic_resolution = _resolve_mult_semantic_category(semantic_segments, categories_df)
            if semantic_resolution:
                resolution_candidates.append((semantic_resolution.confidence, 4, semantic_resolution))
        if exact_resolution:
            resolution_candidates.append((exact_resolution.confidence, 3, exact_resolution))
        normalized_inherited_category_code = _normalize_category_code(inherited_category_code)
        if normalized_inherited_category_code and not _requires_metadata_category_resolution(normalized_inherited_category_code):
            inherited_resolution = MultCategoryResolution(normalized_inherited_category_code, "inherited_group", 80)
            resolution_candidates.append((inherited_resolution.confidence, 2, inherited_resolution))
        if rule_resolution:
            resolution_candidates.append((rule_resolution.confidence, 1, rule_resolution))
        if resolution_candidates:
            _, _, selected_resolution = max(resolution_candidates, key=lambda item: (item[0], item[1]))

    if selected_resolution and selected_resolution.category_code:
        cesta_nombre, categoria_nombre, categoria_nombre_corto = _lookup_category_metadata(selected_resolution.category_code, categories_df)
        metadata = SheetBankMetadata(
            pais_nombre=metadata.pais_nombre,
            cesta_nombre=cesta_nombre,
            categoria_nombre=categoria_nombre,
            categoria_nombre_corto=categoria_nombre_corto,
            categoria_codigo=selected_resolution.category_code,
        )

    rule_config = MULT_METADATA_RULES.get(manufacturer_key, {}) if manufacturer_key else {}
    country_sources: List[object] = []
    # Si no hay regla explicita de pais, se conserva el pais derivado del archivo.
    if rule_config.get("country_source") == "section":
        country_sources.extend([section_title, marca_nombre_limpio])
    else:
        country_sources.extend([marca_nombre_limpio, section_title])
    for source_value in country_sources:
        country_override = _match_metadata_rule(source_value, rule_config.get("country_rules", ()))
        if country_override:
            metadata = SheetBankMetadata(
                pais_nombre=country_override,
                cesta_nombre=metadata.cesta_nombre,
                categoria_nombre=metadata.categoria_nombre,
                categoria_nombre_corto=metadata.categoria_nombre_corto,
                categoria_codigo=metadata.categoria_codigo,
            )
            break

    return metadata


def extract_sheet_subcategory(sheet_name: str) -> str:
    """Obtiene la subcategoria desde el ultimo texto entre parentesis del nombre de hoja."""
    cleaned = _clean_brand_name_from_sheet(sheet_name)
    matches = re.findall(r"\(([^()]+)\)", cleaned)
    if not matches:
        return ""
    raw_subcategory = str(matches[-1]).strip()
    if not raw_subcategory:
        return ""
    return SUBCATEGORY_CANONICAL_BY_NORM.get(_normalize_lookup_text(raw_subcategory), raw_subcategory)

# --- Datos Estaticos cargados en _load_heavy_modules

# --- Función para cargar categorías ---
def load_categories():
    """Carga el catálogo de categorías desde el string embebido."""
    try:
        categories_file = io.StringIO(CATEGORIES_CSV_DATA)
        df = pd.read_csv(categories_file, dtype={'cod': str}).set_index('cod')
        profile_rows = []
        for category_code, row in df.iterrows():
            profile = resolve_category_pipeline_profile(
                category_code,
                row.get("cat", ""),
                row.get("cest", ""),
            )
            profile_rows.append(
                {
                    "pipeline_profile": profile.name,
                    "min_pipeline": profile.min_pipeline,
                    "max_pipeline": profile.max_pipeline,
                    "short_pipeline_bias": profile.short_pipeline_bias,
                    "ultra_fast": profile.ultra_fast,
                }
            )
        if profile_rows:
            profile_df = pd.DataFrame(profile_rows, index=df.index)
            df = pd.concat([df, profile_df], axis=1)
        if os.environ.get('SHOW_CAT_MSG', '1') == '1' and df.index.duplicated().any():
            duplicates = df.index[df.index.duplicated()].unique().tolist()
            print(
                f"{Fore.YELLOW}Advertencia: Se encontraron códigos de categoría duplicados en los datos embebidos: {duplicates}. Se usará la última entrada encontrada para cada código."
            )
        if os.environ.get('SHOW_CAT_MSG', '1') == '1':
            print(Fore.GREEN + "Datos de categorías cargados correctamente desde el script.")
        return df
    except Exception as e:
        print(f"{Fore.RED}{Style.BRIGHT}Error Crítico al cargar datos de categorías desde el string embebido: {e}")
        exit()

# --- Variables Globales y Funciones de Utilidad ---
# --- Variables Globales y Funciones de Utilidad ---
# Nota: SELECTIONS y ROUND_COVERAGE se declaran al inicio del modulo.

def _round_half_up_series(series):
    """Redondea una serie numérica al entero más cercano con umbral .5 (ROUND_HALF_UP).
    Devuelve float con .0 para mantener NaN compatibles.
    """
    # Requiere numpy/pandas cargados; se usa después de esperar la carga pesada
    arr = series.to_numpy(dtype=float)
    # Usar isfinite para evitar afectar NaN/inf
    mask = np.isfinite(arr)
    arr[mask] = np.floor(arr[mask] + 0.5)
    return pd.Series(arr, index=series.index, name=series.name)

def round_coverage_flag():
    """Pregunta/lee si se debe redondear la cobertura (sin decimales, .5 hacia arriba)."""
    env_val = os.environ.get("AUTO_ROUND_COV")
    if env_val is not None:
        env_val_norm = str(env_val).strip().lower()
        do_round = env_val_norm in {"1", "true", "yes", "y", "si", "sí"}
    else:
        print(Fore.CYAN + "\n¿Desea redondear la cobertura (sin decimales, umbral .5)?")
        print(Fore.WHITE + "1 - Sí")
        print(Fore.WHITE + "2 - No")
        opciones = {'1': True, '2': False}
        eleccion = input(Fore.GREEN + "Elija 1 o 2: ")
        do_round = opciones.get(eleccion, False)
    SELECTIONS['Redondeo Cobertura'] = 'Sí' if do_round else 'No'
    clear_and_print_summary()
    return do_round

def _normalize_month_token(token: str) -> str:
    token = (token or "").strip().lower()
    if not token:
        return ""
    normalized = unicodedata.normalize("NFKD", token)
    return "".join(ch for ch in normalized if not unicodedata.combining(ch))

def parse_summary_extra_months(raw_value: Optional[str]) -> List[int]:
    """Convierte una entrada como '8,ago,12' en meses [8, 12]."""
    if raw_value is None:
        return []
    tokens = [t for t in re.split(r"[,\s;/|]+", str(raw_value).strip()) if t]
    if not tokens:
        return []
    months: List[int] = []
    invalid: List[str] = []
    for token in tokens:
        normalized = _normalize_month_token(token)
        if normalized.isdigit():
            month_num = int(normalized)
        else:
            month_num = MONTH_TOKEN_TO_NUMBER.get(normalized, 0)
        if 1 <= month_num <= 12:
            if month_num not in months:
                months.append(month_num)
        else:
            invalid.append(token)
    if invalid:
        raise ValueError(f"Mes(es) inválido(s): {', '.join(invalid)}")
    return months

def parse_summary_extra_months_mode(raw_value: Optional[str]) -> str:
    if raw_value is None:
        return "recent"
    normalized = str(raw_value).strip().lower()
    recent_values = {"recent", "actual", "current", "solo", "ultimo", "último", "ultimo_mes", "single", "one", "1"}
    both_values = {"both", "ambos", "dos", "doble", "dual", "2", "two", "all"}
    if normalized in recent_values:
        return "recent"
    if normalized in both_values:
        return "both"
    raise ValueError(f"Modo de meses extra inválido: {raw_value}")

def format_summary_extra_months(months: Sequence[int]) -> str:
    if not months:
        return "Ninguno"
    return ", ".join(month_abbr[m].capitalize() for m in months if 1 <= m <= 12)

def get_summary_extra_months_from_env() -> List[int]:
    for key in SUMMARY_EXTRA_MONTHS_ENV_KEYS:
        raw = os.environ.get(key)
        if raw is None:
            continue
        try:
            return parse_summary_extra_months(raw)
        except ValueError as exc:
            print(Fore.YELLOW + f"Advertencia: {exc}. Se ignora {key}.")
            return []
    return []

def get_summary_extra_months_mode_from_env() -> Optional[str]:
    for key in SUMMARY_EXTRA_MONTHS_MODE_ENV_KEYS:
        raw = os.environ.get(key)
        if raw is None:
            continue
        try:
            return parse_summary_extra_months_mode(raw)
        except ValueError as exc:
            print(Fore.YELLOW + f"Advertencia: {exc}. Se ignora {key}.")
            return None
    return None

def summary_extra_months_option() -> List[int]:
    """Obtiene meses extra a mostrar en la tabla summary de cobertura."""
    for key in SUMMARY_EXTRA_MONTHS_ENV_KEYS:
        raw = os.environ.get(key)
        if raw is not None:
            months = get_summary_extra_months_from_env()
            SELECTIONS['Meses extra summary'] = format_summary_extra_months(months)
            clear_and_print_summary()
            return months

    print(Fore.CYAN + "\n¿Desea agregar meses extra al summary de cobertura?")
    print(Fore.WHITE + "Ingrese mes(es) (1-12 o nombre, separados por coma). Ej: 8,ago,nov")
    print(Fore.WHITE + "Presione ENTER para continuar sin meses extra.")
    while True:
        raw = input(Fore.GREEN + "Mes(es) extra: ").strip()
        if not raw:
            months = []
            break
        try:
            months = parse_summary_extra_months(raw)
            break
        except ValueError as exc:
            print(Fore.RED + str(exc) + ". Intente nuevamente.")
    SELECTIONS['Meses extra summary'] = format_summary_extra_months(months)
    clear_and_print_summary()
    return months

def summary_extra_months_mode_option(has_extra_months: bool) -> str:
    env_mode = get_summary_extra_months_mode_from_env()
    if env_mode:
        # Evitar confusión: si no hay meses extra, el modo no aplica y no se muestra.
        if has_extra_months:
            SELECTIONS['Modo meses extra summary'] = "Mes más reciente" if env_mode == "recent" else "Año actual y anterior"
            clear_and_print_summary()
        return env_mode
    if not has_extra_months:
        return "recent"
    print(Fore.CYAN + "\n¿Modo de meses extra en summary?")
    print(Fore.WHITE + "1 - Solo mes más reciente (año actual)")
    print(Fore.WHITE + "2 - Dos meses (año actual y año anterior)")
    opciones = {"1": "recent", "2": "both"}
    eleccion = input(Fore.GREEN + "Elija 1 o 2: ").strip()
    modo = opciones.get(eleccion, "recent")
    SELECTIONS['Modo meses extra summary'] = "Mes más reciente" if modo == "recent" else "Año actual y anterior"
    clear_and_print_summary()
    return modo

def clear_and_print_summary():
    """Limpia la terminal y muestra un resumen de las selecciones del usuario."""
    os.system('cls' if os.name == 'nt' else 'clear') # Compatible con Windows y Linux/Mac
    print(Fore.CYAN + Style.BRIGHT + "Resumen de opciones seleccionadas:")

    displayed: Set[str] = set()

    def _get(key: str) -> Optional[object]:
        return SELECTIONS.get(key)

    def _as_text(val: object) -> str:
        if val is None:
            return "-"
        txt = str(val).strip()
        if not txt:
            return "-"
        # Evitar mojibake en consolas Windows (codepages) para el resumen.
        try:
            txt = unicodedata.normalize("NFKD", txt).encode("ascii", "ignore").decode("ascii")
        except Exception:
            pass
        return txt if txt else "-"

    def _line(label: str, key: str, value: object) -> None:
        displayed.add(key)
        print(Fore.BLUE + f"{label}: " + Fore.YELLOW + _as_text(value))

    # --- Archivo / contexto ---
    if _get("Excel") is not None:
        _line("Archivo Excel", "Excel", _get("Excel"))
    if _get("Pais") is not None:
        _line("Pais (detectado)", "Pais", _get("Pais"))
    elif _get("Excel") is not None:
        # El pais se infiere del nombre del archivo al momento de procesarlo.
        _line("Pais (detectado)", "Pais", "Pendiente (se detecta al procesar)")

    # --- Cobertura ---
    cov = _as_text(_get("Cobertura"))
    if cov != "-":
        cov_disp = cov
        scenario_key = normalize_scenario_key(cov)
        if scenario_key == SCENARIO_AUTO:
            cov_disp = "AUTO (usa configuracion predeterminada)"
        elif scenario_key == SCENARIO_AUTO_DUAL_AXIS:
            cov_disp = "AUTO (usa configuracion doble eje)"
        elif scenario_key == SCENARIO_AUTO_OPTIMAL_PIPELINE:
            cov_disp = "AUTO (pipeline optimo por correlacion)"
        elif scenario_key == SCENARIO_PG_GLOBAL_EN:
            cov_disp = "P&G - Global - Ingles"
        elif scenario_key == SCENARIO_NATURA_BR:
            cov_disp = "Natura - Br"
        _line("Tipo de cobertura", "Cobertura", cov_disp)
    if _get("Razón") is not None:
        _line("Razon de cobertura", "Razón", _get("Razón"))
    if _get("Redondeo Cobertura") is not None:
        round_val = str(_get("Redondeo Cobertura")).strip().lower()
        round_disp = "Si (sin decimales)" if round_val in {"si", "sí", "yes", "y", "true", "1"} else "No (1 decimal)"
        _line("Redondeo de cobertura", "Redondeo Cobertura", round_disp)

    # --- Slides ---
    if _get("Slide Cobertura") is not None:
        slide_mode = str(_get("Slide Cobertura")).strip().lower()
        if "complement" in slide_mode:
            slide_disp = "Complementado (Penetracion MAT + Cobertura puntual + Estabilidad)"
        elif "p&g" in slide_mode or slide_mode == "pg":
            slide_disp = "P&G (grafico + tablas editables inferiores)"
        else:
            slide_disp = "Clasico (tabla VAR % MAT)"
        _line("Slide de cobertura", "Slide Cobertura", slide_disp)
    if _get("Slide Evolucion") is not None:
        evo_mode = str(_get("Slide Evolucion")).strip().lower()
        evo_disp = "Simple (lineas de variacion)" if "simple" in evo_mode else "Clasico/avanzado (volumen MAT + barras)"
        _line("Slide evolucion mensual", "Slide Evolucion", evo_disp)
    if _get("Estilo variaciones") is not None:
        var_style = str(_get("Estilo variaciones")).strip().lower()
        var_disp = "Bonito (tarjetas)" if "bonit" in var_style else "Clasico (tabla)"
        _line("Cuadro de variaciones (tendencia)", "Estilo variaciones", var_disp)

    # --- Tendencia ---
    if _get("Eje tendencia") is not None:
        eje = str(_get("Eje tendencia")).strip().lower()
        eje_disp = "Simple (un eje)" if eje == "simple" else ("Doble (2 ejes)" if eje == "doble" else eje)
        _line("Grafico de tendencia", "Eje tendencia", eje_disp)
    if _get("Modo tendencia") is not None:
        _line("Periodicidad tendencia", "Modo tendencia", _as_text(_get("Modo tendencia")))
    if _get("Pipeline PPT") is not None:
        _line("Pipeline PPT/Summary", "Pipeline PPT", _as_text(_get("Pipeline PPT")))

    # --- Idioma ---
    # Mostrar de forma consistente y evitando depender de que el país esté disponible (a veces se define después).
    include_en = str(_get("Inglés") or "").strip().lower() in {"sí", "si", "yes", "y", "true", "1"}
    pais_norm = str(_get("Pais") or "").strip().lower()
    if include_en:
        idioma_disp = "EN (forzado)"
    elif pais_norm in {"brasil", "brazil"}:
        idioma_disp = "PT (por pais)"
    elif pais_norm:
        idioma_disp = "ES (por pais)"
    elif _get("Idioma PPT") is not None:
        # Compatibilidad con el texto legado si existiera.
        idioma_disp = _as_text(_get("Idioma PPT"))
    else:
        idioma_disp = "Auto (por pais)"
    _line("Idioma PPT", "Idioma PPT", idioma_disp)
    displayed.add("Inglés")  # se muestra en Idioma PPT (aunque no exista aún)

    # --- Summary extra months ---
    meses_extra_val = _get("Meses extra summary")
    if meses_extra_val is not None:
        _line("Meses extra (summary)", "Meses extra summary", meses_extra_val)
    modo_val = _get("Modo meses extra summary")
    if modo_val is not None:
        meses_txt = str(meses_extra_val or "").strip().lower()
        no_aplica = (not meses_txt) or (meses_txt in {"ninguno", "-"})
        modo_disp = f"{_as_text(modo_val)}{' (no aplica: no hay meses extra)' if no_aplica else ''}"
        _line("Modo meses extra (summary)", "Modo meses extra summary", modo_disp)

    # Mostrar cualquier otro valor no incluido para evitar "desaparecen opciones".
    remaining = [k for k in SELECTIONS.keys() if k not in displayed]
    if remaining:
        print(Fore.CYAN + "Otros:")
        for k in sorted(remaining):
            print(Fore.BLUE + f"- {k}: " + Fore.YELLOW + _as_text(SELECTIONS.get(k)))

    print("\n" + "-"*50 + "\n")

def print_file_header(idx: int, total: int, filename: str) -> None:
    """Muestra un encabezado visual para la ejecución de un archivo."""
    console.rule(f"[bold cyan]Procesando archivo {idx}/{total}: {filename}")

# --- Función para mostrar resumen de archivos generados ---
def _format_path_for_summary(path_str: str, *, base_dir: Optional[str] = None, max_len: int = 90) -> str:
    """
    Formatea rutas para mostrarlas en consola sin confundir con paths largos:
    - Preferir ruta relativa (a base_dir o al cwd) cuando sea posible.
    - Si sigue siendo muy larga, elidir el medio (mantener inicio y el final).
    """
    if not path_str:
        return ""

    norm = os.path.normpath(str(path_str))
    try:
        abs_path = os.path.abspath(norm)
    except Exception:
        abs_path = norm

    display = abs_path

    def _try_relpath(target: str, base: str) -> Optional[str]:
        try:
            rel = os.path.relpath(target, base)
            # Solo usar relpath si no se va "hacia arriba" (..)
            if not rel.startswith(".."):
                return rel
        except Exception:
            return None
        return None

    if base_dir:
        base_abs = os.path.abspath(os.path.normpath(base_dir))
        rel = _try_relpath(abs_path, base_abs)
        if rel:
            display = rel
    else:
        rel = _try_relpath(abs_path, os.getcwd())
        if rel:
            display = rel

    display = os.path.normpath(display)
    if len(display) <= max_len:
        return display

    # Elidir el medio manteniendo el final (más útil para ubicar el archivo).
    parts = display.split(os.sep)
    if len(parts) <= 2:
        return "..." + display[-(max_len - 3):]

    tail_parts = parts[-3:] if len(parts) >= 3 else parts[-2:]
    tail = os.sep.join(tail_parts)
    head = parts[0]
    candidate = head + os.sep + "..." + os.sep + tail
    if len(candidate) <= max_len:
        return candidate

    head_short = (head[:20] + "...") if len(head) > 23 else (head + "...")
    candidate = head_short + os.sep + "..." + os.sep + tail
    if len(candidate) <= max_len:
        return candidate

    # Fallback: solo el final
    candidate = "..." + os.sep + tail
    if len(candidate) <= max_len:
        return candidate
    return "..." + tail[-(max_len - 3):]

def _format_elapsed(seconds: float) -> str:
    try:
        total = int(round(float(seconds)))
    except Exception:
        return "-"
    if total < 0:
        total = 0
    h = total // 3600
    m = (total % 3600) // 60
    s = total % 60
    return f"{h}:{m:02d}:{s:02d}"


def print_file_locked_error(path_str: str, *, elapsed_seconds: Optional[float] = None) -> None:
    """Muestra un panel rojo cuando no se puede reescribir un archivo por estar en uso (Windows/Excel/PPT abierto)."""
    path_disp = str(path_str or "").strip() or "-"
    try:
        base = os.path.basename(path_disp) if path_disp not in {"-", ""} else "-"
    except Exception:
        base = path_disp

    hora_actual = datetime.now().strftime("%I:%M:%S %p")
    elapsed_line = ""
    if elapsed_seconds is not None:
        elapsed_line = f"\n[white]Tiempo total: [bold]{_format_elapsed(elapsed_seconds)}[/bold][/white]"

    msg = (
        "[bold bright_white]Proceso terminado con error[/bold bright_white]\n\n"
        f"[bold red]Archivo en uso:[/bold red] [white]{base}[/white]\n"
        "[red]No se pudo reescribir porque esta abierto o bloqueado.[/red]\n"
        "[yellow]Cierra el archivo y vuelve a ejecutar.[/yellow]\n\n"
        f"[yellow]Hora de finalizacion:[/yellow] [white]{hora_actual}[/white]"
        f"{elapsed_line}\n\n"
        f"[grey70]{path_disp}[/grey70]"
    )
    console.print()
    console.print(Panel.fit(msg, border_style="red", title="[bold red]Coverages Latam[/bold red]"))
    console.print()


class FileSaveCancelled(Exception):
    """Indica que el usuario cancelo un reintento de guardado."""


def _resolve_elapsed_seconds(
    elapsed_seconds: Optional[float] = None,
    elapsed_seconds_fn: Optional[Callable[[], Optional[float]]] = None,
) -> Optional[float]:
    """Resuelve el tiempo transcurrido usando un valor fijo o una funcion callback."""
    if elapsed_seconds_fn is not None:
        try:
            dynamic_value = elapsed_seconds_fn()
        except Exception:
            dynamic_value = None
        if dynamic_value is not None:
            return dynamic_value
    return elapsed_seconds


def prompt_file_locked_retry(
    path_str: str,
    *,
    action_label: str,
    elapsed_seconds: Optional[float] = None,
    elapsed_seconds_fn: Optional[Callable[[], Optional[float]]] = None,
) -> bool:
    """Pausa el flujo para permitir cerrar el archivo y reintentar sin recomenzar."""
    path_disp = str(path_str or "").strip() or "-"
    try:
        base = os.path.basename(path_disp) if path_disp not in {"-", ""} else "-"
    except Exception:
        base = path_disp
    hora_actual = datetime.now().strftime("%I:%M:%S %p")
    elapsed_value = _resolve_elapsed_seconds(elapsed_seconds, elapsed_seconds_fn)
    elapsed_line = ""
    if elapsed_value is not None:
        elapsed_line = f"\n[white]Tiempo acumulado: [bold]{_format_elapsed(elapsed_value)}[/bold][/white]"
    msg = (
        "[bold bright_white]Proceso en pausa por archivo en uso[/bold bright_white]\n\n"
        f"[bold red]Archivo afectado:[/bold red] [white]{base}[/white]\n"
        f"[bold yellow]Accion pendiente:[/bold yellow] [white]{action_label}[/white]\n"
        "[red]No se puede continuar mientras el archivo siga abierto o bloqueado.[/red]\n"
        "[yellow]Cierra el archivo para habilitar el reintento.[/yellow]\n\n"
        f"[yellow]Hora actual:[/yellow] [white]{hora_actual}[/white]"
        f"{elapsed_line}\n\n"
        f"[grey70]{path_disp}[/grey70]"
    )
    console.print()
    console.print(Panel.fit(msg, border_style="yellow", title="[bold yellow]Coverages Latam[/bold yellow]"))
    console.print()
    try:
        choice = input(f"{Fore.CYAN}Reintentar [Enter] / Cancelar [Q]: {Style.RESET_ALL}").strip().lower()
    except EOFError:
        return False
    return choice not in {"q", "quit", "salir", "cancelar", "n", "no"}


def run_file_write_with_retry(
    target_path: str,
    *,
    action_label: str,
    operation: Callable[[], None],
    elapsed_seconds: Optional[float] = None,
    elapsed_seconds_fn: Optional[Callable[[], Optional[float]]] = None,
) -> None:
    """Ejecuta una operacion de escritura y reintenta si el archivo esta bloqueado."""
    while True:
        try:
            operation()
            return
        except PermissionError as exc:
            locked_path = getattr(exc, "filename", None) or target_path
            should_retry = prompt_file_locked_retry(
                locked_path,
                action_label=action_label,
                elapsed_seconds=elapsed_seconds,
                elapsed_seconds_fn=elapsed_seconds_fn,
            )
            if should_retry:
                continue
            raise FileSaveCancelled(str(locked_path)) from exc


def print_reference_date_detection_warning(
    *,
    file_label: Optional[str] = None,
    sheet_names: Optional[Sequence[str]] = None,
    elapsed_seconds: Optional[float] = None,
) -> None:
    """Muestra un panel amarillo cuando no se detecta una fecha valida para el template."""
    hora_actual = datetime.now().strftime("%I:%M:%S %p")
    elapsed_line = ""
    if elapsed_seconds is not None:
        elapsed_line = f"\n[white]Tiempo total: [bold]{_format_elapsed(elapsed_seconds)}[/bold][/white]"

    file_line = ""
    if file_label:
        file_line = f"\n[white]Archivo: [bold]{file_label}[/bold][/white]"

    sheets_line = ""
    if sheet_names:
        clean_names = [str(name).strip() for name in sheet_names if str(name).strip()]
        if clean_names:
            preview = ", ".join(clean_names[:6])
            if len(clean_names) > 6:
                preview += ", ..."
            sheets_line = f"\n[white]Hojas a revisar: [bold]{preview}[/bold][/white]"

    msg = (
        "[bright_white]Proceso detenido por advertencia[/bright_white]\n\n"
        "[white]No se detectaron correctamente las fechas del archivo Excel.[/white]\n"
        "[white]No se pudo determinar el mes de referencia para nombrar el template.[/white]\n"
        "[white]Revisa la columna de fechas en las hojas procesadas y vuelve a ejecutar.[/white]"
        f"{file_line}"
        f"{sheets_line}\n\n"
        f"[white]Hora de finalizacion: [bold]{hora_actual}[/bold][/white]"
        f"{elapsed_line}"
    )
    console.print()
    console.print(Panel.fit(msg, border_style="yellow", title="Coverages Latam"))
    console.print()


def print_file_summary(
    ruta_excel: str,
    ruta_ppt: str,
    ruta_banco: str,
    ruta_pipeline_report: str = "",
    *,
    elapsed_seconds: Optional[float] = None,
) -> None:
    """Muestra un resumen con las rutas generadas para el archivo."""
    console.print("\n[blue]Resumen de archivos generados:[/blue]")

    items: List[Tuple[str, str]] = [
        ("Excel", ruta_excel),
        ("Presentación", ruta_ppt),
        ("Banco", ruta_banco),
        ("Reporte Pipelines", ruta_pipeline_report),
    ]
    present = [(label, p) for label, p in items if p]

    common_dir = ""
    if present:
        try:
            parents = [os.path.dirname(os.path.abspath(p)) for _, p in present]
            common_dir = os.path.commonpath(parents)
        except Exception:
            common_dir = ""

    if common_dir:
        console.print(f"[cyan]Carpeta:[/] [grey]{_format_path_for_summary(common_dir)}[/grey]")

    for label, p in present:
        filename = os.path.basename(p)
        parent = os.path.dirname(os.path.abspath(p))
        same_parent = False
        if common_dir:
            try:
                same_parent = os.path.normcase(parent) == os.path.normcase(os.path.abspath(common_dir))
            except Exception:
                same_parent = False

        if same_parent:
            console.print(f"[cyan]{label}:[/] [white]{filename}[/white]")
        else:
            parent_disp = _format_path_for_summary(parent, base_dir=common_dir or None)
            console.print(f"[cyan]{label}:[/] [white]{filename}[/white] [grey]({parent_disp})[/grey]")

    # Mostrar panel de proceso completado con hora actual
    hora_actual = datetime.now().strftime("%I:%M:%S %p")
    elapsed_line = ""
    if elapsed_seconds is not None:
        elapsed_line = f"\n[white]Tiempo total: [bold]{_format_elapsed(elapsed_seconds)}[/bold][/white]"
    mensaje = (
        "[bright_white]Proceso completado[/bright_white]\n\n"
        f"[white]Hora de finalizacion: [bold]{hora_actual}[/bold][/white]"
        f"{elapsed_line}"
    )
    console.print()
    console.print(Panel.fit(mensaje, border_style="cyan", title="Coverages Latam"))
    console.print()


DIRTY_MONTH_END_DATE_RE = re.compile(r"^\s*\d+\s+m/e\s+(\d{4})/(\d{2})/(\d{2})\s*$", re.IGNORECASE)


def normalize_input_date_value(value):
    """Normaliza formatos sucios de fecha antes de convertirlos con pandas."""
    if not isinstance(value, str):
        return value
    text = value.strip()
    if not text:
        return value

    if re.match(r"\w{3}-\d{2}$", text):
        try:
            parsed = datetime.strptime(text, "%b-%y")
            return f"{parsed.year:04d}-{parsed.month:02d}-01"
        except ValueError:
            return value

    dirty_month_end_match = DIRTY_MONTH_END_DATE_RE.match(text)
    if dirty_month_end_match:
        try:
            year = int(dirty_month_end_match.group(1))
            month = int(dirty_month_end_match.group(2))
            datetime(year, month, 1)
            return f"{year:04d}-{month:02d}-01"
        except ValueError:
            return value

    return value


def calc_var1(df, coluna, p):
    """
    Calcula variaciones vs período anterior (Y-1) en Python.

    Args:
        df (pd.DataFrame): DataFrame con los datos.
        coluna (str): Nombre de la columna a calcular (e.g., COL_SELL_OUT).
        p (int): Pipeline (shift para Sell_in).

    Returns:
        list: Lista con variaciones [Anual, Semestral, Trimestral].
              Retorna NaN para cálculos imposibles (datos insuficientes).
    """
    n_rows = len(df)
    variations = []

    # Anual (12 vs 12 meses)
    if n_rows >= 24 + p:
        current_sum = df[coluna][n_rows-12-p : n_rows-p].sum() if p != 0 else df[coluna][-12:].sum()
        previous_sum = df[coluna][n_rows-24-p : n_rows-12-p].sum() if p!= 0 else df[coluna][-24:-12].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    # Semestral (6 vs 6 meses)
    if n_rows >= 12 + p:
        current_sum = df[coluna][n_rows-6-p : n_rows-p].sum() if p != 0 else df[coluna][-6:].sum()
        previous_sum = df[coluna][n_rows-12-p : n_rows-6-p].sum() if p!= 0 else df[coluna][-12:-6].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    # Trimestral (3 vs 3 meses)
    if n_rows >= 6 + p:
        current_sum = df[coluna][n_rows-3-p : n_rows-p].sum() if p != 0 else df[coluna][-3:].sum()
        previous_sum = df[coluna][n_rows-6-p : n_rows-3-p].sum() if p!= 0 else df[coluna][-6:-3].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    return variations


def calc_var_same_period_last_year(df, coluna, p):
    """
    Calcula variaciones vs mismo período del año anterior.

    Args:
        df (pd.DataFrame): DataFrame con los datos.
        coluna (str): Nombre de la columna a calcular (e.g., COL_SELL_OUT).
        p (int): Pipeline (shift para Sell_in).

    Returns:
        list: Lista con variaciones [Anual, Semestral, Trimestral].
              Retorna NaN para cálculos imposibles (datos insuficientes).
    """
    n_rows = len(df)
    variations = []

    # Anual (12 meses actuales vs 12 meses del año pasado)
    if n_rows >= 24 + p:
        current_sum = df[coluna][n_rows-12-p : n_rows-p].sum() if p != 0 else df[coluna][-12:].sum()
        previous_sum = df[coluna][n_rows-24-p : n_rows-12-p].sum() if p != 0 else df[coluna][-24:-12].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    # Semestral (6 meses actuales vs 6 meses del año pasado)
    if n_rows >= 18 + p:
        current_sum = df[coluna][n_rows-6-p : n_rows-p].sum() if p != 0 else df[coluna][-6:].sum()
        previous_sum = df[coluna][n_rows-18-p : n_rows-12-p].sum() if p != 0 else df[coluna][-18:-12].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    # Trimestral (3 meses actuales vs 3 meses del año pasado)
    if n_rows >= 15 + p:
        current_sum = df[coluna][n_rows-3-p : n_rows-p].sum() if p != 0 else df[coluna][-3:].sum()
        previous_sum = df[coluna][n_rows-15-p : n_rows-12-p].sum() if p != 0 else df[coluna][-15:-12].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    return variations


def calc_var2(df, coluna, p):
    """
    Calcula variaciones vs período retrasado (Y-2) en Python.

    Args:
        df (pd.DataFrame): DataFrame con los datos.
        coluna (str): Nombre de la columna a calcular (e.g., COL_SELL_OUT).
        p (int): Pipeline (shift para Sell_in).

    Returns:
        list: Lista con variaciones [Anual, Semestral, Trimestral].
              Retorna NaN para cálculos imposibles (datos insuficientes).
    """
    n_rows = len(df)
    variations = []

    # Anual (12 meses actuales vs 12 meses de hace 2 años)
    if n_rows >= 36 + p:
        current_sum = df[coluna][n_rows-12-p : n_rows-p].sum() if p != 0 else df[coluna][-12:].sum()
        previous_sum = df[coluna][n_rows-36-p : n_rows-24-p].sum() if p!= 0 else df[coluna][-36:-24].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    # Semestral (6 meses actuales vs 6 meses de hace 2 años) - CORREGIDO
    if n_rows >= 30 + p: # Necesitamos 6 actuales + 24 para ir 2 años atrás
        current_sum = df[coluna][n_rows-6-p : n_rows-p].sum() if p != 0 else df[coluna][-6:].sum()
        previous_sum = df[coluna][n_rows-30-p : n_rows-24-p].sum() if p!= 0 else df[coluna][-30:-24].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    # Trimestral (3 meses actuales vs 3 meses de hace 2 años) - CORREGIDO
    if n_rows >= 27 + p: # Necesitamos 3 actuales + 24 para ir 2 años atrás
        current_sum = df[coluna][n_rows-3-p : n_rows-p].sum() if p != 0 else df[coluna][-3:].sum()
        previous_sum = df[coluna][n_rows-27-p : n_rows-24-p].sum() if p!= 0 else df[coluna][-27:-24].sum()
        variations.append((current_sum / previous_sum) - 1 if previous_sum else np.nan)
    else:
        variations.append(np.nan)

    return variations


def escalona(df_to_scale):
    """
    Desplaza los valores de cada columna hacia abajo, rellenando con NaN al principio.
    Se utiliza para alinear datos en fórmulas de Excel para cálculos de cobertura.

    Args:
        df_to_scale (pd.DataFrame): DataFrame cuyas columnas serán escalonadas.
    """
    for col in df_to_scale.columns:
        col_idx = df_to_scale.columns.get_loc(col)
        values = list(df_to_scale[col].values)
        # Invierte, trunca desde el inicio según índice, rellena, invierte de nuevo
        scaled_values = (values[::-1][col_idx:] + [np.nan]*col_idx)[::-1]
        df_to_scale[col] = scaled_values

def razao_cov():
    """Devuelve la razón de cobertura elegida o obtenida de las variables de entorno."""
    if os.environ.get("AUTO_RAZON"):
        razon_seleccionada = os.environ["AUTO_RAZON"]
    else:
        print(Fore.CYAN + "\nPregunta: ¿Cuál es la razón de la cobertura?")
        print(Fore.WHITE + "Opciones:")
        print(Fore.WHITE + "1 - Actualización periódica por contrato")
        print(Fore.WHITE + "2 - Conocer nivel de cobertura o pipeline")
        print(Fore.WHITE + "3 - Tendencias Contrarias")
        print(Fore.WHITE + "4 - Renovación de contrato")
        print(Fore.WHITE + "5 - Otras")

        razones = {
            '1': "Actualización periódica por contrato",
            '2': "Conocer nivel de cobertura o pipeline",
            '3': "Tendencias Contrarias",
            '4': "Renovación de contrato",
            '5': "Otras"
        }
        eleccion = input(Fore.GREEN + "Elija el número de la opción (1-5): ")
        razon_seleccionada = razones.get(eleccion, "Otras")  # Default a 'Otras'
    SELECTIONS['Razón'] = razon_seleccionada
    clear_and_print_summary()
    return razon_seleccionada

def tipo_cobertura():
    """Obtiene el tipo de cobertura interactivo o desde las variables de entorno."""
    if os.environ.get("AUTO_COV_TYPE"):
        tipo_seleccionado = os.environ["AUTO_COV_TYPE"]
    else:
        print(Fore.CYAN + "\nPregunta: ¿Qué tipo de cobertura se calculará?")
        print(Fore.WHITE + "Opciones:")
        print(Fore.WHITE + "1 - Cobertura Absoluta (Personalizable)")
        print(Fore.WHITE + "2 - Cobertura Relativa (Personalizable)")
        print(Fore.GREEN + "3 - Template AUTO (usar configuración predeterminada)")
        print(Fore.GREEN + "4 - Template AUTO (usar configuración doble eje)")
        print(
            Fore.WHITE
            + "5 - Template "
            + ansi_truecolor("P&G", SCENARIO_PG_COLOR)
            + Fore.WHITE
            + " - Global - Ingles"
        )
        print(
            Fore.WHITE
            + "6 - Template "
            + ansi_truecolor("Natura", SCENARIO_NATURA_COLOR)
            + Fore.WHITE
            + " - Br"
        )
        print(
            Fore.LIGHTCYAN_EX
            + Style.BRIGHT
            + ">> 7 - Template "
            + Fore.LIGHTMAGENTA_EX
            + "AUTOEXPERIMENTAL"
            + Fore.LIGHTCYAN_EX
            + " (pipeline recomendado por ajuste integral)"
            + Style.RESET_ALL
        )
        tipos = {
            '1': "Absoluta",
            '2': "relativa",
            '3': SCENARIO_AUTO,
            '4': SCENARIO_AUTO_DUAL_AXIS,
            '5': SCENARIO_PG_GLOBAL_EN,
            '6': SCENARIO_NATURA_BR,
            '7': SCENARIO_AUTO_OPTIMAL_PIPELINE,
        }
        eleccion = input(Fore.GREEN + "Elija 1, 2, 3, 4, 5, 6 o 7: ")
        tipo_seleccionado = tipos.get(eleccion, "Absoluta")  # Default a 'Absoluta'
    SELECTIONS['Cobertura'] = tipo_seleccionado
    clear_and_print_summary()
    return tipo_seleccionado

def tipo_eje_tendencia():
    """Elige tipo de gráfico de tendencia de forma interactiva o vía variables de entorno."""
    if os.environ.get("AUTO_EJE"):
        tipo_eje = os.environ["AUTO_EJE"]
    else:
        print(Fore.CYAN + "\n¿Desea el gráfico de tendencia con doble eje?")
        print(Fore.WHITE + "1 - Un solo eje (Sell-in y WP by Numerator juntos)")
        print(Fore.WHITE + "2 - Doble eje (WP by Numerator en eje secundario)")
        opciones = {'1': "simple", '2': "doble"}
        eleccion = input(Fore.GREEN + "Elija 1 o 2: ")
        tipo_eje = opciones.get(eleccion, "simple")
    SELECTIONS['Eje tendencia'] = tipo_eje
    clear_and_print_summary()
    return tipo_eje


def trend_granularity_option() -> str:
    """Elige la periodicidad del gráfico de tendencia."""
    raw_env = next((os.environ.get(k) for k in TREND_GRANULARITY_ENV_KEYS if os.environ.get(k) is not None), None)
    if raw_env is not None:
        granularity = normalize_trend_granularity(raw_env)
    else:
        print(Fore.CYAN + "\n¿Periodicidad del gráfico de tendencia?")
        print(Fore.WHITE + "1 - Mensual")
        print(Fore.WHITE + "2 - Trimestral (agrupado de 3 meses)")
        opciones = {
            "1": "monthly",
            "2": "quarterly",
            "mensual": "monthly",
            "trimestral": "quarterly",
            "trimestre": "quarterly",
        }
        eleccion = input(Fore.GREEN + "Elija 1 o 2: ").strip().lower()
        granularity = opciones.get(eleccion, "monthly")
    SELECTIONS["Modo tendencia"] = trend_granularity_label(granularity)
    clear_and_print_summary()
    return granularity

def variations_box_style_option() -> str:
    """Elige el estilo del cuadro de variaciones (clásico o bonito)."""
    raw_env = next((os.environ.get(k) for k in VARIATIONS_BOX_STYLE_ENV_KEYS if os.environ.get(k) is not None), None)
    if raw_env is not None:
        style = normalize_variations_box_style(raw_env)
    else:
        print(Fore.CYAN + "\n¿Estilo del cuadro de variaciones (en slide de Tendencia)?")
        print(Fore.WHITE + "1 - Clásico (tabla actual)")
        print(Fore.WHITE + "2 - Bonito (tarjetas)")
        opciones = {"1": "classic", "2": "pretty", "clasico": "classic", "clásico": "classic", "bonito": "pretty"}
        eleccion = input(Fore.GREEN + "Elija 1 o 2: ").strip().lower()
        style = opciones.get(eleccion, "classic")
    SELECTIONS["Estilo variaciones"] = "Bonito" if style == "pretty" else "Clasico"
    clear_and_print_summary()
    return style

def coverage_slide_variant_option() -> str:
    """Elige el modo del slide de Cobertura (clásico, complementado o P&G)."""
    raw_env = next((os.environ.get(k) for k in COVERAGE_SLIDE_VARIANT_ENV_KEYS if os.environ.get(k) is not None), None)
    if raw_env is not None:
        variant = normalize_coverage_slide_variant(raw_env)
    else:
        print(Fore.CYAN + "\n¿Modo del slide de Cobertura?")
        print(Fore.WHITE + "1 - Clásico (tabla VAR % MAT)")
        print(Fore.WHITE + "2 - Complementado (Penetración MAT + Cobertura puntual + Estabilidad)")
        print(Fore.WHITE + "3 - P&G (gráfico + tablas editables inferiores)")
        opciones = {
            "1": "classic",
            "2": "complemented",
            "3": "pg",
            "clasico": "classic",
            "clásico": "classic",
            "complementado": "complemented",
            "complemented": "complemented",
            "pg": "pg",
            "p&g": "pg",
        }
        eleccion = input(Fore.GREEN + "Elija 1, 2 o 3: ").strip().lower()
        variant = opciones.get(eleccion, "classic")
    SELECTIONS["Slide Cobertura"] = coverage_slide_variant_label(variant)
    clear_and_print_summary()
    return variant

def evolution_slide_variant_option() -> str:
    """Elige el modo del slide de Evolucion mensual y variacion (simple o clasico/avanzado)."""
    raw_env = next((os.environ.get(k) for k in EVOLUTION_SLIDE_VARIANT_ENV_KEYS if os.environ.get(k) is not None), None)
    if raw_env is not None:
        variant = normalize_evolution_slide_variant(raw_env)
    else:
        print(Fore.CYAN + "\n¿Modo del slide 'Evolucion mensual y variacion'?")
        print(Fore.WHITE + "1 - Simple (solo variacion: lineas, sin volumen mensual)")
        print(Fore.WHITE + "2 - Clasico/avanzado (volumen MAT + barras de variacion)")
        opciones = {
            "1": "simple",
            "2": "classic",
            "simple": "simple",
            "clasico": "classic",
            "clásico": "classic",
            "avanzado": "classic",
            "advanced": "classic",
        }
        eleccion = input(Fore.GREEN + "Elija 1 o 2: ").strip().lower()
        variant = opciones.get(eleccion, "classic")
    SELECTIONS["Slide Evolucion"] = "Simple" if variant == "simple" else "Clasico/Avanzado"
    clear_and_print_summary()
    return variant

def include_english_flag() -> bool:
    """Determina si se deben generar salidas en inglés.

    Usa AUTO_ENGLISH cuando está disponible; de lo contrario, solicita al usuario su preferencia.
    """
    env_val = os.environ.get("AUTO_ENGLISH")
    if env_val is not None:
        include_en = str(env_val).strip().lower() in {"1", "true", "yes", "y", "si", "sí"}
    else:
        print(Fore.CYAN + "\n¿Desea generar la presentación en inglés?")
        print(Fore.WHITE + "1 - Sí (usar bloque ENGLISH de la plantilla)")
        print(Fore.WHITE + "2 - No (usar idioma por país)")
        opciones = {"1": True, "2": False, "si": True, "no": False, "s": True, "n": False}
        while True:
            eleccion = input(Fore.GREEN + "Elija 1 o 2: ").strip().lower()
            if eleccion in opciones:
                include_en = opciones[eleccion]
                break
            if eleccion in {"", "\n"}:
                include_en = False
                break
            print(Fore.RED + "Entrada inválida. Intente nuevamente.")
    SELECTIONS['Inglés'] = 'Sí' if include_en else 'No'
    clear_and_print_summary()
    return include_en


def _get_excel_sheet_load_cache(excel_file_obj) -> Dict[str, SheetLoadCacheEntry]:
    cache = getattr(excel_file_obj, "_coverage_sheet_load_cache", None)
    if cache is None:
        cache = {}
        setattr(excel_file_obj, "_coverage_sheet_load_cache", cache)
    return cache


def _ensure_cached_sheet_metadata_hints(
    cache_entry: SheetLoadCacheEntry,
    sheet_name: str,
) -> None:
    if cache_entry.df_sheet is None:
        return
    if cache_entry.metadata_hints is None:
        metadata_source = cache_entry.metadata_source
        cache_entry.metadata_hints = (
            _extract_sheet_metadata_hints(metadata_source, sheet_name)
            if metadata_source is not None and not metadata_source.empty
            else MultSheetMetadataHints()
        )
        cache_entry.metadata_source = None
    cache_entry.df_sheet.attrs["sheet_metadata_hints"] = cache_entry.metadata_hints


def load_and_preprocess_sheet(excel_file_obj, sheet_name, include_metadata_hints: bool = True):
    """
    Carga una hoja del archivo Excel, la preprocesa (renombra, limpia, fechas)
    y devuelve el DataFrame procesado y la unidad de medida.

    Args:
        excel_file_obj (pd.ExcelFile): Objeto ExcelFile abierto.
        sheet_name (str): Nombre de la hoja a procesar.
        include_metadata_hints (bool): Cuando es True calcula o adjunta metadata
            auxiliar para resolver escenarios MULT.

    Returns:
        tuple: (pd.DataFrame, str) - El DataFrame procesado y la unidad de medida.
               Retorna (None, None) si hay un error al cargar o procesar.
    """
    cache = _get_excel_sheet_load_cache(excel_file_obj)
    cached_entry = cache.get(sheet_name)
    if cached_entry is not None:
        if include_metadata_hints:
            _ensure_cached_sheet_metadata_hints(cached_entry, sheet_name)
        return cached_entry.df_sheet, cached_entry.measure

    try:
        raw_sheet = excel_file_obj.parse(sheet_name, header=None)
        table_row_idx = _find_sheet_table_row_idx(raw_sheet)
        metadata_source = _build_sheet_metadata_source(raw_sheet, table_row_idx)
        sheet_metadata_hints = (
            _extract_sheet_metadata_hints(metadata_source, sheet_name)
            if include_metadata_hints and metadata_source is not None
            else None
        )

        # Detectar inicio real de la tabla buscando "table" en la primera columna
        start_idx = 0
        meta_header_text = None
        try:
            if table_row_idx is not None:
                start_idx = table_row_idx
                meta_header_text = raw_sheet.iloc[start_idx, 0]
        except Exception:
            start_idx = 0
            meta_header_text = None

        df_sheet = raw_sheet.iloc[start_idx:, :].reset_index(drop=True)

        # Validar estructura mínima esperada (al menos 2 filas, 8 columnas)
        rows, cols = df_sheet.shape
        if rows < 2 or cols < 8:
            if cols == 7:
                # Caso específico: 7 columnas – probablemente falta Sell-in del cliente
                print(
                    f"{Fore.RED}{Style.BRIGHT}Error:{Style.RESET_ALL} "
                    f"La hoja '{sheet_name}' no cumple la estructura mínima "
                    f"({rows} filas, {cols} columnas)."
                )
                print(
                    f"{Fore.LIGHTMAGENTA_EX}{Style.BRIGHT}Sugerencia:{Style.RESET_ALL} "
                    f"Probablemente falta la columna de Sell-in del cliente."
                )
                # (Opcional) Ayuda de depuración:
                # print(f"{Fore.LIGHTMAGENTA_EX}Columnas detectadas: {list(df_sheet.columns)}{Style.RESET_ALL}")
            else:
                # Otros casos (<8 columnas o <2 filas)
                print(
                    f"{Fore.RED}{Style.BRIGHT}Error:{Style.RESET_ALL} "
                    f"La hoja '{sheet_name}' tiene una estructura inesperada "
                    f"({rows} filas, {cols} columnas). Se omitirá."
                )
            cache[sheet_name] = SheetLoadCacheEntry(df_sheet=None, measure=None)
            return None, None


        # === Validación temprana adicional: abortar si la columna 8 no tiene datos ===
        # Si existen los 8 encabezados pero no hay datos debajo del encabezado de la columna 8,
        # se omite la hoja para evitar que el programa se rompa más adelante.
        try:
            _col8 = df_sheet.iloc[1:, 7]  # Índice 0-based: 7 es la 8ª columna
            _col8_empty = _col8.isna().all() or (_col8.astype(str).str.strip() == '').all()
        except Exception:
            _col8_empty = True  # si por alguna razón falla, tratamos como vacío

        if _col8_empty:
            print(f"{Fore.RED}Advertencia: La hoja '{sheet_name}' se omitirá porque la columna 8 (Sell-in) no tiene datos debajo del encabezado.")
            cache[sheet_name] = SheetLoadCacheEntry(df_sheet=None, measure=None)
            return None, None
        # === Fin validación adicional ===


        # Obtiene la 'unidad' o 'medida' de la primera fila, columna 2 (índice 1)
        measure = str(df_sheet.iat[0, 1]).replace('Weighted', '').strip()

        # Renombra las columnas al formato estándar
        df_sheet.columns = [COL_DATA, COL_SELL_OUT, COL_PENET, COL_COMPRA_MEDIA, COL_COMPRA_OCA, COL_FREQ, COL_BUYERS, COL_SELL_IN] + list(df_sheet.columns[8:])  # Mantiene columnas extra si existen
        df_sheet = df_sheet.loc[:, [COL_DATA, COL_SELL_IN, COL_SELL_OUT, COL_COMPRA_MEDIA, COL_COMPRA_OCA, COL_FREQ, COL_PENET, COL_BUYERS]]  # Reordena y selecciona

        # Elimina la primera fila (encabezados repetidos) y resetea el índice
        df_sheet = df_sheet.iloc[1:].reset_index(drop=True)

        # Convierte la columna "Data" a tipo datetime
        # Maneja posibles errores de formato o valores nulos
        original_dates = df_sheet[COL_DATA].copy()  # Guardar original por si falla
        try:
            # Normaliza formatos sucios de entrada antes de delegar el parseo general a pandas.
            normalized_dates = df_sheet[COL_DATA].apply(normalize_input_date_value)
            # Convierte el resto (o los ya convertidos) a datetime
            df_sheet[COL_DATA] = pd.to_datetime(normalized_dates, errors='coerce')
        except Exception as e:
            print(f"{Fore.YELLOW}Advertencia: Problema al convertir fechas en hoja '{sheet_name}'. Error: {e}. Se usará la columna original si es posible.")
            fallback_dates = original_dates.apply(normalize_input_date_value)
            df_sheet[COL_DATA] = pd.to_datetime(fallback_dates, errors='coerce')  # Reintentar con la original

        # Eliminar filas donde la fecha no se pudo convertir (NaT)
        initial_rows = len(df_sheet)
        df_sheet.dropna(subset=[COL_DATA], inplace=True)
        if len(df_sheet) < initial_rows:
            print(f"{Fore.YELLOW}Advertencia: Se eliminaron {initial_rows - len(df_sheet)} filas de la hoja '{sheet_name}' por fechas inválidas.")

        if df_sheet.empty:
            print(f"{Fore.RED}Advertencia: La hoja '{sheet_name}' está vacía o no contiene fechas válidas después del preprocesamiento. Se omitirá.")
            cache[sheet_name] = SheetLoadCacheEntry(df_sheet=None, measure=None)
            return None, None

        df_sheet.reset_index(drop=True, inplace=True)

        # Asegurar tipos numéricos (intentar convertir, rellenar NaN con 0 si falla)
        numeric_cols = [COL_SELL_IN, COL_SELL_OUT, COL_COMPRA_MEDIA, COL_COMPRA_OCA, COL_FREQ, COL_PENET, COL_BUYERS]
        for col in numeric_cols:
            df_sheet[col] = pd.to_numeric(df_sheet[col], errors='coerce').fillna(0)

        # Añade columnas de Año, Trimestre, Semestre
        df_sheet[COL_ANO] = df_sheet[COL_DATA].dt.year
        df_sheet[COL_TRI] = df_sheet[COL_DATA].dt.quarter
        df_sheet[COL_SEM] = (df_sheet[COL_DATA].dt.month - 1) // 6 + 1
        df_sheet[COL_DATA] = df_sheet[COL_DATA].dt.date  # Convertir a solo fecha al final
        if sheet_metadata_hints is not None:
            df_sheet.attrs["sheet_metadata_hints"] = sheet_metadata_hints

        cache[sheet_name] = SheetLoadCacheEntry(
            df_sheet=df_sheet,
            measure=measure,
            metadata_source=None if sheet_metadata_hints is not None else metadata_source,
            metadata_hints=sheet_metadata_hints,
        )

        return df_sheet, measure

    except Exception as e:
        print(f"{Fore.RED}Error crítico al cargar o preprocesar la hoja '{sheet_name}': {e}")
        cache[sheet_name] = SheetLoadCacheEntry(df_sheet=None, measure=None)
        return None, None

# --- Funciones de Generación de Gráficos ---

def generar_grafico_evolucion_mensual(
    df_graf,
    pipeline_meses: int = 0,
    lang_idx: int = 2,
    marca_nombre: Optional[str] = None,
    variant: str = "classic",
):
    """
    Genera un gráfico de evolución mensual con series MAT de WP by Numerator vs Sell-in
    y variación interanual.

    Args:
        df_graf (pd.DataFrame): DataFrame con datos mensuales (col 'Data' debe ser datetime).
        pipeline_meses (int): Número de meses de pipeline para desplazar Sell-in.
        lang_idx (int): Identificador de idioma (impacta etiquetas).
        marca_nombre (str, opcional): Nombre de la marca para mensajes de advertencia.

    Returns:
        matplotlib.figure.Figure: Figura de matplotlib con el gráfico, o None si no hay datos.
    """
    if df_graf is None or df_graf.empty or len(df_graf) < 24: # Necesita al menos 24 meses para var YOY
        print(f"{Fore.YELLOW}Advertencia: No se puede generar gráfico de evolución mensual. Datos insuficientes (se requieren >= 24 meses).")
        return None

    variant_norm = normalize_evolution_slide_variant(variant)

    # Usar contexto de estilo para evitar afectar otros gráficos
    with matplotlib.style.context('seaborn-v0_8-whitegrid'):
        df_plot = df_graf.copy()
        df_plot[COL_DATA] = pd.to_datetime(df_plot[COL_DATA]) # Asegurar datetime
        marca_label = (marca_nombre or "N/D").strip() or "N/D"
        needs_exception_warning = False

        # Detectar '-' en valores numéricos y asegurar tipo float
        for col in (COL_SELL_IN, COL_SELL_OUT):
            col_as_str = df_plot[col].astype(str).str.strip()
            dash_mask = col_as_str.eq("-")
            if dash_mask.any():
                needs_exception_warning = True
                df_plot.loc[dash_mask, col] = 0
            df_plot[col] = pd.to_numeric(df_plot[col], errors='coerce').fillna(0)

        # Si hay pipeline, desplazar Sell-in y guardar original si es necesario
        if pipeline_meses > 0:
            # df_plot["Sell_in_original"] = df_plot[COL_SELL_IN].copy() # Descomentar si se necesita el original
            df_plot[COL_SELL_IN] = df_plot[COL_SELL_IN].shift(pipeline_meses)

        # Calcular sumas móviles y variaciones interanuales
        df_plot["Kantar_12m"] = df_plot[COL_SELL_OUT].rolling(12).sum()
        df_plot["Sellin_12m"] = df_plot[COL_SELL_IN].rolling(12).sum()
        kantar_prev = df_plot["Kantar_12m"].shift(12)
        sellin_prev = df_plot["Sellin_12m"].shift(12)
        zero_prev_kantar = kantar_prev == 0
        zero_prev_sellin = sellin_prev == 0
        if zero_prev_kantar.any() or zero_prev_sellin.any():
            needs_exception_warning = True
        safe_kantar_prev = kantar_prev.where(~zero_prev_kantar, 1)
        safe_sellin_prev = sellin_prev.where(~zero_prev_sellin, 1)
        df_plot["Kantar_yoy"] = ((df_plot["Kantar_12m"] / safe_kantar_prev) - 1) * 100
        df_plot["Sellin_yoy"] = ((df_plot["Sellin_12m"] / safe_sellin_prev) - 1) * 100

        # Filtrar NaNs resultantes de rolling/shift
        df_plot = df_plot.dropna(subset=["Kantar_yoy", "Sellin_yoy"]).copy()

        if df_plot.empty:
            print(f"{Fore.YELLOW}Advertencia: No quedan datos para el gráfico de evolución después de calcular YOY.")
            return None

        if needs_exception_warning:
            notify_zero_months_exception(marca_label)

        # Crear figura y ejes con márgenes personalizados
        fig = plt.figure(figsize=(16.5, 8), dpi=100) # Ajustar tamaño si es necesario
        left_margin, right_margin, bottom_margin, top_margin = 0.08, 0.92, 0.18, 0.90
        ax1 = fig.add_axes([left_margin, bottom_margin, right_margin-left_margin, top_margin-bottom_margin])
        ax2 = None
        if variant_norm == "classic":
            ax2 = ax1.twinx()

        var_title = "Variacion Interanual (%)" if lang_idx != 3 else "Year-over-Year Change (%)"
        def _tint_color(color_str: str, mix_with_white: float = 0.78) -> str:
            """Devuelve una version mas clara del color (mezclado con blanco)."""
            try:
                from matplotlib.colors import to_rgb, to_hex
                r, g, b = to_rgb(color_str)
                m = float(mix_with_white)
                if m < 0:
                    m = 0.0
                if m > 1:
                    m = 1.0
                tinted = (r + (1 - r) * m, g + (1 - g) * m, b + (1 - b) * m)
                return to_hex(tinted)
            except Exception:
                return "#E7E6E6"

        def _font_color(col_yoy: str, value: float) -> str:
            if value > 0:
                return COLOR_POS_LABEL if col_yoy == "Kantar_yoy" else COLOR_POS_LABEL_ALT
            if value < 0:
                return COLOR_NEG_LABEL if col_yoy == "Kantar_yoy" else COLOR_NEG_LABEL_ALT
            return "#333333"
        if variant_norm == "simple":
            # Simple: solo variacion (lineas), sin volumen mensual.
            ax1.plot(
                df_plot[COL_DATA],
                df_plot["Kantar_yoy"],
                color=COLOR_KANTAR_LINE,
                marker="o",
                linewidth=2.5,
                markersize=5,
                label="% Var Worldpanel by Numerator",
            )
            ax1.plot(
                df_plot[COL_DATA],
                df_plot["Sellin_yoy"],
                color=COLOR_SELLIN_LINE,
                marker="o",
                linewidth=2.5,
                markersize=5,
                label="% Var Sell-in" + (f" - P:{pipeline_meses}" if pipeline_meses > 0 else ""),
            )
            ax1.set_ylabel(var_title, fontsize=11, labelpad=15)
            ax1.yaxis.set_major_formatter(mtick.PercentFormatter(decimals=0))
            ax1.tick_params(axis='y', labelsize=9)
            ax1.axhline(y=0, color='gray', linestyle='-', alpha=0.5, linewidth=0.8)
            ax1.grid(axis='y', linestyle='--', alpha=0.4)

            offset = 4
            for _, row in df_plot.iterrows():
                for col_yoy, x_offset in [("Kantar_yoy", -offset), ("Sellin_yoy", offset)]:
                    if pd.isna(row[col_yoy]):
                        continue
                    valor = float(row[col_yoy])
                    pos_vert = valor + 1 if valor >= 0 else valor - 1
                    va_align = "bottom" if valor >= 0 else "top"
                    line_color = COLOR_KANTAR_LINE if col_yoy == "Kantar_yoy" else COLOR_SELLIN_LINE
                    bg = _tint_color(line_color, mix_with_white=0.78)
                    ax1.text(
                        row[COL_DATA] + pd.Timedelta(days=x_offset),
                        pos_vert,
                        f"{valor:.1f}%",
                        ha="center",
                        va=va_align,
                        fontsize=7,
                        fontweight="bold",
                        color=_font_color(col_yoy, valor),
                        bbox=dict(boxstyle="round,pad=0.18", facecolor=bg, edgecolor=line_color, linewidth=0.8),
                    )

            # Ajustar limites Y para dar aire a las cajas
            y_min, y_max = ax1.get_ylim()
            pad = max(abs(y_min), abs(y_max)) * 0.18
            ax1.set_ylim(y_min - pad, y_max + pad)
        else:
            # Clasico/avanzado: volumen MAT (lineas) + variacion (barras)
            sellin_label = visible_accum_sell_in_label(lang_idx) + (
                f" - P:{pipeline_meses}" if pipeline_meses > 0 else ""
            )
            ax1.plot(
                df_plot[COL_DATA], df_plot["Kantar_12m"],
                color=COLOR_KANTAR_LINE, marker="o", linewidth=2, markersize=5,
                label=visible_accum_sell_out_label(lang_idx),
            )
            ax1.plot(
                df_plot[COL_DATA],
                df_plot["Sellin_12m"],
                color=COLOR_SELLIN_LINE,
                marker="o",
                linewidth=2,
                markersize=5,
                label=sellin_label,
            )
            ax1.set_ylabel(evolution_mat_axis_label(lang_idx), fontsize=11, labelpad=15)
            ax1.tick_params(axis='y', labelsize=9)
            ax1.set_ylim(bottom=0)
            ax1.grid(axis='y', linestyle='--', alpha=0.4)

            width = 8
            offset = 4
            assert ax2 is not None
            ax2.bar(df_plot[COL_DATA] - pd.DateOffset(days=offset), df_plot["Kantar_yoy"], width=width, color=COLOR_KANTAR_BAR_VAR, edgecolor=COLOR_KANTAR_EDGE_VAR, alpha=0.7, label="% Var Worldpanel by Numerator")
            ax2.bar(df_plot[COL_DATA] + pd.DateOffset(days=offset), df_plot["Sellin_yoy"], width=width, color=COLOR_SELLIN_BAR_VAR, edgecolor=COLOR_SELLIN_EDGE_VAR, alpha=0.7, label="% Var Sell-in")
            ax2.set_ylabel(var_title, fontsize=11, labelpad=15)
            ax2.yaxis.set_major_formatter(mtick.PercentFormatter(decimals=0))
            ax2.tick_params(axis='y', labelsize=9)
            ax2.axhline(y=0, color='gray', linestyle='-', alpha=0.5, linewidth=0.8)

            # Etiquetas en barras con cuadro de fondo segun signo
            for _, row in df_plot.iterrows():
                for col_yoy, x_offset in [("Kantar_yoy", -offset), ("Sellin_yoy", offset)]:
                    if pd.isna(row[col_yoy]):
                        continue
                    valor = float(row[col_yoy])
                    pos_vert = valor + 1 if valor >= 0 else valor - 1
                    va_align = "bottom" if valor >= 0 else "top"
                    bg = "#C6EFCE" if valor > 0 else ("#FFC7CE" if valor < 0 else "#E7E6E6")
                    ax2.text(
                        row[COL_DATA] + pd.Timedelta(days=x_offset),
                        pos_vert,
                        f"{valor:.1f}%",
                        ha="center",
                        va=va_align,
                        fontsize=7,
                        fontweight="bold",
                        color=_font_color(col_yoy, valor),
                        bbox=dict(boxstyle="round,pad=0.18", facecolor=bg, edgecolor="black", linewidth=0.6),
                    )

            y2_min, y2_max = ax2.get_ylim()
            padding = max(abs(y2_min), abs(y2_max)) * 0.15
            ax2.set_ylim(y2_min - padding, y2_max + padding*2)

        # Formato Eje X (Fechas) con extensión de un mes antes y después
        fechas_validas = pd.to_datetime(df_plot[COL_DATA], errors="coerce").dropna()
        if fechas_validas.empty:
            raise ValueError("No hay fechas validas para construir el grafico de evolucion mensual.")
        if int(fechas_validas.min().year) <= 1:
            raise ValueError(
                f"Se detectaron fechas fuera de rango en el grafico de evolucion mensual: min={fechas_validas.min()!s}"
            )
        fecha_min = fechas_validas.min() - pd.DateOffset(months=1)
        fecha_max = fechas_validas.max() + pd.DateOffset(months=1)
        ax1.set_xlim([fecha_min, fecha_max])
        ax1.xaxis.set_major_locator(MonthLocator(interval=1)) # Ajustar intervalo dinámicamente
        ax1.xaxis.set_major_formatter(DateFormatter('%b-%y'))
        ax1.tick_params(axis='x', rotation=45, labelsize=8)

        # Título y Leyenda
        # titulo = "Evolución Mensual y Variación " + (f" (Pipeline: {pipeline_meses})" if pipeline_meses > 0 else "")
        # fig.suptitle(titulo, fontsize=16, fontweight='bold', y=top_margin + 0.05) # Título de la figura
        if variant_norm == "classic":
            lines1, labels1 = ax1.get_legend_handles_labels()
            assert ax2 is not None
            lines2, labels2 = ax2.get_legend_handles_labels()
            ax2.legend(lines1 + lines2, labels1 + labels2, loc="upper left", bbox_to_anchor=(0.01, 0.98), fontsize=9, frameon=True, framealpha=0.8)
        else:
            ax1.legend(loc="upper left", bbox_to_anchor=(0.01, 0.98), fontsize=9, frameon=True, framealpha=0.8)

        # No usar tight_layout con add_axes, márgenes manuales ya aplicados
        # fig.tight_layout(rect=[0, 0, 1, 0.95]) # Ajustar rect si el título se solapa
  
        return fig

def generar_grafico_cobertura(
    slide,
    marca_clean,
    pipeline,
    df_cov_pipe,
    df_pen_pipe,
    lang_idx,
    coverage_label,
    labels_dict,
    *,
    picture_left=None,
    picture_top=None,
    picture_height=None,
):
    """Genera el gráfico de barras de Cobertura vs Penetración y lo añade al slide."""
    cov_series = df_cov_pipe if isinstance(df_cov_pipe, pd.Series) else pd.Series(df_cov_pipe)
    pen_series = df_pen_pipe if isinstance(df_pen_pipe, pd.Series) else pd.Series(df_pen_pipe)
    cov_series = cov_series.rename('coverage')
    pen_series = pen_series.rename('penetracion')
    cov_series = pd.to_numeric(cov_series, errors='coerce')
    pen_series = pd.to_numeric(pen_series, errors='coerce')
    combined = pd.concat([cov_series, pen_series], axis=1, join='inner')
    combined = combined.replace([np.inf, -np.inf], np.nan)
    combined = combined.dropna(subset=['coverage', 'penetracion'])
    if combined.empty:
        print(f"{Fore.YELLOW}Advertencia: No hay datos suficientes para el gráfico de cobertura/penetración (Marca: {marca_clean}, P:{pipeline}).")
        return
    cov_data = combined['coverage'].to_numpy(dtype=float)
    pen_data = combined['penetracion'].to_numpy(dtype=float)
    cov_data = np.where(np.isfinite(cov_data), cov_data, np.nan)
    pen_data = np.where(np.isfinite(pen_data), pen_data, np.nan)
    x_labels = [idx.strftime('%m-%y') if hasattr(idx, 'strftime') else str(idx) for idx in combined.index]
    x_pos = np.arange(len(x_labels))
    fig_cov, ax_cov = plt.subplots(figsize=(12, 4.25), dpi=100)
    bar_width = 0.35
    offset = bar_width / 2
    rects2 = ax_cov.bar(
        x_pos - offset / 1.2,
        pen_data,
        bar_width,
        label=labels_dict.get((lang_idx, 'Graf cob Penet Men'), 'Penetración Mensual'),
        color=COLOR_PENETRACION_BAR,
        edgecolor='black',
        zorder=1,
    )
    rects1 = ax_cov.bar(
        x_pos + offset,
        cov_data,
        bar_width,
        label=coverage_label,
        color=COLOR_COBERTURA_BAR,
        edgecolor='black',
        linewidth=2,
        zorder=2,
        alpha=0.85,
    )
    for rect_group in (rects2, rects1):
        for i, rect in enumerate(rect_group):
            height = rect.get_height()
            if height > 0.1:
                bbox_props = dict(facecolor='#F2F2F2', edgecolor='black', boxstyle='round,pad=0.3')
                if rect_group is rects1:
                    # Resaltar el último mes y cada 12 meses hacia atrás (evita el caso len%12==0).
                    if i % 12 == ((len(rect_group) - 1) % 12):
                        bbox_props['facecolor'] = '#A6A6A6'
                        bbox_props['edgecolor'] = 'black'
                    label_txt = f"{int(np.floor(height + 0.5))}" if globals().get('ROUND_COVERAGE', False) else f"{height:.1f}"
                else:
                    bbox_props['facecolor'] = '#FDEAD9'
                    label_txt = f"{height:.1f}"
                ax_cov.annotate(
                    label_txt,
                    xy=(rect.get_x() + rect.get_width() / 2, height),
                    xytext=(0, 3),
                    textcoords="offset points",
                    ha='center',
                    va='bottom',
                    fontsize=8,
                    bbox=bbox_props,
                )
    ax_cov.set_ylabel(
        f"{coverage_label} | {labels_dict.get((lang_idx, 'Graf cob Penet Men'), 'Penetración Mensual')}",
        fontsize=9,
    )
    title_key = 'Titulo Cob'
    default_title = 'Cobertura Año Móvil' if lang_idx != 3 else 'MOVING YEAR COVERAGE'
    ax_cov.set_title(
        f"{labels_dict.get((lang_idx, title_key), default_title)} | {marca_clean} Pipeline {pipeline}",
        size=16,
    )
    ax_cov.set_xticks(x_pos)
    ax_cov.set_xticklabels(x_labels, rotation=30, ha='right', fontsize=9)
    ax_cov.legend(
        loc='lower center',
        bbox_to_anchor=(0.5, -0.30),
        frameon=False,
        prop={'size': 11},
        ncol=2,
    )
    ax_cov.grid(axis='y', linestyle='--', alpha=0.6)
    ax_cov.set_axisbelow(True)
    ax_cov.spines['top'].set_visible(False)
    ax_cov.spines['right'].set_visible(False)
    ax_cov.spines['left'].set_visible(False)
    max_val = max(np.nanmax(cov_data) if cov_data.size else 0, np.nanmax(pen_data) if pen_data.size else 0)
    ax_cov.set_ylim(bottom=0, top=max_val * 1.15 if max_val else 1)
    ax_cov.margins(x=0)
    plt.tight_layout()
    img_stream = io.BytesIO()
    fig_cov.savefig(img_stream, format='png', bbox_inches='tight', pad_inches=0.1, transparent=True)
    img_stream.seek(0)
    img_pil = Image.open(img_stream)
    bordered = ImageOps.expand(img_pil, border=1, fill='black')
    img_stream_bordered = io.BytesIO()
    bordered.save(img_stream_bordered, format='PNG')
    img_stream_bordered.seek(0)
    slide.shapes.add_picture(
        img_stream_bordered,
        picture_left if picture_left is not None else Inches(0.5),
        picture_top if picture_top is not None else Inches(2.0),
        height=picture_height if picture_height is not None else Inches(4.2),
    )
    plt.close(fig_cov)

def generar_grafico_tendencia(
    slide,
    marca_clean,
    pipeline,
    df_plot,
    lang_idx,
    labels_dict,
    doble_eje: bool = False,
    granularity: str = "monthly",
    box_left=None,
    box_top=None,
    box_width=None,
    box_height=None,
    figsize: Tuple[float, float] = (13, 5),
    legend_y: float = -0.28,
):
    """
    Genera el gráfico de líneas de Tendencia (Sell-in vs Sell-out) y lo añade al slide.
    Si doble_eje=True, WP by Numerator (Sell-out) va en eje secundario.
    """
    if df_plot is None or df_plot.empty or pipeline >= len(df_plot):
         print(f"{Fore.YELLOW}Advertencia: Datos insuficientes para gráfico de Tendencia (Marca: {marca_clean}, P:{pipeline}).")
         return

    fig_trend, ax_trend = plt.subplots(figsize=figsize, dpi=100)
    granularity_norm = normalize_trend_granularity(granularity)

    def _trend_month_abbr_local(month: int) -> str:
        es = ["", "Ene", "Feb", "Mar", "Abr", "May", "Jun", "Jul", "Ago", "Sep", "Oct", "Nov", "Dic"]
        pt = ["", "Jan", "Fev", "Mar", "Abr", "Mai", "Jun", "Jul", "Ago", "Set", "Out", "Nov", "Dez"]
        en = ["", "Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
        table = en if lang_idx == 3 else (pt if lang_idx == 1 else es)
        if 1 <= int(month) <= 12:
            return table[int(month)]
        return "-"

    def _trend_three_month_label(start_dt: "pd.Timestamp", end_dt: "pd.Timestamp") -> str:
        if int(end_dt.month) in {3, 6, 9, 12}:
            quarter = ((int(end_dt.month) - 1) // 3) + 1
            return f"Q{quarter}-{int(end_dt.year) % 100:02d}"
        start_abbr = _trend_month_abbr_local(int(start_dt.month))
        end_abbr = _trend_month_abbr_local(int(end_dt.month))
        start_year = int(start_dt.year) % 100
        end_year = int(end_dt.year) % 100
        if int(start_dt.year) == int(end_dt.year):
            return f"{start_abbr}–{end_abbr} {end_year:02d}"
        return f"{start_abbr} {start_year:02d}–{end_abbr} {end_year:02d}"

    sell_out_series = pd.to_numeric(df_plot[COL_SELL_OUT].iloc[pipeline:], errors="coerce").reset_index(drop=True)
    sell_in_series = pd.to_numeric(df_plot[COL_SELL_IN].iloc[:len(df_plot)-pipeline], errors="coerce").reset_index(drop=True)
    period_tokens = df_plot[COL_DATA].iloc[pipeline:].reset_index(drop=True)

    if len(sell_out_series) != len(sell_in_series):
         print(f"{Fore.RED}Error: Discrepancia de longitud en datos de tendencia para {marca_clean} P:{pipeline}.")
         plt.close(fig_trend)
         return

    if granularity_norm == "quarterly":
        period_dates = pd.to_datetime(period_tokens, format="%m-%y", errors="coerce")
        aggregated_rows: List[Tuple[str, float, float]] = []
        end_idx = len(period_dates) - 1
        while end_idx >= 2:
            start_idx = end_idx - 2
            window_dates = period_dates.iloc[start_idx:end_idx + 1]
            window_sell_in = sell_in_series.iloc[start_idx:end_idx + 1]
            window_sell_out = sell_out_series.iloc[start_idx:end_idx + 1]
            if window_dates.isna().any() or window_sell_in.isna().any() or window_sell_out.isna().any():
                end_idx -= 3
                continue
            label = _trend_three_month_label(window_dates.iloc[0], window_dates.iloc[-1])
            aggregated_rows.append((label, float(window_sell_in.sum()), float(window_sell_out.sum())))
            end_idx -= 3
        aggregated_rows.reverse()
        if not aggregated_rows:
            print(f"{Fore.YELLOW}Advertencia: Datos insuficientes para gráfico de Tendencia trimestral (Marca: {marca_clean}, P:{pipeline}).")
            plt.close(fig_trend)
            return
        x_labels = [row[0] for row in aggregated_rows]
        sell_in_data = np.array([row[1] for row in aggregated_rows], dtype=float)
        sell_out_data = np.array([row[2] for row in aggregated_rows], dtype=float)
        divider_step = 4
    else:
        x_labels = period_tokens.values
        sell_in_data = sell_in_series.to_numpy(dtype=float)
        sell_out_data = sell_out_series.to_numpy(dtype=float)
        divider_step = 12

    sell_in_exponent = trend_axis_magnitude_exponent(sell_in_data)
    sell_out_exponent = trend_axis_magnitude_exponent(sell_out_data)
    sell_in_axis_label = visible_sell_in_label()
    sell_out_axis_label = short_visible_sell_out_axis_label(lang_idx)

    if doble_eje:
        ax2 = ax_trend.twinx()
        lns1 = ax_trend.plot(
            x_labels,
            sell_in_data,
            color=COLOR_SELLIN_TREND_LINE,
            linewidth=4,
            label=f'{visible_sell_in_label()} (P:{pipeline})',
        )
        lns2 = ax2.plot(
            x_labels,
            sell_out_data,
            color=COLOR_SELLOUT_TREND_LINE,
            linewidth=4,
            label=visible_sell_out_label(lang_idx),
        )
        ax_trend.set_ylabel(
            trend_axis_title(sell_in_axis_label, sell_in_exponent, lang_idx),
            color=COLOR_SELLIN_TREND_LINE,
            fontsize=10,
        )
        ax2.set_ylabel(
            trend_axis_title(sell_out_axis_label, sell_out_exponent, lang_idx),
            color=COLOR_SELLOUT_TREND_LINE,
            fontsize=10,
        )
        ax_trend.yaxis.set_major_formatter(build_trend_axis_formatter(lang_idx, sell_in_exponent))
        ax2.yaxis.set_major_formatter(build_trend_axis_formatter(lang_idx, sell_out_exponent))
        # --- CORRECCIÓN: Configurar ambos ejes para empezar desde 0 ---
        ax_trend.set_ylim(bottom=0)
        ax2.set_ylim(bottom=0)
        lns = lns1 + lns2
        labs = [l.get_label() for l in lns]
        ax2.legend(lns, labs, loc='lower center', bbox_to_anchor=(0.5, legend_y), frameon=False, prop={'size': 11}, ncol=2)
    else:
        shared_exponent = trend_axis_magnitude_exponent([*sell_in_data, *sell_out_data])
        lns1 = ax_trend.plot(
            x_labels,
            sell_in_data,
            color=COLOR_SELLIN_TREND_LINE,
            linewidth=4,
            label=f'{visible_sell_in_label()} (P:{pipeline})',
        )
        lns2 = ax_trend.plot(
            x_labels,
            sell_out_data,
            color=COLOR_SELLOUT_TREND_LINE,
            linewidth=4,
            label=visible_sell_out_label(lang_idx),
        )
        ax_trend.set_ylabel(
            trend_axis_title(
                f"{sell_in_axis_label} / {sell_out_axis_label}",
                shared_exponent,
                lang_idx,
            ),
            color='black',
            fontsize=10,
        )
        ax_trend.yaxis.set_major_formatter(build_trend_axis_formatter(lang_idx, shared_exponent))
        ax_trend.set_ylim(bottom=0)
        lns = lns1 + lns2
        labs = [l.get_label() for l in lns]
        ax_trend.legend(lns, labs, loc='lower center', bbox_to_anchor=(0.5, legend_y), frameon=False, prop={'size': 11}, ncol=2)

    # Divisores de ciclo anual: cada 12 meses hacia atrás desde el último dato.
    if len(x_labels) > divider_step:
        for idx in range(len(x_labels) - (divider_step + 1), -1, -divider_step):
            ax_trend.axvline(
                x=x_labels[idx],
                color="#B0B0B0",
                linestyle="--",
                linewidth=1.1,
                alpha=0.35,
                zorder=1,
            )

    x_tick_rotation = 30 if granularity_norm == "quarterly" else 30
    x_tick_size = 7 if granularity_norm == "quarterly" else 9
    ax_trend.tick_params(axis='x', rotation=x_tick_rotation, labelsize=x_tick_size)
    for label in ax_trend.get_xticklabels():
        if granularity_norm == "quarterly":
            label.set_ha('right')
        else:
            label.set_ha('right')
    apply_trend_grid_style(ax_trend, granularity_norm)
    ax_trend.spines['top'].set_visible(False)
    ax_trend.spines['right'].set_visible(False)
    granularity_suffix = {
        "quarterly": {1: " | Trimestral", 2: " | Trimestral", 3: " | Quarterly"},
        "monthly": {1: "", 2: "", 3: ""},
    }[granularity_norm][lang_idx]
    ax_trend.set_title(
        f"{labels_dict.get((lang_idx, 'Titulo Vol'), 'Tendencia en Volumen')}{granularity_suffix} | {marca_clean} P:{pipeline}",
        size=17,
    )

    plt.tight_layout()
    img_stream = io.BytesIO()
    fig_trend.savefig(img_stream, format='png', bbox_inches='tight', pad_inches=0.1, transparent=True)
    img_stream.seek(0)
    img_pil = Image.open(img_stream)
    # Sin contorno: el slide ya maneja el layout, y el borde negro se ve pesado.
    img_stream_bordered = io.BytesIO()
    img_pil.save(img_stream_bordered, format='PNG')
    img_stream_bordered.seek(0)
    # Ubicación por defecto (layout clásico)
    if box_left is None:
        box_left = Inches(0.5)
    if box_top is None:
        box_top = Inches(1.8)

    # Si no se pasa un "box", mantenemos el comportamiento anterior (solo altura).
    if box_width is None and box_height is None:
        slide.shapes.add_picture(img_stream_bordered, box_left, box_top, height=Inches(4.5))
        plt.close(fig_trend)
        return

    # Fit dentro del rectángulo manteniendo aspect ratio.
    if box_width is None:
        box_width = Inches(100)  # sin limite real, se acota por altura
    if box_height is None:
        box_height = Inches(100)  # sin limite real, se acota por ancho

    img_stream_bordered.seek(0)
    try:
        with Image.open(img_stream_bordered) as _img:
            px_w, px_h = _img.size
    except Exception:
        px_w, px_h = (1, 1)
    finally:
        img_stream_bordered.seek(0)

    aspect = (px_w / px_h) if px_h else 1.0
    box_w_in = float(box_width) / 914400.0
    box_h_in = float(box_height) / 914400.0
    placed_w_in = box_w_in
    placed_h_in = placed_w_in / aspect if aspect else box_h_in
    if placed_h_in > box_h_in:
        placed_h_in = box_h_in
        placed_w_in = placed_h_in * aspect

    placed_w = Inches(max(0.1, placed_w_in))
    placed_h = Inches(max(0.1, placed_h_in))
    left = box_left + int((box_width - placed_w) / 2)
    top = box_top + int((box_height - placed_h) / 2)
    slide.shapes.add_picture(img_stream_bordered, left, top, width=placed_w, height=placed_h)
    plt.close(fig_trend)
    

# --- Configuración y estructuras de alto nivel --------------------------------

def normalize_scenario_key(value: object) -> str:
    """Normaliza aliases de escenarios predefinidos del menu de cobertura."""
    normalized = normalize_brand_key(str(value or "")).replace("&", "and")
    normalized = re.sub(r"[^a-z0-9]+", "_", normalized).strip("_")
    if normalized in {"auto", "3"}:
        return SCENARIO_AUTO
    if normalized in {
        "auto_doble_eje",
        "auto_doble",
        "auto_dual_axis",
        "auto_dual",
        "auto_2_ejes",
        "auto_2_eje",
        "auto_2",
        "doble_eje",
        "dual_axis",
        "4",
    }:
        return SCENARIO_AUTO_DUAL_AXIS
    if normalized in {
        "auto_pipeline_optimo",
        "auto_pipeline_optima",
        "auto_optimal_pipeline",
        "auto_optimo",
        "auto_correlacion",
        "auto_correlation",
        "pipeline_optimo",
        "pipeline_optima",
        "correlacion",
        "correlation",
        "7",
    }:
        return SCENARIO_AUTO_OPTIMAL_PIPELINE
    if normalized in {"pg_global_en", "p_g_global_ingles", "p_g_global_english", "pg", "p_g", "5"}:
        return SCENARIO_PG_GLOBAL_EN
    if normalized in {"natura_br", "natura_brasil", "natura", "6"}:
        return SCENARIO_NATURA_BR
    return str(value or "").strip()


@dataclass
class ExecutionOptions:
    coverage_type: str
    coverage_reason: str
    trend_axis: str
    trend_granularity: str
    include_english: bool
    round_coverage: bool
    variations_box_style: str = "classic"
    coverage_slide_variant: str = "classic"
    evolution_slide_variant: str = "classic"
    summary_extra_months: List[int] = field(default_factory=list)
    summary_extra_months_mode: str = "recent"
    variations_include_same_period_last_year: bool = True
    variations_compact_period_labels: bool = False
    optimal_pipeline_mode: bool = False
    auto_mode: bool = False

    @classmethod
    def from_scenario(cls, scenario_key: str) -> Optional["ExecutionOptions"]:
        """Crea opciones predefinidas para escenarios del menu de cobertura."""
        scenario = normalize_scenario_key(scenario_key)
        if scenario == SCENARIO_AUTO:
            return cls(
                coverage_type="Absoluta",
                coverage_reason="Actualización periódica por contrato",
                trend_axis="simple",
                trend_granularity="monthly",
                variations_box_style="classic",
                include_english=False,
                round_coverage=False,
                coverage_slide_variant="classic",
                evolution_slide_variant="classic",
                summary_extra_months=[],
                summary_extra_months_mode="recent",
                variations_include_same_period_last_year=True,
                variations_compact_period_labels=False,
                auto_mode=True,
            )
        if scenario == SCENARIO_AUTO_DUAL_AXIS:
            return cls(
                coverage_type="Absoluta",
                coverage_reason="Actualización periódica por contrato",
                trend_axis="doble",
                trend_granularity="monthly",
                variations_box_style="classic",
                include_english=False,
                round_coverage=False,
                coverage_slide_variant="classic",
                evolution_slide_variant="classic",
                summary_extra_months=[],
                summary_extra_months_mode="recent",
                variations_include_same_period_last_year=True,
                variations_compact_period_labels=False,
                auto_mode=True,
            )
        if scenario == SCENARIO_AUTO_OPTIMAL_PIPELINE:
            return cls(
                coverage_type="Absoluta",
                coverage_reason="Actualización periódica por contrato",
                trend_axis="simple",
                trend_granularity="monthly",
                variations_box_style="classic",
                include_english=False,
                round_coverage=False,
                coverage_slide_variant="classic",
                evolution_slide_variant="classic",
                summary_extra_months=[],
                summary_extra_months_mode="recent",
                variations_include_same_period_last_year=True,
                variations_compact_period_labels=False,
                optimal_pipeline_mode=True,
                auto_mode=True,
            )
        if scenario == SCENARIO_PG_GLOBAL_EN:
            return cls(
                coverage_type="Absoluta",
                coverage_reason="Actualización periódica por contrato",
                trend_axis="simple",
                trend_granularity="monthly",
                variations_box_style="classic",
                include_english=True,
                round_coverage=False,
                coverage_slide_variant="pg",
                evolution_slide_variant="classic",
                summary_extra_months=[],
                summary_extra_months_mode="recent",
                variations_include_same_period_last_year=True,
                variations_compact_period_labels=False,
                auto_mode=True,
            )
        if scenario == SCENARIO_NATURA_BR:
            return cls(
                coverage_type="Absoluta",
                coverage_reason="Actualización periódica por contrato",
                trend_axis="simple",
                trend_granularity="monthly",
                variations_box_style="pretty",
                include_english=False,
                round_coverage=False,
                coverage_slide_variant="complemented",
                evolution_slide_variant="simple",
                summary_extra_months=[],
                summary_extra_months_mode="recent",
                variations_include_same_period_last_year=False,
                variations_compact_period_labels=True,
                auto_mode=True,
            )
        return None

    @classmethod
    def from_environment(cls) -> Optional["ExecutionOptions"]:
        """Crea las opciones cuando se usa la ejecución en modo automático."""
        auto_file = os.environ.get("AUTO_FILE")
        if not auto_file:
            return None
        coverage_type = os.environ.get("AUTO_COV_TYPE", "Absoluta")
        scenario_options = cls.from_scenario(coverage_type)
        if scenario_options:
            return scenario_options
        variations_box_style = normalize_variations_box_style(
            next((os.environ.get(k) for k in VARIATIONS_BOX_STYLE_ENV_KEYS if os.environ.get(k) is not None), None)
        )
        coverage_slide_variant = normalize_coverage_slide_variant(
            next((os.environ.get(k) for k in COVERAGE_SLIDE_VARIANT_ENV_KEYS if os.environ.get(k) is not None), None)
        )
        evolution_slide_variant = normalize_evolution_slide_variant(
            next((os.environ.get(k) for k in EVOLUTION_SLIDE_VARIANT_ENV_KEYS if os.environ.get(k) is not None), None)
        )
        summary_extra_months = get_summary_extra_months_from_env()
        summary_extra_months_mode = get_summary_extra_months_mode_from_env() or "recent"
        auto_mode = coverage_type.strip().lower() == "auto"
        if auto_mode:
            coverage_type = "Absoluta"
            coverage_reason = "Actualización periódica por contrato"
            trend_axis = "simple"
            trend_granularity = "monthly"
            include_english = False
            round_cov = False
        else:
            coverage_reason = os.environ.get("AUTO_RAZON", "Otras")
            trend_axis = os.environ.get("AUTO_EJE", "simple")
            trend_granularity = normalize_trend_granularity(
                next((os.environ.get(k) for k in TREND_GRANULARITY_ENV_KEYS if os.environ.get(k) is not None), None)
            )
            include_english = str(os.environ.get("AUTO_ENGLISH", "0")).strip().lower() in {"1", "true", "yes", "y", "si", "sí"}
            round_cov = str(os.environ.get("AUTO_ROUND_COV", "0")).strip().lower() in {"1", "true", "yes", "y", "si", "sí"}
        return cls(
            coverage_type=coverage_type,
            coverage_reason=coverage_reason,
            trend_axis=trend_axis,
            trend_granularity=trend_granularity,
            variations_box_style=variations_box_style,
            include_english=include_english,
            round_coverage=round_cov,
            coverage_slide_variant=coverage_slide_variant,
            evolution_slide_variant=evolution_slide_variant,
            summary_extra_months=summary_extra_months,
            summary_extra_months_mode=summary_extra_months_mode,
            variations_include_same_period_last_year=True,
            variations_compact_period_labels=False,
            auto_mode=auto_mode,
        )


def apply_execution_options_to_selections(options: ExecutionOptions) -> None:
    """Refleja opciones calculadas en el resumen visible del CLI."""
    SELECTIONS['Cobertura'] = options.coverage_type
    SELECTIONS['Razón'] = options.coverage_reason
    SELECTIONS['Eje tendencia'] = options.trend_axis
    SELECTIONS['Modo tendencia'] = trend_granularity_label(options.trend_granularity)
    SELECTIONS['Pipeline PPT'] = (
        'Recomendado por ajuste integral (P1-P6)' if options.optimal_pipeline_mode else 'Según hoja / todos'
    )
    SELECTIONS['Idioma PPT'] = 'EN (forzado)' if options.include_english else 'ES (por pais)'
    SELECTIONS['Inglés'] = 'Sí' if options.include_english else 'No'
    SELECTIONS['Redondeo Cobertura'] = 'Sí' if options.round_coverage else 'No'
    SELECTIONS["Estilo variaciones"] = (
        "Bonito" if normalize_variations_box_style(options.variations_box_style) == "pretty" else "Clasico"
    )
    SELECTIONS["Slide Cobertura"] = coverage_slide_variant_label(options.coverage_slide_variant)
    SELECTIONS["Slide Evolucion"] = (
        "Simple" if normalize_evolution_slide_variant(options.evolution_slide_variant) == "simple" else "Clasico/Avanzado"
    )
    SELECTIONS['Meses extra summary'] = format_summary_extra_months(options.summary_extra_months)
    if options.summary_extra_months:
        SELECTIONS['Modo meses extra summary'] = (
            "Mes más reciente" if options.summary_extra_months_mode == "recent" else "Año actual y anterior"
        )
    else:
        SELECTIONS.pop('Modo meses extra summary', None)


@dataclass
class PipelineAssets:
    """Recursos calculados para generar las diapositivas de un pipeline."""

    pipeline: int
    marca: str
    coverage_series: "pd.Series"
    penetration_series: "pd.Series"
    variation_table: "pd.DataFrame"
    trend_plot_df: "pd.DataFrame"
    variations_detail: Optional["pd.DataFrame"]
    evolution_figure: Optional["plt.Figure"]
    buyers_mat_actual: Optional[float] = None
    penet_mat_actual: Optional[float] = None
    penet_mat_anterior: Optional[float] = None
    annual_var_cliente_y1: Optional[float] = None
    annual_var_cliente_y2: Optional[float] = None
    annual_var_wp_y1: Optional[float] = None
    annual_var_wp_y2: Optional[float] = None
    measure_unit: Optional[str] = None
    current_year_correlation: Optional[float] = None
    trend_following: Optional[bool] = None


    summary_rows: List[Dict[str, str]] = field(default_factory=list)
    bank_rows: List[Dict[str, object]] = field(default_factory=list)
    lang_index: int = 2

    def as_summary_df(self, labels: Dict[Tuple[int, str], List[str]]) -> "pd.DataFrame":
        return pd.DataFrame(self.summary_rows, columns=labels[(self.lang_index, "Summary")])

    summary_rows: List[Dict[str, str]] = field(default_factory=list)
    bank_rows: List[Dict[str, object]] = field(default_factory=list)

    def as_summary_df(self, labels: Dict[Tuple[int, str], List[str]]) -> "pd.DataFrame":
        return pd.DataFrame(self.summary_rows, columns=labels[(self._lang_index, "Summary")])

    def configure(self, lang_index: int) -> None:
        self._lang_index = lang_index


# --- Pequeñas utilidades -------------------------------------------------------

def compute_coverage_label(coverage_type: str, include_english: bool) -> str:
    """Devuelve el texto de cobertura a mostrar en nombres de archivo y títulos."""
    ctype = coverage_type.strip().lower()
    if ctype == "auto":
        ctype = "absoluta"
    if include_english:
        return "MOVING YEAR COVERAGE" if ctype == "absoluta" else "MOVING YEAR COVERAGE RELATIVE"
    return "Cobertura Absoluta" if ctype == "absoluta" else "Cobertura Relativa"



def determine_language(include_english: bool, pais_nombre: str) -> Tuple[str, int]:
    """Determina el código de idioma y el índice numérico usado por la lógica heredada."""
    if include_english:
        return "EN", 3
    pais_norm = (pais_nombre or "").strip().lower()
    if pais_norm in {"brasil", "brazil"}:
        return "PT", 1
    return "ES", 2


def build_summary_coverage_periods(
    ref_dt: datetime,
    summary_extra_months: Sequence[int],
    extra_months_mode: str,
) -> Tuple[List[datetime], List[datetime], datetime, datetime]:
    """Arma los periodos de cobertura ordenados e identifica los extras."""
    months_to_compare: List[int] = []
    for month_num in summary_extra_months:
        if 1 <= int(month_num) <= 12 and int(month_num) != ref_dt.month and int(month_num) not in months_to_compare:
            months_to_compare.append(int(month_num))

    base_prev = datetime(ref_dt.year - 1, ref_dt.month, 1)
    base_curr = datetime(ref_dt.year, ref_dt.month, 1)

    extras_prev = [datetime(ref_dt.year - 1, month_num, 1) for month_num in months_to_compare] if extra_months_mode == "both" else []
    extras_curr = [datetime(ref_dt.year, month_num, 1) for month_num in months_to_compare]

    if extra_months_mode == "both":
        ordered_periods = extras_prev + [base_prev] + extras_curr + [base_curr]
        extra_periods = extras_prev + extras_curr
    else:
        ordered_periods = [base_prev] + extras_curr + [base_curr]
        extra_periods = extras_curr

    return ordered_periods, extra_periods, base_prev, base_curr

def build_summary_columns(
    lang_index: int,
    fabricante: str,
    ref_dt: datetime,
    summary_extra_months: Sequence[int],
    summary_extra_months_mode: str,
) -> Tuple[List[str], List[datetime], List[str]]:
    coverage_periods, extra_periods, _, _ = build_summary_coverage_periods(
        ref_dt,
        summary_extra_months,
        summary_extra_months_mode,
    )
    summary_base_columns: Dict[int, List[str]] = {
        1: [
            "Fabricante/Marca",
            "Pipeline",
            "Penetração Média Mensal",
            f"% VAR {fabricante}",
            "% VAR Worldpanel by Numerator",
        ],
        2: [
            "Fabricante/Marca",
            "Pipeline",
            "Penetración Media Mensual",
            f"% VAR {fabricante}",
            "% VAR Worldpanel by Numerator",
        ],
        3: [
            "Manufacturer/Brand",
            "Pipeline",
            "Monthly Avg Penetration",
            f"% VAR {fabricante}",
            "% VAR Worldpanel by Numerator",
        ],
    }
    coverage_prefix = "Coverage" if lang_index == 3 else "Cobertura"
    stability_label = {1: "Estabilidade", 2: "Estabilidad", 3: "Stability"}[lang_index]
    summary_columns = list(summary_base_columns[lang_index])
    for period_dt in coverage_periods:
        summary_columns.append(f"{coverage_prefix} {period_dt.strftime('%b-%y')}")
    summary_columns.append(stability_label)
    extra_columns = [f"{coverage_prefix} {period_dt.strftime('%b-%y')}" for period_dt in extra_periods]
    return summary_columns, coverage_periods, extra_columns

def build_labels(
    lang_index: int,
    fabricante: str,
    ref_month_year: str,
    summary_extra_months: Optional[Sequence[int]] = None,
    summary_extra_months_mode: str = "recent",
) -> Dict[Tuple[int, str], List[str] | str]:
    """Reproduce el diccionario de etiquetas usado por el script original."""
    ref_dt = dt.strptime(ref_month_year, "%m-%y")
    extra_months = list(summary_extra_months or [])
    summary_pt, _, extra_cols_pt = build_summary_columns(1, fabricante, ref_dt, extra_months, summary_extra_months_mode)
    summary_es, _, extra_cols_es = build_summary_columns(2, fabricante, ref_dt, extra_months, summary_extra_months_mode)
    summary_en, _, extra_cols_en = build_summary_columns(3, fabricante, ref_dt, extra_months, summary_extra_months_mode)

    return {
        (1, "S1"): " ",
        (1, "Summary"): summary_pt,
        (1, "SummaryExtraCoverageCols"): extra_cols_pt,
        (1, "Graf cob Penet Men"): "Penetração Mensal",
        (1, "Titulo Cob"): "Cobertura em Ano Móvel",
        (1, "Var"): "com",
        (1, "Titulo Vol"): "Tendência em Volumen",
        (2, "S1"): " ",
        (2, "Summary"): summary_es,
        (2, "SummaryExtraCoverageCols"): extra_cols_es,
        (2, "Graf cob Penet Men"): "Penetración Mensual",
        (2, "Titulo Cob"): "Cobertura en Año Móvil",
        (2, "Var"): "con",
        (2, "Titulo Vol"): "Tendencia en Volumen",
        (3, "S1"): " ",
        (3, "Summary"): summary_en,
        (3, "SummaryExtraCoverageCols"): extra_cols_en,
        (3, "Graf cob Penet Men"): "PENETRATION BY PERIOD",
        (3, "Titulo Cob"): "MOVING YEAR COVERAGE",
        (3, "Var"): "with",
        (3, "Titulo Vol"): "TREND IN VOLUME",
        (1, "LowPenFooter"): "Marca de baixa penetração (<200 compradores) - Resultados para uso interno",
        (1, "LowPenFooterPlural"): "Marcas de baixa penetração (<200 compradores) - Resultados para uso interno",
        (1, "LowPenSummarySingular"): "O estudo contém 1 marca de baixa penetração (<200 buyers). Resultados para uso interno",
        (1, "LowPenSummaryPlural"): "O estudo contém {n} marcas de baixa penetração (<200 buyers). Resultados para uso interno",
        (2, "LowPenFooter"): "Marca de baja penetración (<200 compradores) - Resultados para uso interno",
        (2, "LowPenFooterPlural"): "Marcas de baja penetración (<200 compradores) - Resultados para uso interno",
        (2, "LowPenSummarySingular"): "El estudio contiene 1 marca de baja penetración (<200 buyers). Resultados para uso interno",
        (2, "LowPenSummaryPlural"): "El estudio contiene {n} marcas de baja penetración (<200 buyers). Resultados para uso interno",
        (3, "LowPenFooter"): "Low penetration brand (<200 buyers) - For internal use only",
        (3, "LowPenFooterPlural"): "Low penetration brands (<200 buyers) - For internal use only",
        (3, "LowPenSummarySingular"): "This study contains 1 low penetration brand (<200 buyers). For internal use only",
        (3, "LowPenSummaryPlural"): "This study contains {n} low penetration brands (<200 buyers). For internal use only",
    }

def dataframe_to_bordered_stream(
    df: "pd.DataFrame",
    hide_index: bool = True,
    dpi: int = 220,
    styler_fn: Optional[Callable] = None,
) -> io.BytesIO:
    """Convierte un DataFrame en imagen PNG con borde negro.

    Permite aplicar personalizaciones adicionales sobre el Styler mediante ``styler_fn``.
    """
    styler = df.style.set_table_styles(
        [
            {"selector": "*", "props": [("font-size", "10pt"), ("font-family", "Calibri"), ("color", "black"), ("border-style", "solid"), ("border-width", "1px"), ("text-align", "center")]},
            {"selector": "th", "props": [("background-color", "#D9E1F2"), ("font-weight", "bold"), ("padding", "3px 5px")]},
            {"selector": "td", "props": [("padding", "2px 4px")]},
        ]
    )
    if hide_index:
        styler = styler.hide(axis="index")
    if styler_fn is not None:
        styler = styler_fn(styler)
    buffer = io.BytesIO()
    dfi.export(styler, buffer, table_conversion="matplotlib", dpi=dpi)
    buffer.seek(0)
    img = Image.open(buffer)
    bordered = ImageOps.expand(img, border=2, fill="black")
    final_stream = io.BytesIO()
    bordered.save(final_stream, format="PNG")
    final_stream.seek(0)
    return final_stream


def ensure_title_frame(slide: "Presentation"):
    """Garantiza que el slide tenga un cuadro de título y devuelve su text_frame."""
    placeholder = slide.shapes.title
    if placeholder is not None:
        return placeholder.text_frame
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    return textbox.text_frame


class SlideBuilder:
    """Encapsula la lógica de creación de slides para mantener el código ordenado."""

    def __init__(
        self,
        presentation: "Presentation",
        lang_index: int,
        labels: Dict[Tuple[int, str], List[str] | str],
        coverage_label: str,
        coverage_type: str,
        ref_month_year: str,
        manufacturer_name: str,
        country_name: str,
        category_name_display: str,
        tipo_eje_tend: str,
        trend_granularity: str = "monthly",
        variations_box_style: str = "classic",
        coverage_slide_variant: str = "classic",
        variations_compact_period_labels: bool = False,
    ) -> None:
        self.ppt = presentation
        self.lang_index = lang_index
        self.labels = labels
        self.coverage_label = coverage_label
        self.coverage_type = coverage_type
        self.ref_month_year = ref_month_year
        self.manufacturer_name = manufacturer_name
        self.country_name = country_name
        self.category_name_display = category_name_display
        self.tipo_eje_tend = tipo_eje_tend
        self.trend_granularity = normalize_trend_granularity(trend_granularity)
        self.variations_box_style = normalize_variations_box_style(variations_box_style)
        self.coverage_slide_variant = normalize_coverage_slide_variant(coverage_slide_variant)
        self.variations_compact_period_labels = bool(variations_compact_period_labels)

    def _add_picture_fit(
        self,
        slide,
        img_stream: io.BytesIO,
        *,
        left,
        top,
        width,
        height,
        halign: str = "center",  # left|center|right
        valign: str = "center",  # top|center|bottom
    ) -> None:
        """Inserta una imagen en un rectángulo (fit) manteniendo aspect ratio.

        Nota: por defecto centra, pero para el header de cobertura se usa anclaje
        a izquierda/derecha para alinear con el gráfico de coberturas.
        """
        img_stream.seek(0)
        try:
            with Image.open(img_stream) as _img:
                px_w, px_h = _img.size
        except Exception:
            px_w, px_h = (1, 1)
        finally:
            img_stream.seek(0)

        aspect = (px_w / px_h) if px_h else 1.0
        box_w_in = float(width) / 914400.0
        box_h_in = float(height) / 914400.0
        placed_w_in = box_w_in
        placed_h_in = placed_w_in / aspect if aspect else box_h_in
        if placed_h_in > box_h_in:
            placed_h_in = box_h_in
            placed_w_in = placed_h_in * aspect

        placed_w = Inches(max(0.1, placed_w_in))
        placed_h = Inches(max(0.1, placed_h_in))

        # Horizontal alignment inside the box
        _h = (halign or "center").strip().lower()
        if _h == "left":
            left2 = left
        elif _h == "right":
            left2 = left + int(width - placed_w)
        else:
            left2 = left + int((width - placed_w) / 2)

        # Vertical alignment inside the box
        _v = (valign or "center").strip().lower()
        if _v == "top":
            top2 = top
        elif _v == "bottom":
            top2 = top + int(height - placed_h)
        else:
            top2 = top + int((height - placed_h) / 2)
        slide.shapes.add_picture(img_stream, left2, top2, width=placed_w, height=placed_h)

    def _month_abbr(self, month: int) -> str:
        # Abreviaciones locales en mayúsculas (para el cuadro "bonito").
        es = ["", "ENE", "FEB", "MAR", "ABR", "MAY", "JUN", "JUL", "AGO", "SEP", "OCT", "NOV", "DIC"]
        pt = ["", "JAN", "FEV", "MAR", "ABR", "MAI", "JUN", "JUL", "AGO", "SET", "OUT", "NOV", "DEZ"]
        en = ["", "JAN", "FEB", "MAR", "APR", "MAY", "JUN", "JUL", "AUG", "SEP", "OCT", "NOV", "DEC"]
        table = en if self.lang_index == 3 else (pt if self.lang_index == 1 else es)
        if 1 <= int(month) <= 12:
            return table[int(month)]
        return "-"

    def _coverage_metric_title(self) -> str:
        ctype = (self.coverage_type or "").strip().lower()
        if ctype == "auto":
            ctype = "absoluta"
        if self.lang_index == 3:
            return "Absolute Coverage" if ctype == "absoluta" else "Relative Coverage"
        return "Cobertura Absoluta" if ctype == "absoluta" else "Cobertura Relativa"

    def _pen_table_headers(self) -> Tuple[str, str]:
        if self.lang_index == 1:
            return "Ano", "Penetração\nMédia Mensal"
        if self.lang_index == 3:
            return "Year", "Monthly Avg\nPenetration"
        return "Año", "Penetración\nMedia Mensual"

    def _stability_label(self) -> str:
        return {1: "Estabilidade", 2: "Estabilidad", 3: "Stability"}[self.lang_index]

    def _pg_text(self, key: str) -> str:
        texts = {
            1: {
                "var_vol": "% VAR Vol",
                "with_pipeline": "COM PIPELINE={pipeline}",
                "worldpanel": "Worldpanel by Numerator",
                "coverage": "Cobertura",
                "relative_coverage": "Cobertura Relativa",
                "annual_penetration": "Penetração\nAnual",
                "unit_of_measure": "Unidade de medida",
            },
            2: {
                "var_vol": "% VAR Vol",
                "with_pipeline": "CON PIPELINE={pipeline}",
                "worldpanel": "Worldpanel by Numerator",
                "coverage": "Cobertura",
                "relative_coverage": "Cobertura Relativa",
                "annual_penetration": "Penetración\nAnual",
                "unit_of_measure": "Unidad de medida",
            },
            3: {
                "var_vol": "% VAR Vol",
                "with_pipeline": "WITH PIPELINE={pipeline}",
                "worldpanel": "Worldpanel by Numerator",
                "coverage": "Coverage",
                "relative_coverage": "Relative\nCoverage",
                "annual_penetration": "Annual\nPenetration",
                "unit_of_measure": "Unit of measure",
            },
        }
        return texts.get(self.lang_index, texts[2]).get(key, "")

    def _footer_month_abbr(self, month: int) -> str:
        es = ["", "Ene", "Feb", "Mar", "Abr", "May", "Jun", "Jul", "Ago", "Sep", "Oct", "Nov", "Dic"]
        pt = ["", "Jan", "Fev", "Mar", "Abr", "Mai", "Jun", "Jul", "Ago", "Set", "Out", "Nov", "Dez"]
        en = ["", "Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
        table = en if self.lang_index == 3 else (pt if self.lang_index == 1 else es)
        if 1 <= int(month) <= 12:
            return table[int(month)]
        return "-"

    def _resolve_pg_coverage_values(self, coverage_series: "pd.Series") -> Tuple[float, float, float, float]:
        try:
            ref_dt = dt.strptime(self.ref_month_year, "%m-%y")
        except Exception:
            idx = pd.to_datetime(getattr(coverage_series, "index", []), errors="coerce")
            idx = idx[~idx.isna()]
            if len(idx) == 0:
                return (np.nan, np.nan, np.nan, np.nan)
            ref_dt = idx.max().to_pydatetime()

        prev_dt = ref_dt - pd.DateOffset(months=12)
        cov_prev = _coverage_value_for_year_month(coverage_series, int(prev_dt.year), int(prev_dt.month))
        cov_curr = _coverage_value_for_year_month(coverage_series, int(ref_dt.year), int(ref_dt.month))
        pop_val_num = get_population_coverage_percent(self.country_name) / 100.0
        ctype = (self.coverage_type or "").strip().lower()
        if ctype == "auto":
            ctype = "absoluta"

        if ctype == "relativa":
            rel_prev = cov_prev
            rel_curr = cov_curr
            abs_prev = (float(cov_prev) * pop_val_num) if pop_val_num > 0 and pd.notna(cov_prev) else np.nan
            abs_curr = (float(cov_curr) * pop_val_num) if pop_val_num > 0 and pd.notna(cov_curr) else np.nan
        else:
            abs_prev = cov_prev
            abs_curr = cov_curr
            rel_prev = (float(cov_prev) / pop_val_num) if pop_val_num > 0 and pd.notna(cov_prev) else np.nan
            rel_curr = (float(cov_curr) / pop_val_num) if pop_val_num > 0 and pd.notna(cov_curr) else np.nan
        return abs_prev, abs_curr, rel_prev, rel_curr

    def _resolve_pg_bank_coverage_values(self, bank_row: object) -> Tuple[float, float, float, float]:
        try:
            abs_curr = float(bank_row.get('Cobertura Año Mov Actual'))
        except Exception:
            abs_curr = np.nan
        try:
            abs_prev = float(bank_row.get('Cobertura Año Mov Anterior'))
        except Exception:
            abs_prev = np.nan
        try:
            country_name = str(bank_row.get('Pais', self.country_name))
            pop_val_num = get_population_coverage_percent(country_name) / 100.0
        except Exception:
            pop_val_num = 0.0

        ctype = (self.coverage_type or "").strip().lower()
        if ctype == "auto":
            ctype = "absoluta"
        if ctype == "relativa":
            rel_curr = abs_curr
            rel_prev = abs_prev
            abs_curr = (float(abs_curr) * pop_val_num) if pop_val_num > 0 and pd.notna(abs_curr) else np.nan
            abs_prev = (float(abs_prev) * pop_val_num) if pop_val_num > 0 and pd.notna(abs_prev) else np.nan
        else:
            rel_curr = (float(abs_curr) / pop_val_num) if pop_val_num > 0 and pd.notna(abs_curr) else np.nan
            rel_prev = (float(abs_prev) / pop_val_num) if pop_val_num > 0 and pd.notna(abs_prev) else np.nan
        return abs_prev, abs_curr, rel_prev, rel_curr

    def _add_pg_variation_table_shape(self, slide, assets: PipelineAssets, *, left, top, width, height) -> None:
        table_shape = slide.shapes.add_table(3, 3, left, top, width, height)
        table_shape.height = height
        table = table_shape.table
        self._clear_powerpoint_table_style(table)
        table.columns[0].width = int(width * 0.32)
        table.columns[1].width = int(width * 0.34)
        table.columns[2].width = width - table.columns[0].width - table.columns[1].width
        table.rows[0].height = int(height * 0.45)
        table.rows[1].height = int((height - table.rows[0].height) / 2)
        table.rows[2].height = height - table.rows[0].height - table.rows[1].height

        header_bg = self._hex_to_rgb("#8AA6C1")
        white = RGBColor(255, 255, 255)
        white_bg = RGBColor(255, 255, 255)

        try:
            ref_year = int(dt.strptime(self.ref_month_year, "%m-%y").year)
        except Exception:
            ref_year = dt.now().year
        row_labels = [
            f"{ref_year} VS {ref_year - 1}",
            f"{ref_year - 1} VS {ref_year - 2}",
        ]

        self._set_table_cell_text(table.cell(0, 0), self._pg_text("var_vol"), fill_color=header_bg, font_color=white, font_size=10, align=2)
        self._set_table_cell_text(
            table.cell(0, 1),
            f"{self.manufacturer_name.upper()} {self._pg_text('with_pipeline').format(pipeline=assets.pipeline)}",
            fill_color=header_bg,
            font_color=white,
            font_size=10,
            align=2,
        )
        self._set_table_cell_text(table.cell(0, 2), self._pg_text("worldpanel"), fill_color=header_bg, font_color=white, font_size=10, align=2)

        value_specs = [
            (row_labels[0], self._fmt_pct(assets.annual_var_cliente_y1), self._fmt_pct(assets.annual_var_wp_y1)),
            (row_labels[1], self._fmt_pct(assets.annual_var_cliente_y2), self._fmt_pct(assets.annual_var_wp_y2)),
        ]
        for idx, (label, client_val, wp_val) in enumerate(value_specs, start=1):
            client_color = self._summary_variation_font_color(client_val)
            wp_color = self._summary_variation_font_color(wp_val)
            self._set_table_cell_text(table.cell(idx, 0), label, fill_color=header_bg, font_color=white, font_size=10, align=2)
            self._set_table_cell_text(table.cell(idx, 1), client_val, fill_color=white_bg, font_color=client_color, font_size=11, bold=False, align=2)
            self._set_table_cell_text(table.cell(idx, 2), wp_val, fill_color=white_bg, font_color=wp_color, font_size=11, bold=False, align=2)

    def _add_pg_coverage_table_shape(self, slide, assets: PipelineAssets, *, left, top, width, height) -> None:
        table_shape = slide.shapes.add_table(3, 4, left, top, width, height)
        table_shape.height = height
        table = table_shape.table
        self._clear_powerpoint_table_style(table)
        table.columns[0].width = int(width * 0.22)
        table.columns[1].width = int(width * 0.26)
        table.columns[2].width = int(width * 0.26)
        table.columns[3].width = width - table.columns[0].width - table.columns[1].width - table.columns[2].width
        table.rows[0].height = int(height * 0.45)
        table.rows[1].height = int((height - table.rows[0].height) / 2)
        table.rows[2].height = height - table.rows[0].height - table.rows[1].height

        header_bg = self._hex_to_rgb("#8AA6C1")
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
        white_bg = RGBColor(255, 255, 255)

        try:
            ref_dt = dt.strptime(self.ref_month_year, "%m-%y")
            curr_year = int(ref_dt.year)
        except Exception:
            curr_year = dt.now().year
        prev_year = curr_year - 1
        abs_prev, abs_curr, rel_prev, rel_curr = self._resolve_pg_coverage_values(assets.coverage_series)

        def _fmt_cov_cell(value: object) -> str:
            return "-" if value is None or pd.isna(value) else f"{float(value):.1f}"

        self._set_table_cell_text(table.cell(0, 0), "", fill_color=header_bg, font_color=white, font_size=10, align=2)
        self._set_table_cell_text(table.cell(0, 1), self._pg_text("coverage"), fill_color=header_bg, font_color=white, font_size=10, align=2)
        self._set_table_cell_text(table.cell(0, 2), self._pg_text("relative_coverage"), fill_color=header_bg, font_color=white, font_size=10, align=2)
        self._set_table_cell_text(table.cell(0, 3), self._pg_text("annual_penetration"), fill_color=header_bg, font_color=white, font_size=10, align=2)

        rows = [
            (str(curr_year), _fmt_cov_cell(abs_curr), _fmt_cov_cell(rel_curr), ""),
            (str(prev_year), _fmt_cov_cell(abs_prev), _fmt_cov_cell(rel_prev), ""),
        ]
        for idx, row_values in enumerate(rows, start=1):
            self._set_table_cell_text(table.cell(idx, 0), row_values[0], fill_color=header_bg, font_color=white, font_size=10, align=2)
            self._set_table_cell_text(table.cell(idx, 1), row_values[1], fill_color=white_bg, font_color=black, font_size=11, bold=False, align=2)
            self._set_table_cell_text(table.cell(idx, 2), row_values[2], fill_color=white_bg, font_color=black, font_size=11, bold=False, align=2)
            self._set_table_cell_text(table.cell(idx, 3), row_values[3], fill_color=white_bg, font_color=black, font_size=11, bold=False, align=2)

    def _add_pg_footer_text(self, slide, assets: PipelineAssets) -> None:
        try:
            ref_dt = dt.strptime(self.ref_month_year, "%m-%y")
            period_txt = f"{self._footer_month_abbr(ref_dt.month)}'{ref_dt.year % 100:02d}"
        except Exception:
            period_txt = self.ref_month_year
        unit_txt = str(assets.measure_unit or "-").strip() or "-"
        if unit_txt.upper() == "VOLSU":
            unit_txt = "Stats Units"
        footer_txt = f"{self._pg_text('unit_of_measure')}: {unit_txt} | {period_txt} | {self.country_name}"
        tb = slide.shapes.add_textbox(Inches(2.65), Inches(6.95), Inches(8.0), Inches(0.32))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = footer_txt
        p.font.size = Pt(16)
        p.font.bold = False
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = 1

    def _add_pg_subtitle(self, slide, assets: PipelineAssets) -> None:
        tag = "Definition"
        try:
            corr_value = None if assets.current_year_correlation is None or pd.isna(assets.current_year_correlation) else float(assets.current_year_correlation)
            trend_following = bool(assets.trend_following)
            if trend_following and corr_value is not None and (corr_value * 100.0) > 50.0:
                tag = "Correlation"
        except Exception:
            tag = "Definition"
        subtitle = f"Definition of pipeline {int(assets.pipeline)} ({tag})"
        tb = slide.shapes.add_textbox(Inches(0.62), Inches(0.72), Inches(7.2), Inches(0.42))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.bold = False
        p.font.color.rgb = self._hex_to_rgb("#04586C")
        p.alignment = 1

    def _add_pg_category_badge(self, slide) -> None:
        badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.25), Inches(0.0), Inches(1.05), Inches(0.78))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self._hex_to_rgb("#BEF9FF")
        badge.line.fill.background()

        tf = badge.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(3)
        tf.margin_right = Pt(3)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(3)
        p = tf.paragraphs[0]
        p.text = str(self.category_name_display or "").strip() or "-"
        p.font.size = Pt(11)
        p.font.bold = False
        p.font.color.rgb = self._hex_to_rgb("#04586C")
        p.alignment = 1

    def _add_cov_slide_pg_layout(self, slide, assets: PipelineAssets) -> None:
        self._add_pg_subtitle(slide, assets)
        self._add_pg_category_badge(slide)

        tables_top = Inches(5.62)
        tables_height = Inches(0.96)
        left_w = Inches(4.68)
        right_w = Inches(4.68)
        gap = Inches(0.22)
        total_w = left_w + right_w + gap
        left_x = int((self.ppt.slide_width - total_w) / 2)
        right_x = left_x + left_w + gap

        self._add_pg_variation_table_shape(
            slide,
            assets,
            left=left_x,
            top=tables_top,
            width=left_w,
            height=tables_height,
        )
        self._add_pg_coverage_table_shape(
            slide,
            assets,
            left=right_x,
            top=tables_top,
            width=right_w,
            height=tables_height,
        )
        self._add_pg_footer_text(slide, assets)

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> "RGBColor":
        raw = str(hex_color or "").strip().lstrip("#")
        if len(raw) != 6:
            return RGBColor(0, 0, 0)
        try:
            return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
        except Exception:
            return RGBColor(0, 0, 0)

    @staticmethod
    def _set_paragraph_text(
        paragraph,
        text: object,
        *,
        font_size: int,
        font_color: Optional["RGBColor"] = None,
        bold: bool = False,
        align: int = 1,
    ) -> None:
        paragraph.clear()
        paragraph.alignment = align
        paragraph.font.bold = bool(bold)
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = font_color if font_color is not None else RGBColor(0, 0, 0)
        run = paragraph.add_run()
        run.text = "" if text is None else str(text)
        run.font.bold = bool(bold)
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color if font_color is not None else RGBColor(0, 0, 0)

    def _set_table_cell_text(
        self,
        cell,
        text: object,
        *,
        fill_color: Optional["RGBColor"] = None,
        font_color: Optional["RGBColor"] = None,
        font_size: int = 12,
        bold: bool = True,
        align: int = 2,
        word_wrap: bool = True,
    ) -> None:
        if fill_color is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color

        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = bool(word_wrap)
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        try:
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        except Exception:
            pass

        self._set_paragraph_text(
            tf.paragraphs[0],
            text,
            font_size=font_size,
            font_color=font_color if font_color is not None else RGBColor(0, 0, 0),
            bold=bold,
            align=align,
        )

    @staticmethod
    def _clear_table_cell_borders(cell) -> None:
        tc = cell._tc
        tc_pr = tc.get_or_add_tcPr()
        for side in ("lnL", "lnR", "lnT", "lnB"):
            existing = tc_pr.find(qn(f"a:{side}"))
            if existing is not None:
                tc_pr.remove(existing)

    @staticmethod
    def _add_table_cell_border(cell, side: str, color: "RGBColor", width: int = 12700) -> None:
        tc = cell._tc
        tc_pr = tc.get_or_add_tcPr()
        ln = OxmlElement(f"a:{side}")
        ln.set("w", str(int(width)))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")

        solid_fill = OxmlElement("a:solidFill")
        srgb_clr = OxmlElement("a:srgbClr")
        srgb_clr.set("val", f"{int(color[0]):02X}{int(color[1]):02X}{int(color[2]):02X}")
        solid_fill.append(srgb_clr)
        ln.append(solid_fill)

        prst_dash = OxmlElement("a:prstDash")
        prst_dash.set("val", "solid")
        ln.append(prst_dash)

        round_join = OxmlElement("a:round")
        ln.append(round_join)

        head_end = OxmlElement("a:headEnd")
        head_end.set("type", "none")
        head_end.set("w", "med")
        head_end.set("len", "med")
        ln.append(head_end)

        tail_end = OxmlElement("a:tailEnd")
        tail_end.set("type", "none")
        tail_end.set("w", "med")
        tail_end.set("len", "med")
        ln.append(tail_end)

        tc_pr.append(ln)

    @staticmethod
    def _add_transparent_table_cell_border(cell, side: str, width: int = 0) -> None:
        tc = cell._tc
        tc_pr = tc.get_or_add_tcPr()
        ln = OxmlElement(f"a:{side}")
        ln.set("w", str(int(width)))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")

        solid_fill = OxmlElement("a:solidFill")
        prst_clr = OxmlElement("a:prstClr")
        prst_clr.set("val", "black")
        alpha = OxmlElement("a:alpha")
        alpha.set("val", "0")
        prst_clr.append(alpha)
        solid_fill.append(prst_clr)
        ln.append(solid_fill)

        prst_dash = OxmlElement("a:prstDash")
        prst_dash.set("val", "solid")
        ln.append(prst_dash)

        tc_pr.append(ln)

    def _apply_internal_table_borders(
        self,
        slide,
        table,
        *,
        left,
        top,
        table_width,
        table_height,
        color: "RGBColor",
        border_width: int = 12700,
    ) -> None:
        row_count = len(table.rows)
        col_count = len(table.columns)
        for r in range(row_count):
            for c in range(col_count):
                cell = table.cell(r, c)
                self._clear_table_cell_borders(cell)
                if c < col_count - 1:
                    self._add_table_cell_border(cell, "lnR", color, width=border_width)
                if r < row_count - 1:
                    self._add_table_cell_border(cell, "lnB", color, width=border_width)

    @staticmethod
    def _clear_powerpoint_table_style(table) -> None:
        try:
            table.first_row = False
            table.first_col = False
            table.last_row = False
            table.last_col = False
            table.horz_banding = False
            table.vert_banding = False
        except Exception:
            pass
        try:
            tbl_pr = table._tbl.tblPr
            for style_node in list(tbl_pr.findall(qn("a:tableStyleId"))):
                tbl_pr.remove(style_node)
            for attr_name in ("firstRow", "bandRow", "lastRow", "firstCol", "lastCol", "bandCol"):
                if attr_name in tbl_pr.attrib:
                    del tbl_pr.attrib[attr_name]
        except Exception:
            pass

    @staticmethod
    def _normalize_summary_table_value(value: object) -> str:
        if value is None:
            return "-"
        try:
            if "pd" in globals() and pd.isna(value):
                return "-"
        except Exception:
            pass
        txt = str(value).strip()
        return txt if txt else "-"

    @staticmethod
    def _parse_summary_percent_value(value: object) -> Optional[float]:
        txt = SlideBuilder._normalize_summary_table_value(value)
        if txt in {"-", ""}:
            return None
        txt = txt.replace("%", "").replace(",", ".").strip()
        try:
            return float(txt)
        except Exception:
            return None

    @classmethod
    def _summary_variation_font_color(cls, value: object) -> "RGBColor":
        parsed = cls._parse_summary_percent_value(value)
        if parsed is None:
            return RGBColor(0, 0, 0)
        if parsed > 0:
            return RGBColor(0, 97, 0)
        if parsed < 0:
            return RGBColor(156, 0, 6)
        return RGBColor(120, 120, 120)

    @classmethod
    def _variation_table_fill_and_font_colors(cls, value: object) -> Tuple["RGBColor", "RGBColor"]:
        parsed = cls._parse_summary_percent_value(value)
        if parsed is None:
            return RGBColor(255, 255, 255), RGBColor(0, 0, 0)
        if parsed > 0:
            return RGBColor(198, 239, 206), RGBColor(0, 97, 0)
        if parsed < 0:
            return RGBColor(255, 199, 206), RGBColor(156, 0, 6)
        return RGBColor(255, 255, 255), RGBColor(0, 0, 0)

    @staticmethod
    def _normalize_brand_key(value: object) -> str:
        txt = SlideBuilder._normalize_summary_table_value(value)
        if txt == "-":
            return ""
        txt = unicodedata.normalize("NFD", txt)
        txt = "".join(ch for ch in txt if unicodedata.category(ch) != "Mn")
        txt = re.sub(r"\s+", " ", txt).strip().lower()
        return txt

    @staticmethod
    def _format_summary_header_label(col_name: object) -> str:
        """Formatea encabezados largos del summary sin alterar el nombre base de la columna."""
        raw_label = str(col_name).strip()
        match = re.match(r"^(Cobertura|Coverage)\s+([A-Za-z]{3}-\d{2})$", raw_label, flags=re.IGNORECASE)
        if match:
            return f"{match.group(1)}\n{match.group(2)}"
        return raw_label

    def _compute_summary_table_column_widths(self, df_summary: "pd.DataFrame", width: int) -> List[int]:
        """Calcula anchos estables para el summary segun encabezado y columnas visibles."""
        cols = int(len(df_summary.columns))
        if cols <= 0:
            return []
        total_width = int(width)
        if total_width <= 0:
            return []

        col_name_norms = [str(col_name).strip().lower() for col_name in df_summary.columns]

        def _semantic_role(col_name_norm: str, idx: int) -> str:
            if "fabricante/marca" in col_name_norm or "manufacturer/brand" in col_name_norm:
                return "brand"
            if col_name_norm.startswith("pipeline"):
                return "pipeline"
            if "penetr" in col_name_norm:
                return "penetration"
            if "worldpanel by numerator" in col_name_norm:
                return "worldpanel"
            if col_name_norm.startswith("cobertura ") or col_name_norm.startswith("coverage "):
                return "coverage"
            if "estabilidad" in col_name_norm or "estabilidade" in col_name_norm or "stability" in col_name_norm:
                return "stability"
            if idx in (3, 4):
                return "variation"
            return "generic"

        ideal_share_by_role: Dict[str, float] = {
            "brand": 0.19,
            "pipeline": 0.075,
            "penetration": 0.16,
            "variation": 0.12,
            "worldpanel": 0.18,
            "coverage": 0.145,
            "stability": 0.10,
            "generic": 0.11,
        }
        min_share_by_role: Dict[str, float] = {
            "brand": 0.13,
            "pipeline": 0.06,
            "penetration": 0.12,
            "variation": 0.09,
            "worldpanel": 0.13,
            "coverage": 0.10,
            "stability": 0.08,
            "generic": 0.07,
        }

        roles = [_semantic_role(col_name_norm, idx) for idx, col_name_norm in enumerate(col_name_norms)]
        raw_weights: List[float] = []
        min_widths: List[int] = []
        for idx, role in enumerate(roles):
            sample_values = df_summary.iloc[:, idx].head(15).tolist()
            max_cell_len = max((len(self._normalize_summary_table_value(v)) for v in sample_values), default=0)
            header_len = len(self._format_summary_header_label(df_summary.columns[idx]).replace("\n", " "))
            content_factor = max(max_cell_len, header_len)
            weight = ideal_share_by_role.get(role, ideal_share_by_role["generic"]) * (1.0 + min(0.18, max(0, content_factor - 10) * 0.01))
            raw_weights.append(weight)
            min_widths.append(max(int(total_width * min_share_by_role.get(role, min_share_by_role["generic"])), int(total_width * 0.055)))

        min_required = sum(min_widths)
        if min_required >= total_width:
            even_width = total_width // cols
            col_widths = [even_width] * cols
            col_widths[-1] += total_width - sum(col_widths)
            return col_widths

        extra_width = total_width - min_required
        total_weight = sum(raw_weights) or float(cols)
        provisional = [min_widths[idx] + int(round(extra_width * (raw_weights[idx] / total_weight))) for idx in range(cols)]

        width_delta = total_width - sum(provisional)
        if width_delta != 0:
            adjustable_indexes = sorted(range(cols), key=lambda i: raw_weights[i], reverse=(width_delta > 0))
            cursor = 0
            step = 1 if width_delta > 0 else -1
            remaining = abs(width_delta)
            while remaining > 0 and adjustable_indexes:
                idx = adjustable_indexes[cursor % len(adjustable_indexes)]
                next_width = provisional[idx] + step
                if next_width >= min_widths[idx]:
                    provisional[idx] = next_width
                    remaining -= 1
                cursor += 1

        return provisional

    def _add_editable_summary_table(
        self,
        slide,
        df_summary: "pd.DataFrame",
        *,
        left,
        top,
        width,
        max_height,
        low_penetration_brands: Optional[Sequence[str]] = None,
    ) -> Optional[int]:
        if df_summary is None or df_summary.empty:
            return None

        rows = int(len(df_summary.index)) + 1
        cols = int(len(df_summary.columns))
        if rows <= 1 or cols <= 0:
            return None

        max_height = int(max_height)
        if max_height <= 0:
            max_height = int(Inches(4.8))

        header_h = int(Inches(0.34))
        body_rows = max(rows - 1, 0)
        preferred_body_h = int(Inches(0.25))
        min_body_h = int(Inches(0.17))

        if body_rows > 0:
            needed_h = header_h + (body_rows * preferred_body_h)
            if needed_h <= max_height:
                body_h = preferred_body_h
            else:
                body_h = max(min_body_h, int((max_height - header_h) / body_rows))
                needed_h = header_h + (body_rows * body_h)
        else:
            body_h = 0
            needed_h = header_h

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, needed_h)
        table = table_shape.table
        self._clear_powerpoint_table_style(table)

        col_widths = self._compute_summary_table_column_widths(df_summary, width)
        for idx, col_w in enumerate(col_widths):
            table.columns[idx].width = col_w

        table.rows[0].height = header_h
        for r in range(1, rows):
            table.rows[r].height = body_h

        header_bg = RGBColor(217, 225, 242)
        stripe_bg = RGBColor(245, 247, 251)
        white_bg = RGBColor(255, 255, 255)
        soft_red_bg = RGBColor(255, 235, 235)
        black = RGBColor(0, 0, 0)

        if rows <= 9:
            body_font_size = 10
        elif rows <= 14:
            body_font_size = 9
        else:
            body_font_size = 8
        low_penetration_keys: Set[str] = set()
        for brand in (low_penetration_brands or []):
            key = self._normalize_brand_key(brand)
            if key:
                low_penetration_keys.add(key)

        for c, col_name in enumerate(df_summary.columns):
            self._set_table_cell_text(
                table.cell(0, c),
                self._format_summary_header_label(col_name),
                fill_color=header_bg,
                font_color=black,
                font_size=10,
                bold=True,
                align=2,
                word_wrap=True,
            )

        for r, row_values in enumerate(df_summary.itertuples(index=False), start=1):
            brand_key = self._normalize_brand_key(row_values[0]) if cols > 0 else ""
            row_is_low_penetration = brand_key in low_penetration_keys
            for c in range(cols):
                val = self._normalize_summary_table_value(row_values[c])
                align = 1 if c == 0 else 2
                fill = soft_red_bg if row_is_low_penetration else (stripe_bg if r % 2 == 0 else white_bg)
                font_color = self._summary_variation_font_color(row_values[c]) if c in (3, 4) else black
                self._set_table_cell_text(
                    table.cell(r, c),
                    val,
                    fill_color=fill,
                    font_color=font_color,
                    font_size=body_font_size,
                    bold=False,
                    align=align,
                    word_wrap=False,
                )

        separator_color = RGBColor(255, 255, 255)
        separator_width = int(Pt(0.45))
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                self._clear_table_cell_borders(cell)
                self._add_transparent_table_cell_border(cell, "lnL")
                if c < cols - 1:
                    self._add_table_cell_border(cell, "lnR", separator_color, width=separator_width)
                else:
                    self._add_transparent_table_cell_border(cell, "lnR")
                self._add_transparent_table_cell_border(cell, "lnT")
                # No dibujar borde inferior en el header: evita la linea blanca entre header y contenido.
                if r > 0 and r < rows - 1:
                    self._add_table_cell_border(cell, "lnB", separator_color, width=separator_width)
                else:
                    self._add_transparent_table_cell_border(cell, "lnB")
        return needed_h

    def _add_pg_summary_table(
        self,
        slide,
        df_bank: "pd.DataFrame",
        *,
        left,
        top,
        width,
        max_height,
    ) -> Optional[int]:
        if df_bank is None or df_bank.empty:
            return None

        try:
            ref_dt = dt.strptime(self.ref_month_year, "%m-%y")
            prev_year = int(ref_dt.year) - 1
            curr_year = int(ref_dt.year)
            month_label = ref_dt.strftime("%b")
        except Exception:
            curr_year = dt.now().year
            prev_year = curr_year - 1
            month_label = ""
        prev_mat_label = f"MAT {month_label} {prev_year}".strip()
        curr_mat_label = f"MAT {month_label} {curr_year}".strip()

        table_df = df_bank.copy()
        rows = int(len(table_df.index)) + 2
        cols = 11
        if rows <= 2:
            return None

        header_h1 = int(Inches(0.36))
        header_h2 = int(Inches(0.34))
        body_rows = rows - 2
        preferred_body_h = int(Inches(0.28))
        min_body_h = int(Inches(0.20))
        max_height = int(max_height)
        body_h = preferred_body_h
        needed_h = header_h1 + header_h2 + (body_rows * body_h)
        if needed_h > max_height:
            body_h = max(min_body_h, int((max_height - header_h1 - header_h2) / max(body_rows, 1)))
            needed_h = header_h1 + header_h2 + (body_rows * body_h)

        header_bg = self._hex_to_rgb("#8FA9C3")
        header_bg_light = self._hex_to_rgb("#A9CBE8")
        white = RGBColor(255, 255, 255)
        black = RGBColor(35, 35, 35)
        white_bg = RGBColor(255, 255, 255)
        positive_green = RGBColor(0, 176, 80)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, needed_h)
        table = table_shape.table
        self._clear_powerpoint_table_style(table)

        col_ratios = [0.14, 0.085, 0.085, 0.07, 0.11, 0.11, 0.07, 0.07, 0.085, 0.085, 0.09]
        assigned_w = 0
        for idx, ratio in enumerate(col_ratios):
            if idx == cols - 1:
                col_w = int(width - assigned_w)
            else:
                col_w = int(width * ratio)
                assigned_w += col_w
            table.columns[idx].width = col_w

        table.rows[0].height = header_h1
        table.rows[1].height = header_h2
        for row_idx in range(2, rows):
            table.rows[row_idx].height = body_h

        merge_specs = [
            (0, 0, 1, 0),
            (0, 1, 0, 2),
            (0, 3, 1, 3),
            (0, 4, 0, 5),
            (0, 6, 0, 7),
            (0, 8, 0, 9),
            (0, 10, 1, 10),
        ]
        for r1, c1, r2, c2 in merge_specs:
            try:
                table.cell(r1, c1).merge(table.cell(r2, c2))
            except Exception:
                pass

        top_headers = {
            (0, 0): "Brand",
            (0, 1): "Annual Penetration %",
            (0, 3): "PIPELINE",
            (0, 4): f"Var % ({curr_year} vs {prev_year})",
            (0, 6): "Coverage",
            (0, 8): "Relative Coverage",
            (0, 10): "Var. pp",
        }
        bottom_headers = {
            (1, 1): prev_mat_label,
            (1, 2): curr_mat_label,
            (1, 4): f"{str(self.manufacturer_name or 'P&G').upper()} WITH\nPIPELINE",
            (1, 5): "Worldpanel by\nNumerator",
            (1, 6): str(prev_year),
            (1, 7): str(curr_year),
            (1, 8): str(prev_year),
            (1, 9): str(curr_year),
        }

        for (row_idx, col_idx), text in top_headers.items():
            fill = header_bg if col_idx in (0, 3, 4, 10) else header_bg_light
            font_color = white if col_idx in (0, 3, 4) else black
            self._set_table_cell_text(
                table.cell(row_idx, col_idx),
                text,
                fill_color=fill,
                font_color=font_color,
                font_size=9,
                bold=True,
                align=2,
                word_wrap=True,
            )
        for (row_idx, col_idx), text in bottom_headers.items():
            fill = header_bg if col_idx in (4, 5) else header_bg_light
            font_color = white if col_idx in (4, 5) else black
            self._set_table_cell_text(
                table.cell(row_idx, col_idx),
                text,
                fill_color=fill,
                font_color=font_color,
                font_size=8,
                bold=True,
                align=2,
                word_wrap=True,
            )

        def _fmt_num(value: object) -> str:
            try:
                if pd.isna(value):
                    return "-"
                return f"{float(value):.1f}"
            except Exception:
                return "-"

        def _fmt_pct_from_points(value: object) -> str:
            try:
                if pd.isna(value):
                    return "-"
                return f"{float(value):.1f}%"
            except Exception:
                return "-"

        for table_row, (_, bank_row) in enumerate(table_df.iterrows(), start=2):
            abs_prev, abs_curr, rel_prev, rel_curr = self._resolve_pg_bank_coverage_values(bank_row)
            var_pp = (abs_curr - abs_prev) if pd.notna(abs_curr) and pd.notna(abs_prev) else np.nan
            row_values = [
                bank_row.get('Fabricante', '') or bank_row.get('Fabricante/Marca', ''),
                "",
                "",
                bank_row.get('Pipeline', ''),
                _fmt_pct_from_points(bank_row.get('%VAR Cliente', np.nan)),
                _fmt_pct_from_points(bank_row.get('% VAR WP by Numerator', np.nan)),
                _fmt_num(abs_prev),
                _fmt_num(abs_curr),
                _fmt_num(rel_prev),
                _fmt_num(rel_curr),
                _fmt_num(var_pp),
            ]
            for col_idx, value in enumerate(row_values):
                font_color = black
                if col_idx in (4, 5):
                    font_color = self._summary_variation_font_color(value)
                    if self._parse_summary_percent_value(value) and self._parse_summary_percent_value(value) > 0:
                        font_color = positive_green
                self._set_table_cell_text(
                    table.cell(table_row, col_idx),
                    value,
                    fill_color=white_bg,
                    font_color=font_color,
                    font_size=8,
                    bold=False,
                    align=2,
                    word_wrap=False,
                )

        border_color = RGBColor(0, 0, 0)
        border_width = int(Pt(0.75))
        for row_idx in range(rows):
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)
                self._clear_table_cell_borders(cell)
                for side in ("lnL", "lnR", "lnT", "lnB"):
                    self._add_table_cell_border(cell, side, border_color, width=border_width)
        return needed_h

    def _add_editable_variations_table(
        self,
        slide,
        df_variations: "pd.DataFrame",
        *,
        left,
        top,
        width,
        max_height,
    ) -> None:
        if df_variations is None or df_variations.empty:
            return

        compare_lags = (
            pd.to_numeric(df_variations["_CompareLagMonths"], errors="coerce").reset_index(drop=True)
            if "_CompareLagMonths" in df_variations.columns
            else pd.Series(np.nan, index=range(len(df_variations)), dtype=float)
        )
        table_df = df_variations[[col for col in df_variations.columns if not str(col).startswith('_')]].copy()
        table_df.reset_index(drop=True, inplace=True)
        if table_df.empty:
            return

        rows = int(len(table_df.index)) + 1
        cols = int(len(table_df.columns))
        if rows <= 1 or cols <= 0:
            return

        text_columns = [col for col in ("Tipo", "Periodo") if col in table_df.columns]
        value_columns = [col for col in table_df.columns if col not in set(text_columns)]

        max_height = int(max_height)
        if max_height <= 0:
            max_height = int(Inches(1.15))

        header_h = int(Inches(0.22))
        body_rows = max(rows - 1, 0)
        preferred_body_h = int(Inches(0.18))
        min_body_h = int(Inches(0.15))

        if body_rows > 0:
            needed_h = header_h + (body_rows * preferred_body_h)
            if needed_h <= max_height:
                body_h = preferred_body_h
            else:
                body_h = max(min_body_h, int((max_height - header_h) / body_rows))
                needed_h = header_h + (body_rows * body_h)
        else:
            body_h = 0
            needed_h = header_h

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, needed_h)
        table = table_shape.table
        self._clear_powerpoint_table_style(table)

        col_weights: List[int] = []
        for col_name in table_df.columns:
            sample_values = table_df[col_name].head(12).tolist()
            max_cell_len = max((len(self._normalize_summary_table_value(v)) for v in sample_values), default=0)
            col_name_norm = str(col_name).strip().lower()
            header_len = len(str(col_name))
            if col_name_norm == "tipo":
                weight = 8
            elif col_name_norm == "periodo":
                weight = 19
            elif col_name_norm.startswith("wp by numerator"):
                weight = 11
            elif col_name_norm.startswith("cliente p0"):
                weight = 10
            elif col_name_norm.startswith("cliente pipeline"):
                weight = 15
            else:
                # Para columnas numéricas priorizamos el ancho del encabezado,
                # no el contenido, que siempre es corto (porcentajes).
                weight = max(8, min(13, int(round(max(header_len, max_cell_len) * 0.6))))
            col_weights.append(weight)
        weight_sum = sum(col_weights) if col_weights else cols

        width_assigned = 0
        for idx, weight in enumerate(col_weights):
            if idx == cols - 1:
                col_w = int(width - width_assigned)
            else:
                col_w = int(width * (float(weight) / float(weight_sum)))
                width_assigned += col_w
            table.columns[idx].width = max(col_w, int(width * 0.07))

        table.rows[0].height = header_h
        for r in range(1, rows):
            table.rows[r].height = body_h

        white_bg = RGBColor(255, 255, 255)
        year_ago_bg = RGBColor(230, 254, 248)  # #E6FEF8, turquesa muy suave
        black = RGBColor(0, 0, 0)

        header_font_size = 8
        body_font_size = 8 if rows > 5 else 9

        for c, col_name in enumerate(table_df.columns):
            self._set_table_cell_text(
                table.cell(0, c),
                str(col_name),
                fill_color=white_bg,
                font_color=black,
                font_size=header_font_size,
                bold=True,
                align=2,
                word_wrap=True,
            )

        for r, row_values in enumerate(table_df.itertuples(index=False), start=1):
            compare_lag = compare_lags.iloc[r - 1] if (r - 1) < len(compare_lags) else np.nan
            is_year_ago_comparison = pd.notna(compare_lag) and int(compare_lag) == 12
            for c, col_name in enumerate(table_df.columns):
                raw_val = row_values[c]
                if col_name in value_columns:
                    cell_text = "-" if (pd.isna(raw_val) or (isinstance(raw_val, str) and str(raw_val).strip() == "-")) else f"{float(raw_val) * 100:.1f}%"
                    fill_color, font_color = self._variation_table_fill_and_font_colors(raw_val)
                    align = 2
                else:
                    cell_text = self._normalize_summary_table_value(raw_val)
                    fill_color, font_color = white_bg, black
                    align = 1
                if is_year_ago_comparison and col_name in text_columns:
                    fill_color = year_ago_bg
                self._set_table_cell_text(
                    table.cell(r, c),
                    cell_text,
                    fill_color=fill_color,
                    font_color=font_color,
                    font_size=body_font_size,
                    bold=False,
                    align=align,
                    word_wrap=False,
                )
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                self._clear_table_cell_borders(cell)

    def _add_editable_coverage_variation_table(
        self,
        slide,
        variation_table: "pd.DataFrame",
        *,
        left,
        top,
        width,
        height,
    ) -> None:
        if variation_table is None or variation_table.empty:
            return

        table_df = variation_table[[col for col in variation_table.columns if not str(col).startswith("_")]].copy()
        if table_df.empty:
            return

        rows = int(len(table_df.index)) + 1
        cols = int(len(table_df.columns))
        if rows <= 1 or cols <= 0:
            return

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        self._clear_powerpoint_table_style(table)

        if cols == 3:
            width_ratios = [0.28, 0.36, 0.36]
            assigned = 0
            for idx, ratio in enumerate(width_ratios):
                if idx == cols - 1:
                    col_w = int(width - assigned)
                else:
                    col_w = int(width * ratio)
                    assigned += col_w
                table.columns[idx].width = col_w
        else:
            col_w = int(width / max(cols, 1))
            for idx in range(cols):
                table.columns[idx].width = col_w

        header_h = int(height * 0.42)
        body_rows = max(rows - 1, 1)
        body_h = int((height - header_h) / body_rows)
        table.rows[0].height = header_h
        for r in range(1, rows):
            if r == rows - 1:
                table.rows[r].height = height - header_h - (body_h * max(rows - 2, 0))
            else:
                table.rows[r].height = body_h

        white_bg = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)

        for c, col_name in enumerate(table_df.columns):
            header_text = "" if str(col_name).strip() == "" else self._normalize_summary_table_value(col_name)
            self._set_table_cell_text(
                table.cell(0, c),
                header_text,
                fill_color=white_bg,
                font_color=black,
                font_size=9,
                bold=True,
                align=2,
                word_wrap=True,
            )

        for r, row_values in enumerate(table_df.itertuples(index=False), start=1):
            for c, raw_val in enumerate(row_values):
                font_color = black if c == 0 else self._summary_variation_font_color(raw_val)
                self._set_table_cell_text(
                    table.cell(r, c),
                    self._normalize_summary_table_value(raw_val),
                    fill_color=white_bg,
                    font_color=font_color,
                    font_size=11,
                    bold=False,
                    align=2,
                    word_wrap=False,
                )

        self._apply_internal_table_borders(
            slide,
            table,
            left=left,
            top=top,
            table_width=width,
            table_height=height,
            color=black,
            border_width=12700,
        )

    def _add_penetration_header_table_shape(
        self,
        slide,
        *,
        left,
        top,
        width,
        height,
        year_header: str,
        pen_header: str,
        rows: Sequence[Tuple[str, str]],
    ) -> None:
        table_shape = slide.shapes.add_table(3, 2, left, top, width, height)
        table_shape.height = height
        table = table_shape.table
        table.columns[0].width = int(width * 0.40)
        table.columns[1].width = width - table.columns[0].width
        table.rows[0].height = int(height * 0.40)
        table.rows[1].height = int((height - table.rows[0].height) / 2)
        table.rows[2].height = height - table.rows[0].height - table.rows[1].height

        header_bg = RGBColor(0, 0, 0)
        body_bg = self._hex_to_rgb("#D9D9D9")
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)

        row_values: List[Tuple[str, str]] = list(rows[:2])
        while len(row_values) < 2:
            row_values.append(("-", "-"))

        self._set_table_cell_text(table.cell(0, 0), year_header, fill_color=header_bg, font_color=white, font_size=10, align=2)
        self._set_table_cell_text(
            table.cell(0, 1),
            pen_header,
            fill_color=header_bg,
            font_color=white,
            font_size=10,
            align=2,
            word_wrap=False,
        )

        self._set_table_cell_text(table.cell(1, 0), row_values[0][0], fill_color=body_bg, font_color=black, align=1)
        self._set_table_cell_text(table.cell(1, 1), row_values[0][1], fill_color=body_bg, font_color=black, align=3)
        self._set_table_cell_text(table.cell(2, 0), row_values[1][0], fill_color=body_bg, font_color=black, align=1)
        self._set_table_cell_text(table.cell(2, 1), row_values[1][1], fill_color=body_bg, font_color=black, align=3)

    def _add_coverage_stability_header_table_shape(
        self,
        slide,
        *,
        left,
        top,
        width,
        height,
        cov_title: str,
        prev_label: str,
        curr_label: str,
        stability_label: str,
        cov_prev_txt: str,
        cov_curr_txt: str,
        stability_txt: str,
    ) -> None:
        table_shape = slide.shapes.add_table(3, 3, left, top, width, height)
        table_shape.height = height
        table = table_shape.table
        table.columns[0].width = int(width * 0.34)
        table.columns[1].width = int(width * 0.34)
        table.columns[2].width = width - table.columns[0].width - table.columns[1].width
        table.rows[0].height = int(height * 0.34)
        table.rows[1].height = int(height * 0.28)
        table.rows[2].height = height - table.rows[0].height - table.rows[1].height

        header_bg = self._hex_to_rgb("#355D6C")
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
        white_bg = RGBColor(255, 255, 255)

        table.cell(0, 0).merge(table.cell(0, 1))
        table.cell(0, 2).merge(table.cell(1, 2))

        self._set_table_cell_text(table.cell(0, 0), cov_title, fill_color=header_bg, font_color=white, font_size=11, align=2)
        self._set_table_cell_text(table.cell(0, 2), stability_label, fill_color=header_bg, font_color=white, font_size=11, align=2)
        self._set_table_cell_text(table.cell(1, 0), prev_label, fill_color=header_bg, font_color=white, font_size=11, align=2)
        self._set_table_cell_text(table.cell(1, 1), curr_label, fill_color=header_bg, font_color=white, font_size=11, align=2)

        self._set_table_cell_text(table.cell(2, 0), cov_prev_txt, fill_color=white_bg, font_color=black, align=2)
        self._set_table_cell_text(table.cell(2, 1), cov_curr_txt, fill_color=white_bg, font_color=black, align=2)
        self._set_table_cell_text(table.cell(2, 2), stability_txt, fill_color=white_bg, font_color=black, align=2)

    def _add_cov_slide_header_boxes(self, slide, assets: PipelineAssets) -> None:
        """Header del slide de Cobertura en modo 'complemented'."""
        try:
            ref_dt = dt.strptime(self.ref_month_year, "%m-%y")
        except Exception:
            # Fallback: usar el último periodo disponible del índice si hay fechas.
            idx = pd.to_datetime(getattr(assets.coverage_series, "index", []), errors="coerce")
            idx = idx[~idx.isna()]
            if len(idx) == 0:
                return
            ref_dt = idx.max().to_pydatetime()

        prev_dt = ref_dt - pd.DateOffset(months=12)
        curr_label = f"{self._month_abbr(ref_dt.month)}-{ref_dt.year % 100:02d}"
        prev_label = f"{self._month_abbr(prev_dt.month)}-{prev_dt.year % 100:02d}"

        # --- Penetración (MAT actual vs anterior) ---
        year_col, pen_col = self._pen_table_headers()
        mat_curr = f"MAT {curr_label}"
        mat_prev = f"MAT {prev_label}"
        pen_curr = assets.penet_mat_actual
        pen_prev = assets.penet_mat_anterior
        pen_curr_txt = f"{float(pen_curr):.1f}" if (pen_curr is not None and pd.notna(pen_curr)) else "-"
        pen_prev_txt = f"{float(pen_prev):.1f}" if (pen_prev is not None and pd.notna(pen_prev)) else "-"

        # --- Cobertura puntual + estabilidad ---
        cov_title = self._coverage_metric_title()
        stability_label = self._stability_label()
        cov_prev = _coverage_value_for_year_month(assets.coverage_series, int(prev_dt.year), int(prev_dt.month))
        cov_curr = _coverage_value_for_year_month(assets.coverage_series, int(ref_dt.year), int(ref_dt.month))

        def _fmt_cov(v: float) -> str:
            if v is None or pd.isna(v):
                return "-"
            return str(int(np.floor(float(v) + 0.5))) if globals().get("ROUND_COVERAGE", False) else f"{float(v):.1f}"

        stability_txt = "-"
        if cov_prev is not None and cov_curr is not None and pd.notna(cov_prev) and pd.notna(cov_curr):
            if globals().get("ROUND_COVERAGE", False):
                stability_txt = str(int(np.floor(float(cov_curr) + 0.5)) - int(np.floor(float(cov_prev) + 0.5)))
            else:
                stability_txt = f"{(float(cov_curr) - float(cov_prev)):.1f}"

        # Construye tablas nativas de PowerPoint para que el texto sea editable.

        # --- Layout superior: dos cajas en una banda encima del gráfico ---
        # Alinear con el inicio/fin del gráfico de coberturas (que arranca en x=0.5in).
        top = Inches(0.95)
        chart_left = Inches(0.5)
        chart_right = self.ppt.slide_width - Inches(0.5)
        shared_h = Inches(0.90)
        # Cuadros mas angostos, manteniendo posiciones originales:
        # penetracion a la izquierda y cobertura a la derecha.
        total_w = chart_right - chart_left
        left_w = int(total_w * 0.27)
        right_w = int(total_w * 0.37)
        right_left = chart_right - right_w
        pen_top = top - Inches(0.03)

        self._add_penetration_header_table_shape(
            slide,
            left=chart_left,
            top=pen_top,
            width=left_w,
            height=shared_h,
            year_header=year_col,
            pen_header=pen_col,
            rows=[(mat_curr, pen_curr_txt), (mat_prev, pen_prev_txt)],
        )
        self._add_coverage_stability_header_table_shape(
            slide,
            left=right_left,
            top=top,
            width=right_w,
            height=shared_h,
            cov_title=cov_title,
            prev_label=prev_label,
            curr_label=curr_label,
            stability_label=stability_label,
            cov_prev_txt=_fmt_cov(cov_prev),
            cov_curr_txt=_fmt_cov(cov_curr),
            stability_txt=stability_txt,
        )

    @staticmethod
    def _date_minus_months(year: int, month: int, delta: int) -> Tuple[int, int]:
        total = year * 12 + (month - 1) - int(delta)
        y2 = total // 12
        m2 = (total % 12) + 1
        return int(y2), int(m2)

    @staticmethod
    def _safe_float(val: object) -> Optional[float]:
        try:
            if val is None or (isinstance(val, str) and val.strip() == "-"):
                return None
            if "pd" in globals() and pd.isna(val):
                return None
            return float(val)
        except Exception:
            return None

    def _fmt_pct(self, val: object) -> str:
        f = self._safe_float(val)
        if f is None:
            return "-"
        return f"{f * 100:.1f}%"

    def _tipo_label(self, tipo: str) -> str:
        t = (tipo or "").strip().lower()
        if t.startswith("an"):
            return "ANO" if self.lang_index != 3 else "YEAR"
        if t.startswith("sem"):
            return "SEMESTRE" if self.lang_index != 3 else "SEMESTER"
        if t.startswith("tri"):
            return "TRIMESTRE" if self.lang_index != 3 else "QUARTER"
        return (tipo or "").strip().upper()

    def _add_footer_text(self, slide: "Presentation", msg: str) -> None:
        if not msg:
            return
        # Centrar el texto en el espacio "util" a la derecha del logo del template.
        logo_clear = Inches(2.00)
        right = Inches(0.35)
        left = logo_clear
        width = self.ppt.slide_width - left - right
        height = Inches(0.35)
        top = self.ppt.slide_height - height - Inches(0.10)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.text = str(msg)
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 0, 0)
        p.alignment = 1

    def _add_low_penetration_footer(self, slide: "Presentation", buyers_value: Optional[float], threshold: float = 200) -> None:
        """Agrega un aviso al pie del slide cuando buyers promedio < threshold."""
        if buyers_value is None:
            return
        try:
            if "pd" in globals() and pd.isna(buyers_value):
                return
            buyers_num = float(buyers_value)
        except Exception:
            return
        if buyers_num >= float(threshold):
            return

        msg = self.labels.get((self.lang_index, "LowPenFooter")) or "Low penetration brand (<200 buyers) - For internal use only"
        self._add_footer_text(slide, msg)

    def _add_variations_box_pretty(
        self,
        slide: "Presentation",
        variations_detail: "pd.DataFrame",
        pipeline: int,
        trend_plot_df: "pd.DataFrame",
        container_left=None,
        container_top=None,
        container_width=None,
        container_height=None,
    ) -> None:
        """Renderiza el cuadro de variaciones en estilo 'bonito' (shapes, no imagen)."""
        if variations_detail is None or variations_detail.empty:
            return

        wp_col = "WP by Numerator" if "WP by Numerator" in variations_detail.columns else None
        # Sem pipeline (P0) siempre se intenta mostrar cuando existe.
        p0_col = "Cliente P0" if "Cliente P0" in variations_detail.columns else ("Cliente Pipeline (P0)" if "Cliente Pipeline (P0)" in variations_detail.columns else None)
        px_col = f"Cliente Pipeline (P{pipeline})" if f"Cliente Pipeline (P{pipeline})" in variations_detail.columns else (f"Cliente P{pipeline}" if f"Cliente P{pipeline}" in variations_detail.columns else None)
        if wp_col is None and p0_col is None and px_col is None:
            return

        show_pipeline_group = int(pipeline) > 0 and px_col is not None

        # Ubicación: si se pasa un contenedor (layout lado derecho), se usa; si no, fallback.
        if container_left is None:
            container_left = Inches(6.8)
        if container_top is None:
            container_top = Inches(1.15)
        if container_width is None:
            container_width = Inches(6.0)
        if container_height is None:
            container_height = Inches(5.8)

        # Row heights dentro del contenedor.
        row_gap = Inches(0.12)
        row_count = max(1, len(variations_detail))
        row_h = int((container_height - ((row_count - 1) * row_gap)) / row_count)
        if row_h <= 0:
            row_h = Inches(0.50)
        if row_count <= 3:
            row_h = min(row_h, Inches(0.56))

        # Intentar derivar el mes base del gráfico (mm-yy) para construir periodos por pipeline.
        base_year = None
        base_month = None
        if trend_plot_df is not None and not trend_plot_df.empty and COL_DATA in trend_plot_df.columns:
            last_token = str(trend_plot_df[COL_DATA].iloc[-1]).strip()
            try:
                mm_s, yy_s = last_token.split("-")
                base_month = int(mm_s)
                base_year = 2000 + int(yy_s)
            except Exception:
                base_year = None
                base_month = None

        # Colores (aprox. al ejemplo)
        green_border = RGBColor(126, 201, 67)  # #7EC943
        sellin_fill = RGBColor(126, 201, 67)
        kantar_fill = RGBColor(58, 58, 58)     # #3A3A3A
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
        grey = RGBColor(120, 120, 120)
        red = RGBColor(208, 2, 27)

        # Columnas (se escalan para llenar el ancho del contenedor).
        def _emu_to_in(v) -> float:
            return float(v) / 914400.0

        base_cols = {
            "tipo": 1.20,
            "var0": 1.10,
            "wp": 1.10,
            "sell0": 1.20,
        }
        if show_pipeline_group:
            base_cols.update({"varp": 1.10, "sellp": 1.25})
        base_total = sum(base_cols.values())
        scale = _emu_to_in(container_width) / base_total if base_total else 1.0

        col_tipo_w = Inches(base_cols["tipo"] * scale)
        col_var0_w = Inches(base_cols["var0"] * scale)
        col_kantar_w = Inches(base_cols["wp"] * scale)
        col_sell0_w = Inches(base_cols["sell0"] * scale)
        col_varp_w = Inches(base_cols.get("varp", 0.0) * scale)
        col_sellp_w = Inches(base_cols.get("sellp", 0.0) * scale)

        rows_block_h = (row_count * row_h) + ((row_count - 1) * row_gap)
        left = container_left
        top = container_top
        if rows_block_h < container_height:
            top = container_top + int((container_height - rows_block_h) / 2)
        total_w = container_width

        def _add_row_box(y):
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, total_w, row_h)
            box.fill.solid()
            box.fill.fore_color.rgb = white
            box.line.color.rgb = green_border
            box.line.width = Pt(1.5)

        def _add_tipo_text(x, y, text):
            tb = slide.shapes.add_textbox(x, y, col_tipo_w, row_h)
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = False
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            self._set_paragraph_text(tf.paragraphs[0], text, font_size=12, font_color=black, bold=True, align=1)

        def _add_var_text(x, y, w, period_text: str):
            tb = slide.shapes.add_textbox(x, y, w, row_h)
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)

            p1 = tf.paragraphs[0]
            # Mantener texto base, ya que el diseño es un "badge" visual.
            p1_text = "VAR %\nMOVEL" if self.lang_index == 1 else ("YOY %\nCHANGE" if self.lang_index == 3 else "VAR %\nMOVIL")
            self._set_paragraph_text(p1, p1_text, font_size=8, font_color=grey, bold=True, align=1)

            p2 = tf.add_paragraph()
            p2.alignment = 1
            # Acepta "MAT Mar-26 x MAT Mar-25", "MAT Mar-26 vs MAT Mar-25" o variantes.
            period_raw = str(period_text or "").strip()
            parts = re.split(r"\s+(?:vs|x)\s+", period_raw, maxsplit=1, flags=re.IGNORECASE)
            if len(parts) == 1:
                parts = re.split(r"\s*(?:vs|x)\s*", period_raw, maxsplit=1, flags=re.IGNORECASE)
            left_txt = parts[0].strip()
            right_txt = parts[1].strip() if len(parts) > 1 else ""
            if self.variations_compact_period_labels:
                left_txt = re.sub(r"(?i)^(?:MAT|SEM|TRI)\s+", "", left_txt).strip()
                right_txt = re.sub(r"(?i)^(?:MAT|SEM|TRI)\s+", "", right_txt).strip()
            r1 = p2.add_run()
            r1.text = f"{left_txt} " if left_txt else ""
            r1.font.size = Pt(8)
            r1.font.color.rgb = black
            if right_txt:
                rvs = p2.add_run()
                rvs.text = "vs"
                rvs.font.size = Pt(8)
                rvs.font.bold = True
                rvs.font.color.rgb = red
                r2 = p2.add_run()
                r2.text = f" {right_txt}"
                r2.font.size = Pt(8)
                r2.font.color.rgb = black

        def _add_value_card(x, y, w, fill_rgb, title, value):
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, row_h)
            card.fill.solid()
            card.fill.fore_color.rgb = fill_rgb
            card.line.color.rgb = fill_rgb
            card.line.width = Pt(0.5)
            tf = card.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Pt(4)
            tf.margin_right = Pt(4)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)

            p1 = tf.paragraphs[0]
            title_font_size = 7 if len(str(title)) >= 11 else 9
            self._set_paragraph_text(p1, title, font_size=title_font_size, font_color=white, bold=True, align=1)
            # "WP by Numerator" es más largo que SELL-IN; bajamos un poco el tamaño.

            p2 = tf.add_paragraph()
            self._set_paragraph_text(p2, value, font_size=20, font_color=white, bold=True, align=1)

        def _period_label(end_year: Optional[int], end_month: Optional[int], offset: int) -> str:
            if end_year is None or end_month is None:
                return "-"
            prev_y, prev_m = self._date_minus_months(int(end_year), int(end_month), int(offset))
            m1 = f"{self._month_abbr(int(end_month))}-{int(end_year) % 100:02d}"
            m2 = f"{self._month_abbr(int(prev_m))}-{int(prev_y) % 100:02d}"
            return f"{m1} vs {m2}"

        # Encabezados de grupo (Sem pipeline / Pipeline p)
        # Se alinean sobre las tarjetas numéricas (no sobre el bloque completo) para que queden centrados.
        header_h = Inches(0.16)
        header_y = top - Inches(0.36)
        # Permite un pequeño offset negativo para subir un poco más sin mover el cuadro completo.
        if header_y < Inches(-0.06):
            header_y = Inches(-0.06)

        def _add_group_header(x: int, w: int, text: str) -> None:
            tb = slide.shapes.add_textbox(x, header_y, w, header_h)
            tf = tb.text_frame
            tf.clear()
            self._set_paragraph_text(tf.paragraphs[0], text, font_size=10, font_color=grey, bold=False, align=1)

        # Sem pipeline: encima del SELL-IN (verde) del sem pipeline (P0), no sobre el bloque completo.
        x_sem = left + col_tipo_w + col_var0_w + col_kantar_w
        w_sem = col_sell0_w
        _add_group_header(x_sem, w_sem, "Sem pipeline" if self.lang_index != 3 else "No pipeline")

        if show_pipeline_group:
            # Pipeline p: encima del SELL-IN del pipeline (no incluye el badge de periodo).
            x_pip = left + col_tipo_w + col_var0_w + col_kantar_w + col_sell0_w + col_varp_w
            w_pip = col_sellp_w
            _add_group_header(x_pip, w_pip, f"Pipeline {int(pipeline)}")

        rows_to_render = variations_detail.reset_index(drop=True)
        for idx, (_, row) in enumerate(rows_to_render.iterrows()):
            y = top + (idx * (row_h + row_gap))
            _add_row_box(y)

            tipo = row.get("Tipo", "")
            wp_val = row.get(wp_col) if wp_col else None
            p0_val = row.get(p0_col) if p0_col else None
            px_val = row.get(px_col) if px_col else None

            # Periodo sem pipeline (p=0) y pipeline p (p=pipeline), usando el mes base del gráfico.
            sem_end_y, sem_end_m = (base_year, base_month)
            if base_year is not None and base_month is not None:
                pip_end_y, pip_end_m = self._date_minus_months(int(base_year), int(base_month), int(pipeline))
            else:
                pip_end_y, pip_end_m = (None, None)
            sem_period = str(row.get("Periodo", "-"))
            compare_lag = row.get("_CompareLagMonths")
            if pd.notna(compare_lag) and pip_end_y is not None and pip_end_m is not None:
                pip_period = _period_label(pip_end_y, pip_end_m, int(compare_lag))
            else:
                pip_period = sem_period

            x = left
            _add_tipo_text(x, y, self._tipo_label(tipo))
            x += col_tipo_w
            _add_var_text(x, y, col_var0_w, sem_period)
            x += col_var0_w
            if wp_col is not None:
                _add_value_card(x, y, col_kantar_w, kantar_fill, "WP by Numerator", self._fmt_pct(wp_val))
            x += col_kantar_w
            if p0_col is not None:
                _add_value_card(x, y, col_sell0_w, sellin_fill, visible_sell_in_label(), self._fmt_pct(p0_val))
            x += col_sell0_w
            if show_pipeline_group:
                _add_var_text(x, y, col_varp_w, pip_period)
                x += col_varp_w
                _add_value_card(x, y, col_sellp_w, sellin_fill, visible_sell_in_label(), self._fmt_pct(px_val))

    # --- Portada -----------------------------------------------------------------
    def configure_cover(self, pais_nombre: str, fabricante: str, categoria_nombre: str, ref_month_year: str, chosen_lang: str) -> None:
        cover_slide = self.ppt.slides[0]
        line1 = f"{pais_nombre} | {fabricante}"
        try:
            ref_dt = dt.strptime(ref_month_year, "%m-%y")
            meses_es = ["", "enero", "febrero", "marzo", "abril", "mayo", "junio", "julio", "agosto", "septiembre", "octubre", "noviembre", "diciembre"]
            meses_pt = ["", "janeiro", "fevereiro", "março", "abril", "maio", "junho", "julho", "agosto", "setembro", "outubro", "novembro", "dezembro"]
            meses_en = ["", "January", "February", "March", "April", "May", "June", "July", "August", "September", "October", "November", "December"]
            if chosen_lang == "PT":
                month_name = meses_pt[ref_dt.month].capitalize()
                line2 = f"{categoria_nombre} - Corte em {month_name} {ref_dt.year}"
            elif chosen_lang == "EN":
                month_name = meses_en[ref_dt.month]
                line2 = f"{categoria_nombre} - As of {month_name} {ref_dt.year}"
            else:
                month_name = meses_es[ref_dt.month].capitalize()
                line2 = f"{categoria_nombre} - Corte a {month_name} {ref_dt.year}"
        except Exception:
            if chosen_lang == "PT":
                line2 = f"{categoria_nombre} - Corte em {ref_month_year}"
            elif chosen_lang == "EN":
                line2 = f"{categoria_nombre} - As of {ref_month_year}"
            else:
                line2 = f"{categoria_nombre} - Corte a {ref_month_year}"
        textbox = cover_slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(2.5))
        text_frame = textbox.text_frame
        text_frame.clear()
        p1 = text_frame.add_paragraph()
        p1.text = line1
        p1.font.size = Pt(44)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.alignment = 1
        p2 = text_frame.add_paragraph()
        p2.text = line2
        p2.font.size = Pt(36)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = 1

    # --- Pipelines ---------------------------------------------------------------
    def add_pipeline_slides(
        self,
        assets: PipelineAssets,
        marca_nombre_limpio: str,
        lang_index: int,
        coverage_label: str,
        progress: Optional["Progress"] = None,
        task_id: Optional[int] = None,
    ) -> int:
        slides_created = 0
        slide_cov = self.ppt.slides.add_slide(self.ppt.slide_layouts[PPT_LAYOUT_INDEX])
        tx_title_cov = ensure_title_frame(slide_cov)
        p_cov = tx_title_cov.paragraphs[0]
        p_cov.text = f"{marca_nombre_limpio} - Pipeline {assets.pipeline}"
        p_cov.font.bold = True
        p_cov.font.size = Pt(24)
        chart_top = Inches(2.0)
        chart_height = Inches(4.2)
        if self.coverage_slide_variant == "pg":
            chart_top = Inches(1.35)
            chart_height = Inches(3.95)
        generar_grafico_cobertura(
            slide_cov,
            marca_nombre_limpio,
            assets.pipeline,
            assets.coverage_series,
            assets.penetration_series,
            lang_index,
            coverage_label,
            self.labels,
            picture_top=chart_top,
            picture_height=chart_height,
        )
        if self.coverage_slide_variant == "complemented":
            try:
                self._add_cov_slide_header_boxes(slide_cov, assets)
            except Exception as exc:
                print(f"{Fore.YELLOW}Advertencia: No se pudo generar el header complementado (penetración/cobertura) para {marca_nombre_limpio} P{assets.pipeline}. Error: {exc}")
                try:
                    self._add_editable_coverage_variation_table(
                        slide_cov,
                        assets.variation_table,
                        left=Inches(0.5),
                        top=Inches(1.1),
                        width=Inches(6.2),
                        height=Inches(0.62),
                    )
                except Exception as exc2:
                    print(f"{Fore.YELLOW}Advertencia: Tampoco se pudo generar la tabla VAR % MAT (fallback) para {marca_nombre_limpio} P{assets.pipeline}. Error: {exc2}")
        elif self.coverage_slide_variant == "pg":
            try:
                self._add_cov_slide_pg_layout(slide_cov, assets)
            except Exception as exc:
                print(f"{Fore.YELLOW}Advertencia: No se pudo generar el layout P&G para {marca_nombre_limpio} P{assets.pipeline}. Error: {exc}")
                try:
                    self._add_editable_coverage_variation_table(
                        slide_cov,
                        assets.variation_table,
                        left=Inches(0.5),
                        top=Inches(1.1),
                        width=Inches(6.2),
                        height=Inches(0.62),
                    )
                except Exception as exc2:
                    print(f"{Fore.YELLOW}Advertencia: Tampoco se pudo generar la tabla VAR % MAT (fallback) para {marca_nombre_limpio} P{assets.pipeline}. Error: {exc2}")
        else:
            try:
                self._add_editable_coverage_variation_table(
                    slide_cov,
                    assets.variation_table,
                    left=Inches(0.5),
                    top=Inches(1.1),
                    width=Inches(6.2),
                    height=Inches(0.62),
                )
            except Exception as exc:
                print(f"{Fore.YELLOW}Advertencia: No se pudo generar la tabla de variación MAT para {marca_nombre_limpio} P{assets.pipeline}. Error: {exc}")
        self._add_low_penetration_footer(slide_cov, getattr(assets, "buyers_mat_actual", None))
        slides_created += 1
        if progress and task_id is not None:
            progress.update(task_id, advance=1)
        slide_trend = self.ppt.slides.add_slide(self.ppt.slide_layouts[PPT_LAYOUT_INDEX])
        tx_title_trend = ensure_title_frame(slide_trend)
        p_trend = tx_title_trend.paragraphs[0]
        p_trend.text = f"{marca_nombre_limpio} - Pipeline {assets.pipeline}"
        p_trend.font.bold = True
        p_trend.font.size = Pt(24)
        has_variations = assets.variations_detail is not None and not assets.variations_detail.empty
        if self.variations_box_style == "pretty" and has_variations:
            # Layout: gráfico a la izquierda + cuadro bonito a la derecha, mismo alto.
            content_top = Inches(1.15)
            content_bottom = Inches(0.55)
            content_h = self.ppt.slide_height - content_top - content_bottom
            margin_l = Inches(0.35)
            margin_r = Inches(0.25)
            divider_w = Inches(0.06)
            avail_w = self.ppt.slide_width - margin_l - margin_r - divider_w
            left_w = int(avail_w / 2)
            compact_shift = Inches(0.18) if self.variations_compact_period_labels else 0
            chart_left = max(Inches(0.18), margin_l - (Inches(0.12) if self.variations_compact_period_labels else 0))
            chart_top = content_top
            var_top = content_top

            # Divisor vertical (como referencia)
            divider_x = margin_l + left_w - compact_shift
            chart_w = max(Inches(5.8), divider_x - chart_left)
            divider = slide_trend.shapes.add_shape(MSO_SHAPE.RECTANGLE, divider_x, content_top, divider_w, content_h)
            divider.fill.solid()
            # Azul verdoso (segun referencia)
            divider.fill.fore_color.rgb = RGBColor(0, 229, 176)  # #00E5B0
            divider.line.fill.background()

            generar_grafico_tendencia(
                slide_trend,
                marca_nombre_limpio,
                assets.pipeline,
                assets.trend_plot_df,
                lang_index,
                self.labels,
                doble_eje=(self.tipo_eje_tend == "doble"),
                granularity=self.trend_granularity,
                box_left=chart_left,
                box_top=chart_top,
                box_width=chart_w,
                box_height=content_h,
                # Imagen más "alta" para que se aproveche mejor la columna izquierda sin estirar.
                figsize=(9.0, 7.0),
                legend_y=-0.22,
            )
        else:
            generar_grafico_tendencia(
                slide_trend,
                marca_nombre_limpio,
                assets.pipeline,
                assets.trend_plot_df,
                lang_index,
                self.labels,
                doble_eje=(self.tipo_eje_tend == "doble"),
                granularity=self.trend_granularity,
            )
        if has_variations:
            if self.variations_box_style == "pretty":
                # Usa el mismo contenedor del lado derecho que el layout del gráfico.
                content_top = Inches(1.15)
                content_bottom = Inches(0.55)
                content_h = self.ppt.slide_height - content_top - content_bottom
                margin_l = Inches(0.35)
                margin_r = Inches(0.25)
                divider_w = Inches(0.06)
                avail_w = self.ppt.slide_width - margin_l - margin_r - divider_w
                left_w = int(avail_w / 2)
                compact_shift = Inches(0.18) if self.variations_compact_period_labels else 0
                divider_x = margin_l + left_w - compact_shift
                var_left = margin_l + left_w + divider_w
                if self.variations_compact_period_labels:
                    var_left = divider_x + divider_w + Inches(0.08)
                right_w = self.ppt.slide_width - var_left - margin_r
                # Reducir altura a la mitad y centrar verticalmente en el área de contenido.
                var_h = int(content_h * 0.5)
                var_top = content_top + int((content_h - var_h) / 2)
                self._add_variations_box_pretty(
                    slide_trend,
                    assets.variations_detail,
                    assets.pipeline,
                    assets.trend_plot_df,
                    container_left=var_left,
                    container_top=var_top,
                    container_width=right_w,
                    container_height=var_h,
                )
            else:
                table_width = min(int(Inches(6.2)), int(self.ppt.slide_width - Inches(0.6)))
                right_margin = Inches(0.3)
                left_pos = self.ppt.slide_width - table_width - right_margin
                if left_pos < Inches(0.1):
                    left_pos = Inches(0.1)
                top_pos = Inches(0.22)
                self._add_editable_variations_table(
                    slide_trend,
                    assets.variations_detail,
                    left=left_pos,
                    top=top_pos,
                    width=table_width,
                    max_height=Inches(1.15),
                )
        self._add_low_penetration_footer(slide_trend, getattr(assets, "buyers_mat_actual", None))
        slides_created += 1
        if progress and task_id is not None:
            progress.update(task_id, advance=1)
        if assets.evolution_figure is not None:
            slide_evol = self.ppt.slides.add_slide(self.ppt.slide_layouts[PPT_LAYOUT_INDEX])
            tx_title_evol = ensure_title_frame(slide_evol)
            p_evol = tx_title_evol.paragraphs[0]
            if lang_index == 3:
                p_evol.text = f"{marca_nombre_limpio} - Pipeline {assets.pipeline} - Monthly Evolution and YoY Variation"
            elif lang_index == 1:
                p_evol.text = f"{marca_nombre_limpio} - Pipeline {assets.pipeline} - Evolução Mensal e Variação"
            else:
                p_evol.text = f"{marca_nombre_limpio} - Pipeline {assets.pipeline} - Evolución Mensual y Variación"
            p_evol.font.bold = True
            p_evol.font.size = Pt(24)
            buffer = io.BytesIO()
            assets.evolution_figure.savefig(buffer, format="png", dpi=240, bbox_inches="tight", pad_inches=0.08, transparent=True)
            plt.close(assets.evolution_figure)
            buffer.seek(0)
            left = Inches(0.1)
            usable_w = self.ppt.slide_width - 2 * left
            slide_evol.shapes.add_picture(buffer, left, Inches(1.0), width=usable_w)
            self._add_low_penetration_footer(slide_evol, getattr(assets, "buyers_mat_actual", None))
            slides_created += 1
            if progress and task_id is not None:
                progress.update(task_id, advance=1)
        return slides_created

    # --- Resumen -----------------------------------------------------------------
    def add_summary_slide(
        self,
        df_summary: "pd.DataFrame",
        pais_nombre: str,
        categoria_nombre: str,
        low_penetration_brands: Optional[Sequence[str]] = None,
        summary_groups: Optional[Sequence[Tuple[str, "pd.DataFrame"]]] = None,
        df_bank: Optional["pd.DataFrame"] = None,
    ) -> None:
        slide_summary = self.ppt.slides.add_slide(self.ppt.slide_layouts[PPT_LAYOUT_INDEX])
        title_frame = ensure_title_frame(slide_summary)
        p = title_frame.paragraphs[0]
        p.text = f"Summary - {pais_nombre} {categoria_nombre} - {self.coverage_label}"
        p.font.bold = True
        p.font.size = Pt(22 if self.coverage_slide_variant == "pg" else 26)
        tx_s1 = slide_summary.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        s1_frame = tx_s1.text_frame
        s1_frame.text = self.labels.get((self.lang_index, "S1"), "")
        comentarios_box = slide_summary.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(8.5), Inches(0.7))
        comentarios_frame = comentarios_box.text_frame
        comentarios_frame.word_wrap = True
        comentarios_frame.auto_size = True
        comentarios_frame.text = "Comentarios:"

        summary_groups = [
            (period, group_df)
            for period, group_df in (summary_groups or [])
            if group_df is not None and not group_df.empty
        ]
        if not summary_groups and df_summary is not None and not df_summary.empty:
            summary_groups = [(self.ref_month_year, df_summary)]
        if not summary_groups:
            print(f"{Fore.YELLOW}Advertencia: No hay datos para generar la tabla de resumen en el PPT.")
            return

        low_penetration_brands = list(low_penetration_brands or [])

        def _summary_period_title(period_token: str) -> str:
            try:
                period_dt = dt.strptime(str(period_token), "%m-%y")
                month_txt = period_dt.strftime("%b-%y")
            except Exception:
                month_txt = str(period_token)
            if self.lang_index == 3:
                return f"As of {month_txt}"
            if self.lang_index == 1:
                return f"Corte {month_txt}"
            return f"Corte {month_txt}"

        def _add_group_heading(text: str, left, top, width) -> None:
            tb = slide_summary.shapes.add_textbox(left, top, width, Inches(0.22))
            tf = tb.text_frame
            tf.clear()
            p_head = tf.paragraphs[0]
            p_head.text = text
            p_head.font.size = Pt(11)
            p_head.font.bold = True
            p_head.font.color.rgb = RGBColor(0, 0, 0)
            p_head.alignment = 1

        try:
            left = Inches(0.5)
            top = Inches(1.20)
            usable_w = self.ppt.slide_width - 2 * left
            total_h = Inches(4.6)
            if self.coverage_slide_variant == "pg" and df_bank is not None and not df_bank.empty:
                label_box = slide_summary.shapes.add_textbox(Inches(0.18), top + Inches(0.78), Inches(1.22), Inches(0.34))
                label_tf = label_box.text_frame
                label_tf.clear()
                label_tf.word_wrap = True
                label_p = label_tf.paragraphs[0]
                label_p.text = str(categoria_nombre or "").strip()
                label_p.font.size = Pt(12)
                label_p.font.bold = True
                label_p.font.color.rgb = RGBColor(35, 35, 35)
                label_p.alignment = 1
                pg_left = Inches(1.45)
                self._add_pg_summary_table(
                    slide_summary,
                    df_bank,
                    left=pg_left,
                    top=top,
                    width=self.ppt.slide_width - pg_left - Inches(0.35),
                    max_height=total_h,
                )
            elif len(summary_groups) == 1:
                self._add_editable_summary_table(
                    slide_summary,
                    summary_groups[0][1],
                    left=left,
                    top=top,
                    width=usable_w,
                    max_height=total_h,
                    low_penetration_brands=low_penetration_brands,
                )
            else:
                heading_h = Inches(0.22)
                heading_table_gap_h = Inches(0.13)
                gap_h = Inches(0.25)
                available_table_h = int(
                    total_h
                    - (len(summary_groups) * (heading_h + heading_table_gap_h))
                    - ((len(summary_groups) - 1) * gap_h)
                )
                total_body_rows = sum(max(1, len(group_df.index)) for _, group_df in summary_groups)
                y = top
                for idx, (period_token, group_df) in enumerate(summary_groups):
                    group_rows = max(1, len(group_df.index))
                    if idx == len(summary_groups) - 1:
                        table_h = int((top + total_h) - y - heading_h)
                    else:
                        table_h = max(Inches(0.95), int(available_table_h * (group_rows / max(total_body_rows, 1))))
                    _add_group_heading(_summary_period_title(period_token), left, y, usable_w)
                    y += heading_h + heading_table_gap_h
                    used_h = self._add_editable_summary_table(
                        slide_summary,
                        group_df,
                        left=left,
                        top=y,
                        width=usable_w,
                        max_height=table_h,
                        low_penetration_brands=low_penetration_brands,
                    ) or table_h
                    y += int(used_h) + gap_h
        except Exception as exc:
            print(f"{Fore.YELLOW}Advertencia: No se pudo generar la tabla resumen en el PPT. Error: {exc}")
        if low_penetration_brands:
            unique_brands = sorted({str(b).strip() for b in low_penetration_brands if str(b).strip()})
            n = len(unique_brands)
            key = "LowPenSummaryPlural" if n > 1 else "LowPenSummarySingular"
            tpl = self.labels.get((self.lang_index, key))
            if not tpl:
                tpl = self.labels.get((self.lang_index, "LowPenSummaryPlural" if n > 1 else "LowPenSummarySingular"))
            if not tpl:
                tpl = "This study contains {n} low penetration brand(s) (<200 buyers). For internal use only"
            msg = str(tpl).format(n=n)
            self._add_footer_text(slide_summary, msg)

    # --- Post-procesamiento -------------------------------------------------------
    def insert_thanks_text(self, chosen_lang: str) -> None:
        thanks_map = {"ES": "Gracias", "PT": "Obrigado(a)", "EN": "Thank you"}
        thanks_txt = thanks_map.get(chosen_lang, "Gracias")
        if len(self.ppt.slides) <= 6:
            return
        slide7 = self.ppt.slides[6]
        tb = slide7.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(2.5))
        tf7 = tb.text_frame
        tf7.clear()
        p = tf7.add_paragraph()
        p.text = thanks_txt
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = 1

    def reorder_summary_and_credit(self) -> None:
        if len(self.ppt.slides) > 1:
            summary_slide_xml = self.ppt.slides._sldIdLst[-1]
            insert_idx = 7 if len(self.ppt.slides) > 7 else len(self.ppt.slides) - 1
            self.ppt.slides._sldIdLst.insert(insert_idx, summary_slide_xml)
        if len(self.ppt.slides) > 7:
            credit_slide_xml = self.ppt.slides._sldIdLst[6]
            self.ppt.slides._sldIdLst.append(credit_slide_xml)


def parse_file_metadata(excel_file_name: str, categories_df: "pd.DataFrame") -> Tuple[str, str, str, str, str]:
    """Obtiene país, cesta y categoría a partir del nombre del archivo."""
    country_code_str, category_code, fabricante = parse_input_filename_parts(excel_file_name)
    try:
        country_code = int(country_code_str)
    except ValueError as exc:
        raise ValueError(f"El código de país '{country_code_str}' no es numérico") from exc
    try:
        pais_nombre = str(pais.loc[pais.cod == country_code, 'pais'].iloc[0]).strip()
    except Exception as exc:
        raise ValueError(f"No se encontró el país para el código {country_code_str}") from exc
    cesta_nombre, categoria_nombre, categoria_corta = _lookup_category_metadata(category_code, categories_df)
    return pais_nombre, cesta_nombre, categoria_nombre, categoria_corta, fabricante


def ensure_output_folder(root_dir: str, nombre_base_archivo: str) -> str:
    carpeta_salida = os.path.join(root_dir, nombre_base_archivo)
    if not os.path.exists(carpeta_salida):
        os.makedirs(carpeta_salida, exist_ok=True)
    return carpeta_salida



def copy_and_prune_template(root_dir: str, chosen_lang: str) -> Tuple["Presentation", str]:
    """Copia la plantilla base, elimina slides según idioma y devuelve la presentación lista."""
    run_id = os.environ.get('RUN_ID') or datetime.now().strftime('%Y%m%d_%H%M%S')
    tmp_dir = os.path.join(root_dir, 'tmp')
    os.makedirs(tmp_dir, exist_ok=True)
    src_template_path = os.path.join(root_dir, 'Modelo_PPT.pptx')
    if not os.path.exists(src_template_path):
        raise FileNotFoundError(f"No se encontró la plantilla base: {src_template_path}")
    tmp_ppt_name = f"Modelo_PPT_{run_id}_{chosen_lang}.pptx"
    tmp_ppt_path = os.path.join(tmp_dir, tmp_ppt_name)
    shutil.copyfile(src_template_path, tmp_ppt_path)
    ppt = Presentation(tmp_ppt_path)
    keep_indices_by_lang = {
        'ES': {0, 1, 2, 3, 4, 5, 16},
        'PT': {0, 6, 7, 8, 9, 10, 16},
        'EN': {0, 11, 12, 13, 14, 15, 16},
    }
    keep_set = keep_indices_by_lang.get(chosen_lang, keep_indices_by_lang['ES'])
    total_initial = len(ppt.slides)
    delete_list = sorted([i for i in range(total_initial) if i not in keep_set], reverse=True)
    for di in delete_list:
        _delete_slide(ppt, di)
    ppt.save(tmp_ppt_path)
    return Presentation(tmp_ppt_path), tmp_ppt_path


def _delete_slide(pres_obj: "Presentation", idx: int) -> None:
    """Elimina un slide usando la API protegida de python-pptx."""
    sldIdLst = pres_obj.slides._sldIdLst  # type: ignore[attr-defined]
    sldId = sldIdLst[idx]
    rId = sldId.rId
    pres_obj.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def _build_sheet_header_index(ws: "object") -> Dict[str, int]:
    """Devuelve un índice {header: columna} a partir de la fila 1."""
    header_index: Dict[str, int] = {}
    for col in range(1, ws.max_column + 1):
        raw_value = ws.cell(row=1, column=col).value
        if raw_value is None:
            continue
        header = str(raw_value).strip()
        if header and header not in header_index:
            header_index[header] = col
    return header_index


def _find_last_data_row(ws: "object", data_col: int, start_row: int = 2) -> int:
    """Encuentra la última fila contigua con datos en la columna de fecha."""
    row = start_row
    last_valid = start_row - 1
    while row <= ws.max_row:
        value = ws.cell(row=row, column=data_col).value
        if value is None or str(value).strip() == "":
            break
        last_valid = row
        row += 1
    return last_valid


def _find_first_nonempty_row(ws: "object", col: int, start_row: int, end_row: int) -> Optional[int]:
    """Devuelve la primera fila con valor no vacío en un rango vertical."""
    for row in range(start_row, end_row + 1):
        value = ws.cell(row=row, column=col).value
        if value is None:
            continue
        if isinstance(value, str) and value.strip() == "":
            continue
        return row
    return None


def autofit_worksheet_columns(
    ws: "object",
    *,
    min_width: float = 10.0,
    max_width: float = 36.0,
    padding: float = 2.0,
) -> None:
    """Ajusta el ancho de columnas segun encabezado y contenido visible."""
    for col_cells in ws.iter_cols(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
        max_len = 0
        col_letter = None
        for cell in col_cells:
            if col_letter is None:
                col_letter = cell.column_letter
            value = cell.value
            if value is None:
                continue
            if isinstance(value, datetime):
                display_text = value.strftime("%b-%y")
            else:
                display_text = str(value).strip()
            if display_text.startswith("="):
                # En templates con formulas, el texto de la formula no representa
                # el ancho visible del valor calculado.
                display_text = "0000.0%"
            if not display_text:
                continue
            max_len = max(max_len, len(display_text))
        if not col_letter:
            continue
        target_width = min(max_width, max(min_width, float(max_len) + padding))
        ws.column_dimensions[col_letter].width = target_width


def _excel_lang_code(include_english: bool, pais_nombre: str) -> str:
    if include_english:
        return "EN"
    return "PT" if (pais_nombre or "").strip().lower() in {"brasil", "brazil"} else "ES"


def _parse_pipeline_from_sheet_name(sheet_name: str) -> int:
    match = re.match(r"(?i)^p([0-6])_", str(sheet_name or "").strip())
    return int(match.group(1)) if match else 0


def _clean_brand_name_from_sheet(sheet_name: str) -> str:
    cleaned = re.sub(r"(?i)^p[0-6]_", "", str(sheet_name or "")).strip()
    return cleaned or str(sheet_name or "N/D")


def _safe_hex(color_value: str) -> str:
    return str(color_value or "").strip().replace("#", "")[:6] or "000000"


def _set_line_series_color(series_obj: "object", color_value: str, width: int = 28575) -> None:
    color_hex = _safe_hex(color_value)
    try:
        series_obj.graphicalProperties.line.solidFill = color_hex
        series_obj.graphicalProperties.line.width = width
    except Exception:
        pass
    try:
        series_obj.graphicalProperties.solidFill = color_hex
    except Exception:
        pass


def _set_bar_series_color(series_obj: "object", color_value: str, line_color: str = "000000") -> None:
    fill_hex = _safe_hex(color_value)
    line_hex = _safe_hex(line_color)
    try:
        series_obj.graphicalProperties.solidFill = fill_hex
    except Exception:
        pass
    try:
        series_obj.graphicalProperties.line.solidFill = line_hex
    except Exception:
        pass


def _find_excel_header_col(headers: Dict[str, int], canonical_header: str) -> Optional[int]:
    """Busca una columna por encabezado interno o por su etiqueta visible."""
    direct = headers.get(canonical_header)
    if direct is not None:
        return direct
    visible = VISIBLE_EXCEL_HEADER_MAP.get(canonical_header)
    if visible:
        return headers.get(visible)
    return None


def add_native_excel_charts(
    xlsx_path: str,
    *,
    coverage_label: str,
    trend_axis: str,
    evolution_slide_variant: str,
    include_english: bool,
    pais_nombre: str,
) -> None:
    """
    Inserta graficos nativos de Excel (editables) en cada hoja de marca,
    replicando los 3 graficos principales del flujo PPT:
    - Cobertura vs penetracion mensual (barras).
    - Tendencia de volumen Sell-in vs Sell-out (lineas).
    - Evolucion mensual y variacion interanual (simple o clasico).
    """
    from openpyxl import load_workbook as _load_wb_chart
    from openpyxl.chart import (
        BarChart as _BarChart,
        LineChart as _LineChart,
        Reference as _Reference,
    )
    from openpyxl.chart.label import DataLabelList as _DataLabelList
    from openpyxl.chart.axis import DisplayUnitsLabelList as _DisplayUnitsLabelList
    from openpyxl.chart.series import SeriesLabel as _SeriesLabel
    from openpyxl.chart.shapes import GraphicalProperties as _GraphicalProperties
    from openpyxl.utils import get_column_letter as _get_col_letter

    lang_code = _excel_lang_code(include_english, pais_nombre)
    lang_index = {"PT": 1, "ES": 2, "EN": 3}[lang_code]
    trend_axis_norm = str(trend_axis or "").strip().lower()
    evolution_variant_norm = normalize_evolution_slide_variant(evolution_slide_variant)

    def _numeric_column_values(ws, col: int, start_row: int, end_row: int) -> List[object]:
        values: List[object] = []
        for row in range(start_row, end_row + 1):
            value = ws.cell(row=row, column=col).value
            if isinstance(value, str) and re.fullmatch(r"=[A-Z]{1,3}[1-9][0-9]*", value.strip(), re.IGNORECASE):
                value = ws[value.strip()[1:]].value
            values.append(value)
        return values

    def _format_trend_excel_axis(axis, title: str, values: Iterable[object]) -> None:
        exponent = trend_axis_magnitude_exponent(values)
        axis.title = trend_axis_title(title, exponent, lang_index)
        if exponent:
            axis.dispUnits = _DisplayUnitsLabelList(custUnit=float(10 ** exponent))
            abbreviation = trend_axis_magnitude_abbreviation(exponent)
            axis.numFmt = f'0.##"{abbreviation}"'
        else:
            axis.numFmt = "#,##0.##"

    chart_titles = {
        "ES": {
            "coverage_title": "Cobertura en Año Móvil",
            "penetration_label": "Penetracion Mensual",
            "trend_title": "Tendencia en Volumen",
            "evolution_title": "Evolucion Mensual y Variacion",
            "evolution_var_axis": "Variacion Interanual",
            "evolution_monthly_axis": "Volumen Mensual",
        },
        "PT": {
            "coverage_title": "Cobertura em Ano Movel",
            "penetration_label": "Penetracao Mensal",
            "trend_title": "Tendencia em Volumen",
            "evolution_title": "Evolucao Mensal e Variacao",
            "evolution_var_axis": "Variacao Interanual",
            "evolution_monthly_axis": "Volumen Mensual",
        },
        "EN": {
            "coverage_title": "MOVING YEAR COVERAGE",
            "penetration_label": "PENETRATION BY PERIOD",
            "trend_title": "TREND IN VOLUME",
            "evolution_title": "Monthly Evolution and YoY Variation",
            "evolution_var_axis": "YoY Variation",
            "evolution_monthly_axis": "Monthly Volume",
        },
    }[lang_code]
    chart_scale = 1.2
    chart_anchor_col = "AA"

    def _tint_hex_color(color_value: str, mix_with_white: float = 0.78) -> str:
        """Aclara un color HEX mezclándolo con blanco."""
        hex_color = _safe_hex(color_value)
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
        except Exception:
            return "E7E6E6"
        m = max(0.0, min(1.0, float(mix_with_white)))
        r2 = int(round(r + (255 - r) * m))
        g2 = int(round(g + (255 - g) * m))
        b2 = int(round(b + (255 - b) * m))
        return f"{r2:02X}{g2:02X}{b2:02X}"

    def _apply_variation_labels(series_obj: "object", line_color: str) -> None:
        """Muestra valor puntual con fondo difuminado del color de línea y color por signo."""
        dlabels = _DataLabelList()
        dlabels.showVal = True
        dlabels.showSerName = False
        dlabels.showCatName = False
        dlabels.showLegendKey = False
        dlabels.showPercent = False
        dlabels.separator = " "
        # Color de fuente por signo (Excel evalúa el formato en tiempo de cálculo).
        dlabels.numFmt = "[Green]0.0%;[Red]-0.0%;0.0%"
        series_obj.dLbls = dlabels
        # Fondo difuminado con color de la línea de la serie.
        try:
            dlabels.spPr = _GraphicalProperties(solidFill=_tint_hex_color(line_color, mix_with_white=0.78))
            if getattr(dlabels.spPr, "line", None) is not None:
                dlabels.spPr.line.solidFill = _safe_hex(line_color)
        except Exception:
            pass

    wb = _load_wb_chart(xlsx_path)
    wb.calculation.calcMode = "auto"
    wb.calculation.fullCalcOnLoad = True
    wb.calculation.forceFullCalc = True

    for ws in wb.worksheets:
        headers = _build_sheet_header_index(ws)
        pipeline = _parse_pipeline_from_sheet_name(ws.title)
        cov_header = f"P{pipeline}"

        data_col = _find_excel_header_col(headers, COL_DATA)
        sell_out_col = _find_excel_header_col(headers, COL_SELL_OUT)
        sell_in_sim_col = _find_excel_header_col(headers, COL_SELL_IN_SIM)
        pen_col = _find_excel_header_col(headers, COL_PENET)
        cov_col = headers.get(cov_header)
        evo_kantar_col = _find_excel_header_col(headers, COL_EVO_KANTAR_YOY)
        evo_sellin_col = _find_excel_header_col(headers, COL_EVO_SELLIN_YOY)
        if any(col is None for col in (data_col, sell_out_col, sell_in_sim_col, pen_col, cov_col)):
            continue
        data_col = int(data_col)
        sell_out_col = int(sell_out_col)
        sell_in_sim_col = int(sell_in_sim_col)
        pen_col = int(pen_col)
        cov_col = int(cov_col)

        last_data_row = _find_last_data_row(ws, data_col=data_col, start_row=2)
        if last_data_row < 3:
            continue
        n_data_rows = last_data_row - 1
        if n_data_rows < 12:
            continue

        # Evita duplicados sin borrar graficos ajenos a esta rutina.
        def _chart_anchor_cell(ch: "object") -> str:
            try:
                anchor = ch.anchor
                if isinstance(anchor, str):
                    return anchor.upper()
                if hasattr(anchor, "_from"):
                    return f"{_get_col_letter(anchor._from.col + 1)}{anchor._from.row + 1}"
            except Exception:
                return ""
            return ""

        if hasattr(ws, "_charts"):
            target_anchors = {
                f"{chart_anchor_col}2",
                f"{chart_anchor_col}22",
                f"{chart_anchor_col}42",
                # Limpieza de anclas anteriores para evitar duplicados al regenerar.
                "W2",
                "W22",
                "W42",
            }
            ws._charts = [c for c in ws._charts if _chart_anchor_cell(c) not in target_anchors]  # type: ignore[attr-defined]

        brand_name = _clean_brand_name_from_sheet(ws.title)
        trend_start = 2 + pipeline + EXCEL_TREND_INITIAL_GAP_MONTHS
        trend_end = last_data_row
        sell_in_start = 2 + EXCEL_TREND_INITIAL_GAP_MONTHS
        sell_in_end = last_data_row - pipeline

        # 1) Cobertura vs Penetracion (rangos directos de columnas originales).
        cov_start = _find_first_nonempty_row(ws, cov_col, start_row=2, end_row=last_data_row)
        if cov_start is not None and cov_start <= last_data_row:
            coverage_chart = _BarChart()
            coverage_chart.type = "col"
            coverage_chart.grouping = "clustered"
            coverage_chart.overlap = 0
            coverage_chart.gapWidth = 85
            coverage_chart.height = 7.1 * chart_scale
            coverage_chart.width = 16.2 * chart_scale
            coverage_chart.title = f"{chart_titles['coverage_title']} | {brand_name} Pipeline {pipeline}"
            coverage_chart.y_axis.title = f"{coverage_label} | {chart_titles['penetration_label']}"
            coverage_chart.y_axis.scaling.min = 0
            coverage_chart.y_axis.numFmt = "0.0"
            coverage_chart.x_axis.number_format = "yyyy-mm"
            coverage_chart.x_axis.numFmt = "yyyy-mm"
            coverage_chart.x_axis.tickLblPos = "low"
            coverage_chart.x_axis.tickLblSkip = 1
            coverage_chart.x_axis.tickMarkSkip = 1
            coverage_chart.x_axis.delete = False
            coverage_chart.legend.position = "b"
            coverage_chart.legend.overlay = False

            coverage_chart.add_data(
                _Reference(ws, min_col=pen_col, min_row=cov_start, max_row=last_data_row),
                titles_from_data=False,
            )
            coverage_chart.series[-1].title = _SeriesLabel(v=chart_titles["penetration_label"])
            coverage_chart.add_data(
                _Reference(ws, min_col=cov_col, min_row=cov_start, max_row=last_data_row),
                titles_from_data=False,
            )
            coverage_chart.series[-1].title = _SeriesLabel(v=coverage_label)
            coverage_chart.set_categories(
                _Reference(ws, min_col=data_col, min_row=cov_start, max_row=last_data_row)
            )
            coverage_chart.dataLabels = _DataLabelList()
            coverage_chart.dataLabels.showVal = True
            coverage_chart.dataLabels.showSerName = False
            coverage_chart.dataLabels.showCatName = False
            coverage_chart.dataLabels.showLegendKey = False
            coverage_chart.dataLabels.showPercent = False
            coverage_chart.dataLabels.numFmt = "0.0"
            coverage_chart.dataLabels.separator = " "
            _set_bar_series_color(coverage_chart.series[0], COLOR_PENETRACION_BAR)
            _set_bar_series_color(coverage_chart.series[1], COLOR_COBERTURA_BAR)
            ws.add_chart(coverage_chart, f"{chart_anchor_col}2")

        # 2) Tendencia (rangos directos de columnas reales; omite los primeros 6 meses).
        if trend_start <= trend_end and sell_in_start <= sell_in_end:
            trend_categories = _Reference(
                ws,
                min_col=data_col,
                min_row=trend_start,
                max_row=trend_end,
            )
            trend_chart = _LineChart()
            trend_chart.style = 2
            trend_chart.height = 7.1 * chart_scale
            trend_chart.width = 16.2 * chart_scale
            trend_chart.title = f"{chart_titles['trend_title']} | {brand_name} P:{pipeline}"
            trend_chart.x_axis.number_format = "yyyy-mm"
            trend_chart.x_axis.numFmt = "yyyy-mm"
            trend_chart.x_axis.tickLblPos = "low"
            trend_chart.x_axis.tickLblSkip = 1
            trend_chart.x_axis.tickMarkSkip = 1
            trend_chart.x_axis.delete = False
            trend_chart.legend.position = "b"
            trend_chart.legend.overlay = False
            trend_chart.y_axis.scaling.min = 0

            if trend_axis_norm == "doble":
                sell_in_values = _numeric_column_values(ws, sell_in_sim_col, sell_in_start, sell_in_end)
                _format_trend_excel_axis(trend_chart.y_axis, visible_sell_in_label(), sell_in_values)
                trend_chart.add_data(
                    _Reference(ws, min_col=sell_in_sim_col, min_row=sell_in_start, max_row=sell_in_end),
                    titles_from_data=False,
                )
                trend_chart.series[-1].title = _SeriesLabel(v=f"{visible_sell_in_label()} (P:{pipeline})")
                trend_chart.set_categories(trend_categories)

                trend_chart2 = _LineChart()
                trend_chart2.y_axis.axId = 200
                trend_chart2.y_axis.crosses = "max"
                sell_out_values = _numeric_column_values(ws, sell_out_col, trend_start, trend_end)
                _format_trend_excel_axis(
                    trend_chart2.y_axis,
                    short_visible_sell_out_axis_label(lang_index),
                    sell_out_values,
                )
                trend_chart2.add_data(
                    _Reference(ws, min_col=sell_out_col, min_row=trend_start, max_row=trend_end),
                    titles_from_data=False,
                )
                trend_chart2.series[-1].title = _SeriesLabel(v=visible_sell_out_label(lang_index))
                trend_chart += trend_chart2
            else:
                sell_in_values = _numeric_column_values(ws, sell_in_sim_col, sell_in_start, sell_in_end)
                sell_out_values = _numeric_column_values(ws, sell_out_col, trend_start, trend_end)
                _format_trend_excel_axis(
                    trend_chart.y_axis,
                    f"{visible_sell_in_label()} / {short_visible_sell_out_axis_label(lang_index)}",
                    [*sell_in_values, *sell_out_values],
                )
                trend_chart.add_data(
                    _Reference(ws, min_col=sell_in_sim_col, min_row=sell_in_start, max_row=sell_in_end),
                    titles_from_data=False,
                )
                trend_chart.series[-1].title = _SeriesLabel(v=f"{visible_sell_in_label()} (P:{pipeline})")
                trend_chart.add_data(
                    _Reference(ws, min_col=sell_out_col, min_row=trend_start, max_row=trend_end),
                    titles_from_data=False,
                )
                trend_chart.series[-1].title = _SeriesLabel(v=visible_sell_out_label(lang_index))
                trend_chart.set_categories(trend_categories)

            if len(trend_chart.series) >= 1:
                _set_line_series_color(trend_chart.series[0], COLOR_SELLIN_TREND_LINE)
            if len(trend_chart.series) >= 2:
                _set_line_series_color(trend_chart.series[1], COLOR_SELLOUT_TREND_LINE)
            ws.add_chart(trend_chart, f"{chart_anchor_col}22")

        # 3) Evolucion mensual y variacion interanual (nutrida por columnas V/W del Excel).
        if (
            n_data_rows >= 24
            and trend_start <= trend_end
            and sell_in_start <= sell_in_end
            and evo_kantar_col is not None
            and evo_sellin_col is not None
        ):
            evo_start_k = _find_first_nonempty_row(ws, evo_kantar_col, start_row=2, end_row=last_data_row)
            evo_start_s = _find_first_nonempty_row(ws, evo_sellin_col, start_row=2, end_row=last_data_row)
            evo_start_s_shifted = (evo_start_s + pipeline) if evo_start_s is not None else None
            evo_candidates = [r for r in (evo_start_k, evo_start_s_shifted) if r is not None]
            if not evo_candidates:
                continue
            evo_start = max(evo_candidates)
            if evo_start > last_data_row:
                continue

            sellin_var_start = max(2, evo_start - pipeline)
            sellin_var_end = max(2, last_data_row - pipeline)
            if sellin_var_start > sellin_var_end:
                continue

            evo_categories = _Reference(
                ws,
                min_col=data_col,
                min_row=evo_start,
                max_row=last_data_row,
            )

            if evolution_variant_norm == "simple":
                evol_chart = _LineChart()
                evol_chart.style = 13
                evol_chart.height = 7.5 * chart_scale
                evol_chart.width = 16.2 * chart_scale
                evol_chart.title = f"{chart_titles['evolution_title']} | {brand_name} P:{pipeline}"
                evol_chart.y_axis.title = chart_titles["evolution_var_axis"]
                evol_chart.y_axis.numFmt = "0.0%"
                evol_chart.x_axis.number_format = "yyyy-mm"
                evol_chart.x_axis.numFmt = "yyyy-mm"
                evol_chart.x_axis.tickLblPos = "low"
                evol_chart.x_axis.tickLblSkip = 1
                evol_chart.x_axis.tickMarkSkip = 1
                evol_chart.x_axis.delete = False
                evol_chart.legend.position = "b"
                evol_chart.legend.overlay = False
                evol_chart.add_data(
                    _Reference(
                        ws,
                        min_col=evo_kantar_col,
                        min_row=evo_start,
                        max_row=last_data_row,
                    ),
                    titles_from_data=False,
                )
                evol_chart.series[-1].title = _SeriesLabel(v=COL_EVO_KANTAR_YOY)
                evol_chart.add_data(
                    _Reference(
                        ws,
                        min_col=evo_sellin_col,
                        min_row=sellin_var_start,
                        max_row=sellin_var_end,
                    ),
                    titles_from_data=False,
                )
                evol_chart.series[-1].title = _SeriesLabel(v=COL_EVO_SELLIN_YOY)
                evol_chart.set_categories(evo_categories)
                if len(evol_chart.series) >= 1:
                    _set_line_series_color(evol_chart.series[0], COLOR_KANTAR_LINE)
                    _apply_variation_labels(evol_chart.series[0], COLOR_KANTAR_LINE)
                if len(evol_chart.series) >= 2:
                    _set_line_series_color(evol_chart.series[1], COLOR_SELLIN_LINE)
                    _apply_variation_labels(evol_chart.series[1], COLOR_SELLIN_LINE)
                ws.add_chart(evol_chart, f"{chart_anchor_col}42")
            else:
                # Modo clasico: solo tendencia de variaciones (sin volumen), con etiquetas puntuales.
                evol_var_chart = _LineChart()
                evol_var_chart.style = 2
                evol_var_chart.height = 7.5 * chart_scale
                evol_var_chart.width = 16.2 * chart_scale
                evol_var_chart.title = f"{chart_titles['evolution_title']} | {brand_name} P:{pipeline}"
                evol_var_chart.y_axis.title = chart_titles["evolution_var_axis"]
                evol_var_chart.y_axis.numFmt = "0.0%"
                evol_var_chart.x_axis.number_format = "yyyy-mm"
                evol_var_chart.x_axis.numFmt = "yyyy-mm"
                evol_var_chart.x_axis.tickLblPos = "low"
                evol_var_chart.x_axis.tickLblSkip = 1
                evol_var_chart.x_axis.tickMarkSkip = 1
                evol_var_chart.x_axis.delete = False
                evol_var_chart.legend.position = "b"
                evol_var_chart.legend.overlay = False
                evol_var_chart.add_data(
                    _Reference(
                        ws,
                        min_col=evo_kantar_col,
                        min_row=evo_start,
                        max_row=last_data_row,
                    ),
                    titles_from_data=False,
                )
                evol_var_chart.series[-1].title = _SeriesLabel(v=COL_EVO_KANTAR_YOY)
                evol_var_chart.add_data(
                    _Reference(
                        ws,
                        min_col=evo_sellin_col,
                        min_row=sellin_var_start,
                        max_row=sellin_var_end,
                    ),
                    titles_from_data=False,
                )
                evol_var_chart.series[-1].title = _SeriesLabel(v=COL_EVO_SELLIN_YOY)
                evol_var_chart.set_categories(evo_categories)
                if len(evol_var_chart.series) >= 1:
                    _set_line_series_color(evol_var_chart.series[0], COLOR_KANTAR_BAR_VAR)
                    _apply_variation_labels(evol_var_chart.series[0], COLOR_KANTAR_BAR_VAR)
                if len(evol_var_chart.series) >= 2:
                    _set_line_series_color(evol_var_chart.series[1], COLOR_SELLIN_BAR_VAR)
                    _apply_variation_labels(evol_var_chart.series[1], COLOR_SELLIN_BAR_VAR)
                ws.add_chart(evol_var_chart, f"{chart_anchor_col}42")

    wb.save(xlsx_path)


def generate_excel_template(
    root_dir: str,
    excel_file_obj: 'pd.ExcelFile',
    marcas: Sequence[str],
    pais_nombre: str,
    categoria_nombre: str,
    categoria_nombre_corto: str,
    fabricante: str,
    coverage_label: str,
    coverage_type: str,
    coverage_reason: str,
    trend_axis: str,
    evolution_slide_variant: str,
    include_english: bool,
    output_descriptor: str = "",
    elapsed_seconds_fn: Optional[Callable[[], Optional[float]]] = None,
) -> Tuple[str, str, str, str]:
    """Genera el archivo Excel temporal y devuelve datos clave."""
    try:
        console.print("\n[bold cyan]Generando archivo Excel temporal...[/bold cyan]")
    except Exception:
        print(Fore.CYAN + "\nGenerando archivo Excel temporal...")
    excel_temp_path = os.path.join(root_dir, EXCEL_TEMP_FILENAME)
    ref_month_year: Optional[str] = None
    sheets_without_reference_date: List[str] = []
    try:
        with pd.ExcelWriter(excel_temp_path) as writer:
            # Recorrer cada hoja (marca) del archivo
            total_sheets = len(marcas) if hasattr(marcas, "__len__") else 0
            status = console.status("Procesando hojas Excel...", spinner="line")
            status.start()
            try:
                for idx_sheet, marca_sheet_name in enumerate(marcas, start=1):
                    status.update(f"Procesando hoja {idx_sheet}/{total_sheets}: {marca_sheet_name}")

                    # 1.1) Carga y preprocesa la hoja usando la función refactorizada
                    df_marca, measure_unit = load_and_preprocess_sheet(
                        excel_file_obj,
                        marca_sheet_name,
                        include_metadata_hints=False,
                    )

                    # Si la carga falló, continuar con la siguiente hoja
                    if df_marca is None:
                        continue

                    # Guardar número original de filas de datos para fórmulas Excel
                    original_data_rows = len(df_marca)
                    if original_data_rows < 12:
                        console.print(
                            f"[yellow]Advertencia:[/] Hoja '{marca_sheet_name}' tiene < 12 meses de datos ({original_data_rows}). "
                            "Algunos calculos de Excel pueden fallar o dar NaN."
                        )
                        # Continuar de todos modos, pero con precaución

                    # Actualizar fecha de referencia global usando la ultima fecha valida detectada.
                    ref_date_value = pd.to_datetime(df_marca[COL_DATA].iloc[-1], errors="coerce")
                    if pd.isna(ref_date_value):
                        sheets_without_reference_date.append(marca_sheet_name)
                        console.print(
                            f"[yellow]Advertencia:[/] No se detecto una fecha valida en la ultima fila de "
                            f"'{marca_sheet_name}'. Se omitira para calcular el mes de referencia del template."
                        )
                    else:
                        ref_month_year = ref_date_value.strftime('%m-%y')

                    # --- 1.5) Creación de columnas con fórmulas Excel ---
                    df_excel = df_marca.copy() # Trabajar sobre una copia para Excel
                    # Hacer los índices basados en 1 y añadir offset de header (fila 1)
                    excel_row_offset = 2

                    # Sell_in_sim (Ejemplo - ajustable manualmente en Excel si se necesita)
                    # La fórmula asume que Sell_in está en la columna B
                    df_excel[COL_SELL_IN_SIM] = [f"=B{r}" for r in range(excel_row_offset, original_data_rows + excel_row_offset)] + [np.nan] * (len(df_excel) - original_data_rows)

                    # Acumulados (MAT - Moving Annual Total) - comienzan desde la fila 12 de datos
                    # Las fórmulas asumen Sell_out en C y Sell_in_sim en L
                    for i in range(11, original_data_rows):
                        row_excel = i + excel_row_offset
                        df_excel.loc[i, COL_ACUM_SELL_OUT] = f"=SUM(C{row_excel - 11}:C{row_excel})"
                        df_excel.loc[i, COL_ACUM_SELL_IN] = f"=SUM(L{row_excel - 11}:L{row_excel})" # Usa Sell_in_sim (L)

                    # --- 1.6) Cálculo de coberturas (pipeline 0 a 6) en Excel ---
                    pop_value_decimal = get_population_coverage_percent(pais_nombre) / 100.0
                    cov_formulas_list = []
                    max_rows_excel = original_data_rows + excel_row_offset -1 # Última fila con datos en Excel

                    for r_idx in range(original_data_rows): # Iterar sobre índices de datos (0 a N-1)
                        excel_current_row = r_idx + excel_row_offset
                        row_formulas = {}
                        if r_idx >= 11: # Cobertura solo se calcula desde el mes 12
                            for p in range(7): # Pipelines P0 a P6
                                # Fila de Excel para el numerador (Acum_Sell_in) - con pipeline
                                num_row_excel = excel_current_row + p
                                # Fila de Excel para el denominador (Acum_Sell_out) - sin pipeline
                                den_row_excel = excel_current_row

                                # Verificar que las filas referenciadas existan
                                if num_row_excel <= max_rows_excel and den_row_excel <= max_rows_excel:
                                    # La fórmula asume Acum_Sell_in en N y Acum_Sell_out en M
                                    #anterior  m{den_row_excel}/n{num_row_excel}*100
                                    base_formula = f"M{num_row_excel}/N{den_row_excel}*100"
                                    if coverage_type == "relativa":
                                        formula = f"=IFERROR(({base_formula})/{pop_value_decimal},NA())"
                                    else:
                                        formula = f"=IFERROR({base_formula},NA())"
                                    row_formulas[f'P{p}'] = formula
                                else:
                                     row_formulas[f'P{p}'] = np.nan # O "" o NA()
                        else:
                            # Rellenar con NaN para las primeras 11 filas
                            for p in range(7):
                                 row_formulas[f'P{p}'] = np.nan

                        cov_formulas_list.append(row_formulas)

                    df_cov_excel = pd.DataFrame(cov_formulas_list, index=df_excel.index[:original_data_rows])

                    # Escalonar las columnas de cobertura
                    df_cov_excel_scaled = df_cov_excel.copy()
                    escalona(df_cov_excel_scaled) # Escalonar la copia




        # -------------------------------------------------------
                    # 1.7 & 1.8) Cálculo de variaciones (Y-1 e Y-2) en Excel
                    # -------------------------------------------------------

                    # ► VARIABLES EXTRA que tu código “heredado” sigue ocupando
                    n_data          = original_data_rows                            # filas con datos
                    last_row_excel  = n_data + excel_row_offset - 1                 # última fila real en Excel
                    min_periods_for_layout = 42                                    # 36 meses + 6 pipelines para mantener layout
                    missing_periods = max(0, min_periods_for_layout - original_data_rows)

                    def build_if_no_zero(num_range: str, den_range: str, formula_body: str) -> str:
                        """
                        Envuelve una fórmula con validación de ceros: si hay 0 en alguna de las ventanas, devuelve "-".
                        'formula_body' no debe llevar '=' al inicio.
                        """
                        return f"=IF(OR(COUNTIF({num_range},0)>0,COUNTIF({den_range},0)>0),\"-\",{formula_body})"

                    def format_period_label(prefix: str, lag_meses: int) -> str:
                        """
                        Construye etiquetas tipo 'MAT ene-26 x MAT ene-25' de forma segura.
                        Si no hay historia suficiente, devuelve '-'.
                        """
                        idx_curr = n_data - 1
                        idx_prev = idx_curr - lag_meses
                        if idx_curr < 0 or idx_prev < 0 or idx_curr >= n_data or idx_prev >= n_data:
                            return "-"
                        fecha_curr = pd.to_datetime(df_excel.iloc[idx_curr][COL_DATA], errors="coerce")
                        fecha_prev = pd.to_datetime(df_excel.iloc[idx_prev][COL_DATA], errors="coerce")
                        if pd.isna(fecha_curr) or pd.isna(fecha_prev):
                            return "-"
                        return f"{prefix} {fecha_curr.strftime('%b-%y')} x {prefix} {fecha_prev.strftime('%b-%y')}"

                    def format_period_label_between_offsets(prefix: str, curr_end_offset: int, prev_end_offset: int) -> str:
                        """
                        Construye etiquetas usando dos cortes explícitos desde la fecha más reciente.
                        Ejemplo:
                        - curr_end_offset=12, prev_end_offset=24
                        - devuelve 'MAT Dec-24 x MAT Dec-23'
                        """
                        idx_curr = (n_data - 1) - curr_end_offset
                        idx_prev = (n_data - 1) - prev_end_offset
                        if idx_curr < 0 or idx_prev < 0 or idx_curr >= n_data or idx_prev >= n_data:
                            return "-"
                        fecha_curr = pd.to_datetime(df_excel.iloc[idx_curr][COL_DATA], errors="coerce")
                        fecha_prev = pd.to_datetime(df_excel.iloc[idx_prev][COL_DATA], errors="coerce")
                        if pd.isna(fecha_curr) or pd.isna(fecha_prev):
                            return "-"
                        return f"{prefix} {fecha_curr.strftime('%b-%y')} x {prefix} {fecha_prev.strftime('%b-%y')}"

                    def rango_excel(end_row: int, meses: int) -> tuple[int, int]:
                        """Devuelve (inicio, fin) inclusivo para una ventana de 'meses' que termina en 'end_row'."""
                        return end_row - (meses - 1), end_row

                    def formula_yoy_excel(col: str, end_row: int, meses: int, lag_meses: int) -> str:
                        """
                        = SUM( col[num_ini:num_fin] ) / SUM( col[den_ini:den_fin] ) - 1
                        donde el denominador termina en end_row - lag_meses y tiene el mismo tamaño 'meses'.
                        """
                        # Numerador: ventana actual (tamaño 'meses') que termina en end_row
                        num_ini, num_fin = rango_excel(end_row, meses)
                        # Denominador: misma ventana 'meses', pero que termina 'lag_meses' antes
                        den_fin = end_row - lag_meses
                        den_ini, den_fin = rango_excel(den_fin, meses)
                        num_range = f"{col}{num_ini}:{col}{num_fin}"
                        den_range = f"{col}{den_ini}:{col}{den_fin}"
                        formula_body = f"SUM({num_range})/SUM({den_range})-1"
                        return build_if_no_zero(num_range, den_range, formula_body)

                    # ---------- Y-1 -------------------------------------------------
                    # Ventanas: MAT=12, SEM=6, TRI=3 comparadas contra su misma ventana 12/6/3 meses antes.
                    y1_periods = [
                        ("Anual", "MAT", 12, 12),       # (tipo, etiqueta, meses_ventana, lag_meses)
                        ("Semestral", "SEM", 6, 6),
                        ("Trimestral", "TRI", 3, 3),
                    ]
                    var = pd.DataFrame(
                        [
                            [tipo, format_period_label(etiqueta, lag)]
                            for tipo, etiqueta, _, lag in y1_periods
                        ],
                        columns=['Tipo', 'Periodo']
                    )

                    # Variaciones WP by Numerator
                    var_wp = []
                    for _, _, meses, lag in y1_periods:
                        required = meses + lag
                        if n_data >= required:
                            var_wp.append(formula_yoy_excel("C", last_row_excel, meses, lag))
                        else:
                            var_wp.append("-")
                    var['WP by Numerator'] = var_wp

                    # Variaciones Cliente
                    for p in range(7):
                        end_row_p = last_row_excel - p
                        cli_var = []
                        for _, _, meses, lag in y1_periods:
                            required = meses + lag
                            if (n_data - p) >= required:
                                cli_var.append(formula_yoy_excel("L", end_row_p, meses, lag))
                            else:
                                cli_var.append("-")
                        var[f'Cliente P{p}'] = cli_var

                    # ---------- Mismo período del año pasado (solo SEM / TRI) -----
                    # Complementa la tabla del template Excel con las mismas filas
                    # extra que se agregaron al cuadro de la PPT.
                    same_period_last_year_periods = [
                        ("Semestral", "SEM", 6, 12),
                        ("Trimestral", "TRI", 3, 12),
                    ]
                    var_same_period_last_year = pd.DataFrame(
                        [
                            [tipo, format_period_label(etiqueta, lag)]
                            for tipo, etiqueta, _, lag in same_period_last_year_periods
                        ],
                        columns=['Tipo', 'Periodo']
                    )

                    wp_same_period_last_year = []
                    for _, _, meses, lag in same_period_last_year_periods:
                        required = meses + lag
                        if n_data >= required:
                            wp_same_period_last_year.append(formula_yoy_excel("C", last_row_excel, meses, lag))
                        else:
                            wp_same_period_last_year.append("-")
                    var_same_period_last_year['WP by Numerator'] = wp_same_period_last_year

                    for p in range(7):
                        end_row_p = last_row_excel - p
                        cli_same_period_last_year = []
                        for _, _, meses, lag in same_period_last_year_periods:
                            required = meses + lag
                            if (n_data - p) >= required:
                                cli_same_period_last_year.append(formula_yoy_excel("L", end_row_p, meses, lag))
                            else:
                                cli_same_period_last_year.append("-")
                        var_same_period_last_year[f'Cliente P{p}'] = cli_same_period_last_year

                    # ---------- Y-2 -------------------------------------------------
                    # Ventanas: MAT=12, SEM=6, TRI=3  (todas comparadas contra el mismo tamaño W, 24 meses antes)
                    periods = [
                        ('Anual',      'MAT', 12,  24),   # (tipo, etiqueta, meses_ventana, lag_meses)
                        ('Semestral',  'SEM', 6,   24),
                        ('Trimestral', 'TRI', 3,   24),
                    ]

                    # Texto de periodo (formato robusto para evitar KeyError con pocos meses)
                    aux = pd.DataFrame(
                        [
                            [tipo, format_period_label(etiqueta, lag)]
                            for tipo, etiqueta, _, lag in periods
                        ],
                        columns=['Tipo', 'Periodo']
                    )

                    # Reglas de suficiencia de datos por ventana para Y-2:
                    #  - MAT (12): requiere >= 12 + 24 = 36 meses
                    #  - SEM (6):  requiere >= 6  + 24 = 30 meses
                    #  - TRI (3):  requiere >= 3  + 24 = 27 meses

                    # ► WP by Numerator (columna C)
                    wp_y2_formulas = []
                    for _, _, meses, lag in periods:
                        required = meses + lag
                        if n_data >= required:
                            wp_y2_formulas.append(formula_yoy_excel("C", last_row_excel, meses, lag))
                        else:
                            wp_y2_formulas.append("-")
                    aux['WP by Numerator'] = wp_y2_formulas

                    # ► Clientes P0..P6 (columna L), ajustando el fin por 'p'
                    for p in range(7):
                        end_row_p = last_row_excel - p
                        cli_y2 = []
                        for _, _, meses, lag in periods:
                            required = meses + lag
                            # Suficiencia: descontamos 'p' del total disponible para ese cliente
                            if (n_data - p) >= required:
                                cli_y2.append(formula_yoy_excel("L", end_row_p, meses, lag))
                            else:
                                cli_y2.append("-")
                        aux[f'Cliente P{p}'] = cli_y2

                    # Limpiar variaciones sin sentido sin crear columnas negativas
                    if missing_periods > 0:
                        console.print(
                            f"[yellow]Advertencia:[/] Correlaciones/variaciones con periodos incompletos para "
                            f"[green]{marca_sheet_name}[/green] ({original_data_rows}/{min_periods_for_layout}); "
                            "se calculan correlaciones posibles; faltantes='-'."
                        )


                    # ---------- Unir Y-1 y Y-2 --------------------------------------
                    current_variations_block = pd.concat(
                        [
                            var.iloc[[0]].copy(),
                            var.iloc[[1]].copy(),
                            var_same_period_last_year.iloc[[0]].copy(),
                            var.iloc[[2]].copy(),
                            var_same_period_last_year.iloc[[1]].copy(),
                        ],
                        ignore_index=True,
                    )
                    df_variations_excel = pd.concat([current_variations_block, aux], ignore_index=True)

                    # ---------- Fila auxiliar de validación para P&G ---------------
                    # Expone explícitamente el cálculo anual del año previo vs su año anterior:
                    # con corte Dec-25 => MAT Dec-24 x MAT Dec-23.
                    prev_year_validation = {
                        'Tipo': 'Anual',
                        'Periodo': format_period_label_between_offsets('MAT', 12, 24) if n_data >= 36 else '-',
                        'WP by Numerator': "-",
                    }
                    if n_data >= 36:
                        prev_year_validation['WP by Numerator'] = formula_yoy_excel("C", last_row_excel - 12, 12, 12)
                    for p in range(7):
                        end_row_prev_year_p = last_row_excel - 12 - p
                        if (n_data - p) >= 36:
                            prev_year_validation[f'Cliente P{p}'] = formula_yoy_excel("L", end_row_prev_year_p, 12, 12)
                        else:
                            prev_year_validation[f'Cliente P{p}'] = "-"
                    df_variations_excel.loc[len(df_variations_excel)] = prev_year_validation



                    # --- 1.9) Cálculo de correlaciones en Excel (MAT) ---
                    # Se genera un diccionario con fórmulas de correlación para cada pipeline (P0 a P6)
                    # Se construyen fórmulas Excel que calculan la correlación Pearson entre dos rangos de 12 filas:
                    #   uno en la columna M y otro en la columna N, considerando el desplazamiento (pipeline).
                    # Los índices son base 1 y se garantiza que cada rango tenga exactamente 12 filas; de lo contrario, se asigna '-'.
            
                    # ---------- Correlaciones: 12m, 2 años antes (12m terminando hace 24m), 2 años (ventana 24m) ----------

                    series_sell_out = pd.to_numeric(df_marca[COL_SELL_OUT], errors="coerce")
                    series_sell_in = pd.to_numeric(df_marca[COL_SELL_IN], errors="coerce")

                    def _window_invalid(series: "pd.Series", start_row_excel: int, end_row_excel: int) -> bool:
                        """
                        Valida que la ventana exista y no tenga NaN.
                        Devuelve True si la ventana es inválida.
                        """
                        start_idx = start_row_excel - excel_row_offset
                        end_idx = end_row_excel - excel_row_offset
                        if start_idx < 0 or end_idx >= len(series):
                            return True
                        window = series.iloc[start_idx:end_idx+1]
                        return window.isna().any()

                    def _build_correl_row(label: str, window: int, end_offset: int = 0) -> dict:
                        """
                        Genera una fila de correlaciones entre M y N para:
                        - window: tamaño de ventana (12 o 24)
                        - end_offset: 0 = ventana termina en last_row_excel (reciente)
                                        24 = ventana termina 24 meses antes (para '2 años antes')
                        N se alinea con M desplazando p filas hacia arriba (n_start = m_start - p).
                        Si no hay suficientes datos para esa p y esa ventana, devuelve '-'.
                        """
                        row = {'Correlacion': label}

                        # ¿Hay datos suficientes para esta ventana y desplazamiento?
                        if n_data >= window + end_offset:
                            # Ventana base en M
                            row_ini = last_row_excel - end_offset - (window - 1)
                            row_fin = last_row_excel - end_offset

                            # Respetar que la fila 1 es encabezado
                            m_start = max(row_ini, 2)
                            m_end   = max(row_fin, 2)

                            for p in range(0, 7):  # P0..P6
                                # Cada pipeline consume filas adicionales; valida que haya datos suficientes
                                if (n_data - p) < (window + end_offset):
                                    row[f'P{p}'] = '-'
                                    continue

                                n_start = max(row_ini - p, 2)
                                n_end   = max(row_fin - p, 2)

                                # Ambas ventanas deben tener exactamente 'window' filas
                                if (m_end - m_start + 1 == window) and (n_end - n_start + 1 == window):
                                    # Si hay NaN en alguna ventana, se considera incompleto y se marca con '-'.
                                    # Los ceros son datos validos para correlacion; CORREL solo falla si no hay varianza.
                                    if _window_invalid(series_sell_out, m_start, m_end) or _window_invalid(series_sell_in, n_start, n_end):
                                        row[f'P{p}'] = "-"
                                    else:
                                        # Usa coma ',' en argumentos; función en inglés 'CORREL' como en tu flujo actual
                                        m_range = f"M{m_start}:M{m_end}"
                                        n_range = f"N{n_start}:N{n_end}"
                                        row[f'P{p}'] = (
                                            f"=IFERROR(IF(OR(COUNTBLANK({m_range})>0,COUNTBLANK({n_range})>0),\"-\","
                                            f"CORREL({m_range},{n_range})),\"-\")"
                                        )
                                else:
                                    row[f'P{p}'] = '-'
                        else:
                            for p in range(0, 7):
                                row[f'P{p}'] = '-'

                        return row

                    # Construye las 3 filas en el orden solicitado
                    rows = [
                        _build_correl_row('Año Actual', 12, end_offset=0),                   # últimos 12 meses
                        _build_correl_row('1 año antes', 12, end_offset=12),           # 12 meses que terminaron hace 12 meses (Año anterior)
                        _build_correl_row('2 años (ventana de 24 meses)', 24, 0),       # últimos 24 meses
                    ]

                    # Ordenar columnas: Correlacion, P0..P6 (incluye P6)
                    cols = ['Correlacion'] + [f'P{i}' for i in range(7)]
                    df_correlations_excel = pd.DataFrame(rows)[cols]




                    # --- 1.10) Promedio de Penetración y Buyers (MAT) en Excel ---
                    avg_formulas = []
                    # MAT Actual
                    if n_data >= 12:
                         start_avg_curr = last_row_excel - 11
                         end_avg_curr = last_row_excel
                         # Asume Penet en G, Buyers en H
                         avg_formulas.append({'Media': 'Penet MAT Actual', 'Valor': f"=AVERAGE(G{start_avg_curr}:G{end_avg_curr})"})
                         avg_formulas.append({'Media': 'Buyers MAT Actual', 'Valor': f"=AVERAGE(H{start_avg_curr}:H{end_avg_curr})"})
                    else:
                         avg_formulas.append({'Media': 'Penet MAT Actual', 'Valor': f"=AVERAGE(G{excel_row_offset}:G{last_row_excel})"}) # Promedio de lo disponible
                         avg_formulas.append({'Media': 'Buyers MAT Actual', 'Valor': f"=AVERAGE(H{excel_row_offset}:H{last_row_excel})"})

                    # MAT Anterior
                    if n_data >= 24:
                         start_avg_prev = last_row_excel - 23
                         end_avg_prev = last_row_excel - 12
                         avg_formulas.append({'Media': 'Penet MAT Anterior', 'Valor': f"=AVERAGE(G{start_avg_prev}:G{end_avg_prev})"})
                         avg_formulas.append({'Media': 'Buyers MAT Anterior', 'Valor': f"=AVERAGE(H{start_avg_prev}:H{end_avg_prev})"})
                    else:
                         avg_formulas.append({'Media': 'Penet MAT Anterior', 'Valor': np.nan}) # O NA()
                         avg_formulas.append({'Media': 'Buyers MAT Anterior', 'Valor': np.nan})

                    df_averages_excel = pd.DataFrame(avg_formulas)


                    # --- 1.10-bis) Variaciones YoY 12m para gráfico de evolución (columnas V y W) ---
                    evo_kantar_formulas: List[object] = []
                    evo_sellin_formulas: List[object] = []
                    for r_idx in range(original_data_rows):
                        row_excel = r_idx + excel_row_offset

                        # WP by Numerator YoY 12m: (SUM últimos 12) / (SUM 12 previos) - 1
                        if row_excel >= (excel_row_offset + 23):
                            num_range_c = f"C{row_excel - 11}:C{row_excel}"
                            den_range_c = f"C{row_excel - 23}:C{row_excel - 12}"
                            evo_kantar_formulas.append(
                                f"=IFERROR(SUM({num_range_c})/IF(SUM({den_range_c})=0,1,SUM({den_range_c}))-1,NA())"
                            )
                        else:
                            evo_kantar_formulas.append(np.nan)

                        # Sell-in YoY 12m SIN pipeline (pipeline se aplica al graficar)
                        if row_excel >= (excel_row_offset + 23):
                            sell_end = row_excel
                            sell_start = sell_end - 11
                            sell_prev_end = sell_end - 12
                            sell_prev_start = sell_prev_end - 11
                            num_range_l = f"L{sell_start}:L{sell_end}"
                            den_range_l = f"L{sell_prev_start}:L{sell_prev_end}"
                            evo_sellin_formulas.append(
                                f"=IFERROR(SUM({num_range_l})/IF(SUM({den_range_l})=0,1,SUM({den_range_l}))-1,NA())"
                            )
                        else:
                            evo_sellin_formulas.append(np.nan)

                    df_evolution_excel = pd.DataFrame(
                        {
                            COL_EVO_KANTAR_YOY: evo_kantar_formulas,
                            COL_EVO_SELLIN_YOY: evo_sellin_formulas,
                        },
                        index=df_excel.index[:original_data_rows],
                    )

                    # --- 1.11) Ensamblar DataFrame final para Excel ---
                    # Unir datos originales + coberturas + variaciones de evolución (V, W)
                    df_excel_final = pd.concat([df_excel, df_cov_excel_scaled, df_evolution_excel], axis=1)

                    # Crear la sección de resumen (Variaciones, Promedios, Correlación + Estabilidad)
                    # Añadir filas vacías y reorganizar
                    df_variations_excel['spacer1'] = np.nan
                    # df_averages_excel['spacer2'] = np.nan
                    df_correlations_excel['spacer3'] = np.nan

                    # Aplanar las tablas de resumen para concatenarlas horizontalmente
                    summary_part1 = df_variations_excel.T.reset_index().T # Variaciones
                    summary_part2 = df_averages_excel.T.reset_index().T   # Promedios
                    summary_part3 = df_correlations_excel.T.reset_index().T # Correlaciones

                    # Crear un DataFrame vacío con el número correcto de columnas para alinear
                    max_cols = df_excel_final.shape[1]
                    summary_placeholder = pd.DataFrame(np.nan, index=range(max(len(summary_part1), len(summary_part2), len(summary_part3))), columns=df_excel_final.columns)

                    # Rellenar el placeholder (esto requiere manejo cuidadoso de índices y columnas)
                    # Simplificación: Crear el df_excel_summary_part como antes y concatenar al final
                    df_excel_summary_part = pd.concat([df_variations_excel.reset_index(drop=True),
                                                      df_averages_excel.reset_index(drop=True),
                                                      df_correlations_excel.reset_index(drop=True)], axis=1)

                    # Añadir fila vacía de separación
                    df_excel_final.loc[len(df_excel_final)] = [np.nan] * len(df_excel_final.columns)

                    # Poner "Estabilidad" 2 filas arriba de la fila de encabezado "Correlacion":
                    #   Estabilidad
                    #   (fila en blanco)
                    #   Correlacion / P0..P6 (encabezado)
                    stab_row = {c: np.nan for c in df_excel_summary_part.columns}
                    if "Correlacion" in stab_row:
                        stab_row["Correlacion"] = "Estabilidad"
                        # Asume Cobertura P0-P6 en columnas O a U (después de escalonar)
                        coverage_start_col_idx = 15  # Col O es la 15 (1-based)
                        for p in range(7):
                            key = f"P{p}"
                            if key not in stab_row:
                                continue
                            col_letter = get_column_letter(coverage_start_col_idx + p)
                            # OJO: estos valores ya vienen "escalonados" (pipeline aplicado) por `escalona()`,
                            # así que la estabilidad se calcula con la misma fila para todos los pipelines (como P0).
                            row_last_cov = last_row_excel
                            row_prev_cov = last_row_excel - 12
                            # Requiere 12 meses hacia atrás y suficiente historia para que ambas coberturas existan
                            if row_last_cov >= excel_row_offset and row_prev_cov >= excel_row_offset and (original_data_rows >= (24 + p)):
                                stab_row[key] = f"=IFERROR({col_letter}{row_last_cov}-{col_letter}{row_prev_cov},NA())"
                            else:
                                stab_row[key] = "-"

                    df_stability_above = pd.DataFrame([stab_row], columns=df_excel_summary_part.columns)

                    # Añadir nombres de columnas del resumen como cabecera
                    summary_header = pd.DataFrame([df_excel_summary_part.columns], columns=df_excel_summary_part.columns)
                    df_excel_summary_part_with_header = pd.concat(
                        [df_stability_above, summary_header, df_excel_summary_part],
                        ignore_index=True,
                    )

                    # Ajustar columnas del resumen para que coincidan con el df principal y concatenar
                    # --- INICIO CAMBIO ---
                    # Si el número de columnas no coincide, agrega columnas vacías
                    n_main_cols = df_excel_final.shape[1]
                    n_summary_cols = df_excel_summary_part_with_header.shape[1]
                    if n_summary_cols < n_main_cols:
                        # Agrega columnas vacías al resumen
                        for i in range(n_summary_cols, n_main_cols):
                            df_excel_summary_part_with_header[f'empty_{i}'] = np.nan
                    elif n_summary_cols > n_main_cols:
                        # Si el resumen tiene más columnas, recórtalas
                        df_excel_summary_part_with_header = df_excel_summary_part_with_header.iloc[:, :n_main_cols]
                    # Ahora reasigna los nombres de columnas
                    df_excel_summary_part_with_header.columns = df_excel_final.columns
                    # --- FIN CAMBIO ---

                    df_excel_final = pd.concat([df_excel_final, df_excel_summary_part_with_header], ignore_index=True)

                    # --- 1.13) Exportar a la hoja de Excel ---
                    df_excel_final.to_excel(writer, sheet_name=marca_sheet_name, index=False)

            finally:
                status.stop()

        print(Fore.GREEN + f"Archivo Excel temporal '{EXCEL_TEMP_FILENAME}' generado.")

        # Aplicar formato de color y porcentaje a la sección de Correlaciones (como en excel_color.py)
        try:
            def apply_correlation_formatting(xlsx_path: str) -> None:
                from openpyxl import load_workbook as _load_wb
                from openpyxl.formatting.rule import ColorScaleRule as _ColorScaleRule
                from openpyxl.utils import get_column_letter as _get_col_letter
                wb = _load_wb(xlsx_path)
                for ws in wb.worksheets:
                    found = False
                    for row in ws.iter_rows(values_only=False):
                        for cell in row:
                            if isinstance(cell.value, str) and str(cell.value).strip().lower() == 'correlacion':
                                header_row = cell.row
                                header_col = cell.column
                                start_col = header_col + 1  # P0
                                end_col = start_col + 6     # P6

                                # Detectar largo dinámico: desde la fila siguiente hasta que la fila esté vacía en P0..P6
                                r = header_row + 1
                                last_row = r - 1
                                while True:
                                    vals = [ws.cell(row=r, column=c).value for c in range(start_col, end_col + 1)]
                                    if all(v is None for v in vals):
                                        break
                                    last_row = r
                                    r += 1

                                if last_row >= header_row + 1:
                                    # Formato de porcentaje 0.0% en P0..P6
                                    for rr in range(header_row + 1, last_row + 1):
                                        for cc in range(start_col, end_col + 1):
                                            ws.cell(row=rr, column=cc).number_format = '0.0%'

                                    # Regla de escala de color 3-colores (rojo-amarillo-verde)
                                    rng = f"{_get_col_letter(start_col)}{header_row + 1}:{_get_col_letter(end_col)}{last_row}"
                                    color_scale = _ColorScaleRule(
                                        start_type='min', start_color='F8696B',
                                        mid_type='percentile', mid_value=50, mid_color='FFEB84',
                                        end_type='max', end_color='63BE7B'
                                    )
                                    ws.conditional_formatting.add(rng, color_scale)
                                found = True
                                break
                        if found:
                            break
                wb.save(xlsx_path)

            apply_correlation_formatting(excel_temp_path)
            print(Fore.GREEN + "Formato de correlaciones aplicado (colores y porcentaje).")

            # Variaciones: formato porcentaje y reglas de color rojo(<0)/verde(>0)
            def apply_variations_formatting(xlsx_path: str) -> None:
                from openpyxl import load_workbook as _load_wb2
                from openpyxl.utils import get_column_letter as _col_letter
                from openpyxl.formatting.rule import Rule as _Rule
                from openpyxl.styles import PatternFill as _PatternFill, Font as _Font
                from openpyxl.styles.differential import DifferentialStyle as _Diff
                wb2 = _load_wb2(xlsx_path)
                for ws in wb2.worksheets:
                    header_row = None
                    wp_col = None
                    # Buscar el encabezado 'WP by Numerator'
                    for row in ws.iter_rows(values_only=False):
                        for cell in row:
                            if isinstance(cell.value, str) and str(cell.value).strip().lower() == 'wp by numerator':
                                header_row = cell.row
                                wp_col = cell.column
                                break
                        if header_row:
                            break
                    if not header_row:
                        continue
                    # Detectar columnas de Cliente P0..P6 consecutivas hacia la derecha
                    end_col = wp_col
                    p = 0
                    while True:
                        header_cell = ws.cell(row=header_row, column=wp_col + 1 + p)
                        val = header_cell.value
                        if isinstance(val, str) and val.strip().lower() == f'cliente p{p}':
                            end_col = wp_col + 1 + p
                            p += 1
                            if p > 20:  # seguridad
                                break
                        else:
                            break
                    # Si no se detectaron clientes, por defecto tomar WP + 7 clientes
                    if end_col == wp_col:
                        end_col = wp_col + 7
                    # Determinar rango de filas con datos (hasta que todas las columnas estén vacías)
                    r = header_row + 1
                    last_row = r - 1
                    while True:
                        vals = [ws.cell(row=r, column=c).value for c in range(wp_col, end_col + 1)]
                        if all(v is None for v in vals):
                            break
                        last_row = r
                        r += 1
                    if last_row < header_row + 1:
                        continue
                    # Aplicar formato porcentaje
                    for rr in range(header_row + 1, last_row + 1):
                        for cc in range(wp_col, end_col + 1):
                            ws.cell(row=rr, column=cc).number_format = '0.0%'
                    data_range = f"{_col_letter(wp_col)}{header_row + 1}:{_col_letter(end_col)}{last_row}"
                    # Regla < 0%: relleno rojo claro (#FFC7CE), texto rojo oscuro (#9C0006)
                    dxf_red = _Diff(
                        fill=_PatternFill(start_color='FFC7CE', end_color='FFC7CE', fill_type='solid'),
                        font=_Font(color='9C0006')
                    )
                    rule_red = _Rule(type='cellIs', operator='lessThan', formula=['0'], dxf=dxf_red)
                    ws.conditional_formatting.add(data_range, rule_red)

                    # Regla > 0%: relleno verde claro (#C6EFCE), texto verde oscuro (#006100)
                    dxf_green = _Diff(
                        fill=_PatternFill(start_color='C6EFCE', end_color='C6EFCE', fill_type='solid'),
                        font=_Font(color='006100')
                    )
                    rule_green = _Rule(type='cellIs', operator='greaterThan', formula=['0'], dxf=dxf_green)
                    ws.conditional_formatting.add(data_range, rule_green)

                    # Estabilidad (2 decimales + mismos colores rojo/verde que variaciones)
                    # La fila "Estabilidad" queda justo arriba del header "Correlacion".
                    stab_cell = None
                    for row in ws.iter_rows(values_only=False):
                        for cell in row:
                            if isinstance(cell.value, str) and cell.value.strip().lower() == "estabilidad":
                                # Validar que debajo esté el header 'Correlacion' para evitar falsos positivos
                                below = ws.cell(row=cell.row + 1, column=cell.column).value
                                if isinstance(below, str) and below.strip().lower() == "correlacion":
                                    stab_cell = cell
                                    break
                        if stab_cell:
                            break

                    if stab_cell:
                        stab_row = stab_cell.row
                        start_col = stab_cell.column + 1  # P0
                        end_col = start_col + 6           # P6
                        for cc in range(start_col, end_col + 1):
                            ws.cell(row=stab_row, column=cc).number_format = "0.00"
                        stab_range = f"{_col_letter(start_col)}{stab_row}:{_col_letter(end_col)}{stab_row}"
                        ws.conditional_formatting.add(
                            stab_range,
                            _Rule(type="cellIs", operator="lessThan", formula=["0"], dxf=dxf_red),
                        )
                        ws.conditional_formatting.add(
                            stab_range,
                            _Rule(type="cellIs", operator="greaterThan", formula=["0"], dxf=dxf_green),
                        )

                wb2.save(xlsx_path)
            apply_variations_formatting(excel_temp_path)
            print(Fore.GREEN + "Formato de variaciones aplicado (0.0% + rojo/verde).")

            def apply_coverage_values_formatting(xlsx_path: str) -> None:
                """Formatea coberturas (P0..P6), variaciones (V,W) y resalta cortes clave."""
                from openpyxl import load_workbook as _load_wb3
                from openpyxl.styles import PatternFill as _PatternFill
                wb3 = _load_wb3(xlsx_path)
                current_cov_fill = _PatternFill(start_color="F8CBAD", end_color="F8CBAD", fill_type="solid")
                prev12_cov_fill = _PatternFill(start_color="DDEBF7", end_color="DDEBF7", fill_type="solid")
                for ws in wb3.worksheets:
                    data_col = None
                    coverage_cols = []
                    evolution_var_cols = []
                    for col in range(1, ws.max_column + 1):
                        header_value = ws.cell(row=1, column=col).value
                        if header_value is None:
                            continue
                        header_text = str(header_value).strip()
                        if header_text == COL_DATA:
                            data_col = col
                        if header_text in {f"P{i}" for i in range(7)}:
                            coverage_cols.append(col)
                        if header_text in {COL_EVO_KANTAR_YOY, COL_EVO_SELLIN_YOY}:
                            evolution_var_cols.append(col)
                    if data_col is None or (not coverage_cols and not evolution_var_cols):
                        continue

                    row = 2
                    last_data_row = 1
                    while row <= ws.max_row:
                        value = ws.cell(row=row, column=data_col).value
                        if value is None or (isinstance(value, str) and value.strip() == ""):
                            break
                        last_data_row = row
                        row += 1
                    if last_data_row < 2:
                        continue

                    for rr in range(2, last_data_row + 1):
                        for cc in coverage_cols:
                            ws.cell(row=rr, column=cc).number_format = "0.0"
                        for cc in evolution_var_cols:
                            ws.cell(row=rr, column=cc).number_format = "0.0%"

                    # Resaltar coberturas del corte actual y de hace 12 meses para lectura rápida.
                    row_current = last_data_row
                    row_prev12 = last_data_row - 12 if (last_data_row - 12) >= 2 else None
                    for cc in coverage_cols:
                        ws.cell(row=row_current, column=cc).fill = current_cov_fill
                        if row_prev12 is not None:
                            ws.cell(row=row_prev12, column=cc).fill = prev12_cov_fill

                wb3.save(xlsx_path)

            apply_coverage_values_formatting(excel_temp_path)
            print(Fore.GREEN + "Formato aplicado: coberturas (1 decimal) y YoY evolución (0.0%).")

            add_native_excel_charts(
                excel_temp_path,
                coverage_label=coverage_label,
                trend_axis=trend_axis,
                evolution_slide_variant=evolution_slide_variant,
                include_english=include_english,
                pais_nombre=pais_nombre,
            )
            print(Fore.GREEN + "Graficos nativos de Excel insertados (editables).")

            def apply_visible_excel_labels(xlsx_path: str) -> None:
                from openpyxl import load_workbook as _load_wb4

                wb4 = _load_wb4(xlsx_path)
                for ws in wb4.worksheets:
                    for col in range(1, ws.max_column + 1):
                        header_cell = ws.cell(row=1, column=col)
                        if header_cell.value in VISIBLE_EXCEL_HEADER_MAP:
                            header_cell.value = VISIBLE_EXCEL_HEADER_MAP[header_cell.value]
                wb4.save(xlsx_path)

            apply_visible_excel_labels(excel_temp_path)
            print(Fore.GREEN + "Etiquetas visibles de Sell-in/Sell-out actualizadas en Excel.")

            # Colorear pestañas por marca para identificar rápidamente cada grupo de hojas.
            apply_template_tab_colors(excel_temp_path, marcas)
            print(Fore.GREEN + "Color de pestañas aplicado por marca en el template.")

            def apply_template_autofit(xlsx_path: str) -> None:
                from openpyxl import load_workbook as _load_wb5
                wb4 = _load_wb5(xlsx_path)
                for ws in wb4.worksheets:
                    autofit_worksheet_columns(ws, min_width=10.0, max_width=30.0, padding=2.0)
                wb4.save(xlsx_path)

            apply_template_autofit(excel_temp_path)
            print(Fore.GREEN + "Ancho de columnas ajustado automaticamente en el template.")
        except Exception as e:
            print(Fore.YELLOW + f"No se pudo completar el postproceso del template Excel: {e}")

    except PermissionError:
        # Usualmente pasa cuando el Excel de salida esta abierto/bloqueado.
        if os.path.exists(excel_temp_path):
            try:
                os.remove(excel_temp_path)
            except Exception:
                pass
        raise
    except Exception as e:
        print(f"{Fore.RED}{Style.BRIGHT}Error crítico durante la generación del archivo Excel: {e}")
        if os.path.exists(excel_temp_path):
             os.remove(excel_temp_path) # Limpiar si falla
        exit()

    # --- 1.14) Renombrar y mover archivo Excel final ---
    if not ref_month_year:
         if os.path.exists(excel_temp_path):
              os.remove(excel_temp_path)
         source_excel = getattr(excel_file_obj, "io", None) or getattr(excel_file_obj, "_io", None)
         source_label = os.path.basename(str(source_excel)) if source_excel else None
         print_reference_date_detection_warning(
             file_label=source_label,
             sheet_names=sheets_without_reference_date or list(marcas),
         )
         exit()

    categoria_salida = build_output_category_segment(categoria_nombre_corto, output_descriptor)
    nombre_base_archivo = f"{pais_nombre}-{categoria_salida}-{fabricante}-{ref_month_year}_{coverage_label}"
    carpeta_salida = os.path.join(root_dir, nombre_base_archivo) # Carpeta con el mismo nombre base

    if not os.path.exists(carpeta_salida):
        try:
            os.makedirs(carpeta_salida)
            print(Fore.BLUE + "Carpeta de salida creada")
        except OSError as e:
            print(f"{Fore.RED}Error al crear carpeta de salida '{carpeta_salida}': {e}")
            if os.path.exists(excel_temp_path): os.remove(excel_temp_path)
            exit()
    else:
        print(Fore.YELLOW + "Carpeta de salida ya existe, no se creara de nuevo")

    nombre_template_final = build_bounded_output_filename(carpeta_salida, f"Template_{nombre_base_archivo}.xlsx")
    ruta_template_final = os.path.join(carpeta_salida, nombre_template_final)

    def write_final_template() -> None:
        if os.path.exists(ruta_template_final):
            print(Fore.YELLOW + f"Archivo Excel ya existe. Se sobrescribirá.")
            os.remove(ruta_template_final)
        os.rename(excel_temp_path, ruta_template_final)

    try:
        run_file_write_with_retry(
            ruta_template_final,
            action_label="guardar el template Excel final",
            operation=write_final_template,
            elapsed_seconds_fn=elapsed_seconds_fn,
        )
        print(Fore.GREEN + "Archivo Excel final guardado")
    except FileSaveCancelled:
        if os.path.exists(excel_temp_path):
            try:
                os.remove(excel_temp_path)
            except Exception:
                pass
        raise
    except Exception as e:
        print(f"{Fore.RED}Error al mover/renombrar archivo Excel final: {e}")
        if os.path.exists(excel_temp_path): os.remove(excel_temp_path) # Limpiar temporal si falla el renombrado
        exit()


    return ref_month_year, carpeta_salida, nombre_base_archivo, ruta_template_final

def compute_coverage_dataframe(
    df_marca: "pd.DataFrame",
    pais_nombre: str,
    coverage_type: str,
    round_coverage: bool,
    marca_nombre: Optional[str] = None,
) -> "pd.DataFrame":
    """Calcula la cobertura rolling de 12 meses para cada pipeline."""
    acum_sell_out_py = df_marca[COL_SELL_OUT].rolling(window=12, min_periods=12).sum()
    acum_sell_out_py.index = df_marca[COL_DATA]
    df_coverage = pd.DataFrame(index=acum_sell_out_py.index)
    marca_label = (marca_nombre or "N/D").strip() or "N/D"
    needs_exception_warning = False
    for p in range(7):
        sell_in_shifted = df_marca[COL_SELL_IN].shift(p)
        acum_sell_in_shifted = sell_in_shifted.rolling(window=12, min_periods=12).sum()
        acum_sell_in_shifted.index = df_marca[COL_DATA]
        zero_mask = acum_sell_in_shifted == 0
        if zero_mask.any():
            needs_exception_warning = True
            acum_sell_in_shifted = acum_sell_in_shifted.copy()
            acum_sell_in_shifted.loc[zero_mask] = 1
        coverage_p = (acum_sell_out_py / acum_sell_in_shifted) * 100
        coverage_p = coverage_p.replace([np.inf, -np.inf], np.nan)
        df_coverage[f'P{p}'] = coverage_p
    pop_val_num = get_population_coverage_percent(pais_nombre) / 100.0
    if coverage_type.lower() == "relativa" and pop_val_num > 0:
        df_coverage = df_coverage / pop_val_num
    if round_coverage:
        df_coverage = df_coverage.apply(_round_half_up_series)
    else:
        df_coverage = df_coverage.round(1)
    if needs_exception_warning:
        notify_zero_months_exception(marca_label)
    return df_coverage


def compute_variations_dataframe(df_marca: "pd.DataFrame") -> "pd.DataFrame":
    period_types = ["Anual", "Semestral", "Trimestral"]
    df_variations = pd.DataFrame(columns=['Tipo', 'Periodo', 'WP by Numerator'] + [f'Cliente P{p}' for p in range(7)])
    kantar_vars_y1 = calc_var1(df_marca, COL_SELL_OUT, 0)
    cliente_vars_y1 = {p: calc_var1(df_marca, COL_SELL_IN, p) for p in range(7)}
    for i, p_type in enumerate(period_types):
        row = {'Tipo': p_type, 'Periodo': f'{p_type} vs Y-1', 'WP by Numerator': kantar_vars_y1[i]}
        for p in range(7):
            row[f'Cliente P{p}'] = cliente_vars_y1[p][i]
        df_variations.loc[len(df_variations)] = row
    kantar_vars_y2 = calc_var2(df_marca, COL_SELL_OUT, 0)
    cliente_vars_y2 = {p: calc_var2(df_marca, COL_SELL_IN, p) for p in range(7)}
    for i, p_type in enumerate(period_types):
        row = {'Tipo': p_type, 'Periodo': f'{p_type} vs Y-2', 'WP by Numerator': kantar_vars_y2[i]}
        for p in range(7):
            row[f'Cliente P{p}'] = cliente_vars_y2[p][i]
        df_variations.loc[len(df_variations)] = row
    return df_variations


def compute_averages(df_marca: "pd.DataFrame") -> Dict[str, float]:
    averages = {}
    n_data = len(df_marca)
    if n_data >= 12:
        averages['Penet_MAT_Actual'] = df_marca[COL_PENET].iloc[-12:].mean()
        averages['Buyers_MAT_Actual'] = df_marca[COL_BUYERS].iloc[-12:].mean()
        averages['Freq_MAT_Actual'] = df_marca[COL_FREQ].iloc[-12:].mean()
    else:
        averages['Penet_MAT_Actual'] = df_marca[COL_PENET].mean()
        averages['Buyers_MAT_Actual'] = df_marca[COL_BUYERS].mean()
        averages['Freq_MAT_Actual'] = df_marca[COL_FREQ].mean()
    if n_data >= 24:
        averages['Penet_MAT_Anterior'] = df_marca[COL_PENET].iloc[-24:-12].mean()
        averages['Buyers_MAT_Anterior'] = df_marca[COL_BUYERS].iloc[-24:-12].mean()
    else:
        averages['Penet_MAT_Anterior'] = np.nan
        averages['Buyers_MAT_Anterior'] = np.nan
    return averages


def compute_pipeline_mat_correlation(
    df_marca: "pd.DataFrame",
    pipeline: int,
    window: int = 12,
    end_offset: int = 0,
) -> Optional[float]:
    """Calcula correlación Pearson entre acumulados MAT de WP y Sell-in.

    Replica la pestaña de correlación del Excel: compara la serie MAT de
    WP by Numerator contra la serie MAT de Sell-in desplazada por pipeline.
    """
    if df_marca is None or df_marca.empty:
        return np.nan

    try:
        p = int(pipeline)
    except Exception:
        return np.nan

    try:
        window_int = int(window)
        end_offset_int = int(end_offset)
    except Exception:
        return np.nan
    if window_int <= 0 or end_offset_int < 0:
        return np.nan

    series_sell_out = pd.to_numeric(df_marca[COL_SELL_OUT], errors="coerce")
    series_sell_in = pd.to_numeric(df_marca[COL_SELL_IN], errors="coerce")
    acum_sell_out = series_sell_out.rolling(window=12, min_periods=12).sum()
    acum_sell_in = series_sell_in.shift(max(p, 0)).rolling(window=12, min_periods=12).sum()

    end_pos = len(df_marca) - end_offset_int
    start_pos = end_pos - window_int
    if start_pos < 0 or end_pos > len(df_marca):
        return np.nan
    sell_out_window = acum_sell_out.iloc[start_pos:end_pos]
    sell_in_window = acum_sell_in.iloc[start_pos:end_pos]

    if len(sell_out_window) != window_int or len(sell_in_window) != window_int:
        return np.nan
    if sell_out_window.isna().any() or sell_in_window.isna().any():
        return np.nan
    try:
        if float(sell_out_window.std()) == 0.0 or float(sell_in_window.std()) == 0.0:
            return np.nan
    except Exception:
        return np.nan

    try:
        corr, _ = pearsonr(sell_out_window.to_numpy(dtype=float), sell_in_window.to_numpy(dtype=float))
        return float(corr) if np.isfinite(corr) else np.nan
    except Exception:
        return np.nan


def compute_pipeline_current_year_correlation(df_marca: "pd.DataFrame", pipeline: int) -> Optional[float]:
    """Calcula la correlación MAT del Año Actual para un pipeline dado."""
    return compute_pipeline_mat_correlation(df_marca, pipeline, window=12, end_offset=0)


def compute_previous_year_annual_variation(df_marca: "pd.DataFrame", coluna: str, pipeline: int = 0) -> Optional[float]:
    """Calcula la variación anual del año previo vs su año anterior.

    Ejemplo con referencia dic-25:
    - compara MAT dic-24 vs MAT dic-23
    - para sell-in respeta el desplazamiento del pipeline
    """
    if df_marca is None or df_marca.empty:
        return np.nan

    try:
        p = int(pipeline)
    except Exception:
        return np.nan

    n_rows = len(df_marca)
    if n_rows < 36 + max(p, 0):
        return np.nan

    try:
        series = pd.to_numeric(df_marca[coluna], errors="coerce")
        if p != 0:
            current_sum = series.iloc[n_rows - 24 - p:n_rows - 12 - p].sum()
            previous_sum = series.iloc[n_rows - 36 - p:n_rows - 24 - p].sum()
        else:
            current_sum = series.iloc[n_rows - 24:n_rows - 12].sum()
            previous_sum = series.iloc[n_rows - 36:n_rows - 24].sum()
        if pd.isna(previous_sum) or float(previous_sum) == 0.0:
            return np.nan
        return (float(current_sum) / float(previous_sum)) - 1.0
    except Exception:
        return np.nan


@dataclass(frozen=True)
class OptimalPipelineCandidate:
    pipeline: int
    current_correlation: Optional[float]
    current_variation: Optional[float]
    wp_current_variation: Optional[float]
    variation_distance_points: Optional[float]
    current_trend_match: bool
    previous_year_correlation: Optional[float]
    two_year_correlation: Optional[float]
    previous_year_variation: Optional[float]
    wp_previous_year_variation: Optional[float]
    historical_trend_match: bool
    recent_shipment_outlier: bool = False
    forced_by_sheet: bool = False


@dataclass(frozen=True)
class OptimalPipelineSelection:
    pipeline: int
    reason: str
    candidates: Tuple[OptimalPipelineCandidate, ...]


@dataclass(frozen=True)
class AutoPipelineComparison:
    """Resultados paralelos del modo experimental.

    ``correlation`` conserva el candidato con mayor correlación MAT de Año
    Actual. ``balanced`` es la recomendación operativa que también considera
    variación, categoría, longitud, historia y outliers.
    """

    correlation: OptimalPipelineSelection
    balanced: OptimalPipelineSelection


@dataclass(frozen=True)
class CategoryPipelineProfile:
    name: str = "normal"
    min_pipeline: int = 1
    max_pipeline: int = 6
    short_pipeline_bias: float = 0.0
    ultra_fast: bool = False


OPTIMAL_PIPELINE_CORR_TIE_TOLERANCE = 0.02
OPTIMAL_PIPELINE_DISTANCE_TOLERANCE_POINTS = 5.0
OPTIMAL_PIPELINE_STRONG_CORR_MIN = 0.50
OPTIMAL_PIPELINE_HISTORICAL_CORR_MIN = 0.50
OPTIMAL_PIPELINE_BALANCED_CORR_MIN = 0.25
OPTIMAL_PIPELINE_EXACT_VARIATION_DISTANCE_POINTS = 0.10
OPTIMAL_PIPELINE_VISUAL_VARIATION_DISTANCE_POINTS = 0.50
OPTIMAL_PIPELINE_VISUAL_MAX_PIPELINE = 4
OPTIMAL_PIPELINE_MATERIAL_VARIATION_IMPROVEMENT_POINTS = 1.00
OPTIMAL_PIPELINE_MATERIAL_VARIATION_RATIO_MAX = 0.85
OPTIMAL_PIPELINE_MATERIAL_CORR_SACRIFICE_MAX = 0.15
OPTIMAL_PIPELINE_MATERIAL_CORR_MIN = 0.50
OPTIMAL_PIPELINE_VARIATION_DISTANCE_WEIGHT = 30.0
OPTIMAL_PIPELINE_NORMAL_LENGTH_PENALTY = 2.0
OPTIMAL_PIPELINE_OUTLIER_PENALTY = 6.0
AUTO_PIPELINE_CONFLICT_MEDIUM_CORR_LOSS = 0.10
AUTO_PIPELINE_CONFLICT_HIGH_CORR_LOSS = 0.25
HIGH_ROTATION_CATEGORY_CODES: Set[str] = frozenset({
    "MAYO",
    "KETC",
    "TOMA",
    "BOUI",
    "SOUP",
    "SNAC",
    "BISC",
    "CRSN",
    "PORK",
    "MEAT",
    "CHCK",
    "SAUS",
    "CRBR",
    "CRDF",
    "CRFO",
    "CRCU",
    "MXEV",
    "CRSA",
    "COCP",
    "PETF",
})
ULTRA_FAST_CATEGORY_CODES: Set[str] = frozenset({
    "PORK",
    "MEAT",
    "CHCK",
    "FISH",
    "SAUS",
    "HAMS",
})


def extract_forced_pipeline_from_sheet_name(sheet_name: str) -> Optional[int]:
    match = re.match(r"(?i)^p([0-6])_", str(sheet_name or ""))
    if match:
        try:
            return int(match.group(1))
        except Exception:
            return None
    return None


def _is_finite_number(value: object) -> bool:
    try:
        return value is not None and not pd.isna(value) and np.isfinite(float(value))
    except Exception:
        return False


def _same_variation_direction(left: object, right: object) -> bool:
    if not _is_finite_number(left) or not _is_finite_number(right):
        return False
    left_f = float(left)
    right_f = float(right)
    return (left_f * right_f) > 0 or (left_f == 0.0 and right_f == 0.0)


def _variation_distance_points(left: object, right: object) -> Optional[float]:
    if not _is_finite_number(left) or not _is_finite_number(right):
        return np.nan
    return abs(float(left) - float(right)) * 100.0


def _robust_outlier_score(series: "pd.Series", window: int = 18) -> float:
    values = pd.to_numeric(series, errors="coerce").dropna()
    if len(values) < 8:
        return 0.0
    tail = values.iloc[-min(int(window), len(values)):]
    if len(tail) < 8:
        return 0.0

    def _mad_score(sample: "pd.Series") -> float:
        sample = pd.to_numeric(sample, errors="coerce").dropna()
        if len(sample) < 5:
            return 0.0
        median = float(sample.median())
        mad = float((sample - median).abs().median())
        if mad <= 0:
            std = float(sample.std())
            if std <= 0:
                return 0.0
            return float((sample - median).abs().max() / std)
        return float(((sample - median).abs() / mad).max())

    level_score = _mad_score(tail)
    diff_score = _mad_score(tail.diff())
    return max(level_score, diff_score)


def detect_recent_shipment_outlier(df_marca: "pd.DataFrame", pipeline: int) -> bool:
    if df_marca is None or df_marca.empty:
        return False
    try:
        p = max(0, int(pipeline))
    except Exception:
        p = 0
    try:
        sell_in_series = pd.to_numeric(df_marca[COL_SELL_IN], errors="coerce").shift(p)
        sell_out_series = pd.to_numeric(df_marca[COL_SELL_OUT], errors="coerce")
        sell_in_score = _robust_outlier_score(sell_in_series, window=18)
        sell_out_score = _robust_outlier_score(sell_out_series, window=18)
        return sell_in_score >= 4.0 and sell_in_score >= (sell_out_score + 1.0)
    except Exception:
        return False


def resolve_category_pipeline_profile(
    category_code: object = None,
    category_name: object = None,
    basket_name: object = None,
) -> CategoryPipelineProfile:
    code = _normalize_category_code(category_code)
    basket_norm = _normalize_metadata_match_text(basket_name)
    text_norm = _normalize_metadata_match_text(f"{category_name or ''} {category_code or ''}")

    if code in ULTRA_FAST_CATEGORY_CODES or any(
        term in text_norm
        for term in {
            "carne",
            "carnicos",
            "porcina",
            "pollo",
            "pescado",
            "cecinas",
            "yoghurt",
            "yogurt",
            "yogur",
            "leche liquida",
            "leche fermentada",
            "queso fresco",
            "petit suisse",
            "crema de leche",
        }
    ):
        return CategoryPipelineProfile(
            name="ultra_alta_rotacion",
            min_pipeline=1,
            max_pipeline=2,
            short_pipeline_bias=12.0,
            ultra_fast=True,
        )

    if (
        "alimento" in basket_norm
        or "bebida" in basket_norm
        or "lacteo" in basket_norm
        or "lacteos" in basket_norm
        or code in HIGH_ROTATION_CATEGORY_CODES
        or "diet y light" in text_norm
    ):
        return CategoryPipelineProfile(
            name="alta_rotacion",
            min_pipeline=1,
            max_pipeline=3,
            short_pipeline_bias=8.0,
            ultra_fast=False,
        )

    return CategoryPipelineProfile()


def build_optimal_pipeline_candidates(
    df_marca: "pd.DataFrame",
    df_variations: "pd.DataFrame",
    forced_pipeline: Optional[int] = None,
) -> Tuple[OptimalPipelineCandidate, ...]:
    try:
        wp_current = df_variations.loc[df_variations["Tipo"] == "Anual", "WP by Numerator"].iloc[0]
    except Exception:
        wp_current = np.nan
    wp_previous = compute_previous_year_annual_variation(df_marca, COL_SELL_OUT, 0)

    candidates: List[OptimalPipelineCandidate] = []
    for pipeline in range(1, 7):
        try:
            current_variation = df_variations.loc[
                df_variations["Tipo"] == "Anual",
                f"Cliente P{pipeline}",
            ].iloc[0]
        except Exception:
            current_variation = np.nan
        previous_variation = compute_previous_year_annual_variation(df_marca, COL_SELL_IN, pipeline)
        candidates.append(
            OptimalPipelineCandidate(
                pipeline=pipeline,
                current_correlation=compute_pipeline_current_year_correlation(df_marca, pipeline),
                current_variation=current_variation,
                wp_current_variation=wp_current,
                variation_distance_points=_variation_distance_points(current_variation, wp_current),
                current_trend_match=_same_variation_direction(current_variation, wp_current),
                previous_year_correlation=compute_pipeline_mat_correlation(df_marca, pipeline, window=12, end_offset=12),
                two_year_correlation=compute_pipeline_mat_correlation(df_marca, pipeline, window=24, end_offset=0),
                previous_year_variation=previous_variation,
                wp_previous_year_variation=wp_previous,
                historical_trend_match=_same_variation_direction(previous_variation, wp_previous),
                recent_shipment_outlier=detect_recent_shipment_outlier(df_marca, pipeline),
                forced_by_sheet=(forced_pipeline == pipeline),
            )
        )
    return tuple(candidates)


def select_correlation_pipeline(
    candidates: Sequence[OptimalPipelineCandidate],
    forced_pipeline: Optional[int] = None,
) -> OptimalPipelineSelection:
    """Selecciona el candidato AUTO Correlación sin reglas de balanceo.

    El prefijo de hoja no altera el ranking cuando hay correlaciones finitas;
    solo sirve como fallback cuando no existe evidencia calculable. Esto hace
    que el resultado sea un contrafactual limpio frente a AUTO Balanceado.
    """

    candidates_tuple = tuple(candidates)
    finite = [
        candidate for candidate in candidates_tuple
        if _is_finite_number(candidate.current_correlation)
    ]
    if finite:
        chosen = max(
            finite,
            key=lambda candidate: (
                float(candidate.current_correlation),
                -candidate.pipeline,
            ),
        )
        corr = float(chosen.current_correlation)
        reason = (
            "máxima correlación MAT de Año Actual entre P1-P6"
            if corr > 0
            else "mejor correlación MAT disponible, pero la señal no es positiva; requiere revisión"
        )
        return OptimalPipelineSelection(chosen.pipeline, reason, candidates_tuple)

    forced_pipeline = forced_pipeline if forced_pipeline in range(1, 7) else None
    fallback_pipeline = forced_pipeline or 1
    reason = (
        "sin correlaciones calculables; se muestra el pipeline indicado en la hoja como fallback"
        if forced_pipeline is not None
        else "sin correlaciones calculables; se muestra P1 como fallback informativo"
    )
    return OptimalPipelineSelection(fallback_pipeline, reason, candidates_tuple)


def _choose_near_top_current_candidate(
    candidates: Sequence[OptimalPipelineCandidate],
) -> Optional[OptimalPipelineCandidate]:
    usable = [
        candidate for candidate in candidates
        if _is_finite_number(candidate.current_correlation)
        and float(candidate.current_correlation) > 0
        and candidate.current_trend_match
    ]
    if not usable:
        return None
    top_corr = max(float(candidate.current_correlation) for candidate in usable)
    near_top = [
        candidate for candidate in usable
        if float(candidate.current_correlation) >= top_corr - OPTIMAL_PIPELINE_CORR_TIE_TOLERANCE
    ]
    finite_distances = [
        float(candidate.variation_distance_points)
        for candidate in near_top
        if _is_finite_number(candidate.variation_distance_points)
    ]
    if finite_distances:
        best_distance = min(finite_distances)
        distance_cap = best_distance + OPTIMAL_PIPELINE_DISTANCE_TOLERANCE_POINTS
        near_top = [
            candidate for candidate in near_top
            if not _is_finite_number(candidate.variation_distance_points)
            or float(candidate.variation_distance_points) <= distance_cap
        ]
    forced_near_top = [candidate for candidate in near_top if candidate.forced_by_sheet]
    if forced_near_top:
        return sorted(forced_near_top, key=lambda candidate: candidate.pipeline)[0]
    return sorted(near_top, key=lambda candidate: candidate.pipeline)[0] if near_top else None


def _historical_support_score(candidate: OptimalPipelineCandidate) -> float:
    score = 0.0
    if _is_finite_number(candidate.previous_year_correlation):
        score = max(score, float(candidate.previous_year_correlation))
    if _is_finite_number(candidate.two_year_correlation):
        score = max(score, float(candidate.two_year_correlation) * 0.85)
    if candidate.historical_trend_match:
        score += 0.10
    return score


def _current_fit_score(
    candidate: OptimalPipelineCandidate,
    min_pipeline: int = 1,
    length_penalty: float = OPTIMAL_PIPELINE_NORMAL_LENGTH_PENALTY,
) -> float:
    if (
        not _is_finite_number(candidate.current_correlation)
        or not _is_finite_number(candidate.variation_distance_points)
        or not candidate.current_trend_match
    ):
        return float("-inf")
    corr = float(candidate.current_correlation)
    if corr <= 0:
        return float("-inf")
    distance = max(0.0, float(candidate.variation_distance_points))
    pipeline_penalty = max(0, candidate.pipeline - max(1, int(min_pipeline))) * float(length_penalty)
    outlier_penalty = OPTIMAL_PIPELINE_OUTLIER_PENALTY if candidate.recent_shipment_outlier else 0.0
    historical_bonus = 3.0 if candidate.historical_trend_match else 0.0
    forced_bonus = 1.0 if candidate.forced_by_sheet else 0.0
    return (
        corr * 100.0
        - distance * OPTIMAL_PIPELINE_VARIATION_DISTANCE_WEIGHT
        - pipeline_penalty
        - outlier_penalty
        + historical_bonus
        + forced_bonus
    )


def _format_pp(value: object) -> str:
    if not _is_finite_number(value):
        return "N/D"
    return f"{float(value):.2f}pp"


def _format_corr(value: object) -> str:
    if not _is_finite_number(value):
        return "N/D"
    return f"{float(value):.3f}"


def _choose_balanced_current_candidate(
    candidates: Sequence[OptimalPipelineCandidate],
    min_pipeline: int = 1,
    length_penalty: float = OPTIMAL_PIPELINE_NORMAL_LENGTH_PENALTY,
) -> Optional[OptimalPipelineCandidate]:
    usable = [
        candidate for candidate in candidates
        if candidate.current_trend_match
        and _is_finite_number(candidate.current_correlation)
        and float(candidate.current_correlation) >= OPTIMAL_PIPELINE_BALANCED_CORR_MIN
        and _is_finite_number(candidate.variation_distance_points)
    ]
    if not usable:
        return None
    exact_variation = [
        candidate for candidate in usable
        if float(candidate.variation_distance_points) <= OPTIMAL_PIPELINE_EXACT_VARIATION_DISTANCE_POINTS
    ]
    if exact_variation:
        forced_exact = [candidate for candidate in exact_variation if candidate.forced_by_sheet]
        if forced_exact:
            return sorted(forced_exact, key=lambda candidate: candidate.pipeline)[0]
        return sorted(
            exact_variation,
            key=lambda candidate: (
                candidate.pipeline,
                -float(candidate.current_correlation),
            ),
        )[0]

    visual_variation = [
        candidate for candidate in usable
        if float(candidate.variation_distance_points) <= OPTIMAL_PIPELINE_VISUAL_VARIATION_DISTANCE_POINTS
        and candidate.pipeline <= OPTIMAL_PIPELINE_VISUAL_MAX_PIPELINE
    ]
    if visual_variation:
        best_visual = sorted(
            visual_variation,
            key=lambda candidate: (
                -_current_fit_score(
                    candidate,
                    min_pipeline=min_pipeline,
                    length_penalty=length_penalty,
                ),
                candidate.pipeline,
            ),
        )[0]
        return best_visual
    return None


def _choose_material_variation_candidate(
    candidates: Sequence[OptimalPipelineCandidate],
    min_pipeline: int = 1,
    length_penalty: float = OPTIMAL_PIPELINE_NORMAL_LENGTH_PENALTY,
) -> Optional[OptimalPipelineCandidate]:
    """Permite un balance intermedio cuando mejora materialmente la variación.

    Esta banda solo compite contra el candidato de máxima correlación actual y
    exige conservar una correlación sólida, reducir el gap tanto en términos
    absolutos como relativos y usar un pipeline más corto.
    """
    usable = [
        candidate for candidate in candidates
        if candidate.current_trend_match
        and _is_finite_number(candidate.current_correlation)
        and float(candidate.current_correlation) >= OPTIMAL_PIPELINE_MATERIAL_CORR_MIN
        and _is_finite_number(candidate.variation_distance_points)
    ]
    if not usable:
        return None

    top_corr = max(usable, key=lambda candidate: float(candidate.current_correlation))
    top_gap = float(top_corr.variation_distance_points)
    if top_gap <= 0:
        return None

    material_candidates = []
    for candidate in usable:
        if candidate.pipeline >= top_corr.pipeline:
            continue
        corr_sacrifice = float(top_corr.current_correlation) - float(candidate.current_correlation)
        candidate_gap = float(candidate.variation_distance_points)
        gap_improvement = top_gap - candidate_gap
        gap_ratio = candidate_gap / top_gap
        if (
            corr_sacrifice <= OPTIMAL_PIPELINE_MATERIAL_CORR_SACRIFICE_MAX
            and gap_improvement >= OPTIMAL_PIPELINE_MATERIAL_VARIATION_IMPROVEMENT_POINTS
            and gap_ratio <= OPTIMAL_PIPELINE_MATERIAL_VARIATION_RATIO_MAX
        ):
            material_candidates.append(candidate)

    if not material_candidates:
        return None
    return sorted(
        material_candidates,
        key=lambda candidate: (
            -_current_fit_score(
                candidate,
                min_pipeline=min_pipeline,
                length_penalty=length_penalty,
            ),
            candidate.pipeline,
        ),
    )[0]


def _material_variation_reason(
    chosen: OptimalPipelineCandidate,
    candidates: Sequence[OptimalPipelineCandidate],
) -> str:
    top_corr = max(
        (
            candidate for candidate in candidates
            if candidate.current_trend_match
            and _is_finite_number(candidate.current_correlation)
            and _is_finite_number(candidate.variation_distance_points)
        ),
        key=lambda candidate: float(candidate.current_correlation),
    )
    gap_improvement = (
        float(top_corr.variation_distance_points)
        - float(chosen.variation_distance_points)
    )
    corr_sacrifice = (
        float(top_corr.current_correlation)
        - float(chosen.current_correlation)
    )
    months_shorter = top_corr.pipeline - chosen.pipeline
    month_label = "mes" if months_shorter == 1 else "meses"
    return (
        "ajuste balanceado material de variación anual: "
        f"P{chosen.pipeline} corr={_format_corr(chosen.current_correlation)}, "
        f"gap={_format_pp(chosen.variation_distance_points)}; "
        f"top correlación P{top_corr.pipeline} corr={_format_corr(top_corr.current_correlation)}, "
        f"gap={_format_pp(top_corr.variation_distance_points)}; "
        f"mejora gap={_format_pp(gap_improvement)}, "
        f"sacrificio correlación={corr_sacrifice:.3f} y "
        f"pipeline {months_shorter} {month_label} más corto"
    )


def _balanced_current_reason(
    chosen: OptimalPipelineCandidate,
    candidates: Sequence[OptimalPipelineCandidate],
    min_pipeline: int = 1,
    length_penalty: float = OPTIMAL_PIPELINE_NORMAL_LENGTH_PENALTY,
) -> str:
    chosen_score = _current_fit_score(
        chosen,
        min_pipeline=min_pipeline,
        length_penalty=length_penalty,
    )
    reason = (
        "score balanceado de Año Actual: "
        f"P{chosen.pipeline} corr={_format_corr(chosen.current_correlation)}, "
        f"gap variación={_format_pp(chosen.variation_distance_points)}, "
        f"score={chosen_score:.1f}"
    )
    if (
        _is_finite_number(chosen.variation_distance_points)
        and float(chosen.variation_distance_points) <= OPTIMAL_PIPELINE_EXACT_VARIATION_DISTANCE_POINTS
    ):
        reason = "ajuste casi exacto de variación anual; " + reason
    elif (
        _is_finite_number(chosen.variation_distance_points)
        and float(chosen.variation_distance_points) <= OPTIMAL_PIPELINE_VISUAL_VARIATION_DISTANCE_POINTS
    ):
        reason = "alineación visual de variación anual; " + reason
    top_corr_candidates = [
        candidate for candidate in candidates
        if candidate.current_trend_match
        and _is_finite_number(candidate.current_correlation)
        and float(candidate.current_correlation) > 0
    ]
    if top_corr_candidates:
        top_corr = max(top_corr_candidates, key=lambda candidate: float(candidate.current_correlation))
        if top_corr.pipeline != chosen.pipeline:
            top_score = _current_fit_score(
                top_corr,
                min_pipeline=min_pipeline,
                length_penalty=length_penalty,
            )
            reason += (
                f"; top correlación P{top_corr.pipeline} "
                f"corr={_format_corr(top_corr.current_correlation)}, "
                f"gap variación={_format_pp(top_corr.variation_distance_points)}, "
                f"score={top_score:.1f}"
            )
    scored_candidates = [
        candidate for candidate in top_corr_candidates
        if _is_finite_number(candidate.variation_distance_points)
        and np.isfinite(
            _current_fit_score(
                candidate,
                min_pipeline=min_pipeline,
                length_penalty=length_penalty,
            )
        )
    ]
    if scored_candidates:
        top_score_candidate = max(
            scored_candidates,
            key=lambda candidate: _current_fit_score(
                candidate,
                min_pipeline=min_pipeline,
                length_penalty=length_penalty,
            ),
        )
        if top_score_candidate.pipeline != chosen.pipeline:
            reason += (
                f"; top score P{top_score_candidate.pipeline} "
                f"corr={_format_corr(top_score_candidate.current_correlation)}, "
                f"gap variación={_format_pp(top_score_candidate.variation_distance_points)}, "
                f"score={_current_fit_score(top_score_candidate, min_pipeline=min_pipeline, length_penalty=length_penalty):.1f}"
            )
    return reason


def _is_fast_consumable_category(
    category_code: object = None,
    category_name: object = None,
    basket_name: object = None,
) -> bool:
    return resolve_category_pipeline_profile(category_code, category_name, basket_name).name != "normal"


def _current_shipment_delta(candidate: OptimalPipelineCandidate) -> Optional[float]:
    if not _is_finite_number(candidate.current_variation) or not _is_finite_number(candidate.wp_current_variation):
        return np.nan
    return float(candidate.current_variation) - float(candidate.wp_current_variation)


def _choose_fast_consumable_candidate(
    candidates: Sequence[OptimalPipelineCandidate],
    profile: CategoryPipelineProfile,
) -> Optional[OptimalPipelineCandidate]:
    min_pipeline = max(1, int(profile.min_pipeline))
    max_pipeline = min(6, max(min_pipeline, int(profile.max_pipeline)))
    short_candidates = [candidate for candidate in candidates if min_pipeline <= candidate.pipeline <= max_pipeline]
    usable = [
        candidate for candidate in short_candidates
        if _is_finite_number(candidate.current_correlation)
        and float(candidate.current_correlation) > 0
        and candidate.current_trend_match
    ]
    if usable:
        forced_usable = [candidate for candidate in usable if candidate.forced_by_sheet]
        if forced_usable:
            forced = sorted(forced_usable, key=lambda candidate: candidate.pipeline)[0]
            best_corr = max(float(candidate.current_correlation) for candidate in usable)
            best_distance = min(
                float(candidate.variation_distance_points)
                for candidate in usable
                if _is_finite_number(candidate.variation_distance_points)
            )
            forced_distance = (
                float(forced.variation_distance_points)
                if _is_finite_number(forced.variation_distance_points)
                else best_distance
            )
            if (
                float(forced.current_correlation) >= best_corr - 0.15
                and forced_distance <= best_distance + 10.0
            ):
                return forced

        def _score(candidate: OptimalPipelineCandidate) -> float:
            corr_score = float(candidate.current_correlation) * 100.0
            distance_penalty = (
                float(candidate.variation_distance_points) * 1.2
                if _is_finite_number(candidate.variation_distance_points)
                else 15.0
            )
            length_penalty = (candidate.pipeline - min_pipeline) * float(profile.short_pipeline_bias)
            delta = _current_shipment_delta(candidate)
            shipment_adjustment = 0.0
            if _is_finite_number(delta):
                # Si el sell-in crece por encima de WP, permitimos algo más de pipeline.
                # Si cae más que WP, reforzamos pipelines cortos.
                shipment_adjustment = min(8.0, max(-8.0, float(delta) * 100.0 * 0.4)) * (candidate.pipeline - 1)
            forced_bonus = 3.0 if candidate.forced_by_sheet else 0.0
            return corr_score - distance_penalty - length_penalty + shipment_adjustment + forced_bonus

        return sorted(usable, key=lambda candidate: (-_score(candidate), candidate.pipeline))[0]

    historical = [
        candidate for candidate in short_candidates
        if candidate.historical_trend_match
        and _historical_support_score(candidate) >= OPTIMAL_PIPELINE_HISTORICAL_CORR_MIN
    ]
    if historical:
        forced_historical = [candidate for candidate in historical if candidate.forced_by_sheet]
        if forced_historical:
            return sorted(forced_historical, key=lambda candidate: candidate.pipeline)[0]
        best_score = max(_historical_support_score(candidate) for candidate in historical)
        near_best = [
            candidate for candidate in historical
            if _historical_support_score(candidate) >= best_score - OPTIMAL_PIPELINE_CORR_TIE_TOLERANCE
        ]
        return sorted(near_best, key=lambda candidate: candidate.pipeline)[0]

    forced_short = [candidate for candidate in short_candidates if candidate.forced_by_sheet]
    if forced_short:
        return sorted(forced_short, key=lambda candidate: candidate.pipeline)[0]
    return short_candidates[0] if short_candidates else None


def _fast_consumable_reason(candidate: OptimalPipelineCandidate) -> str:
    if (
        _is_finite_number(candidate.current_correlation)
        and float(candidate.current_correlation) > 0
        and candidate.current_trend_match
    ):
        suffix = " pese a outlier reciente de shipments" if candidate.recent_shipment_outlier else ""
        return f"categoría con perfil de pipeline corto: se prioriza pipeline corto con correlación y variación actuales razonables{suffix}"
    if candidate.historical_trend_match and _historical_support_score(candidate) >= OPTIMAL_PIPELINE_HISTORICAL_CORR_MIN:
        suffix = " y outlier reciente de shipments" if candidate.recent_shipment_outlier else ""
        return f"categoría con perfil de pipeline corto: el Año Actual no domina y se conserva pipeline corto con respaldo histórico{suffix}"
    if candidate.recent_shipment_outlier:
        return "categoría con perfil de pipeline corto: se conserva pipeline corto porque hay outlier reciente de shipments"
    if candidate.forced_by_sheet:
        return "categoría con perfil de pipeline corto: sin señal concluyente, se conserva el pipeline corto indicado en la hoja"
    return "categoría con perfil de pipeline corto: sin señal concluyente, se usa el pipeline corto disponible"


def select_optimal_pipeline(
    df_marca: "pd.DataFrame",
    df_variations: "pd.DataFrame",
    forced_pipeline: Optional[int] = None,
    category_code: object = None,
    category_name: object = None,
    basket_name: object = None,
) -> OptimalPipelineSelection:
    """Elige la recomendación AUTO Balanceado para P1-P6.

    Esta decisión combina correlación, variación, perfil de categoría,
    longitud, historia y outliers. El contrafactual puro de correlación se
    obtiene por separado con :func:`select_correlation_pipeline`.
    """
    forced_pipeline = forced_pipeline if forced_pipeline in range(1, 7) else None
    candidates = build_optimal_pipeline_candidates(df_marca, df_variations, forced_pipeline)
    category_profile = resolve_category_pipeline_profile(category_code, category_name, basket_name)

    if category_profile.name != "normal":
        chosen_fast = _choose_fast_consumable_candidate(candidates, category_profile)
        if chosen_fast is not None:
            return OptimalPipelineSelection(
                chosen_fast.pipeline,
                _fast_consumable_reason(chosen_fast),
                candidates,
            )

    chosen_balanced = _choose_balanced_current_candidate(candidates)
    if chosen_balanced is not None:
        return OptimalPipelineSelection(
            chosen_balanced.pipeline,
            _balanced_current_reason(chosen_balanced, candidates),
            candidates,
        )

    chosen_material = _choose_material_variation_candidate(candidates)
    if chosen_material is not None:
        return OptimalPipelineSelection(
            chosen_material.pipeline,
            _material_variation_reason(chosen_material, candidates),
            candidates,
        )

    strong_candidates = [
        candidate for candidate in candidates
        if _is_finite_number(candidate.current_correlation)
        and float(candidate.current_correlation) >= OPTIMAL_PIPELINE_STRONG_CORR_MIN
        and candidate.current_trend_match
    ]
    chosen = _choose_near_top_current_candidate(strong_candidates)
    if chosen is not None:
        return OptimalPipelineSelection(
            chosen.pipeline,
            "correlación MAT de Año Actual positiva y variación anual alineada",
            candidates,
        )

    chosen = _choose_near_top_current_candidate(candidates)
    if chosen is not None:
        return OptimalPipelineSelection(
            chosen.pipeline,
            "correlación MAT de Año Actual positiva, aunque débil, y variación anual alineada",
            candidates,
        )

    historical_candidates = [
        candidate for candidate in candidates
        if candidate.historical_trend_match
        and _historical_support_score(candidate) >= OPTIMAL_PIPELINE_HISTORICAL_CORR_MIN
    ]
    if historical_candidates:
        forced_historical = [candidate for candidate in historical_candidates if candidate.forced_by_sheet]
        if forced_historical:
            chosen = sorted(forced_historical, key=lambda candidate: candidate.pipeline)[0]
            reason = "respaldo histórico: el Año Actual se rompe, pero la hoja y la historia sostienen el pipeline"
            if chosen.recent_shipment_outlier:
                reason = "respaldo histórico con outlier reciente de shipments: se conserva el pipeline indicado"
            return OptimalPipelineSelection(
                chosen.pipeline,
                reason,
                candidates,
            )
        best_score = max(_historical_support_score(candidate) for candidate in historical_candidates)
        near_best = [
            candidate for candidate in historical_candidates
            if _historical_support_score(candidate) >= best_score - OPTIMAL_PIPELINE_CORR_TIE_TOLERANCE
        ]
        chosen = sorted(near_best, key=lambda candidate: candidate.pipeline)[0]
        reason = "respaldo histórico por correlación y tendencia anual previa"
        if chosen.recent_shipment_outlier:
            reason = "respaldo histórico con outlier reciente de shipments"
        return OptimalPipelineSelection(
            chosen.pipeline,
            reason,
            candidates,
        )

    if forced_pipeline is not None:
        forced_candidates = [candidate for candidate in candidates if candidate.pipeline == forced_pipeline]
        if forced_candidates and forced_candidates[0].recent_shipment_outlier:
            return OptimalPipelineSelection(
                forced_pipeline,
                "outlier reciente de shipments rompe la señal; se conserva el pipeline indicado en la hoja",
                candidates,
            )
        return OptimalPipelineSelection(
            forced_pipeline,
            "sin señal concluyente; se conserva el pipeline indicado en la hoja",
            candidates,
        )

    return OptimalPipelineSelection(
        1,
        "sin señal concluyente; se usa P1 como fallback operativo",
        candidates,
    )


def compute_trend_plot_df(df_marca: "pd.DataFrame") -> "pd.DataFrame":
    df_trend_plot = df_marca[[COL_DATA, COL_SELL_IN, COL_SELL_OUT]].copy()
    df_trend_plot[COL_DATA] = df_trend_plot[COL_DATA].apply(lambda x: x.strftime('%m-%y'))
    return df_trend_plot


def build_variation_table(
    fabricante: str,
    labels: Dict[Tuple[int, str], List[str] | str],
    lang_index: int,
    pipeline: int,
    ref_month_year: str,
    var_cliente_mat: Optional[float],
    var_kantar_mat: Optional[float],
) -> "pd.DataFrame":
    label_var = labels[(lang_index, 'Var')]
    data = {
        " ": [f"VAR % MAT ({ref_month_year})"],
        f"{fabricante} {label_var} Pipeline {pipeline}": [f"{var_cliente_mat*100:.1f}%" if pd.notna(var_cliente_mat) else "-"],
        "Worldpanel by Numerator": [f"{var_kantar_mat*100:.1f}%" if pd.notna(var_kantar_mat) else "-"],
    }
    return pd.DataFrame(data)

def build_variations_detail_table(
    df_variations: "pd.DataFrame",
    pipeline: int,
    df_marca: "pd.DataFrame",
    include_same_period_last_year: bool = True,
) -> "pd.DataFrame":
    """Construye la tabla de variaciones utilizada en el slide de tendencia."""
    if df_variations is None or df_variations.empty:
        return pd.DataFrame()
    filtered = df_variations[df_variations['Periodo'].astype(str).str.contains('Y-1', na=False)].copy()
    if filtered.empty:
        return pd.DataFrame()

    base_columns = [col for col in ['Tipo', 'Periodo', 'WP by Numerator', 'Cliente P0'] if col in filtered.columns]
    if not base_columns:
        return pd.DataFrame()
    detail_df = filtered[base_columns].copy()

    pipeline_col = f'Cliente P{pipeline}'
    pipeline_detail_col = f'Cliente Pipeline (P{pipeline})' if int(pipeline) > 0 else None
    if pipeline_detail_col and pipeline_col in df_variations.columns:
        detail_df[pipeline_detail_col] = df_variations.loc[detail_df.index, pipeline_col].values

    def _format_month(dt: "pd.Timestamp") -> str:
        if pd.isna(dt):
            return "-"
        dt = pd.to_datetime(dt)
        return f"{month_abbr[dt.month]}-{dt.year % 100:02d}"

    def _build_period_text(label: str, end_dt: "pd.Timestamp", compare_lag: int) -> str:
        previous_dt = end_dt - pd.DateOffset(months=compare_lag)
        return f"{label} {_format_month(end_dt)} x {label} {_format_month(previous_dt)}"

    if df_marca is not None and not df_marca.empty:
        try:
            current_dt = pd.to_datetime(df_marca[COL_DATA].iloc[-1])
            period_specs = {
                'Anual': {'label': 'MAT', 'current_lag': 12, 'yoy_lag': None, 'show_yoy': False},
                'Semestral': {'label': 'SEM', 'current_lag': 6, 'yoy_lag': 12, 'show_yoy': include_same_period_last_year},
                'Trimestral': {'label': 'TRI', 'current_lag': 3, 'yoy_lag': 12, 'show_yoy': include_same_period_last_year},
            }
            wp_yoy_vars = calc_var_same_period_last_year(df_marca, COL_SELL_OUT, 0)
            p0_yoy_vars = calc_var_same_period_last_year(df_marca, COL_SELL_IN, 0)
            px_yoy_vars = calc_var_same_period_last_year(df_marca, COL_SELL_IN, pipeline) if int(pipeline) > 0 else None
            period_idx = {'Anual': 0, 'Semestral': 1, 'Trimestral': 2}
            vertical_rows: List[Dict[str, object]] = []
            for _, row in detail_df.iterrows():
                tipo = row.get('Tipo')
                spec = period_specs.get(tipo)
                if spec is None or pd.isna(current_dt):
                    row_dict = row.to_dict()
                    row_dict['_CompareLagMonths'] = np.nan
                    vertical_rows.append(row_dict)
                    continue

                label = spec['label']
                current_row = row.to_dict()
                current_row['Periodo'] = _build_period_text(label, current_dt, int(spec['current_lag']))
                current_row['_CompareLagMonths'] = int(spec['current_lag'])
                vertical_rows.append(current_row)

                idx = period_idx.get(tipo)
                if spec['show_yoy'] and spec['yoy_lag'] is not None and idx is not None:
                    yoy_row = row.to_dict()
                    yoy_row['Periodo'] = _build_period_text(label, current_dt, int(spec['yoy_lag']))
                    yoy_row['_CompareLagMonths'] = int(spec['yoy_lag'])
                    yoy_row['WP by Numerator'] = wp_yoy_vars[idx] if idx < len(wp_yoy_vars) else np.nan
                    yoy_row['Cliente P0'] = p0_yoy_vars[idx] if idx < len(p0_yoy_vars) else np.nan
                    if pipeline_detail_col:
                        if px_yoy_vars is not None and idx < len(px_yoy_vars):
                            yoy_row[pipeline_detail_col] = px_yoy_vars[idx]
                        else:
                            yoy_row[pipeline_detail_col] = np.nan
                    vertical_rows.append(yoy_row)
            detail_df = pd.DataFrame(vertical_rows)
        except Exception:
            pass

    if '_CompareLagMonths' not in detail_df.columns:
        detail_df['_CompareLagMonths'] = np.nan

    ordered_columns = ['Tipo', 'Periodo', 'WP by Numerator', 'Cliente P0']
    if pipeline_detail_col:
        ordered_columns.append(pipeline_detail_col)
    ordered_columns.append('_CompareLagMonths')
    detail_df = detail_df[[col for col in ordered_columns if col in detail_df.columns]]
    detail_df.reset_index(drop=True, inplace=True)
    return detail_df


def build_evolution_figure(
    df_marca: "pd.DataFrame",
    pipeline: int,
    lang_index: int,
    marca_nombre: str,
    variant: str = "classic",
) -> Optional["plt.Figure"]:
    if len(df_marca) < 24:
        return None
    df_evol = df_marca[[COL_DATA, COL_SELL_IN, COL_SELL_OUT]].copy()
    df_evol[COL_DATA] = pd.to_datetime(df_evol[COL_DATA])
    return generar_grafico_evolucion_mensual(
        df_evol,
        pipeline,
        lang_index,
        marca_nombre=marca_nombre,
        variant=variant,
    )


def _coverage_value_for_year_month(coverage_series: "pd.Series", year: int, month: int) -> float:
    if coverage_series is None or coverage_series.empty:
        return np.nan
    idx = pd.to_datetime(coverage_series.index, errors="coerce")
    values = pd.to_numeric(coverage_series, errors="coerce")
    clean_series = pd.Series(values.to_numpy(dtype=float), index=idx).dropna()
    clean_series = clean_series[~clean_series.index.isna()]
    if clean_series.empty:
        return np.nan
    matched = clean_series[(clean_series.index.year == year) & (clean_series.index.month == month)]
    if matched.empty:
        return np.nan
    return float(matched.iloc[-1])


def _coverage_value_to_number(value: float, round_coverage: bool) -> float | int:
    if pd.notna(value):
        return int(np.floor(float(value) + 0.5)) if round_coverage else round(float(value), 1)
    return 0


def _coverage_value_to_text(value: float, round_coverage: bool) -> str:
    if pd.notna(value):
        return str(int(np.floor(float(value) + 0.5))) if round_coverage else f"{float(value):.1f}"
    return "0" if round_coverage else "0.0"



def build_summary_and_bank_rows(
    pipeline: int,
    marca_nombre_limpio: str,
    subcategoria_nombre: str,
    coverage_series: "pd.Series",
    df_variations: "pd.DataFrame",
    averages: Dict[str, float],
    labels: Dict[Tuple[int, str], List[str] | str],
    lang_index: int,
    fabricante: str,
    pais_nombre: str,
    categoria_nombre: str,
    cesta_nombre: str,
    coverage_reason: str,
    measure_unit: str,
    coverage_type: str,
    ref_month_year: str,
    round_coverage: bool,
    summary_extra_months: Sequence[int],
    summary_extra_months_mode: str,
) -> Tuple[Dict[str, str], Dict[str, object], float, float, float, str]:
    ref_dt = dt.strptime(ref_month_year, '%m-%y')
    summary_columns, coverage_periods, _ = build_summary_columns(
        lang_index=lang_index,
        fabricante=fabricante,
        ref_dt=ref_dt,
        summary_extra_months=summary_extra_months,
        summary_extra_months_mode=summary_extra_months_mode,
    )
    _, _, base_prev, base_curr = build_summary_coverage_periods(
        ref_dt,
        summary_extra_months,
        summary_extra_months_mode,
    )
    coverage_anterior = _coverage_value_for_year_month(coverage_series, base_prev.year, base_prev.month)
    coverage_actual = _coverage_value_for_year_month(coverage_series, base_curr.year, base_curr.month)

    var_cliente_anual_y1 = df_variations.loc[df_variations['Tipo'] == 'Anual', f'Cliente P{pipeline}'].iloc[0]
    var_kantar_anual_y1 = df_variations.loc[df_variations['Tipo'] == 'Anual', 'WP by Numerator'].iloc[0]
    tendencia_alineada = "NO"
    if pd.notna(var_cliente_anual_y1) and pd.notna(var_kantar_anual_y1):
        if (var_cliente_anual_y1 * var_kantar_anual_y1) > 0:
            tendencia_alineada = "SI"
        elif var_cliente_anual_y1 == 0 and var_kantar_anual_y1 == 0:
            tendencia_alineada = "SI"

    cov_actual_val = _coverage_value_to_number(coverage_actual, round_coverage)
    cov_anterior_val = _coverage_value_to_number(coverage_anterior, round_coverage)
    estabilidad = (cov_actual_val - cov_anterior_val) if round_coverage else round(cov_actual_val - cov_anterior_val, 1)

    summary_row = {
        summary_columns[0]: marca_nombre_limpio,
        summary_columns[1]: pipeline,
        summary_columns[2]: f"{averages.get('Penet_MAT_Actual', 0):.1f}%",
        summary_columns[3]: f"{var_cliente_anual_y1*100:.1f}%" if pd.notna(var_cliente_anual_y1) else "0.0%",
        summary_columns[4]: f"{var_kantar_anual_y1*100:.1f}%" if pd.notna(var_kantar_anual_y1) else "0.0%",
    }

    coverage_col_idx = 5
    for period_dt in coverage_periods:
        cov_value = _coverage_value_for_year_month(coverage_series, period_dt.year, period_dt.month)
        summary_row[summary_columns[coverage_col_idx]] = _coverage_value_to_text(cov_value, round_coverage)
        coverage_col_idx += 1
    summary_row[summary_columns[-1]] = str(estabilidad) if round_coverage else f"{estabilidad:.1f}"

    banco_row = {
        'Periodo': ref_dt.date(),
        'Fabricante': fabricante,
        'Categoria': categoria_nombre,
        'Fabricante/Marca': marca_nombre_limpio,
        'Cesta': cesta_nombre,
        'Subcategoria': subcategoria_nombre,
        'Panel': 'PNC',
        'Unidad': measure_unit,
        'Razon': coverage_reason,
        'Pais': pais_nombre,
        'Ampliacion': 'SI',
        'Penet Media Ano Mov Atual': round(averages.get('Penet_MAT_Actual', 0), 1),
        'Penet Media Ano Mov Anterior': round(averages.get('Penet_MAT_Anterior', 0), 1),
        'Raw Buyers Media Ano Mov Atual': round(averages.get('Buyers_MAT_Actual', 0), 1),
        'Frecuencia Media Mensual': round(averages.get('Freq_MAT_Actual', 0), 1),
        'Pipeline': pipeline,
        'Cobertura Año Mov Actual': cov_actual_val,
        'Cobertura Año Mov Anterior': cov_anterior_val,
        '%VAR Cliente': round(var_cliente_anual_y1 * 100, 1) if pd.notna(var_cliente_anual_y1) else 0,
        '% VAR WP by Numerator': round(var_kantar_anual_y1 * 100, 1) if pd.notna(var_kantar_anual_y1) else 0,
        'Misma Tendencia': tendencia_alineada,
        'Estabilidad': estabilidad,
    }
    return summary_row, banco_row, cov_actual_val, cov_anterior_val, estabilidad, tendencia_alineada


COVERAGE_BANK_COLUMNS = [
    'Periodo', 'Fabricante', 'Categoria', 'Fabricante/Marca', 'Cesta', 'Subcategoria', 'Panel', 'Unidad',
    'Razon', 'Pais', 'Ampliacion', 'Penet Media Ano Mov Atual', 'Penet Media Ano Mov Anterior',
    'Raw Buyers Media Ano Mov Atual', 'Frecuencia Media Mensual', 'Pipeline', 'Cobertura Año Mov Actual',
    'Cobertura Año Mov Anterior', '%VAR Cliente', '% VAR WP by Numerator', 'Misma Tendencia', 'Estabilidad'
]


PIPELINE_REPORT_BASE_COLUMNS = [
    'Fabricante/Marca',
    'Cesta',
    'Penet Media Ano Mov Atual',
    'Raw Buyers Media Ano Mov Atual',
    'Frecuencia Media Mensual',
    'Pipeline',
    '%VAR Cliente',
    '% VAR WP by Numerator',
    'Estabilidad',
]
PIPELINE_REPORT_VARIATION_COLUMNS = [f'Variación Cliente P{p}' for p in range(1, 7)]
PIPELINE_REPORT_CORRELATION_COLUMNS = [f'Correlación P{p}' for p in range(1, 7)]
PIPELINE_REPORT_DIAGNOSTIC_COLUMNS = [
    'Confianza selección',
    'Gap variación seleccionado',
    'Correlación seleccionada',
    'Misma dirección selección',
    'Pipeline top correlación',
    'Correlación top',
    'Gap variación top correlación',
    'Pipeline top variación',
    'Gap variación top',
    'Correlación top variación',
]
PIPELINE_REPORT_AUTO_MODE_COLUMNS = [
    'Pipeline AUTO Correlación',
    'Correlación AUTO Correlación',
    'Gap variación AUTO Correlación',
    'Motivo AUTO Correlación',
    'Pipeline AUTO Balanceado',
    'Correlación AUTO Balanceado',
    'Gap variación AUTO Balanceado',
    'Motivo AUTO Balanceado',
    'Tipo de decisión AUTO Balanceado',
    'Conflicto AUTO Correlación vs Balanceado',
    'Pérdida de correlación Balanceado',
    'Mejora gap de variación Balanceado',
    'Revisión requerida',
]
PIPELINE_REPORT_REASON_COLUMN = 'Motivo de selección de pipeline'


def build_pipeline_report_columns(ref_month_year: str) -> List[str]:
    try:
        ref_dt = dt.strptime(ref_month_year, '%m-%y')
        prev_dt = ref_dt.replace(year=ref_dt.year - 1)
        prev_coverage_col = f"Cobertura {prev_dt.strftime('%m-%y')}"
        curr_coverage_col = f"Cobertura {ref_dt.strftime('%m-%y')}"
    except Exception:
        prev_coverage_col = 'Cobertura Año Mov Anterior'
        curr_coverage_col = 'Cobertura Año Mov Actual'
    return (
        PIPELINE_REPORT_BASE_COLUMNS[:8]
        + [prev_coverage_col, curr_coverage_col]
        + PIPELINE_REPORT_BASE_COLUMNS[8:]
        + PIPELINE_REPORT_VARIATION_COLUMNS
        + PIPELINE_REPORT_CORRELATION_COLUMNS
        + PIPELINE_REPORT_DIAGNOSTIC_COLUMNS
        + PIPELINE_REPORT_AUTO_MODE_COLUMNS
        + [PIPELINE_REPORT_REASON_COLUMN]
    )


def _round_report_number(value: object, digits: int = 1) -> object:
    return round(float(value), digits) if _is_finite_number(value) else np.nan


def _pipeline_selection_confidence(
    selected_candidate: Optional[OptimalPipelineCandidate],
    selection_reason: str,
    candidates: Sequence[OptimalPipelineCandidate] = (),
) -> str:
    reason_norm = _normalize_metadata_match_text(selection_reason)
    if selected_candidate is None:
        return "Baja"
    corr = selected_candidate.current_correlation
    gap = selected_candidate.variation_distance_points
    top_current_corr = max(
        (
            float(candidate.current_correlation)
            for candidate in candidates
            if _is_finite_number(candidate.current_correlation)
            and float(candidate.current_correlation) > 0
        ),
        default=np.nan,
    )
    corr_loss = (
        float(top_current_corr) - float(corr)
        if _is_finite_number(top_current_corr) and _is_finite_number(corr)
        else np.nan
    )
    if (
        selected_candidate.current_trend_match
        and _is_finite_number(corr)
        and _is_finite_number(gap)
        and float(corr) >= OPTIMAL_PIPELINE_STRONG_CORR_MIN
        and float(gap) <= OPTIMAL_PIPELINE_VISUAL_VARIATION_DISTANCE_POINTS
        and (not _is_finite_number(corr_loss) or float(corr_loss) <= 0.05)
    ):
        return "Alta"
    if selected_candidate.recent_shipment_outlier and "outlier" in reason_norm:
        return "Media"
    if (
        selected_candidate.current_trend_match
        and _is_finite_number(corr)
        and float(corr) >= OPTIMAL_PIPELINE_BALANCED_CORR_MIN
        and (not _is_finite_number(corr_loss) or float(corr_loss) <= AUTO_PIPELINE_CONFLICT_HIGH_CORR_LOSS)
    ):
        return "Media"
    if (
        selected_candidate.historical_trend_match
        and _historical_support_score(selected_candidate) >= OPTIMAL_PIPELINE_HISTORICAL_CORR_MIN
    ):
        return "Media"
    return "Baja"


def build_pipeline_selection_diagnostics(
    selected_pipeline: object,
    candidates: Sequence[OptimalPipelineCandidate],
    selection_reason: str,
) -> Dict[str, object]:
    try:
        selected_pipeline_int = int(selected_pipeline)
    except Exception:
        selected_pipeline_int = None

    candidate_by_pipeline = {candidate.pipeline: candidate for candidate in candidates}
    selected_candidate = (
        candidate_by_pipeline.get(selected_pipeline_int)
        if selected_pipeline_int is not None
        else None
    )
    correlation_candidates = [
        candidate for candidate in candidates
        if _is_finite_number(candidate.current_correlation)
    ]
    trend_candidates = [
        candidate for candidate in candidates
        if candidate.current_trend_match
        and _is_finite_number(candidate.current_correlation)
        and float(candidate.current_correlation) > 0
    ]
    variation_candidates = [
        candidate for candidate in trend_candidates
        if _is_finite_number(candidate.variation_distance_points)
    ]
    top_corr = (
        max(correlation_candidates, key=lambda candidate: float(candidate.current_correlation))
        if correlation_candidates
        else None
    )
    top_variation = (
        sorted(
            variation_candidates,
            key=lambda candidate: (
                float(candidate.variation_distance_points),
                candidate.pipeline,
            ),
        )[0]
        if variation_candidates
        else None
    )

    diagnostics = {
        'Confianza selección': _pipeline_selection_confidence(selected_candidate, selection_reason, candidates),
        'Gap variación seleccionado': _round_report_number(
            selected_candidate.variation_distance_points if selected_candidate is not None else np.nan,
            2,
        ),
        'Correlación seleccionada': _round_report_number(
            selected_candidate.current_correlation if selected_candidate is not None else np.nan,
            3,
        ),
        'Misma dirección selección': (
            "SI" if selected_candidate is not None and selected_candidate.current_trend_match else "NO"
        ),
        'Pipeline top correlación': top_corr.pipeline if top_corr is not None else "",
        'Correlación top': _round_report_number(
            top_corr.current_correlation if top_corr is not None else np.nan,
            3,
        ),
        'Gap variación top correlación': _round_report_number(
            top_corr.variation_distance_points if top_corr is not None else np.nan,
            2,
        ),
        'Pipeline top variación': top_variation.pipeline if top_variation is not None else "",
        'Gap variación top': _round_report_number(
            top_variation.variation_distance_points if top_variation is not None else np.nan,
            2,
        ),
        'Correlación top variación': _round_report_number(
            top_variation.current_correlation if top_variation is not None else np.nan,
            3,
        ),
    }
    return diagnostics


def _pipeline_decision_type(selection_reason: str) -> str:
    reason_norm = _normalize_metadata_match_text(selection_reason)
    if "categoria con perfil" in reason_norm:
        return "Restricción/prior de categoría"
    if "ajuste casi exacto" in reason_norm or "alineacion visual" in reason_norm:
        return "Override por alineación de variación"
    if "ajuste balanceado material" in reason_norm:
        return "Balance material de variación"
    if "respaldo historico" in reason_norm:
        return "Respaldo histórico"
    if "outlier" in reason_norm:
        return "Fallback por outlier"
    if "indicado en la hoja" in reason_norm or "pipeline indicado" in reason_norm:
        return "Prefijo de hoja"
    if "sin senal concluyente" in reason_norm or "fallback" in reason_norm:
        return "Fallback operativo"
    return "Correlación y variación actuales"


def build_auto_pipeline_comparison_diagnostics(
    comparison: AutoPipelineComparison,
) -> Dict[str, object]:
    candidates = comparison.balanced.candidates or comparison.correlation.candidates
    candidate_by_pipeline = {candidate.pipeline: candidate for candidate in candidates}
    correlation_candidate = candidate_by_pipeline.get(comparison.correlation.pipeline)
    balanced_candidate = candidate_by_pipeline.get(comparison.balanced.pipeline)

    correlation_corr = (
        correlation_candidate.current_correlation if correlation_candidate is not None else np.nan
    )
    balanced_corr = balanced_candidate.current_correlation if balanced_candidate is not None else np.nan
    correlation_gap = (
        correlation_candidate.variation_distance_points if correlation_candidate is not None else np.nan
    )
    balanced_gap = balanced_candidate.variation_distance_points if balanced_candidate is not None else np.nan
    corr_loss = (
        max(0.0, float(correlation_corr) - float(balanced_corr))
        if _is_finite_number(correlation_corr) and _is_finite_number(balanced_corr)
        else np.nan
    )
    gap_improvement = (
        float(correlation_gap) - float(balanced_gap)
        if _is_finite_number(correlation_gap) and _is_finite_number(balanced_gap)
        else np.nan
    )

    if comparison.correlation.pipeline == comparison.balanced.pipeline:
        conflict = "Sin conflicto"
    elif not _is_finite_number(corr_loss):
        conflict = "No evaluable"
    elif float(corr_loss) >= AUTO_PIPELINE_CONFLICT_HIGH_CORR_LOSS:
        conflict = "Alto"
    elif float(corr_loss) >= AUTO_PIPELINE_CONFLICT_MEDIUM_CORR_LOSS:
        conflict = "Medio"
    else:
        conflict = "Bajo"

    confidence = _pipeline_selection_confidence(
        balanced_candidate,
        comparison.balanced.reason,
        candidates,
    )
    correlation_reason_norm = _normalize_metadata_match_text(comparison.correlation.reason)
    review_required = (
        conflict in {"Alto", "No evaluable"}
        or confidence == "Baja"
        or "requiere revision" in correlation_reason_norm
        or "sin correlaciones calculables" in correlation_reason_norm
    )

    return {
        'Pipeline AUTO Correlación': comparison.correlation.pipeline,
        'Correlación AUTO Correlación': _round_report_number(correlation_corr, 3),
        'Gap variación AUTO Correlación': _round_report_number(correlation_gap, 2),
        'Motivo AUTO Correlación': comparison.correlation.reason,
        'Pipeline AUTO Balanceado': comparison.balanced.pipeline,
        'Correlación AUTO Balanceado': _round_report_number(balanced_corr, 3),
        'Gap variación AUTO Balanceado': _round_report_number(balanced_gap, 2),
        'Motivo AUTO Balanceado': comparison.balanced.reason,
        'Tipo de decisión AUTO Balanceado': _pipeline_decision_type(comparison.balanced.reason),
        'Conflicto AUTO Correlación vs Balanceado': conflict,
        'Pérdida de correlación Balanceado': _round_report_number(corr_loss, 3),
        'Mejora gap de variación Balanceado': _round_report_number(gap_improvement, 2),
        'Revisión requerida': "SI" if review_required else "NO",
    }


def build_pipeline_report_row(
    bank_row: Dict[str, object],
    df_variations: "pd.DataFrame",
    candidates: Sequence[OptimalPipelineCandidate],
    selection_reason: str,
    ref_month_year: str,
    auto_comparison: Optional[AutoPipelineComparison] = None,
) -> Dict[str, object]:
    report_columns = build_pipeline_report_columns(ref_month_year)
    try:
        ref_dt = dt.strptime(ref_month_year, '%m-%y')
        prev_dt = ref_dt.replace(year=ref_dt.year - 1)
        prev_coverage_col = f"Cobertura {prev_dt.strftime('%m-%y')}"
        curr_coverage_col = f"Cobertura {ref_dt.strftime('%m-%y')}"
    except Exception:
        prev_coverage_col = 'Cobertura Año Mov Anterior'
        curr_coverage_col = 'Cobertura Año Mov Actual'

    row: Dict[str, object] = {col: "" for col in report_columns}
    for col in [
        'Fabricante/Marca',
        'Cesta',
        'Penet Media Ano Mov Atual',
        'Raw Buyers Media Ano Mov Atual',
        'Frecuencia Media Mensual',
        'Pipeline',
        '%VAR Cliente',
        '% VAR WP by Numerator',
        'Estabilidad',
    ]:
        row[col] = bank_row.get(col, "")
    row[prev_coverage_col] = bank_row.get('Cobertura Año Mov Anterior', "")
    row[curr_coverage_col] = bank_row.get('Cobertura Año Mov Actual', "")

    try:
        annual_row = df_variations.loc[df_variations['Tipo'] == 'Anual'].iloc[0]
    except Exception:
        annual_row = {}
    for p in range(1, 7):
        value = annual_row.get(f'Cliente P{p}', np.nan) if hasattr(annual_row, "get") else np.nan
        row[f'Variación Cliente P{p}'] = round(float(value) * 100, 1) if pd.notna(value) else np.nan

    candidate_by_pipeline = {candidate.pipeline: candidate for candidate in candidates}
    for p in range(1, 7):
        candidate = candidate_by_pipeline.get(p)
        corr = candidate.current_correlation if candidate is not None else np.nan
        row[f'Correlación P{p}'] = round(float(corr), 3) if _is_finite_number(corr) else np.nan
    row.update(
        build_pipeline_selection_diagnostics(
            bank_row.get('Pipeline', ""),
            candidates,
            selection_reason,
        )
    )
    if auto_comparison is None:
        try:
            balanced_pipeline = int(bank_row.get('Pipeline', 1))
        except Exception:
            balanced_pipeline = 1
        balanced_selection = OptimalPipelineSelection(
            balanced_pipeline,
            selection_reason,
            tuple(candidates),
        )
        auto_comparison = AutoPipelineComparison(
            correlation=select_correlation_pipeline(candidates),
            balanced=balanced_selection,
        )
    row.update(build_auto_pipeline_comparison_diagnostics(auto_comparison))
    row[PIPELINE_REPORT_REASON_COLUMN] = selection_reason or ""
    return row


def generate_presentation_and_bank(
    root_dir: str,
    excel_file_obj: "pd.ExcelFile",
    marcas: Sequence[str],
    pais_nombre: str,
    category_code: str,
    categories_df: "pd.DataFrame",
    categoria_nombre: str,
    categoria_nombre_corto: str,
    fabricante: str,
    cesta_nombre: str,
    coverage_label: str,
    coverage_type: str,
    coverage_reason: str,
    ref_month_year: str,
    carpeta_salida: str,
    nombre_base_archivo: str,
    include_english: bool,
    trend_axis: str,
    trend_granularity: str,
    variations_box_style: str,
    coverage_slide_variant: str,
    evolution_slide_variant: str,
    round_coverage: bool,
    summary_extra_months: Sequence[int],
    summary_extra_months_mode: str,
    variations_include_same_period_last_year: bool = True,
    variations_compact_period_labels: bool = False,
    optimal_pipeline_mode: bool = False,
    elapsed_seconds_fn: Optional[Callable[[], Optional[float]]] = None,
) -> Tuple[str, "pd.DataFrame", "pd.DataFrame", "pd.DataFrame"]:
    chosen_lang, lang_index = determine_language(include_english, pais_nombre)
    ppt, tmp_ppt_path = copy_and_prune_template(root_dir, chosen_lang)
    labels = build_labels(lang_index, fabricante, ref_month_year, summary_extra_months, summary_extra_months_mode)
    builder = SlideBuilder(
        ppt,
        lang_index,
        labels,
        coverage_label,
        coverage_type=coverage_type,
        ref_month_year=ref_month_year,
        manufacturer_name=fabricante,
        country_name=pais_nombre,
        category_name_display=(categoria_nombre_corto or categoria_nombre),
        tipo_eje_tend=trend_axis,
        trend_granularity=trend_granularity,
        variations_box_style=variations_box_style,
        coverage_slide_variant=coverage_slide_variant,
        variations_compact_period_labels=variations_compact_period_labels,
    )
    builder.configure_cover(pais_nombre, fabricante, categoria_nombre, ref_month_year, chosen_lang)

    summary_rows: List[Dict[str, str]] = []
    summary_rows_by_period: "OrderedDict[str, List[Dict[str, str]]]" = OrderedDict()
    bank_rows: List[Dict[str, object]] = []
    pipeline_report_rows: List[Dict[str, object]] = []
    low_penetration_brands: List[str] = []
    brand_section_map: Dict[str, List[int]] = {}
    current_section_title: Optional[str] = None
    current_metadata_group: Optional[str] = None
    current_metadata_category_code: Optional[str] = None

    total_slides_to_generate = 0
    needs_mult_metadata_hints = _requires_metadata_category_resolution(category_code)
    for marca_sheet_name in marcas:
        df_marca_ppt, _ = load_and_preprocess_sheet(
            excel_file_obj,
            marca_sheet_name,
            include_metadata_hints=False,
        )
        if df_marca_ppt is None:
            continue
        forced_pipeline = extract_forced_pipeline_from_sheet_name(marca_sheet_name)
        if optimal_pipeline_mode:
            pipelines_to_run = [1]
        else:
            pipelines_to_run = [int(forced_pipeline)] if forced_pipeline is not None else list(range(7))
        n_slides_marca = len(pipelines_to_run) * (2 + (1 if len(df_marca_ppt) >= 24 else 0))
        total_slides_to_generate += n_slides_marca

    progress = Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        BarColumn(),
        TextColumn("[progress.percentage]{task.percentage:>3.0f}%"),
        TextColumn("{task.completed}/{task.total}"),
        TimeElapsedColumn(),
        TimeRemainingColumn(),
        transient=True,
    )

    with progress:
        task_id = progress.add_task("Creando Diapositivas PPT", total=total_slides_to_generate + 1)
        for marca_sheet_name in marcas:
            df_marca_ppt, measure_unit = load_and_preprocess_sheet(
                excel_file_obj,
                marca_sheet_name,
                include_metadata_hints=needs_mult_metadata_hints,
            )
            if df_marca_ppt is None:
                continue
            marca_nombre_limpio = re.sub(r"(?i)^p[0-6]_", "", marca_sheet_name)
            subcategoria_nombre = extract_sheet_subcategory(marca_sheet_name)
            forced_pipeline = extract_forced_pipeline_from_sheet_name(marca_sheet_name)
            metadata_group_title, next_metadata_group = build_metadata_group_for_sheet(
                marca_sheet_name,
                current_metadata_group,
                fabricante,
            )
            section_title, next_section_title = build_section_title_for_sheet(
                marca_sheet_name,
                current_section_title,
            )
            sheet_bank_metadata = resolve_sheet_bank_metadata(
                category_code=category_code,
                fabricante=fabricante,
                marca_nombre_limpio=marca_nombre_limpio,
                section_title=metadata_group_title,
                categories_df=categories_df,
                default_pais_nombre=pais_nombre,
                default_cesta_nombre=cesta_nombre,
                default_categoria_nombre=categoria_nombre,
                default_categoria_nombre_corto=categoria_nombre_corto,
                sheet_metadata_hints=df_marca_ppt.attrs.get("sheet_metadata_hints"),
                inherited_category_code=current_metadata_category_code,
            )
            current_metadata_group = next_metadata_group
            if _requires_metadata_category_resolution(category_code) and is_total_group_sheet(marca_sheet_name):
                next_code = _normalize_category_code(sheet_bank_metadata.categoria_codigo)
                current_metadata_category_code = (
                    next_code
                    if next_code and not _requires_metadata_category_resolution(next_code)
                    else None
                )
            issues_detected = detect_brand_data_issues(df_marca_ppt, window=0)
            if issues_detected:
                for issue in issues_detected:
                    if issue == "zero_dash":
                        notify_zero_months_exception(marca_nombre_limpio)
                    elif issue == "negative":
                        notify_negative_values_exception(marca_nombre_limpio)
            df_coverage = compute_coverage_dataframe(
                df_marca_ppt,
                pais_nombre,
                coverage_type,
                round_coverage,
                marca_nombre=marca_nombre_limpio,
            )
            df_variations = compute_variations_dataframe(df_marca_ppt)
            optimal_selection: Optional[OptimalPipelineSelection] = None
            auto_comparison: Optional[AutoPipelineComparison] = None
            selection_reason_by_pipeline: Dict[int, str] = {}
            pipeline_candidates = build_optimal_pipeline_candidates(
                df_marca_ppt,
                df_variations,
                forced_pipeline,
            )
            if optimal_pipeline_mode:
                optimal_selection = select_optimal_pipeline(
                    df_marca_ppt,
                    df_variations,
                    forced_pipeline=forced_pipeline,
                    category_code=sheet_bank_metadata.categoria_codigo,
                    category_name=sheet_bank_metadata.categoria_nombre,
                    basket_name=sheet_bank_metadata.cesta_nombre,
                )
                pipelines_to_run = [optimal_selection.pipeline]
                pipeline_candidates = optimal_selection.candidates
                auto_comparison = AutoPipelineComparison(
                    correlation=select_correlation_pipeline(
                        pipeline_candidates,
                        forced_pipeline=forced_pipeline,
                    ),
                    balanced=optimal_selection,
                )
                selection_reason_by_pipeline[optimal_selection.pipeline] = optimal_selection.reason
                print(
                    Fore.CYAN
                    + f"Pipeline AUTO Balanceado para {marca_nombre_limpio}: P{optimal_selection.pipeline} "
                    + f"({optimal_selection.reason}). "
                    + f"AUTO Correlación: P{auto_comparison.correlation.pipeline}."
                )
            else:
                pipelines_to_run = [int(forced_pipeline)] if forced_pipeline is not None else list(range(7))
                for pipeline_option in pipelines_to_run:
                    if forced_pipeline is not None:
                        selection_reason_by_pipeline[int(pipeline_option)] = "Pipeline indicado en el nombre de la hoja"
                    else:
                        selection_reason_by_pipeline[int(pipeline_option)] = "Pipeline generado por configuración no automática"
            averages = compute_averages(df_marca_ppt)
            sheet_ref_value = pd.to_datetime(df_marca_ppt[COL_DATA].iloc[-1], errors="coerce")
            sheet_ref_month_year = (
                sheet_ref_value.strftime("%m-%y")
                if not pd.isna(sheet_ref_value)
                else ref_month_year
            )
            notify_buyers_threshold(marca_nombre_limpio, averages.get('Buyers_MAT_Actual'))
            try:
                buyers_val = averages.get('Buyers_MAT_Actual')
                if buyers_val is not None and not pd.isna(buyers_val) and float(buyers_val) < 200:
                    if marca_nombre_limpio not in low_penetration_brands:
                        low_penetration_brands.append(marca_nombre_limpio)
            except Exception:
                pass
            df_trend_plot = compute_trend_plot_df(df_marca_ppt)
            for pipeline in pipelines_to_run:
                coverage_series = df_coverage[f'P{pipeline}']
                var_cliente_mat = df_variations.loc[df_variations['Tipo'] == 'Anual', f'Cliente P{pipeline}'].iloc[0]
                var_kantar_mat = df_variations.loc[df_variations['Tipo'] == 'Anual', 'WP by Numerator'].iloc[0]
                trend_following = False
                if pd.notna(var_cliente_mat) and pd.notna(var_kantar_mat):
                    if (var_cliente_mat * var_kantar_mat) > 0 or (var_cliente_mat == 0 and var_kantar_mat == 0):
                        trend_following = True
                current_year_correlation = compute_pipeline_current_year_correlation(df_marca_ppt, pipeline)
                annual_var_cliente_y2 = compute_previous_year_annual_variation(df_marca_ppt, COL_SELL_IN, pipeline)
                annual_var_wp_y2 = compute_previous_year_annual_variation(df_marca_ppt, COL_SELL_OUT, 0)
                variation_table = build_variation_table(
                    fabricante,
                    labels,
                    lang_index,
                    pipeline,
                    sheet_ref_month_year,
                    var_cliente_mat,
                    var_kantar_mat,
                )
                variations_detail = build_variations_detail_table(
                    df_variations,
                    pipeline,
                    df_marca_ppt,
                    include_same_period_last_year=variations_include_same_period_last_year,
                )
                evolution_figure = build_evolution_figure(
                    df_marca_ppt,
                    pipeline,
                    lang_index,
                    marca_nombre_limpio,
                    variant=evolution_slide_variant,
                )
                assets = PipelineAssets(
                    pipeline=pipeline,
                    marca=marca_nombre_limpio,
                    coverage_series=coverage_series,
                    penetration_series=df_marca_ppt.set_index(COL_DATA)[COL_PENET].loc[coverage_series.dropna().index],
                    variation_table=variation_table,
                    trend_plot_df=df_trend_plot,
                    variations_detail=variations_detail,
                    evolution_figure=evolution_figure,
                    buyers_mat_actual=averages.get('Buyers_MAT_Actual'),
                    penet_mat_actual=averages.get('Penet_MAT_Actual'),
                    penet_mat_anterior=averages.get('Penet_MAT_Anterior'),
                    annual_var_cliente_y1=var_cliente_mat,
                    annual_var_cliente_y2=annual_var_cliente_y2,
                    annual_var_wp_y1=var_kantar_mat,
                    annual_var_wp_y2=annual_var_wp_y2,
                    measure_unit=measure_unit,
                    current_year_correlation=current_year_correlation,
                    trend_following=trend_following,
                )
                slide_start_idx = len(ppt.slides)
                slides_created = builder.add_pipeline_slides(
                    assets,
                    marca_nombre_limpio=marca_nombre_limpio,
                    lang_index=lang_index,
                    coverage_label=builder.coverage_label,
                    progress=progress,
                    task_id=task_id,
                )
                current_section_title = next_section_title
                register_section_slide_range(
                    brand_section_map,
                    section_title,
                    slide_start_idx,
                    slides_created,
                )
                summary_row, bank_row, _, _, _, _ = build_summary_and_bank_rows(
                    pipeline=pipeline,
                    marca_nombre_limpio=marca_nombre_limpio,
                    subcategoria_nombre=subcategoria_nombre,
                    coverage_series=coverage_series,
                    df_variations=df_variations,
                    averages=averages,
                    labels=labels,
                    lang_index=lang_index,
                    fabricante=fabricante,
                    pais_nombre=sheet_bank_metadata.pais_nombre,
                    categoria_nombre=sheet_bank_metadata.categoria_nombre,
                    cesta_nombre=sheet_bank_metadata.cesta_nombre,
                    coverage_reason=coverage_reason,
                    measure_unit=measure_unit,
                    coverage_type=coverage_type,
                    ref_month_year=sheet_ref_month_year,
                    round_coverage=round_coverage,
                    summary_extra_months=summary_extra_months,
                    summary_extra_months_mode=summary_extra_months_mode,
                )
                summary_rows.append(summary_row)
                summary_rows_by_period.setdefault(sheet_ref_month_year, []).append(summary_row)
                bank_rows.append(bank_row)
                if optimal_pipeline_mode:
                    pipeline_report_rows.append(
                        build_pipeline_report_row(
                            bank_row=bank_row,
                            df_variations=df_variations,
                            candidates=pipeline_candidates,
                            selection_reason=selection_reason_by_pipeline.get(
                                int(pipeline),
                                "Pipeline generado por configuración actual",
                            ),
                            ref_month_year=sheet_ref_month_year,
                            auto_comparison=auto_comparison,
                        )
                    )
        progress.update(task_id, advance=1)

    summary_groups: List[Tuple[str, "pd.DataFrame"]] = []
    for period_token, rows_for_period in summary_rows_by_period.items():
        try:
            period_dt = dt.strptime(period_token, "%m-%y")
            period_columns, _, _ = build_summary_columns(
                lang_index=lang_index,
                fabricante=fabricante,
                ref_dt=period_dt,
                summary_extra_months=summary_extra_months,
                summary_extra_months_mode=summary_extra_months_mode,
            )
        except Exception:
            period_columns = labels[(lang_index, 'Summary')]
        group_df = pd.DataFrame(rows_for_period)
        if not group_df.empty:
            group_df = group_df.reindex(columns=period_columns)
        summary_groups.append((period_token, group_df))

    if summary_groups:
        df_summary = pd.concat([group_df for _, group_df in summary_groups], ignore_index=True, sort=False)
    else:
        df_summary = pd.DataFrame(summary_rows)
        if not df_summary.empty:
            df_summary = df_summary.reindex(columns=labels[(lang_index, 'Summary')])
    df_bank = pd.DataFrame(bank_rows, columns=COVERAGE_BANK_COLUMNS)
    df_pipeline_report = pd.DataFrame(pipeline_report_rows)
    if not df_pipeline_report.empty:
        df_pipeline_report = df_pipeline_report.reindex(columns=build_pipeline_report_columns(ref_month_year))

    builder.add_summary_slide(
        df_summary,
        pais_nombre,
        categoria_nombre,
        low_penetration_brands=low_penetration_brands,
        summary_groups=summary_groups,
        df_bank=df_bank,
    )
    builder.insert_thanks_text(chosen_lang)
    builder.reorder_summary_and_credit()

    section_slide_map: Dict[str, List[int]] = {}
    intro_title = "Intro" if chosen_lang == "EN" else ("Inicio" if chosen_lang == "ES" else "Inicio")
    summary_title = "Summary" if chosen_lang == "EN" else ("Resumen" if chosen_lang == "ES" else "Resumo")
    closing_title = "Closing" if chosen_lang == "EN" else ("Cierre" if chosen_lang == "ES" else "Fechamento")
    intro_count = min(6, len(ppt.slides))
    if intro_count > 0:
        register_section_slide_range(section_slide_map, intro_title, 0, intro_count)
    if len(ppt.slides) > 6:
        register_section_slide_range(section_slide_map, summary_title, 6, 1)
    for title, indexes in brand_section_map.items():
        for slide_idx in indexes:
            register_section_slide_range(section_slide_map, title, slide_idx, 1)
    if len(ppt.slides) > 7:
        register_section_slide_range(section_slide_map, closing_title, len(ppt.slides) - 1, 1)

    nombre_ppt_final = build_bounded_output_filename(carpeta_salida, f"{nombre_base_archivo}.pptx")
    ruta_ppt_final = os.path.join(carpeta_salida, nombre_ppt_final)
    summary_slide_index = 6 if len(ppt.slides) > 7 else max(0, len(ppt.slides) - 1)
    summary_table_specs = [
        (
            [builder._format_summary_header_label(col_name) for col_name in group_df.columns],
            builder._compute_summary_table_column_widths(group_df, ppt.slide_width - (2 * Inches(0.5))),
        )
        for _, group_df in summary_groups
        if group_df is not None and not group_df.empty
    ]

    def write_final_presentation() -> None:
        if os.path.exists(ruta_ppt_final):
            os.remove(ruta_ppt_final)
        ppt.save(ruta_ppt_final)
        for summary_table_headers, summary_table_widths in summary_table_specs:
            if not summary_table_headers or not summary_table_widths:
                continue
            apply_table_grid_widths_in_pptx(
                ruta_ppt_final,
                slide_index=summary_slide_index,
                header_row=summary_table_headers,
                column_widths=summary_table_widths,
            )
        apply_summary_table_border_style_in_pptx(ruta_ppt_final)
        apply_powerpoint_sections(ruta_ppt_final, section_slide_map)
        apply_variation_table_internal_borders_in_pptx(ruta_ppt_final)
        apply_trend_variation_table_transparent_style_in_pptx(ruta_ppt_final)

    run_file_write_with_retry(
        ruta_ppt_final,
        action_label="guardar la presentación final",
        operation=write_final_presentation,
        elapsed_seconds_fn=elapsed_seconds_fn,
    )

    return ruta_ppt_final, df_summary, df_bank, df_pipeline_report


def save_coverage_bank(
    df_bank: "pd.DataFrame",
    carpeta_salida: str,
    nombre_base_archivo: str,
    fabricante: str,
    categoria_nombre: str,
    categoria_nombre_corto: str,
    pais_nombre: str,
    ref_month_year: str,
    coverage_label: str,
    coverage_type: str,
    coverage_slide_variant: str,
    output_descriptor: str = "",
    elapsed_seconds_fn: Optional[Callable[[], Optional[float]]] = None,
) -> str:
    df_bank = df_bank.copy()
    try:
        mes_ejecucion_dt = datetime.now().date().replace(day=1)
        if 'Mes_Ejecucion' not in df_bank.columns:
            df_bank.insert(0, 'Mes_Ejecucion', mes_ejecucion_dt)
        else:
            df_bank['Mes_Ejecucion'] = mes_ejecucion_dt
    except Exception as exc:
        print(f"{Fore.YELLOW}Advertencia: No se pudo agregar la columna 'Mes_Ejecucion': {exc}")
    categoria_para_banco = build_output_category_segment(categoria_nombre_corto or categoria_nombre, output_descriptor)
    nombre_banco_final = build_bounded_output_filename(
        carpeta_salida,
        f"Banco_{fabricante}_{categoria_para_banco}_{pais_nombre}_{ref_month_year}_{coverage_label}.xlsx",
    )
    ruta_banco_final = os.path.join(carpeta_salida, nombre_banco_final)

    def write_bank_file() -> None:
        df_bank.to_excel(ruta_banco_final, index=False)
        from openpyxl import load_workbook as _wb_load
        from openpyxl.styles import PatternFill as _PatternFill, Font as _Font, Alignment as _Alignment, Border as _Border, Side as _Side
        wb_bank = _wb_load(ruta_banco_final)
        for ws in wb_bank.worksheets:
            header_map = {}
            for cell in ws[1]:
                if cell.value is not None:
                    header_map[str(cell.value).strip().lower()] = cell.column
            for header_name in ['periodo', 'mes_ejecucion']:
                col_idx = header_map.get(header_name)
                if col_idx is None:
                    continue
                for r in range(2, ws.max_row + 1):
                    c = ws.cell(row=r, column=col_idx)
                    c.number_format = 'mmm-yy'
            if ws.title.lower() != "summary":
                autofit_worksheet_columns(ws, min_width=11.0, max_width=34.0, padding=2.0)

        if normalize_coverage_slide_variant(coverage_slide_variant) == "pg" and not df_bank.empty:
            if "summary" in wb_bank.sheetnames:
                del wb_bank["summary"]
            ws_sum = wb_bank.create_sheet("summary", 0)

            try:
                ref_dt = dt.strptime(ref_month_year, "%m-%y")
                month_label = ref_dt.strftime("%b")
                prev_year_label = f"{int(ref_dt.year) - 1}"
                curr_year_label = f"{int(ref_dt.year)}"
                prev_mat_label = f"MAT {month_label}{str(int(ref_dt.year) - 1)[-2:]}"
                curr_mat_label = f"MAT {month_label}{str(int(ref_dt.year))[-2:]}"
            except Exception:
                prev_year_label = "Y-1"
                curr_year_label = "Y"
                prev_mat_label = "MAT Y-1"
                curr_mat_label = "MAT Y"

            header_fill = _PatternFill(fill_type="solid", fgColor="8FA9C3")
            header_font = _Font(color="FFFFFF", bold=True)
            body_font = _Font(color="000000", bold=False)
            body_align = _Alignment(horizontal="center", vertical="center")
            header_align = _Alignment(horizontal="center", vertical="center", wrap_text=True)
            border = _Border(
                left=_Side(style="thin", color="000000"),
                right=_Side(style="thin", color="000000"),
                top=_Side(style="thin", color="000000"),
                bottom=_Side(style="thin", color="000000"),
            )

            headers_top = {
                "A1": "Brand",
                "B1": "Annual Penetration %",
                "D1": "PIPELINE",
                "E1": f"Var % ({curr_year_label} vs {prev_year_label})",
                "G1": "Coverage",
                "I1": "Relative Coverage",
                "K1": "Var. pp",
            }
            headers_bottom = {
                "B2": prev_mat_label,
                "C2": curr_mat_label,
                "E2": f"{fabricante.upper()} WITH\nPIPELINE",
                "F2": "Worldpanel by\nNumerator",
                "G2": prev_year_label,
                "H2": curr_year_label,
                "I2": prev_year_label,
                "J2": curr_year_label,
            }

            ws_sum.merge_cells("A1:A2")
            ws_sum.merge_cells("B1:C1")
            ws_sum.merge_cells("D1:D2")
            ws_sum.merge_cells("E1:F1")
            ws_sum.merge_cells("G1:H1")
            ws_sum.merge_cells("I1:J1")
            ws_sum.merge_cells("K1:K2")

            for cell_ref, value in headers_top.items():
                cell = ws_sum[cell_ref]
                cell.value = value
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = header_align
                cell.border = border
            for cell_ref, value in headers_bottom.items():
                cell = ws_sum[cell_ref]
                cell.value = value
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = header_align
                cell.border = border

            for row_idx in (1, 2):
                for col_idx in range(1, 12):
                    ws_sum.cell(row=row_idx, column=col_idx).border = border

            def _derive_pg_cov_values(bank_row: object) -> Tuple[Optional[float], Optional[float], Optional[float], Optional[float]]:
                try:
                    abs_curr = float(bank_row.get('Cobertura Año Mov Actual'))
                except Exception:
                    abs_curr = np.nan
                try:
                    abs_prev = float(bank_row.get('Cobertura Año Mov Anterior'))
                except Exception:
                    abs_prev = np.nan
                try:
                    pop_val = get_population_coverage_percent(str(bank_row.get('Pais', pais_nombre))) / 100.0
                except Exception:
                    pop_val = 0.0

                ctype = (coverage_type or "").strip().lower()
                if ctype == "auto":
                    ctype = "absoluta"
                if ctype == "relativa":
                    rel_curr = abs_curr
                    rel_prev = abs_prev
                    abs_curr = (abs_curr * pop_val) if pop_val > 0 and pd.notna(abs_curr) else np.nan
                    abs_prev = (abs_prev * pop_val) if pop_val > 0 and pd.notna(abs_prev) else np.nan
                else:
                    rel_curr = (abs_curr / pop_val) if pop_val > 0 and pd.notna(abs_curr) else np.nan
                    rel_prev = (abs_prev / pop_val) if pop_val > 0 and pd.notna(abs_prev) else np.nan
                return abs_prev, abs_curr, rel_prev, rel_curr

            summary_start_row = 3
            for idx, (_, bank_row) in enumerate(df_bank.iterrows(), start=summary_start_row):
                abs_prev, abs_curr, rel_prev, rel_curr = _derive_pg_cov_values(bank_row)
                row_values = [
                    bank_row.get('Fabricante/Marca', ''),
                    None,
                    None,
                    bank_row.get('Pipeline', ''),
                    (float(bank_row.get('%VAR Cliente', 0)) / 100.0) if pd.notna(bank_row.get('%VAR Cliente', np.nan)) else None,
                    (float(bank_row.get('% VAR WP by Numerator', 0)) / 100.0) if pd.notna(bank_row.get('% VAR WP by Numerator', np.nan)) else None,
                    abs_prev,
                    abs_curr,
                    rel_prev,
                    rel_curr,
                    (abs_curr - abs_prev) if pd.notna(abs_curr) and pd.notna(abs_prev) else None,
                ]
                for col_idx, value in enumerate(row_values, start=1):
                    cell = ws_sum.cell(row=idx, column=col_idx)
                    cell.value = value
                    cell.font = body_font
                    cell.alignment = body_align
                    cell.border = border
                    if col_idx in (5, 6) and value is not None:
                        cell.number_format = '0.0%'
                    elif col_idx in (7, 8, 9, 10, 11) and value is not None:
                        cell.number_format = '0.0'

            col_widths = {
                "A": 18,
                "B": 14,
                "C": 14,
                "D": 10,
                "E": 14,
                "F": 16,
                "G": 11,
                "H": 11,
                "I": 14,
                "J": 11,
                "K": 10,
            }
            for col_letter, width in col_widths.items():
                ws_sum.column_dimensions[col_letter].width = width
            ws_sum.row_dimensions[1].height = 26
            ws_sum.row_dimensions[2].height = 28

            # Evitar que Excel abra el libro con hojas agrupadas/seleccionadas.
            for ws in wb_bank.worksheets:
                ws.sheet_view.tabSelected = False
            ws_sum.sheet_view.tabSelected = True
            try:
                summary_idx = wb_bank.sheetnames.index("summary")
                wb_bank.active = summary_idx
                if getattr(wb_bank, "views", None):
                    wb_bank.views[0].activeTab = summary_idx
                    wb_bank.views[0].firstSheet = summary_idx
            except Exception:
                pass

        wb_bank.save(ruta_banco_final)

    try:
        run_file_write_with_retry(
            ruta_banco_final,
            action_label="guardar el banco de coberturas",
            operation=write_bank_file,
            elapsed_seconds_fn=elapsed_seconds_fn,
        )
    except FileSaveCancelled:
        raise
    except Exception as exc:
        print(f"{Fore.YELLOW}Advertencia: No se pudo aplicar formato mmm-yy en Banco: {exc}")
    print(Fore.MAGENTA + "-> Banco de coberturas guardado")
    return ruta_banco_final


def save_pipeline_report(
    df_pipeline_report: "pd.DataFrame",
    carpeta_salida: str,
    fabricante: str,
    categoria_nombre: str,
    categoria_nombre_corto: str,
    pais_nombre: str,
    ref_month_year: str,
    coverage_label: str,
    output_descriptor: str = "",
    elapsed_seconds_fn: Optional[Callable[[], Optional[float]]] = None,
) -> str:
    categoria_para_reporte = build_output_category_segment(categoria_nombre_corto or categoria_nombre, output_descriptor)
    nombre_reporte_final = build_bounded_output_filename(
        carpeta_salida,
        f"Reporte de Pipelines {fabricante}_{categoria_para_reporte}_{pais_nombre}_{ref_month_year}_{coverage_label}.xlsx",
    )
    ruta_reporte_final = os.path.join(carpeta_salida, nombre_reporte_final)

    def write_report_file() -> None:
        report_df = df_pipeline_report.copy()
        report_df.to_excel(ruta_reporte_final, index=False, sheet_name="Reporte Pipelines")

        from openpyxl import load_workbook as _wb_load
        from openpyxl.formatting.rule import ColorScaleRule as _ColorScaleRule
        from openpyxl.formatting.rule import Rule as _Rule
        from openpyxl.styles import PatternFill as _PatternFill, Font as _Font, Alignment as _Alignment, Border as _Border, Side as _Side
        from openpyxl.styles.differential import DifferentialStyle as _Diff
        from openpyxl.utils import get_column_letter as _get_col_letter

        wb_report = _wb_load(ruta_reporte_final)
        ws = wb_report.active
        ws.title = "Reporte Pipelines"

        header_fill = _PatternFill(fill_type="solid", fgColor="404040")
        soft_header_fill = _PatternFill(fill_type="solid", fgColor="D9EAF7")
        correlation_header_fill = _PatternFill(fill_type="solid", fgColor="1F4E78")
        balanced_header_fill = _PatternFill(fill_type="solid", fgColor="548235")
        comparison_header_fill = _PatternFill(fill_type="solid", fgColor="BF9000")
        review_header_fill = _PatternFill(fill_type="solid", fgColor="C00000")
        header_font = _Font(color="FFFFFF", bold=True)
        soft_header_font = _Font(color="1F4E78", bold=True)
        header_alignment = _Alignment(horizontal="center", vertical="center", wrap_text=True)
        body_alignment = _Alignment(horizontal="center", vertical="center", wrap_text=True)
        thin_side = _Side(style="thin", color="D9D9D9")
        thin_border = _Border(left=thin_side, right=thin_side, top=thin_side, bottom=thin_side)
        selected_side = _Side(style="thick", color="000000")
        selected_border = _Border(
            left=selected_side,
            right=selected_side,
            top=selected_side,
            bottom=selected_side,
        )
        soft_red_fill = _PatternFill(fill_type="solid", fgColor="FFEBEB")
        soft_yellow_fill = _PatternFill(fill_type="solid", fgColor="FFF2CC")
        soft_green_fill = _PatternFill(fill_type="solid", fgColor="E2F0D9")
        dxf_red = _Diff(
            fill=_PatternFill(start_color='FFC7CE', end_color='FFC7CE', fill_type='solid'),
            font=_Font(color='9C0006'),
        )
        dxf_green = _Diff(
            fill=_PatternFill(start_color='C6EFCE', end_color='C6EFCE', fill_type='solid'),
            font=_Font(color='006100'),
        )
        rule_red = _Rule(type='cellIs', operator='lessThan', formula=['0'], dxf=dxf_red)
        rule_green = _Rule(type='cellIs', operator='greaterThan', formula=['0'], dxf=dxf_green)
        critical_headers = {
            'Fabricante/Marca',
            'Cesta',
            'Pipeline',
            '%VAR Cliente',
            '% VAR WP by Numerator',
            PIPELINE_REPORT_REASON_COLUMN,
            *PIPELINE_REPORT_VARIATION_COLUMNS,
            *PIPELINE_REPORT_CORRELATION_COLUMNS,
            *PIPELINE_REPORT_DIAGNOSTIC_COLUMNS,
            *PIPELINE_REPORT_AUTO_MODE_COLUMNS,
        }

        header_map: Dict[str, int] = {}
        for cell in ws[1]:
            if cell.value is None:
                continue
            header = str(cell.value).strip()
            header_map[header] = cell.column
            if header == 'Revisión requerida':
                cell.fill = review_header_fill
                cell.font = header_font
            elif header in {
                'Conflicto AUTO Correlación vs Balanceado',
                'Pérdida de correlación Balanceado',
                'Mejora gap de variación Balanceado',
            }:
                cell.fill = comparison_header_fill
                cell.font = header_font
            elif 'AUTO Correlación' in header and 'Balanceado' not in header:
                cell.fill = correlation_header_fill
                cell.font = header_font
            elif 'AUTO Balanceado' in header:
                cell.fill = balanced_header_fill
                cell.font = header_font
            elif header in critical_headers:
                cell.fill = header_fill
                cell.font = header_font
            else:
                cell.fill = soft_header_fill
                cell.font = soft_header_font
            cell.alignment = header_alignment
            cell.border = selected_border

        percent_columns = {
            '%VAR Cliente',
            '% VAR WP by Numerator',
            *PIPELINE_REPORT_VARIATION_COLUMNS,
        }
        decimal_columns = {
            'Penet Media Ano Mov Atual',
            'Raw Buyers Media Ano Mov Atual',
            'Frecuencia Media Mensual',
            'Estabilidad',
            'Gap variación seleccionado',
            'Gap variación top correlación',
            'Gap variación top',
            'Gap variación AUTO Correlación',
            'Gap variación AUTO Balanceado',
            'Mejora gap de variación Balanceado',
        }
        correlation_columns = {
            *PIPELINE_REPORT_CORRELATION_COLUMNS,
            'Correlación seleccionada',
            'Correlación top',
            'Correlación top variación',
            'Correlación AUTO Correlación',
            'Correlación AUTO Balanceado',
            'Pérdida de correlación Balanceado',
        }

        buyers_col = header_map.get('Raw Buyers Media Ano Mov Atual')
        for row_idx in range(2, ws.max_row + 1):
            row_is_low_buyers = False
            if buyers_col is not None:
                try:
                    buyers_value = ws.cell(row=row_idx, column=buyers_col).value
                    row_is_low_buyers = pd.notna(buyers_value) and float(buyers_value) < 200
                except Exception:
                    row_is_low_buyers = False
            for col_idx in range(1, ws.max_column + 1):
                cell = ws.cell(row=row_idx, column=col_idx)
                cell.alignment = body_alignment
                cell.border = thin_border
                if row_is_low_buyers:
                    cell.fill = soft_red_fill
                header = ws.cell(row=1, column=col_idx).value
                if header in percent_columns:
                    cell.number_format = '0.0'
                elif header in decimal_columns:
                    cell.number_format = '0.0'
                elif header in correlation_columns:
                    cell.number_format = '0.000'

            pipeline_col = header_map.get('Pipeline')
            if pipeline_col is None:
                continue
            try:
                selected_pipeline = int(ws.cell(row=row_idx, column=pipeline_col).value)
            except Exception:
                selected_pipeline = None
            if selected_pipeline is None:
                continue

            ws.cell(row=row_idx, column=pipeline_col).border = selected_border
            for selected_header in (
                f'Variación Cliente P{selected_pipeline}',
                f'Correlación P{selected_pipeline}',
                'Pipeline AUTO Balanceado',
            ):
                selected_col = header_map.get(selected_header)
                if selected_col is not None:
                    ws.cell(row=row_idx, column=selected_col).border = selected_border

            review_col = header_map.get('Revisión requerida')
            if review_col is not None:
                review_cell = ws.cell(row=row_idx, column=review_col)
                review_cell.fill = soft_red_fill if str(review_cell.value).strip().upper() == "SI" else soft_green_fill
            conflict_col = header_map.get('Conflicto AUTO Correlación vs Balanceado')
            if conflict_col is not None:
                conflict_cell = ws.cell(row=row_idx, column=conflict_col)
                conflict_value = str(conflict_cell.value or "").strip().lower()
                if conflict_value in {"alto", "no evaluable"}:
                    conflict_cell.fill = soft_red_fill
                elif conflict_value in {"medio", "bajo"}:
                    conflict_cell.fill = soft_yellow_fill
                elif conflict_value == "sin conflicto":
                    conflict_cell.fill = soft_green_fill

        for col_name in ['%VAR Cliente', '% VAR WP by Numerator', *PIPELINE_REPORT_VARIATION_COLUMNS, 'Estabilidad']:
            col_idx = header_map.get(col_name)
            if col_idx is None or ws.max_row < 2:
                continue
            data_range = f"{_get_col_letter(col_idx)}2:{_get_col_letter(col_idx)}{ws.max_row}"
            ws.conditional_formatting.add(data_range, rule_red)
            ws.conditional_formatting.add(data_range, rule_green)

        corr_cols = [header_map.get(col_name) for col_name in PIPELINE_REPORT_CORRELATION_COLUMNS]
        corr_cols = [col_idx for col_idx in corr_cols if col_idx is not None]
        if corr_cols and ws.max_row >= 2:
            first_corr_col = min(corr_cols)
            last_corr_col = max(corr_cols)
            corr_range = f"{_get_col_letter(first_corr_col)}2:{_get_col_letter(last_corr_col)}{ws.max_row}"
            ws.conditional_formatting.add(
                corr_range,
                _ColorScaleRule(
                    start_type='min', start_color='F8696B',
                    mid_type='percentile', mid_value=50, mid_color='FFEB84',
                    end_type='max', end_color='63BE7B',
                ),
            )

        ws.freeze_panes = "A2"
        ws.auto_filter.ref = ws.dimensions
        autofit_worksheet_columns(ws, min_width=11.0, max_width=44.0, padding=2.0)
        for reason_header in (
            PIPELINE_REPORT_REASON_COLUMN,
            'Motivo AUTO Correlación',
            'Motivo AUTO Balanceado',
        ):
            reason_col = header_map.get(reason_header)
            if reason_col is not None:
                ws.column_dimensions[ws.cell(row=1, column=reason_col).column_letter].width = 72
        wb_report.save(ruta_reporte_final)

    try:
        run_file_write_with_retry(
            ruta_reporte_final,
            action_label="guardar el reporte de pipelines",
            operation=write_report_file,
            elapsed_seconds_fn=elapsed_seconds_fn,
        )
    except FileSaveCancelled:
        raise
    except Exception as exc:
        print(f"{Fore.YELLOW}Advertencia: No se pudo guardar el reporte de pipelines: {exc}")
    print(Fore.MAGENTA + "-> Reporte de pipelines guardado")
    return ruta_reporte_final


def cleanup_temp_dir(root_dir: str) -> None:
    tmp_dir = os.path.join(root_dir, 'tmp')
    if os.path.isdir(tmp_dir):
        shutil.rmtree(tmp_dir)
        print(Fore.BLUE + "Carpeta temporal ./tmp eliminada")


class CoverageStudioUltraApp:
    def __init__(self) -> None:
        self.root_dir = os.path.dirname(os.path.abspath(__file__))
        os.chdir(self.root_dir)
        self.categories: Optional["pd.DataFrame"] = None
        self._script_start_monotonic: Optional[float] = None
        # Mantiene una asignación estable de fabricante->color durante toda la corrida.
        self._brand_color_cache: Dict[str, Tuple[int, int, int]] = {}

    def list_excel_files(self) -> List[str]:
        return [f for f in os.listdir(self.root_dir) if f.endswith('.xlsx') and not f.startswith('~$') and f != EXCEL_TEMP_FILENAME]

    def _brand_color_rgb(self, brand: str) -> Tuple[int, int, int]:
        """Asigna un color consistente y distinguible por fabricante."""
        key = normalize_brand_key(brand)
        if not key:
            return (255, 235, 59)
        cached = self._brand_color_cache.get(key)
        if cached:
            return cached
        used_colors = set(self._brand_color_cache.values())
        for candidate in TERMINAL_BRAND_COLOR_SEQUENCE:
            readable_candidate = lift_color_to_min_luminance(candidate)
            if readable_candidate not in used_colors:
                self._brand_color_cache[key] = readable_candidate
                return readable_candidate
        fallback = TERMINAL_BRAND_COLOR_SEQUENCE[len(self._brand_color_cache) % len(TERMINAL_BRAND_COLOR_SEQUENCE)]
        fallback = lift_color_to_min_luminance(fallback)
        self._brand_color_cache[key] = fallback
        return fallback

    def _colorize_filename_brand(self, filename: str) -> str:
        """Colorea solo el fabricante dentro del nombre del .xlsx."""
        prefix, brand, suffix, ext = parse_filename_brand(filename)
        if not brand:
            return filename
        rgb = self._brand_color_rgb(brand)
        # Se usa ANSI truecolor (24-bit) para tener una paleta amplia.
        brand_colored = ansi_truecolor(brand, rgb)
        return f"{prefix}{brand_colored}{Fore.BLUE}{suffix}{ext}"

    def ensure_categories_loaded(self) -> None:
        if self.categories is None:
            wait_for_heavy_modules()
            self.categories = load_categories()

    def select_files(self, excel_list: Sequence[str]) -> List[str]:
        print(Fore.CYAN + "Archivos Excel (.xlsx) encontrados:")
        max_name_len = max((len(archivo) for archivo in excel_list), default=0)
        idx_width = len(str(len(excel_list))) if excel_list else 1
        for i, archivo in enumerate(excel_list, start=1):
            meta = quick_file_metadata(archivo)
            archivo_coloreado = self._colorize_filename_brand(archivo)
            prefix = Fore.BLUE + f"{i:>{idx_width}}. "
            pad = " " * (max_name_len - len(archivo) + 3)
            if meta:
                print(prefix + f"{archivo_coloreado}{pad}" + Fore.YELLOW + f"| {meta}")
            else:
                print(prefix + f"{archivo_coloreado}")
        while True:
            opcion = input(
                Fore.WHITE
                + f"Seleccione el número de archivo a procesar (1-{len(excel_list)}).\n"
                + "Puede separar varios con comas o escribir 'all': "
            )
            opcion = opcion.strip().lower()
            if opcion in {"all", "todos", "*"}:
                selected_indices = list(range(1, len(excel_list) + 1))
            else:
                try:
                    selected_indices = [int(x) for x in opcion.split(',') if x]
                except ValueError:
                    print(Fore.RED + Style.BRIGHT + "Entrada inválida. Ingrese números separados por coma o 'all'.")
                    continue
                if not all(1 <= idx <= len(excel_list) for idx in selected_indices):
                    print(Fore.RED + "Uno o más números están fuera de rango. Intente nuevamente.")
                    continue
            selected_files = [excel_list[idx - 1] for idx in selected_indices]
            SELECTIONS['Excel'] = ", ".join(selected_files)
            clear_and_print_summary()
            return selected_files

    def gather_interactive_options(self) -> ExecutionOptions:
        coverage_type = tipo_cobertura()
        scenario_options = ExecutionOptions.from_scenario(coverage_type)
        if scenario_options:
            apply_execution_options_to_selections(scenario_options)
            clear_and_print_summary()
            return scenario_options
        else:
            coverage_type_value = coverage_type
            coverage_reason = razao_cov()
            trend_axis = tipo_eje_tendencia()
            trend_granularity = trend_granularity_option()
            variations_box_style = variations_box_style_option()
            coverage_slide_variant = coverage_slide_variant_option()
            evolution_slide_variant = evolution_slide_variant_option()
            include_english = include_english_flag()
            round_cov = round_coverage_flag()
            summary_extra_months = summary_extra_months_option()
            summary_extra_months_mode = summary_extra_months_mode_option(bool(summary_extra_months))
            auto_mode = False
        return ExecutionOptions(
            coverage_type=coverage_type_value,
            coverage_reason=coverage_reason,
            trend_axis=trend_axis,
            trend_granularity=trend_granularity,
            variations_box_style=variations_box_style,
            include_english=include_english,
            round_coverage=round_cov,
            coverage_slide_variant=coverage_slide_variant,
            evolution_slide_variant=evolution_slide_variant,
            summary_extra_months=summary_extra_months,
            summary_extra_months_mode=summary_extra_months_mode,
            variations_include_same_period_last_year=True,
            variations_compact_period_labels=False,
            auto_mode=auto_mode,
        )


    def process_file(self, excel_file_name: str, options: ExecutionOptions, idx: int, total: int) -> None:
        global ROUND_COVERAGE
        ROUND_COVERAGE = options.round_coverage
        self.ensure_categories_loaded()
        excel_file_path = os.path.join(self.root_dir, excel_file_name)
        elapsed = None
        if self._script_start_monotonic is not None:
            try:
                elapsed = time.monotonic() - float(self._script_start_monotonic)
            except Exception:
                elapsed = None
        def get_elapsed() -> Optional[float]:
            if self._script_start_monotonic is None:
                return elapsed
            try:
                return time.monotonic() - float(self._script_start_monotonic)
            except Exception:
                return elapsed
        try:
            try:
                excel_file_obj = pd.ExcelFile(excel_file_path)
                marcas = excel_file_obj.sheet_names
            except FileNotFoundError:
                print(f"{Fore.RED}{Style.BRIGHT}Error: No se encontró el archivo seleccionado: {excel_file_path}")
                return
            except PermissionError:
                # Input bloqueado (raro) o sin permisos.
                print_file_locked_error(excel_file_path, elapsed_seconds=elapsed)
                return
            except Exception as exc:
                print(f"{Fore.RED}{Style.BRIGHT}Error al abrir el archivo Excel '{excel_file_name}': {exc}")
                return

            try:
                _, category_code, _ = parse_input_filename_parts(excel_file_name)
                output_descriptor = extract_input_filename_descriptor(excel_file_name)
                pais_nombre, cesta_nombre, categoria_nombre, categoria_nombre_corto, fabricante = parse_file_metadata(excel_file_name, self.categories)
            except ValueError as exc:
                print(f"{Fore.RED}{Style.BRIGHT}{exc}")
                return

            # Asegurar que el resumen refleje opciones también en modo AUTO (sin selección interactiva).
            SELECTIONS['Excel'] = excel_file_name
            apply_execution_options_to_selections(options)

            SELECTIONS['Pais'] = pais_nombre
            clear_and_print_summary()
            print_file_header(idx, total, excel_file_name)

            coverage_label = compute_coverage_label(options.coverage_type, options.include_english)
            ref_month_year, carpeta_salida, nombre_base_archivo, ruta_template_final = generate_excel_template(
                self.root_dir,
                excel_file_obj,
                marcas,
                pais_nombre,
                categoria_nombre,
                categoria_nombre_corto,
                fabricante,
                coverage_label,
                options.coverage_type,
                options.coverage_reason,
                options.trend_axis,
                options.evolution_slide_variant,
                options.include_english,
                output_descriptor=output_descriptor,
                elapsed_seconds_fn=get_elapsed,
            )
            ruta_ppt_final, df_summary, df_bank, df_pipeline_report = generate_presentation_and_bank(
                root_dir=self.root_dir,
                excel_file_obj=excel_file_obj,
                marcas=marcas,
                pais_nombre=pais_nombre,
                category_code=category_code,
                categories_df=self.categories,
                categoria_nombre=categoria_nombre,
                categoria_nombre_corto=categoria_nombre_corto,
                fabricante=fabricante,
                cesta_nombre=cesta_nombre,
                coverage_label=coverage_label,
                coverage_type=options.coverage_type,
                coverage_reason=options.coverage_reason,
                ref_month_year=ref_month_year,
                carpeta_salida=carpeta_salida,
                nombre_base_archivo=nombre_base_archivo,
                include_english=options.include_english,
                trend_axis=options.trend_axis,
                trend_granularity=options.trend_granularity,
                variations_box_style=options.variations_box_style,
                coverage_slide_variant=options.coverage_slide_variant,
                evolution_slide_variant=options.evolution_slide_variant,
                round_coverage=options.round_coverage,
                summary_extra_months=options.summary_extra_months,
                summary_extra_months_mode=options.summary_extra_months_mode,
                variations_include_same_period_last_year=options.variations_include_same_period_last_year,
                variations_compact_period_labels=options.variations_compact_period_labels,
                optimal_pipeline_mode=options.optimal_pipeline_mode,
                elapsed_seconds_fn=get_elapsed,
            )
            ruta_banco_final = save_coverage_bank(
                df_bank=df_bank,
                carpeta_salida=carpeta_salida,
                nombre_base_archivo=nombre_base_archivo,
                fabricante=fabricante,
                categoria_nombre=categoria_nombre,
                categoria_nombre_corto=categoria_nombre_corto,
                pais_nombre=pais_nombre,
                ref_month_year=ref_month_year,
                coverage_label=coverage_label,
                coverage_type=options.coverage_type,
                coverage_slide_variant=options.coverage_slide_variant,
                output_descriptor=output_descriptor,
                elapsed_seconds_fn=get_elapsed,
            )
            ruta_pipeline_report_final = ""
            if options.optimal_pipeline_mode:
                ruta_pipeline_report_final = save_pipeline_report(
                    df_pipeline_report=df_pipeline_report,
                    carpeta_salida=carpeta_salida,
                    fabricante=fabricante,
                    categoria_nombre=categoria_nombre,
                    categoria_nombre_corto=categoria_nombre_corto,
                    pais_nombre=pais_nombre,
                    ref_month_year=ref_month_year,
                    coverage_label=coverage_label,
                    output_descriptor=output_descriptor,
                    elapsed_seconds_fn=get_elapsed,
                )
            print_file_summary(
                ruta_template_final,
                ruta_ppt_final,
                ruta_banco_final,
                ruta_pipeline_report_final,
                elapsed_seconds=elapsed,
            )
            report_zero_months_exceptions()
        except FileSaveCancelled:
            print(Fore.YELLOW + "Guardado cancelado por el usuario.")
            return
        except PermissionError as exc:
            locked_path = getattr(exc, "filename", None) or str(exc)
            print_file_locked_error(locked_path, elapsed_seconds=elapsed)
            return

    def run(self) -> None:
        if self._script_start_monotonic is None:
            self._script_start_monotonic = time.monotonic()
        excel_list = self.list_excel_files()
        if not excel_list:
            print(f"{Fore.RED}{Style.BRIGHT}Error: No se encontraron archivos .xlsx en la carpeta: {self.root_dir}")
            return
        env_options = ExecutionOptions.from_environment()
        if env_options:
            excel_file_name = os.environ['AUTO_FILE']
            idx = int(os.environ.get('AUTO_INDEX', '1'))
            total = int(os.environ.get('AUTO_TOTAL', '1'))
            self.process_file(excel_file_name, env_options, idx, total)
            cleanup_temp_dir(self.root_dir)
            return
        selected_files = self.select_files(excel_list)
        options = self.gather_interactive_options()
        total = len(selected_files)
        for idx, excel_file_name in enumerate(selected_files, start=1):
            self.process_file(excel_file_name, options, idx, total)
        cleanup_temp_dir(self.root_dir)



def main() -> None:
    app = CoverageStudioUltraApp()
    start_mono = time.monotonic()
    try:
        app._script_start_monotonic = start_mono
        app.run()
    except PermissionError as exc:
        locked_path = getattr(exc, "filename", None) or str(exc)
        print_file_locked_error(locked_path, elapsed_seconds=(time.monotonic() - start_mono))
        try:
            cleanup_temp_dir(app.root_dir)
        except Exception:
            pass
    except KeyboardInterrupt:
        end_time = datetime.now().strftime("%I:%M:%S %p")
        elapsed = _format_elapsed(time.monotonic() - start_mono)
        msg = (
            "[bright_white]Programa terminado por el usuario[/bright_white]\n\n"
            f"[white]Hora de finalizacion: [bold]{end_time}[/bold][/white]\n"
            f"[white]Tiempo total: [bold]{elapsed}[/bold][/white]\n\n"
            "[grey]Hasta luego.[/grey]"
        )
        console.print()
        console.print(Panel.fit(msg, border_style="yellow", title="Coverages Latam"))
        console.print()
        try:
            cleanup_temp_dir(app.root_dir)
        except Exception:
            pass


if __name__ == "__main__":
    main()
